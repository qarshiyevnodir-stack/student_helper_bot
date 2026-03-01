import os
import requests
import json
import logging
import random
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()

# Initialize OpenAI client (Manus pre-configured)
client = OpenAI()

def search_image(query):
    """Search for an image using Unsplash Source API with random seed to avoid duplicates."""
    try:
        # Adding a random seed to get different images for different slides
        seed = random.randint(1, 1000)
        url = f"https://source.unsplash.com/featured/?{query.replace(' ', ',')}&sig={seed}"
        response = requests.get(url, allow_redirects=True, timeout=15)
        if response.status_code == 200:
            image_path = f"temp_{hash(query)}_{seed}.jpg"
            with open(image_path, 'wb') as f:
                f.write(response.content)
            return image_path
    except Exception as e:
        logging.error(f"Error searching image for '{query}': {e}")
    return None

def generate_presentation(topic, slide_count, template_path):
    """Generate a PowerPoint presentation by strictly modifying a template."""
    
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")
        
    prs = Presentation(template_path)
    
    # 1. STICK TO SLIDE COUNT: Remove extra slides or add if needed
    # First, remove extra slides from the end
    while len(prs.slides) > slide_count:
        # Python-pptx doesn't have a direct slide removal, we have to use internal XML manipulation
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]
    
    # If template has fewer slides than requested, add new ones using the second layout (Title and Content)
    while len(prs.slides) < slide_count:
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        prs.slides.add_slide(slide_layout)

    # 2. GENERATE CONTENT using GPT
    prompt = f"""Create a professional presentation outline for the topic '{topic}' in Uzbek language. 
    Total slides: {slide_count}.
    For each slide, provide:
    1. 'title': A concise title.
    2. 'content': 3-4 bullet points of detailed information.
    3. 'image_query': 2-3 English keywords for a relevant image.
    
    Respond ONLY with a JSON object in this format:
    {{
      "slides": [
        {{
          "title": "Slayd sarlavhasi",
          "content": ["Ma'lumot 1", "Ma'lumot 2", "Ma'lumot 3"],
          "image_query": "technology computer"
        }},
        ...
      ]
    }}"""
    
    try:
        response = client.chat.completions.create(
            model="gpt-4.1-mini",
            messages=[{"role": "system", "content": "You are a professional presentation creator. You write in Uzbek language."},
                      {"role": "user", "content": prompt}],
            response_format={"type": "json_object"}
        )
        
        data = json.loads(response.choices[0].message.content)
        slides_data = data.get('slides', [])
    except Exception as e:
        logging.error(f"GPT content generation failed: {e}")
        slides_data = [{"title": topic, "content": ["Ma'lumot topilmadi."], "image_query": topic}] * slide_count

    # 3. FILL THE PRESENTATION AND REPLACE IMAGES
    for i, slide_info in enumerate(slides_data[:slide_count]):
        slide = prs.slides[i]
        
        # A. Update Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_info.get('title', '')
        
        # B. Replace ALL text in all text boxes with new content
        content_points = slide_info.get('content', [])
        content_text = "\n".join(content_points)
        
        # Find the main body text box (usually the one with most area that isn't the title)
        body_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                if body_shape is None or (shape.width * shape.height > body_shape.width * body_shape.height):
                    body_shape = shape
            
            # Clear any other text boxes to avoid "old content" showing up
            if shape.has_text_frame and shape != slide.shapes.title and shape != body_shape:
                # If it's a small text box, just clear it
                if shape.width * shape.height < (prs.slide_width * prs.slide_height * 0.1):
                    shape.text = ""

        if body_shape:
            body_shape.text = content_text

        # C. Replace Images
        image_query = slide_info.get('image_query', topic)
        new_image_path = search_image(image_query)
        
        if new_image_path:
            try:
                # Find existing pictures on the slide
                pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
                
                if pics:
                    # Replace the first existing picture's position and size
                    first_pic = pics[0]
                    left, top, width, height = first_pic.left, first_pic.top, first_pic.width, first_pic.height
                    
                    # Remove the old picture
                    sp = first_pic._element
                    sp.getparent().remove(sp)
                    
                    # Add new picture in the same spot
                    slide.shapes.add_picture(new_image_path, left, top, width=width, height=height)
                else:
                    # If no picture exists, add it to a default position (right side)
                    slide.shapes.add_picture(new_image_path, Inches(6), Inches(1.5), width=Inches(3.5))
                
                os.remove(new_image_path)
            except Exception as e:
                logging.error(f"Error replacing image on slide {i}: {e}")
                if os.path.exists(new_image_path):
                    os.remove(new_image_path)
                
    output_filename = f"temp_output_{hash(topic)}.pptx"
    prs.save(output_filename)
    return output_filename
