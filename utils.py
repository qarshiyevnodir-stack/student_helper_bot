import json
import logging
import os
import copy
import random
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

from openai import OpenAI
import traceback
from PIL import Image as PILImage

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
    base_url=os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1"),
)

PIXABAY_API_KEY = os.getenv("PIXABAY_API_KEY")

# ─────────────────────────────────────────────
# Template slide structure (1.pptx):
#   Index 0  → Slide 1: TITLE            (sarlavha)
#   Index 1  → Slide 2: TITLE_AND_BODY   (reja)
#   Index 2  → Slide 3: BLANK_1_1_1_1_1_1  (kontent 1 — 3 ustun)
#   Index 3  → Slide 4: TITLE_AND_TWO_COLUMNS_1_1 (kontent 2 — 2 ustun)
#   Index 4  → Slide 5: ONE_COLUMN_TEXT  (kontent 3 — chapda rasm, o'ngda matn)
#   Index 5  → Slide 6: BLANK_1_1        (kontent 4 — katta iqtibos)
#   Index 6  → Slide 7: CUSTOM           (kontent 5 — o'ngda rasm, chapda matn)
#   Index 7  → Slide 8: TITLE_AND_BODY_1 (xulosa)
#
# Takrorlash qoidasi:
#   5 ta slayd  → 1 marta  (3-7 slaydlar)
#   10 ta slayd → 2 marta
#   15 ta slayd → 3 marta
#   20 ta slayd → 4 marta
#   25 ta slayd → 5 marta
#   30 ta slayd → 6 marta
# ─────────────────────────────────────────────

CONTENT_SLIDE_TEMPLATE_INDICES = [2, 3, 4, 5, 6]  # Shablondagi 3-7 slaydlar indekslari


# ═══════════════════════════════════════════════════════════════
# YORDAMCHI FUNKSIYALAR
# ═══════════════════════════════════════════════════════════════

def duplicate_slide(prs, template_slide_index):
    """Shablondagi slaydni nusxalab, taqdimot oxiriga qo'shadi."""
    template_slide = prs.slides[template_slide_index]
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Yangi slayddagi standart shakllarni o'chirish
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # Shablon slayddagi barcha shakllarni nusxalash
    for shape in template_slide.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shape._element))

    return new_slide


def move_slide(prs, old_index, new_index):
    """Slaydni old_index dan new_index ga ko'chiradi."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide_elem = slides[old_index]
    xml_slides.remove(slide_elem)
    xml_slides.insert(new_index, slide_elem)


def build_slide_structure(prs, requested_content_count):
    """
    Taqdimot tuzilmasini quradi:
      - 1-2 slaydlar: har doim boshida (shablon)
      - 3-7 slaydlar: requested_content_count / 5 marta takrorlanadi
      - 8-slayd: har doim oxirida

    Qaytaradi: kontent slaydlari indekslari ro'yxati (final taqdimotda)
    """
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5

    logging.info(f"Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")

    extra_sets_needed = full_repeats - 1  # Birinchi to'plam allaqachon bor

    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")

    # Xulosa slaydini (index 7) oxiriga ko'chirish
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)

    logging.info(f"Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def find_placeholder_by_idx(slide, idx):
    """Placeholder ni indeks bo'yicha topadi."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def find_textbox_by_position(slide, min_left_cm, max_left_cm):
    """
    Berilgan chapdan oraliq koordinatada joylashgan TEXT_BOX ni topadi.
    3-slayddagi o'rta va o'ng ustunlar uchun ishlatiladi.
    """
    for shape in slide.shapes:
        left_cm = shape.left / 914400 * 2.54
        if min_left_cm <= left_cm <= max_left_cm:
            if hasattr(shape, "text_frame"):
                return shape
    return None


def auto_shrink_text(shape, text, base_font_pt, min_font_pt=10, bold=False):
    """
    Matnni shape ga yozadi. Agar matn uzun bo'lsa, shriftni kichraytiradi.
    base_font_pt dan boshlaydi, min_font_pt gacha kamaytiradi.
    """
    if shape is None:
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    # Matn uzunligiga qarab shrift o'lchamini hisoblash
    char_count = len(text)
    if char_count <= 60:
        font_pt = base_font_pt
    elif char_count <= 100:
        font_pt = max(min_font_pt, base_font_pt - 4)
    elif char_count <= 150:
        font_pt = max(min_font_pt, base_font_pt - 8)
    elif char_count <= 200:
        font_pt = max(min_font_pt, base_font_pt - 12)
    else:
        font_pt = max(min_font_pt, base_font_pt - 16)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    if bold:
        run.font.bold = True


def calc_body_font_pt(total_chars, base_pt=16, min_pt=10, max_pt=20):
    """
    Matn uzunligiga qarab shrift o'lchamini hisoblaydi.
    Matn qisqa bo'lsa kattaroq, uzun bo'lsa kichikroq shrift.
    total_chars: barcha matn belgilarining umumiy soni
    """
    if total_chars <= 150:
        return min(max_pt, base_pt + 4)   # 20pt
    elif total_chars <= 300:
        return min(max_pt, base_pt + 2)   # 18pt
    elif total_chars <= 500:
        return base_pt                    # 16pt
    elif total_chars <= 800:
        return max(min_pt, base_pt - 2)   # 14pt
    elif total_chars <= 1200:
        return max(min_pt, base_pt - 4)   # 12pt
    else:
        return max(min_pt, base_pt - 6)   # 10pt


def set_text_list_auto(shape, items, base_font_pt=18, min_font_pt=10):
    """
    Ro'yxat matnini yozadi. Elementlar ko'p yoki uzun bo'lsa shriftni kichraytiradi.
    """
    if shape is None or not items:
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    # Umumiy belgilar soniga qarab shrift o'lchamini hisoblash
    total_chars = sum(len(str(item)) for item in items)
    item_count = len(items)

    if total_chars <= 100 and item_count <= 3:
        font_pt = base_font_pt
    elif total_chars <= 200 and item_count <= 5:
        font_pt = max(min_font_pt, base_font_pt - 2)
    elif total_chars <= 350:
        font_pt = max(min_font_pt, base_font_pt - 4)
    elif total_chars <= 500:
        font_pt = max(min_font_pt, base_font_pt - 6)
    else:
        font_pt = max(min_font_pt, base_font_pt - 8)

    from pptx.enum.text import PP_ALIGN
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(item)
        run.font.size = Pt(font_pt)


def fetch_image_preview_urls(image_query, count=3):
    """
    Pixabay dan rasm URL larini qaytaradi (yuklab olmaydi).
    Foydalanuvchiga ko'rsatish uchun ishlatiladi.
    Qaytaradi: list of cdn_url yoki []
    """
    import re
    if not PIXABAY_API_KEY:
        return []
    try:
        url = (
            f"https://pixabay.com/api/"
            f"?key={PIXABAY_API_KEY}"
            f"&q={requests.utils.quote(image_query)}"
            f"&image_type=photo&orientation=horizontal"
            f"&per_page={max(count + 2, 5)}&safesearch=true"
        )
        resp = requests.get(url, timeout=10)
        resp.raise_for_status()
        hits = resp.json().get("hits", [])
        results = []
        for hit in hits:
            preview_url = hit.get("previewURL", "")
            if not preview_url:
                continue
            if preview_url.lower().endswith(".jpg"):
                cdn_url = re.sub(r'_\d+\.jpg$', '_640.jpg', preview_url)
            else:
                cdn_url = preview_url
            results.append(cdn_url)
            if len(results) >= count:
                break
        return results
    except Exception as e:
        logging.error(f"fetch_image_preview_urls xatolik ({image_query}): {e}")
        return []


def split_text_into_blocks(body_text, n):
    """
    Matnni n ta alohida gapga bo'lish.
    Har bir blok alohida gap bo'lishi kerak - bir-birining davomi emas.
    """
    import re
    if not body_text or n <= 1:
        if n == 1:
            return [body_text]
        return [body_text] + [''] * (n - 1)
    # Avval gaplarga bo'lish (. ! ? bilan tugagan)
    sentences = re.split(r'(?<=[.!?])\s+', body_text.strip())
    sentences = [s.strip() for s in sentences if s.strip()]
    if len(sentences) >= n:
        # Gaplarni n ta blokga teng taqsimlash
        blocks = []
        per_block = max(1, len(sentences) // n)
        for i in range(n):
            start = i * per_block
            end = start + per_block if i < n - 1 else len(sentences)
            block = " ".join(sentences[start:end])
            if block and not block[-1] in '.!?':
                block += '.'
            blocks.append(block)
        return blocks
    else:
        # Gaplar yetarli emas - so'zlarga bo'lish, har blok bosh harf + nuqta
        words = body_text.split()
        blocks = []
        per_block = max(1, len(words) // n)
        for i in range(n):
            start = i * per_block
            end = start + per_block if i < n - 1 else len(words)
            block = " ".join(words[start:end])
            if block:
                block = block[0].upper() + block[1:]
                if not block[-1] in '.!?':
                    block += '.'
            blocks.append(block)
        return blocks


def fetch_image(image_query):
    """
    Pixabay orqali rasm yuklab oladi.
    CDN previewURL dan _640.jpg o'lchamli rasm oladi (rate limit yo'q).
    Qaytaradi: lokal fayl yo'li yoki None.
    """
    import re
    if not PIXABAY_API_KEY:
        logging.warning("PIXABAY_API_KEY yo'q. Rasm o'tkazib yuborildi.")
        return None
    try:
        url = (f"https://pixabay.com/api/"
               f"?key={PIXABAY_API_KEY}"
               f"&q={requests.utils.quote(image_query)}"
               f"&image_type=photo&orientation=horizontal"
               f"&per_page=5&safesearch=true")
        resp = requests.get(url, timeout=10)
        resp.raise_for_status()
        hits = resp.json().get("hits", [])
        if not hits:
            logging.warning(f"Rasm topilmadi: {image_query}")
            return None

        img_data = None
        img_ext = "jpg"

        for hit in hits:
            preview_url = hit.get("previewURL", "")
            if not preview_url:
                continue

            # Faqat .jpg previewURL larni ishlatamiz (PNG uchun 640px ishlamaydi)
            if not preview_url.lower().endswith(".jpg"):
                continue

            # _150.jpg -> _640.jpg (CDN URL, rate limit yo'q)
            cdn_url = re.sub(r'_\d+\.jpg$', '_640.jpg', preview_url)
            img_resp = requests.get(cdn_url, timeout=15)
            if img_resp.status_code == 200 and "image" in img_resp.headers.get("Content-Type", ""):
                img_data = img_resp.content
                logging.info(f"Rasm yuklandi (640px): {image_query}")
                break

            # Fallback: previewURL (150px)
            img_resp2 = requests.get(preview_url, timeout=15)
            if img_resp2.status_code == 200 and "image" in img_resp2.headers.get("Content-Type", ""):
                img_data = img_resp2.content
                logging.info(f"Rasm yuklandi (150px fallback): {image_query}")
                break

        # Agar .jpg topilmasa, istalgan previewURL dan yuklab olish
        if not img_data:
            for hit in hits:
                preview_url = hit.get("previewURL", "")
                if not preview_url:
                    continue
                img_resp = requests.get(preview_url, timeout=15)
                if img_resp.status_code == 200 and "image" in img_resp.headers.get("Content-Type", ""):
                    img_data = img_resp.content
                    ct = img_resp.headers.get("Content-Type", "")
                    img_ext = "png" if "png" in ct else "jpg"
                    logging.info(f"Rasm yuklandi (preview fallback): {image_query}")
                    break

        if not img_data:
            logging.warning(f"Rasm yuklab bo'lmadi: {image_query}")
            return None

        img_path = f"/tmp/slide_img_{random.randint(0, 99999)}.{img_ext}"
        with open(img_path, "wb") as f:
            f.write(img_data)
        logging.info(f"Rasm saqlandi: {img_path}")
        return img_path
    except Exception as e:
        logging.error(f"Rasm yuklashda xatolik ({image_query}): {e}")
        return None


def place_image_in_placeholder(slide, ph_idx, img_path):
    """
    Rasmni placeholder o'lchamida va koordinatasida joylashtiradi.
    Placeholder ni o'zi esa ko'rinmas qilinadi (rasm ustiga qo'yiladi).
    """
    if not img_path:
        return
    ph = find_placeholder_by_idx(slide, ph_idx)
    if ph is None:
        logging.warning(f"Placeholder idx={ph_idx} topilmadi.")
        return
    try:
        slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
        os.remove(img_path)
        logging.info(f"Rasm placeholder idx={ph_idx} ga joylashtirildi.")
    except Exception as e:
        logging.error(f"Rasm joylashtirish xatolik: {e}")


def save_user_image_to_tmp(image_bytes):
    """
    Foydalanuvchi yuborgan rasm bytes ni /tmp ga saqlaydi.
    Pillow bilan JPEG ga convert qiladi va sifatini saqlaydi.
    Qaytaradi: lokal fayl yo'li yoki None.
    """
    if not image_bytes:
        return None
    try:
        img_path = f"/tmp/user_slide_img_{random.randint(0, 999999)}.jpg"
        img = PILImage.open(BytesIO(image_bytes))
        # RGBA -> RGB (JPEG RGBA ni qo'llab-quvvatlamaydi)
        if img.mode in ('RGBA', 'P', 'LA'):
            img = img.convert('RGB')
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        img.save(img_path, 'JPEG', quality=92)
        logging.info(f"Foydalanuvchi rasmi saqlandi: {img_path}")
        return img_path
    except Exception as e:
        logging.error(f"Foydalanuvchi rasmini saqlashda xatolik: {e}")
        return None


# ═══════════════════════════════════════════════════════════════
# GPT KONTENT YARATISH
# ═══════════════════════════════════════════════════════════════

def generate_slide_content(topic, slide_number, total_slides, language,
                           is_plan=False, is_conclusion=False,
                           slide_type=None):
    """GPT orqali slayd uchun kontent yaratadi."""

    if is_plan:
        prompt = (
            f"Mavzu: '{topic}'. Taqdimot rejasini (plan) yarat. "
            f"Faqat 3 yoki 4 ta asosiy nuqta yoz (na ko'proq, na kamroq). Til: {language}. "
            f"JSON formatida qaytarish shart: {{\"title\": \"Reja\", \"content\": [\"1-nuqta\", \"2-nuqta\", \"3-nuqta\"]}}. "
            f"content massivida faqat 3 yoki 4 ta element bo'lsin."
        )
    elif is_conclusion:
        prompt = (
            f"Mavzu: '{topic}'. Xulosa slayd uchun matn yarat. "
            f"Asosiy xulosalar. Til: {language}. "
            f"JSON formatida: {{\"title\": \"Xulosa\", \"content\": [\"...\", \"...\"]}}"
        )
    elif slide_type == "three_columns":
        # 3-slayd: sarlavha + 3 ta alohida ustun matni
        prompt = (
            f"Mavzu: '{topic}'. Bu {total_slides} ta slaydli taqdimotning {slide_number}-slaydiga kontent yarat. "
            f"Til: {language}. "
            f"Slaydda 3 ta alohida ustun bor. Har bir ustun uchun 3-5 jumlali batafsil matn yoz. "
            f"JSON formatida: {{\"title\": \"...\", \"col1\": \"...\", \"col2\": \"...\", \"col3\": \"...\", \"image_query\": \"...\"}}"
        )
    elif slide_type == "two_columns":
        # 4-slayd: sarlavha + 2 ta ustun matni (faqat 2 punkt)
        prompt = (
            f"Mavzu: '{topic}'. Bu {total_slides} ta slaydli taqdimotning {slide_number}-slaydiga kontent yarat. "
            f"Til: {language}. "
            f"Slaydda 2 ta ustun bor, har biriga 4-6 jumlali batafsil paragraf yoz. "
            f"Faqat 2 ta ustun matni kerak, ko'proq emas. "
            f"JSON formatida: {{\"title\": \"...\", \"col1\": \"...\", \"col2\": \"...\", \"image_query\": \"...\"}}"
        )
    else:
        # Oddiy slayd: sarlavha + 2 punkt matn + rasm so'rovi
        prompt = (
            f"Mavzu: '{topic}'. Bu {total_slides} ta slaydli taqdimotning {slide_number}-slaydiga kontent yarat. "
            f"Til: {language}. "
            f"2 ta punkt matn yoz, har biri 3-5 jumlali batafsil bo'lsin. "
            f"JSON formatida: {{\"title\": \"...\", \"content\": [\"...\", \"...\"], \"image_query\": \"...\"}}"
        )

    try:
        response = client.chat.completions.create(
            model="gpt-4.1-mini",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Siz taqdimot slaydlari uchun kontent yaratuvchi yordamchisiz. "
                        "Faqat JSON formatida javob bering. Matnlar berilgan tilda bo'lsin. "
                        "Matnlar qisqa va aniq bo'lsin."
                    )
                },
                {"role": "user", "content": prompt}
            ],
            response_format={"type": "json_object"}
        )
        content = response.choices[0].message.content
        return json.loads(content)
    except Exception as e:
        logging.error(f"GPT xatolik (slayd {slide_number}): {e}")
        return None


# ═══════════════════════════════════════════════════════════════
# SLAYD TO'LDIRISH FUNKSIYALARI
# ═══════════════════════════════════════════════════════════════

def fill_slide_1_title(slide, topic, name_surname):
    """
    Slayd 1 — Muqova (TITLE).
    Sarlavha: katta shrift, matn uzun bo'lsa kichrayadi.
    Subtitle: ism-familiya yoki taqdimotchi nomi.
    """
    title_ph = find_placeholder_by_idx(slide, 0)   # CENTER_TITLE
    subtitle_ph = find_placeholder_by_idx(slide, 1) # SUBTITLE

    if title_ph:
        auto_shrink_text(title_ph, topic.upper(), base_font_pt=40, min_font_pt=20, bold=True)

    if subtitle_ph:
        if name_surname and name_surname.strip():
            auto_shrink_text(subtitle_ph, name_surname, base_font_pt=24, min_font_pt=14)
        else:
            # Ism kiritilmasa, placeholder ni bo'sh qilish (shablon matni ko'rinmasin)
            subtitle_ph.text_frame.clear()
            subtitle_ph.text_frame.paragraphs[0].text = ""


def fill_slide_2_plan(slide, plan_data):
    """
    Slayd 2 — Reja (TITLE_AND_BODY).
    Sarlavha har doim 'Reja', asosiy matn raqamli ro'yxat (1. 2. 3. 4.).
    """
    title_ph = find_placeholder_by_idx(slide, 0)  # TITLE
    body_ph  = find_placeholder_by_idx(slide, 1)  # BODY

    # Sarlavha har doim "Reja"
    if title_ph:
        auto_shrink_text(title_ph, "Reja", base_font_pt=32, bold=True)

    if body_ph:
        import re
        from pptx.oxml.ns import qn
        from lxml import etree
        from pptx.enum.text import PP_ALIGN
        # Maksimal 4 ta nuqta
        content = plan_data.get("content", [])[:4]
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, item in enumerate(content):
            # Mavjud N. yoki N.M. prefiksini olib tashlash
            text = re.sub(r'^[\d]+[\d\.]*\.?\s*', '', str(item)).strip()
            label = f"{i+1}. {text}"
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            # Shablon bullet va tartib raqamini o'chirish
            pPr = p._p.get_or_add_pPr()
            # Mavjud buChar, buAutoNum, buNone larni olib tashlash
            for child in list(pPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('buChar', 'buAutoNum', 'buNone', 'buClr', 'buFont', 'buSzPct'):
                    pPr.remove(child)
            # buNone qo'shish — hech qanday belgi yo'q
            etree.SubElement(pPr, qn('a:buNone'))
            run = p.add_run()
            run.text = label
            run.font.size = Pt(20)


def fill_slide_3_three_columns(slide, content_data):
    """
    Slayd 3 — Uch ustunli kontent (BLANK_1_1_1_1_1_1).
    Sarlavha (idx=0) + 3 ta ustun (idx=1, idx=10, idx=11).
    Har bir ustun uchun alohida matn.
    """
    title_ph = find_placeholder_by_idx(slide, 0)   # TITLE
    col1_ph  = find_placeholder_by_idx(slide, 1)   # Chap ustun
    col2_ph  = find_placeholder_by_idx(slide, 10)  # O'rta ustun
    col3_ph  = find_placeholder_by_idx(slide, 11)  # O'ng ustun

    # Agar yangi placeholder topilmasa, koordinata bo'yicha topishga urinish
    if col2_ph is None:
        col2_ph = find_textbox_by_position(slide, 8.5, 11.5)
    if col3_ph is None:
        col3_ph = find_textbox_by_position(slide, 16.0, 20.0)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=28, bold=True)

    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")
    col3_text = content_data.get("col3", "")

    # Agar col1/col2/col3 yo'q bo'lsa, content ro'yxatidan foydalanish
    if not col1_text:
        items = content_data.get("content", ["", "", ""])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""
        col3_text = items[2] if len(items) > 2 else ""

    if col1_ph:
        auto_shrink_text(col1_ph, col1_text, base_font_pt=20, min_font_pt=12)
    if col2_ph:
        auto_shrink_text(col2_ph, col2_text, base_font_pt=20, min_font_pt=12)
    if col3_ph:
        auto_shrink_text(col3_ph, col3_text, base_font_pt=20, min_font_pt=12)


def fill_slide_4_two_columns(slide, content_data):
    """
    Slayd 4 — Ikki ustunli kontent (TITLE_AND_TWO_COLUMNS_1_1).
    Sarlavha + 2 ta ustun (faqat 2 ta paragraf, ko'proq bo'lsa qolganini o'tkazib yuboradi).
    """
    title_ph = find_placeholder_by_idx(slide, 0)  # TITLE
    col1_ph  = find_placeholder_by_idx(slide, 2)  # Chap ustun
    col2_ph  = find_placeholder_by_idx(slide, 1)  # O'ng ustun

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=28, bold=True)

    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")

    # Agar col1/col2 yo'q bo'lsa, content ro'yxatidan foydalanish (faqat 2 ta)
    if not col1_text:
        items = content_data.get("content", ["", ""])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""

    if col1_ph:
        auto_shrink_text(col1_ph, col1_text, base_font_pt=16, min_font_pt=10)
    if col2_ph:
        auto_shrink_text(col2_ph, col2_text, base_font_pt=16, min_font_pt=10)


def fill_slide_5_image_left(slide, content_data, image_query):
    """
    Slayd 5 — Chapda rasm, o'ngda matn (ONE_COLUMN_TEXT).
    Sarlavha (idx=0) + Asosiy matn (idx=1) + Rasm (idx=2).
    Shrift: sarlavha kichrayadi, asosiy matn 13pt.
    """
    title_ph = find_placeholder_by_idx(slide, 0)  # TITLE (o'ngda)
    body_ph  = find_placeholder_by_idx(slide, 1)  # SUBTITLE (o'ngda)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=26, min_font_pt=14, bold=True)

    if body_ph:
        items = content_data.get("content", [])
        # Faqat 2 ta paragraf, tartiblanmagan avzasiz, uzunroq matn
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(str(i)) for i in items[:2])
        if total_chars <= 150:
            font_pt = 18
        elif total_chars <= 300:
            font_pt = 16
        elif total_chars <= 500:
            font_pt = 14
        elif total_chars <= 700:
            font_pt = 12
        else:
            font_pt = 10
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from lxml import etree
        from pptx.util import Emu
        for idx_p, item in enumerate(items[:2]):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            pPr = p._p.get_or_add_pPr()
            # Barcha bullet/indent elementlarini o'chirish
            for child in list(pPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('buChar', 'buAutoNum', 'buNone', 'buClr', 'buFont', 'buSzPct', 'indent', 'marL'):
                    pPr.remove(child)
            etree.SubElement(pPr, qn('a:buNone'))
            # Indent va margin ni nolga tushirish
            pPr.set('indent', '0')
            pPr.set('marL', '0')
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_pt)

    # Rasm yuklab olish va joylashtirish
    # Agar image_query fayl yo'li bo'lsa (user_images), to'g'ridan-to'g'ri ishlatish
    if image_query and os.path.isfile(image_query):
        place_image_in_placeholder(slide, 2, image_query)
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "nature"))
        img_path = fetch_image(query)
        place_image_in_placeholder(slide, 2, img_path)


def fill_slide_6_quote(slide, content_data):
    """
    Slayd 6 — Katta iqtibos / diqqat slayd (BLANK_1_1).
    Sarlavha (idx=0) tepada, Asosiy matn (idx=1) pastda.
    Matn ko'p bo'lsa shrift kichrayadi. Tekislanish tepadan.
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    from lxml import etree
    from pptx.oxml.ns import qn
    title_ph = find_placeholder_by_idx(slide, 0)  # TITLE (tepada)
    body_ph  = find_placeholder_by_idx(slide, 1)  # SUBTITLE (pastda, katta)

    # idx=6 bo'sh qo'shimcha placeholder ni o'chirish
    extra_ph = find_placeholder_by_idx(slide, 6)
    if extra_ph is not None:
        sp = extra_ph._element
        sp.getparent().remove(sp)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=26, min_font_pt=14, bold=True)

    if body_ph:
        items = content_data.get("content", [])
        if items:
            # 4 paragraf (2 ta asosiy + har birini 2 ga bo'lish), 16pt, tepadan tekislash
            expanded = []
            for item in items[:2]:
                text = str(item)
                # Matnni taxminan 2 ga bo'lish (birinchi gap va qolgan qism)
                sentences = text.replace('! ', '!|').replace('. ', '.|').replace('? ', '?|').split('|')
                mid = max(1, len(sentences) // 2)
                part1 = ' '.join(sentences[:mid]).strip()
                part2 = ' '.join(sentences[mid:]).strip()
                if part1:
                    expanded.append(part1)
                if part2:
                    expanded.append(part2)
            full_text = "\n\n".join(expanded) if expanded else "\n\n".join(str(i) for i in items)
            auto_shrink_text(body_ph, full_text, base_font_pt=16, min_font_pt=11)
            from pptx.enum.text import MSO_ANCHOR
            body_ph.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def fill_slide_7_image_right(slide, content_data, image_query):
    """
    Slayd 7 — O'ngda rasm, chapda matn (CUSTOM).
    Sarlavha (idx=0) + Asosiy matn (idx=1) + Rasm (idx=2).
    Sarlavha tepadan 3sm, matn uzunroq, 13pt.
    """
    from pptx.util import Cm
    title_ph = find_placeholder_by_idx(slide, 0)  # TITLE (chapda)
    body_ph  = find_placeholder_by_idx(slide, 1)  # SUBTITLE (chapda)

    if title_ph:
        # Sarlavha tepadan 1.5sm, so'zlar ko'p bo'lsa shrift kichrayadi
        title_ph.top = Cm(1.5)
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=32, min_font_pt=14, bold=True)

    if body_ph:
        items = content_data.get("content", [])
        # Tartiblanmagan avzasiz, kattaroq shrift
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(str(i)) for i in items[:2])
        if total_chars <= 200:
            font_pt = 18
        elif total_chars <= 350:
            font_pt = 16
        elif total_chars <= 500:
            font_pt = 14
        else:
            font_pt = 12
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from lxml import etree
        for idx_p, item in enumerate(items[:2]):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            pPr = p._p.get_or_add_pPr()
            # Barcha bullet/indent elementlarini o'chirish
            for child in list(pPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('buChar', 'buAutoNum', 'buNone', 'buClr', 'buFont', 'buSzPct', 'indent', 'marL'):
                    pPr.remove(child)
            etree.SubElement(pPr, qn('a:buNone'))
            # Indent va margin ni nolga tushirish
            pPr.set('indent', '0')
            pPr.set('marL', '0')
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_pt)

    # Rasm yuklab olish va joylashtirish
    # Agar image_query fayl yo'li bo'lsa (user_images), to'g'ridan-to'g'ri ishlatish
    if image_query and os.path.isfile(image_query):
        place_image_in_placeholder(slide, 2, image_query)
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "nature"))
        img_path = fetch_image(query)
        place_image_in_placeholder(slide, 2, img_path)


def fill_slide_8_conclusion(slide, conclusion_data):
    """
    Slayd 8 — Xulosa / Rahmat (TITLE_AND_BODY_1).
    Sarlavha (idx=0) + Asosiy matn (idx=1).
    """
    title_ph = find_placeholder_by_idx(slide, 0)  # TITLE
    body_ph  = find_placeholder_by_idx(slide, 1)  # BODY

    if title_ph:
        auto_shrink_text(title_ph, conclusion_data.get("title", "Xulosa"), base_font_pt=32, bold=True)
    if body_ph:
        set_text_list_auto(body_ph, conclusion_data.get("content", []), base_font_pt=20)


# ═══════════════════════════════════════════════════════════════
# 2 BOSQICHLI TIZIM
# ═══════════════════════════════════════════════════════════════

# Slayd turi xaritasi (0-indexed kontent slayd pozitsiyasi → tur nomi)
SLIDE_TYPE_NAMES = {
    0: "three_columns",   # 3-slayd: 3 ustun
    1: "two_columns",     # 4-slayd: 2 ustun
    2: "image_left",      # 5-slayd: chapda rasm
    3: "quote",           # 6-slayd: iqtibos
    4: "image_right",     # 7-slayd: o'ngda rasm
}

# 3-shablon uchun slayd turi xaritasi
SLIDE_TYPE_NAMES_T3 = {
    0: "one_column",      # 3-slayd: bir ustunli matn
    1: "two_columns",     # 4-slayd: 2 ustun
    2: "three_columns",   # 5-slayd: 3 ustun
    3: "image_left",      # 6-slayd: chapda rasm
    4: "quote",           # 7-slayd: chapda matn
}


def generate_plan_with_titles(topic, slide_count, language, template_num=None):
    """
    1-BOSQICH: Mavzu bo'yicha reja va har bir kontent slayd sarlavhasini yaratadi.

    Qaytaradi:
    {
      "plan": ["1. ...", "2. ...", "3. ..."],   # 3-5 ta, slayd soniga mos
      "slide_titles": ["Sarlavha 1", "Sarlavha 2", ...]  # kontent slaydlar uchun
    }
    """
    # Slayd soniga mos reja punktlari soni
    plan_count_map = {5: 3, 10: 4, 15: 4, 20: 5, 25: 5, 30: 5}
    # Platinum1 (37) va Gamma2 (38) uchun doim 3 ta reja
    if template_num in (37, 38):
        plan_count = 3
    else:
        plan_count = plan_count_map.get(slide_count, 4)

    # Kontent slaydlar soni (5 ta shablon × takrorlash)
    content_count = slide_count  # 5, 10, 15, 20, 25, 30

    prompt = (
        f"Mavzu: '{topic}'. Taqdimot uchun reja va slayd sarlavhalarini yarat.\n"
        f"Slaydlar soni: {slide_count}. Til: {language}.\n\n"
        f"QOIDALAR:\n"
        f"1. 'plan' massivida AYNAN {plan_count} ta nuqta bo'lsin. "
        f"Har bir nuqta mavzuni chuqur yorituvchi, aniq va o'ziga xos bo'lsin. "
        f"'Kirish', 'Asosiy qism', 'Xulosa' kabi umumiy so'zlar ISHLATILMASIN. "
        f"Har bir nuqta FAQAT oddiy arab raqami bilan boshlansin: '1. matn', '2. matn' — "
        f"ichki raqamlash (1.1, 1.2) MUTLAQO ISHLATILMASIN.\n"
        f"2. 'slide_titles' massivida AYNAN {content_count} ta sarlavha bo'lsin. "
        f"Har bir sarlavha qisqa (3-6 so'z), aniq va mavzuga oid bo'lsin. "
        f"Sarlavhalar reja nuqtalariga mos va izchil bo'lsin.\n\n"
        f"JSON formatida qaytarish shart:\n"
        f"{{\"plan\": [\"1. ...\", \"2. ...\"], \"slide_titles\": [\"...\", \"...\", ...]}}"
    )

    try:
        response = client.chat.completions.create(
            model="gpt-4.1-mini",
            max_tokens=1500,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Siz taqdimot reja va sarlavhalar yaratuvchi mutaxasssissiz. "
                        "Faqat JSON formatida javob bering. "
                        "Reja nuqtalari mavzudan kelib chiqib, aniq va o'ziga xos bo'lsin. "
                        "Umumiy iboralar ishlatmang. "
                        "Reja nuqtalarida HECH QACHON ichki raqamlash (1.1, 1.2) ishlatma."
                    )
                },
                {"role": "user", "content": prompt}
            ],
            response_format={"type": "json_object"}
        )
        result = json.loads(response.choices[0].message.content)
        # Reja punktlarini cheklash va tozalash
        raw_plan = result.get("plan", [])[:plan_count]
        # Ichki raqamlashni tozalash: '1.1. matn' → '1. matn'
        import re
        clean_plan = []
        for idx, item in enumerate(raw_plan):
            # Har qanday "N.M. " yoki "N.M " prefiksini olib tashlash
            cleaned = re.sub(r'^\d+\.\d+\.?\s*', '', str(item)).strip()
            # Agar raqam bilan boshlanmasa, qo'shish
            if not re.match(r'^\d+\.', cleaned):
                cleaned = f"{idx+1}. {cleaned}"
            clean_plan.append(cleaned)
        result["plan"] = clean_plan
        result["slide_titles"] = result.get("slide_titles", [])[:content_count]
        logging.info(f"1-bosqich tayyor: {len(result['plan'])} reja, {len(result['slide_titles'])} sarlavha")
        return result
    except Exception as e:
        logging.error(f"generate_plan_with_titles xatolik: {e}")
        return None


def generate_all_content(topic, slide_count, language, slide_titles, slide_type_names=None):
    """
    2-BOSQICH: Tasdiqlangan sarlavhalar bo'yicha barcha kontent slaydlar uchun
    matnlarni BITTA GPT so'rovida yaratadi.

    Qaytaradi: list of dicts, har biri bir kontent slayd uchun.
    """
    if slide_type_names is None:
        slide_type_names = SLIDE_TYPE_NAMES
    content_count = len(slide_titles)
    slides_info = []
    for i, title in enumerate(slide_titles):
        stype = slide_type_names.get(i % 5, "image_left")
        if stype == "one_column":
            fmt = '{"title": "' + title + '", "content": ["...", "...", "...", "..."], "image_query": "..."}'
            desc = "4-6 ta paragraf, har biri 2-3 jumla, matn blokini to'liq to'ldirsin"
        elif stype == "four_columns":
            fmt = '{"title": "' + title + '", "content": ["...", "...", "...", "..."], "image_query": "..."}'
            desc = "4 ta ALOHIDA paragraf (content massivida 4 ta element - HAMMASI MAJBURIY), har biri kamida 3-5 jumla, hech biri bo'sh qolmasin"
        elif stype == "four_blocks":
            fmt = '{"title": "' + title + '", "content": ["...", "...", "...", "..."], "image_query": "..."}'
            desc = "4 ta ALOHIDA blok (content massivida 4 ta element - HAMMASI MAJBURIY), har biri kamida 3-5 jumla, hech biri bo'sh qolmasin"
        elif stype == "single_body":
            fmt = '{"title": "' + title + '", "content": ["...", "...", "..."], "image_query": "..."}'
            desc = "3 ta paragraf, har biri 3-5 jumla, image_query inglizcha"
        elif stype == "four_blocks_2x2":
            fmt = '{"title": "' + title + '", "content": ["...", "...", "...", "..."], "image_query": "..."}'
            desc = "4 ta ALOHIDA blok (content massivida 4 ta element - HAMMASI MAJBURIY), har biri 3-5 jumla, hech biri bo'sh qolmasin"
        elif stype == "numbered_list_4":
            fmt = '{"title": "' + title + '", "content": ["...", "...", "...", "..."], "image_query": "..."}'
            desc = "4 ta raqamli element (content massivida 4 ta element - HAMMASI MAJBURIY), har biri 3-5 jumla"
        elif stype in ("icon_list_3", "icon_list_3_large"):
            fmt = '{"title": "' + title + '", "content": ["...", "...", "..."], "image_query": "..."}'
            desc = "3 ta element (content massivida 3 ta element - UCHALA MAJBURIY), har biri 3-5 jumla"
        elif stype == "two_plus_one":
            fmt = '{"title": "' + title + '", "content": ["...", "...", "..."], "image_query": "..."}'
            desc = "3 ta blok (content massivida 3 ta element - UCHALA MAJBURIY), har biri 3-5 jumla"
        elif stype == "three_columns":
            fmt = '{"title": "' + title + '", "content": ["...", "...", "..."], "image_query": "..."}'
            desc = "3 ta ALOHIDA paragraf (content massivida 3 ta element - UCHALA MAJBURIY), har biri kamida 4-6 jumla, hech biri bo'sh qolmasin"
        elif stype == "two_columns":
            fmt = '{"title": "' + title + '", "content": ["...", "..."], "image_query": "..."}'
            desc = "2 ta ustun, har biri kamida 6-8 jumla, matn blokini to'ldirsin"
        elif stype in ("image_left", "image_right"):
            fmt = '{"title": "' + title + '", "content": ["...", "...", "..."], "image_query": "..."}'
            desc = "3 ta punkt, har biri 3-5 jumla, image_query inglizcha"
        else:  # quote
            fmt = '{"title": "' + title + '", "content": ["...", "..."], "image_query": "..."}'
            desc = "2 ta paragraf, har biri 3-5 jumla"
        slides_info.append(f"  Slayd {i+1} ('{title}', format: {desc}): {fmt}")

    slides_list_str = "\n".join(slides_info)

    prompt = (
        f"Mavzu: '{topic}'. Til: {language}.\n"
        f"Quyidagi {content_count} ta slayd uchun kontent yoz.\n"
        f"Har bir slayd uchun ko'rsatilgan formatda JSON ob'ekt qaytarish kerak.\n\n"
        f"{slides_list_str}\n\n"
        f"Barcha slaydlarni bitta JSON massivida qaytaring:\n"
        f"{{\"slides\": [{{...}}, {{...}}, ...]}}"
    )

    # Slayd soniga qarab max_tokens hisoblash
    max_tok_map = {5: 2000, 10: 3500, 15: 5000, 20: 6500, 25: 8000, 30: 10000}
    max_tok = max_tok_map.get(len(slide_titles), 4000)

    try:
        response = client.chat.completions.create(
            model="gpt-4.1-mini",
            max_tokens=max_tok,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Siz taqdimot slaydlari uchun kontent yaratuvchi mutaxasssissiz. "
                        "Faqat JSON formatida javob bering. "
                        "Har bir slayd uchun batafsil, ma'lumotli, to'liq matn yozing — "
                        "matnlarni qisqartirma, har bir ustun/paragraf uchun kamida 3-5 jumla yoz. "
                        "image_query har doim ingliz tilida bo'lsin."
                    )
                },
                {"role": "user", "content": prompt}
            ],
            response_format={"type": "json_object"}
        )
        result = json.loads(response.choices[0].message.content)
        slides = result.get("slides", [])
        logging.info(f"2-bosqich tayyor: {len(slides)} slayd kontent")
        return slides
    except Exception as e:
        logging.error(f"generate_all_content xatolik: {type(e).__name__}: {e}")
        logging.error(traceback.format_exc())
        return None


# ═══════════════════════════════════════════════════════════════
# ASOSIY FUNKSIYA
# ═══════════════════════════════════════════════════════════════

def generate_template_1_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None,
                                       content_data_list=None, user_images=None):
    """
    8 slaydli shablon asosida taqdimot yaratadi.

    2-bosqichli tizim bilan ishlash uchun:
      - plan: {"content": ["1. ...", "2. ...", ...]}  (1-bosqichdan)
      - content_data_list: list of dicts (2-bosqichdan generate_all_content())

    Agar content_data_list berilmasa, eski usulda har bir slayd uchun alohida GPT so'rovi yuboriladi.
    """
    logging.info(f"Taqdimot yaratilmoqda: mavzu='{topic}', slaydlar={requested_slide_count}, til={language}")

    # ── 1. Shablon tuzilmasini qurish ──
    total_content_slides = build_slide_structure(prs, requested_slide_count)
    total_slides = len(prs.slides)

    # ── 2. Kontent ma'lumotlarini tayyorlash ──

    # Reja
    if plan is None or not isinstance(plan, dict) or not plan.get("content"):
        plan = {"title": "Reja", "content": ["Kirish", "Asosiy qism", "Xulosa"]}

    # Kontent slaydlari
    if content_data_list is None:
        # Eski usul: har bir slayd uchun alohida GPT so'rovi
        slide_type_map = {
            0: "three_columns",
            1: "two_columns",
            2: None,
            3: None,
            4: None,
        }
        content_data_list = []
        for i in range(total_content_slides):
            stype = slide_type_map.get(i % 5, None)
            data = generate_slide_content(topic, i + 3, total_slides, language, slide_type=stype)
            if not data:
                if stype == "three_columns":
                    data = {"title": f"{topic} — {i+1}", "col1": "Birinchi ustun", "col2": "Ikkinchi ustun", "col3": "Uchinchi ustun", "image_query": topic}
                elif stype == "two_columns":
                    data = {"title": f"{topic} — {i+1}", "col1": "Birinchi ustun", "col2": "Ikkinchi ustun", "image_query": topic}
                else:
                    data = {"title": f"{topic} — {i+1}", "content": ["Asosiy ma'lumot", "Qo'shimcha tafsilotlar"], "image_query": topic}
            content_data_list.append(data)
    else:
        # 2-bosqich: content_data_list tayyor, faqat uzunligini tekshirish
        while len(content_data_list) < total_content_slides:
            content_data_list.append({"title": topic, "content": ["Ma'lumot"], "image_query": topic})

    # Xulosa
    conclusion = generate_slide_content(topic, total_slides, total_slides, language, is_conclusion=True)
    if not conclusion:
        conclusion = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}

    # ── 3. Slaydlarni to'ldirish ──

    # Slayd 1: Sarlavha
    fill_slide_1_title(prs.slides[0], topic, name_surname)

    # Slayd 2: Reja
    fill_slide_2_plan(prs.slides[1], plan)

    # Slaydlar 3 dan (total_slides - 2) gacha: Kontent
    user_img_idx = 0  # foydalanuvchi rasmlari indeksi
    for i in range(total_content_slides):
        slide_index = i + 2  # 0-indexed
        slide = prs.slides[slide_index]
        data = content_data_list[i]
        image_query = data.get("image_query", topic)

        slide_type = i % 5  # 0=3-slayd, 1=4-slayd, 2=5-slayd, 3=6-slayd, 4=7-slayd

        if slide_type == 0:
            fill_slide_3_three_columns(slide, data)
        elif slide_type == 1:
            fill_slide_4_two_columns(slide, data)
        elif slide_type == 2:
            # Foydalanuvchi rasmi bormi?
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_slide_5_image_left(slide, data, img_path if img_path else image_query)
            else:
                fill_slide_5_image_left(slide, data, image_query)
        elif slide_type == 3:
            fill_slide_6_quote(slide, data)
        elif slide_type == 4:
            # Foydalanuvchi rasmi bormi?
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_slide_7_image_right(slide, data, img_path if img_path else image_query)
            else:
                fill_slide_7_image_right(slide, data, image_query)

        logging.info(f"  Slayd {slide_index + 1} to'ldirildi (tur {slide_type + 3}): {data.get('title', '')}")

    # Oxirgi slayd: Xulosa
    fill_slide_8_conclusion(prs.slides[-1], conclusion)

    # ── 4. Faylni xotiraga saqlash ──
    prs_bytes = BytesIO()
    prs.save(prs_bytes)
    prs_bytes.seek(0)

    logging.info(f"Taqdimot tayyor: {total_slides} ta slayd")
    return prs_bytes


# Eski funksiya (zaxira sifatida saqlanadi)
def generate_presentation(prs, topic, slide_count, language, name_surname=""):
    return generate_template_1_presentation(prs, topic, slide_count, language, name_surname)


# ═══════════════════════════════════════════════════════════════
# 2-SHABLON TUZILMASI (2.pptx)
# ═══════════════════════════════════════════════════════════════
#
# Slayd indekslari (0-based):
#   0 → Slayd 1: TITLE            — Muqova (sarlavha + subtitle)
#   1 → Slayd 2: TITLE            — Reja (sarlavha + subtitle)
#   2 → Slayd 3: SECTION_TITLE    — Bo'lim sarlavhasi (sarlavha + tavsif)
#   3 → Slayd 4: TWO_COLUMNS      — Ikki ustunli (sarlavha + 2 ustun)
#   4 → Slayd 5: CUSTOM_1         — Uch ustunli (sarlavha + 3 ustun)
#   5 → Slayd 6: CUSTOM_1_1       — Sarlavha + asosiy matn + qo'shimcha
#   6 → Slayd 7: CUSTOM_4_1_1_1_1 — Rasm + matn (o'ngda rasm)
#   7 → Slayd 8: TITLE_AND_BODY   — Xulosa / Thanks
#
# Takrorlash: 3-7 slaydlar (index 2-6) takrorlanadi
# ═══════════════════════════════════════════════════════════════

CONTENT_SLIDE_TEMPLATE_INDICES_2 = [2, 3, 4, 5, 6]  # 2-shablondagi 3-7 slaydlar


def build_slide_structure_2(prs, requested_content_count):
    """
    2-shablon uchun tuzilma quradi.
    Xuddi 1-shablon kabi: 3-7 slaydlar takrorlanadi, 8-slayd oxirida.
    """
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5

    logging.info(f"[T2] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")

    extra_sets_needed = full_repeats - 1

    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_2:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T2] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")

    # Xulosa slaydini (index 7) oxiriga ko'chirish
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)

    logging.info(f"[T2] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


# ─── 2-Shablon slayd to'ldirish funksiyalari ───────────────────

def fill_t2_slide_1_cover(slide, topic, name_surname):
    """
    2-Shablon Slayd 1 — Muqova (TITLE).
    idx=0: Sarlavha (CENTER_TITLE, 55pt, markazda)
    idx=1: Subtitle (ism-familiya yoki sana, markazda)
    """
    title_ph    = find_placeholder_by_idx(slide, 0)
    subtitle_ph = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, topic.upper(), base_font_pt=48, min_font_pt=22, bold=True)

    if subtitle_ph:
        if name_surname and name_surname.strip():
            auto_shrink_text(subtitle_ph, name_surname, base_font_pt=22, min_font_pt=14)
        else:
            # Bo'sh qoldirish — shablon matnini o'chirish
            tf = subtitle_ph.text_frame
            tf.clear()
            if tf.paragraphs:
                tf.paragraphs[0].text = ""


def fill_t2_slide_2_plan(slide, plan_data):
    """
    2-Shablon Slayd 2 — Reja (TITLE layout).
    idx=0: Sarlavha — har doim "Reja"
    idx=1: Subtitle — raqamli ro'yxat (1. 2. 3. 4.)
    """
    import re
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, "Reja", base_font_pt=40, bold=True)

    if body_ph:
        content = plan_data.get("content", [])[:4]
        numbered = []
        for i, item in enumerate(content):
            text = re.sub(r'^[\d]+[\d\.]*\.?\s*', '', str(item)).strip()
            numbered.append(f"{i+1}. {text}")

        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True

        total_chars = sum(len(s) for s in numbered)
        if total_chars <= 150:
            font_pt = 22
        elif total_chars <= 280:
            font_pt = 18
        else:
            font_pt = 15

        for i, item in enumerate(numbered):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            # Bullet o'chirish
            pPr = p._p.get_or_add_pPr()
            for tag in ['a:buChar', 'a:buAutoNum', 'a:buFont', 'a:buSzPct']:
                for el in pPr.findall(qn(tag)):
                    pPr.remove(el)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)


def fill_t2_slide_3_section(slide, content_data):
    """
    2-Shablon Slayd 3 — Bo'lim sarlavhasi (SECTION_TITLE_AND_DESCRIPTION).
    idx=0: Sarlavha (markazda, katta)
    idx=1: Tavsif (sarlavha ostida, kichikroq)
    col1/col2/col3 yoki content ro'yxatidan matn oladi.
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=32, min_font_pt=18, bold=True)

    if body_ph:
        # col1/col2/col3 formatidan yoki content ro'yxatidan matn olish
        col1 = content_data.get("col1", "")
        col2 = content_data.get("col2", "")
        col3 = content_data.get("col3", "")
        if col1 or col2 or col3:
            # col1+col2+col3 ni bitta matn sifatida birlashtirish
            combined = " ".join(filter(None, [col1, col2, col3]))
            items = [combined]
        else:
            items = content_data.get("content", [])

        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True

        total_chars = sum(len(str(i)) for i in items[:3])
        font_pt = 18 if total_chars <= 200 else 15 if total_chars <= 400 else 12

        for idx_p, item in enumerate(items[:3]):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            for tag in ['a:buChar', 'a:buAutoNum']:
                for el in pPr.findall(qn(tag)):
                    pPr.remove(el)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_pt)

    # Shape[4] (teskari text_box, rot=180) shablon dan o'chirilgan — hech narsa qilmasa ham bo'ladi


def fill_t2_slide_4_two_columns(slide, content_data):
    """
    2-Shablon Slayd 4 — Ikki ustunli (TITLE_AND_TWO_COLUMNS).
    idx=0: Sarlavha (yuqorida, keng)
    idx=1: Chap ustun (chapdan 2.30sm)
    idx=2: O'ng ustun (chapdan 13.34sm)
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    col1_ph  = find_placeholder_by_idx(slide, 1)
    col2_ph  = find_placeholder_by_idx(slide, 2)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=28, min_font_pt=16, bold=True)

    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")
    if not col1_text:
        items = content_data.get("content", ["", ""])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""

    def write_col(ph, text):
        if not ph or not text:
            return
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True
        font_pt = 16 if len(text) <= 300 else 13 if len(text) <= 500 else 11
        p = tf.paragraphs[0]
        pPr = p._p.get_or_add_pPr()
        for tag in ['a:buChar', 'a:buAutoNum']:
            for el in pPr.findall(qn(tag)):
                pPr.remove(el)
        etree.SubElement(pPr, qn('a:buNone'))
        pPr.set('marL', '0')
        pPr.set('indent', '0')
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_pt)

    write_col(col1_ph, col1_text)
    write_col(col2_ph, col2_text)


def fill_t2_slide_5_three_columns(slide, content_data):
    """
    2-Shablon Slayd 5 — Uch ustunli (CUSTOM_1).
    idx=0: Sarlavha (yuqorida, keng)
    idx=1: Chap ustun (placeholder, chapdan 2.01sm)
    shape[2]: O'rta ustun (TEXT_BOX, chapdan 9.38sm)
    shape[3]: O'ng ustun (TEXT_BOX, chapdan 17.02sm)
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    col1_ph  = find_placeholder_by_idx(slide, 1)
    # O'rta va o'ng ustunlarni to'g'ridan-to'g'ri shape indeksi bilan topish
    col2_ph = None
    col3_ph = None
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        if not hasattr(shape, 'text_frame'):
            continue
        left_cm = shape.left / 914400 * 2.54 if shape.left else 0
        if 8.0 <= left_cm <= 12.0:
            col2_ph = shape
        elif 16.0 <= left_cm <= 21.0:
            col3_ph = shape

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=26, min_font_pt=16, bold=True)

    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")
    col3_text = content_data.get("col3", "")
    if not col1_text:
        items = content_data.get("content", ["", "", ""])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""
        col3_text = items[2] if len(items) > 2 else ""
    # Fallback: agar col3 bo'sh bo'lsa, col1 yoki col2 dan olinsin
    if not col3_text and col2_text:
        # col2 ni ikkiga bo'lib col3 ga berish
        words = col2_text.split('. ')
        half = len(words) // 2
        if half > 0:
            col3_text = '. '.join(words[half:]).strip()
            col2_text = '. '.join(words[:half]).strip()
            if col3_text and not col3_text.endswith('.'):
                col3_text += '.'
    if not col3_text and col1_text:
        col3_text = col1_text

    def write_col(ph, text):
        if not ph:
            return
        if not text:
            text = ""
        tf = ph.text_frame
        # XML darajasida to'liq tozalash — barcha paragraflarni o'chirish
        txBody = tf._txBody
        # Barcha <a:p> elementlarini o'chirib, bitta yangi qo'shish
        for p_el in txBody.findall(qn('a:p')):
            txBody.remove(p_el)
        # Yangi paragraf qo'shish
        from lxml import etree as _etree
        new_p = _etree.SubElement(txBody, qn('a:p'))
        pPr = _etree.SubElement(new_p, qn('a:pPr'))
        _etree.SubElement(pPr, qn('a:buNone'))
        pPr.set('marL', '0')
        pPr.set('indent', '0')
        # Alignment
        pPr.set('algn', 'l')
        # Run qo'shish
        new_r = _etree.SubElement(new_p, qn('a:r'))
        rPr = _etree.SubElement(new_r, qn('a:rPr'))
        rPr.set('lang', 'uz-UZ')
        rPr.set('dirty', '0')
        font_pt = 16 if len(text) <= 200 else 13 if len(text) <= 400 else 11
        rPr.set('sz', str(font_pt * 100))
        t_el = _etree.SubElement(new_r, qn('a:t'))
        t_el.text = text
        tf.word_wrap = True

    write_col(col1_ph, col1_text)
    write_col(col2_ph, col2_text)
    write_col(col3_ph, col3_text)


def fill_t2_slide_6_text(slide, content_data):
    """
    2-Shablon Slayd 6 — Sarlavha + asosiy matn + qo'shimcha (CUSTOM_1_1).
    idx=0: Sarlavha (yuqorida)
    idx=1: Asosiy matn (o'rtada, katta)
    idx=6: Qo'shimcha matn (pastda)
    """
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)
    extra_ph = find_placeholder_by_idx(slide, 6)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=26, min_font_pt=14, bold=True)

    if body_ph:
        items = content_data.get("content", [])
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        # Matnni 4 paragrafga kengaytirish
        expanded = []
        for item in items[:2]:
            text = str(item)
            sentences = text.replace('! ', '!|').replace('. ', '.|').replace('? ', '?|').split('|')
            mid = max(1, len(sentences) // 2)
            part1 = ' '.join(sentences[:mid]).strip()
            part2 = ' '.join(sentences[mid:]).strip()
            if part1:
                expanded.append(part1)
            if part2:
                expanded.append(part2)

        total_chars = sum(len(s) for s in expanded)
        font_pt = 16 if total_chars <= 300 else 14 if total_chars <= 500 else 12

        for idx_p, item in enumerate(expanded):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            for tag in ['a:buChar', 'a:buAutoNum']:
                for el in pPr.findall(qn(tag)):
                    pPr.remove(el)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)
        body_ph.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # idx=6 placeholder shablon dan o'chirilgan — hech narsa qilmasa ham bo'ladi


def fill_t2_slide_7_image(slide, content_data, image_query):
    """
    2-Shablon Slayd 7 — Rasm va matn (CUSTOM_4_1_1_1_1).
    idx=0: Sarlavha (chapda, yuqorida)
    idx=1: Matn bloki (chapda, pastroq)
    Rasm: o'ng tomonda katta ko'k to'rtburchak (GROUP shape, chapdan 14.27sm)
    """
    from pptx.util import Cm
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=28, min_font_pt=14, bold=True)

    if body_ph:
        items = content_data.get("content", [])
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(str(i)) for i in items[:2])
        font_pt = 16 if total_chars <= 200 else 13 if total_chars <= 400 else 11

        for idx_p, item in enumerate(items[:2]):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            for tag in ['a:buChar', 'a:buAutoNum']:
                for el in pPr.findall(qn(tag)):
                    pPr.remove(el)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_pt)

    # Rasm: o'ng tomondagi katta GROUP shape o'rniga rasm qo'yish
    # GROUP shape koordinatalari: left=14.27sm, top=1.97sm, width=9.44sm, height=10.35sm
    from pptx.util import Cm
    # Agar image_query fayl yo'li bo'lsa (user_images), to'g'ridan-to'g'ri ishlatish
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "nature"))
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(14.27)
            top  = Cm(1.97)
            width = Cm(9.44)
            height = Cm(10.35)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T2] Rasm slayd 7 ga joylashtirildi.")
        except Exception as e:
            logging.error(f"[T2] Rasm joylashtirish xatolik: {e}")


def fill_t2_slide_8_conclusion(slide, conclusion_data):
    """
    2-Shablon Slayd 8 — Xulosa / Thanks (TITLE_AND_BODY).
    idx=0: Sarlavha (yuqorida, chapda)
    idx=1: Asosiy matn (katta maydon)
    """
    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)

    if title_ph:
        # Xulosa sarlavhasi — mavzuga mos til bilan
        title_text = conclusion_data.get("title", "Xulosa")
        auto_shrink_text(title_ph, title_text, base_font_pt=28, min_font_pt=16, bold=True)

    if body_ph:
        set_text_list_auto(body_ph, conclusion_data.get("content", []), base_font_pt=18)


# ═══════════════════════════════════════════════════════════════
# 2-SHABLON ASOSIY FUNKSIYA
# ═══════════════════════════════════════════════════════════════

def generate_template_2_presentation(prs, topic, requested_slide_count, language,
                                      name_surname="", plan=None,
                                      content_data_list=None, user_images=None):
    """
    2-shablon (2.pptx) asosida taqdimot yaratadi.
    1-shablon bilan bir xil 2-bosqichli tizim bilan ishlaydi.
    """
    logging.info(f"[T2] Taqdimot yaratilmoqda: mavzu='{topic}', slaydlar={requested_slide_count}, til={language}")

    # ── 1. Shablon tuzilmasini qurish ──
    total_content_slides = build_slide_structure_2(prs, requested_slide_count)

    # ── 2. Kontent ma'lumotlarini tayyorlash ──

    # Reja
    if plan is None or not isinstance(plan, dict) or not plan.get("content"):
        plan = {"title": "Reja", "content": ["Kirish", "Asosiy qism", "Xulosa"]}

    # Kontent slaydlari
    if content_data_list is None:
        # Zaxira: har bir slayd uchun alohida GPT so'rovi
        slide_type_map = {
            0: "three_columns",
            1: "two_columns",
            2: None,
            3: None,
            4: None,
        }
        content_data_list = []
        for i in range(total_content_slides):
            stype = slide_type_map.get(i % 5, None)
            data = generate_slide_content(topic, i + 3, len(prs.slides), language, slide_type=stype)
            if not data:
                data = {"title": f"{topic} — {i+1}", "content": ["Ma'lumot"], "image_query": topic}
            content_data_list.append(data)
    else:
        while len(content_data_list) < total_content_slides:
            content_data_list.append({"title": topic, "content": ["Ma'lumot"], "image_query": topic})

    # Xulosa
    conclusion = generate_slide_content(topic, len(prs.slides), len(prs.slides), language, is_conclusion=True)
    if not conclusion:
        conclusion = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}

    # ── 3. Slaydlarni to'ldirish ──

    # Slayd 1: Muqova
    fill_t2_slide_1_cover(prs.slides[0], topic, name_surname)

    # Slayd 2: Reja
    fill_t2_slide_2_plan(prs.slides[1], plan)

    # Kontent slaydlar (3-7 takrorlanadi)
    # 2-shablon slayd turlari (0-indexed kontent pozitsiyasi):
    # 0 → slayd 3: SECTION_TITLE (bo'lim sarlavhasi + tavsif)
    # 1 → slayd 4: TWO_COLUMNS (2 ustun)
    # 2 → slayd 5: THREE_COLUMNS (3 ustun)
    # 3 → slayd 6: TEXT (sarlavha + katta matn)
    # 4 → slayd 7: IMAGE (rasm + matn)

    user_img_idx = 0
    for i in range(total_content_slides):
        slide_index = i + 2
        slide = prs.slides[slide_index]
        data = content_data_list[i]
        image_query = data.get("image_query", topic)
        slide_type = i % 5

        if slide_type == 0:
            fill_t2_slide_3_section(slide, data)
        elif slide_type == 1:
            fill_t2_slide_4_two_columns(slide, data)
        elif slide_type == 2:
            fill_t2_slide_5_three_columns(slide, data)
        elif slide_type == 3:
            fill_t2_slide_6_text(slide, data)
        elif slide_type == 4:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t2_slide_7_image(slide, data, img_path if img_path else image_query)
            else:
                fill_t2_slide_7_image(slide, data, image_query)

        logging.info(f"  [T2] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")

    # Oxirgi slayd: Xulosa
    fill_t2_slide_8_conclusion(prs.slides[-1], conclusion)

    # ── 4. Faylni xotiraga saqlash ──
    prs_bytes = BytesIO()
    prs.save(prs_bytes)
    prs_bytes.seek(0)

    logging.info(f"[T2] Taqdimot tayyor: {len(prs.slides)} ta slayd")
    return prs_bytes


# ═══════════════════════════════════════════════════════════════
# 3-SHABLON TUZILMASI (3.pptx)
# ═══════════════════════════════════════════════════════════════
#
# Slayd indekslari (0-based):
#   0 → Slayd 1: TITLE            — Muqova (sarlavha + subtitle)
#   1 → Slayd 2: SECTION_HEADER   — Reja (idx=0 ro'yxat, idx=2 "REJA" fixed)
#   2 → Slayd 3: ONE_COLUMN_TEXT  — Bir ustunli matn
#   3 → Slayd 4: TWO_COLUMNS      — Ikki ustunli
#   4 → Slayd 5: THREE_COLUMNS    — Uch ustunli
#   5 → Slayd 6: IMAGE_LEFT       — Chapda rasm, o'ngda matn
#   6 → Slayd 7: QUOTE            — Chapda matn, o'ngda dekorativ
#   7 → Slayd 8: CONCLUSION       — Xulosa / Thanks
#
# Takrorlash: 3-7 slaydlar (index 2-6) takrorlanadi
# ═══════════════════════════════════════════════════════════════

CONTENT_SLIDE_TEMPLATE_INDICES_3 = [2, 3, 4, 5, 6]  # 3-shablondagi 3-7 slaydlar


def build_slide_structure_3(prs, requested_content_count):
    """
    3-shablon uchun tuzilma quradi.
    3-7 slaydlar takrorlanadi, 8-slayd oxirida.
    """
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5

    logging.info(f"[T3] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")

    extra_sets_needed = full_repeats - 1

    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_3:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T3] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")

    # Xulosa slaydini (index 7) oxiriga ko'chirish
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)

    logging.info(f"[T3] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


# ─── 3-Shablon slayd to'ldirish funksiyalari ───────────────────

def fill_t3_slide_1_cover(slide, topic, name_surname):
    """
    3-Shablon Slayd 1 — Muqova (TITLE).
    idx=0: CENTER_TITLE — asosiy sarlavha (katta, oq)
    idx=1: SUBTITLE — ism-familiya yoki sana (kichik)
    """
    title_ph    = find_placeholder_by_idx(slide, 0)
    subtitle_ph = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, topic.upper(), base_font_pt=48, min_font_pt=22, bold=False)

    if subtitle_ph:
        if name_surname and name_surname.strip():
            auto_shrink_text(subtitle_ph, name_surname, base_font_pt=20, min_font_pt=12)
        else:
            tf = subtitle_ph.text_frame
            tf.clear()
            if tf.paragraphs:
                tf.paragraphs[0].text = ""


def fill_t3_slide_2_plan(slide, plan_data):
    """
    3-Shablon Slayd 2 — Reja (SECTION_HEADER).
    idx=0: TITLE — reja ro'yxati (katta maydon, chapda)
    idx=2: TITLE — "REJA" yozuvi (o'ngda, FIXED — o'zgartirmaslik)
    """
    import re
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    body_ph = find_placeholder_by_idx(slide, 0)   # Ro'yxat maydoni
    # idx=2 ni o'zgartirmaymiz — "REJA" yozuvi fixed

    if body_ph:
        content = plan_data.get("content", [])[:5]
        numbered = []
        for i, item in enumerate(content):
            text = re.sub(r'^[\d]+[\d\.]*\.?\s*', '', str(item)).strip()
            numbered.append(f"{i+1}. {text}")

        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True

        total_chars = sum(len(s) for s in numbered)
        if total_chars <= 150:
            font_pt = 24
        elif total_chars <= 300:
            font_pt = 20
        else:
            font_pt = 16

        for i, item in enumerate(numbered):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            for child in list(pPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('buChar', 'buAutoNum', 'buNone', 'buClr', 'buFont', 'buSzPct', 'indent', 'marL', 'algn'):
                    pPr.remove(child)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            pPr.set('algn', 'l')  # chapga tekislash
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)


def fill_t3_slide_3_one_column(slide, content_data):
    """
    3-Shablon Slayd 3 — Bir ustunli matn (ONE_COLUMN_TEXT).
    idx=0: TITLE — sarlavha (yuqorida)
    idx=1: SUBTITLE — katta matn bloki (markazda)
    """
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=28, min_font_pt=14, bold=True)

    if body_ph:
        from pptx.util import Inches, Emu
        from pptx.dml.color import RGBColor
        # Placeholder ni kattalashtirish: sarlavha pastidan slayd pastiga qadar
        body_ph.left   = Inches(1.35)
        body_ph.top    = Inches(1.50)
        body_ph.width  = Inches(7.21)
        body_ph.height = Inches(3.70)

        items = content_data.get("content", [])
        # Matn uzunligiga qarab shrift o'lchamini hisoblash
        total_chars = sum(len(str(i)) for i in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=10, max_pt=20)

        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True

        for idx_p, item in enumerate(items):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            for child in list(pPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('buChar', 'buAutoNum', 'buNone', 'buClr', 'buFont', 'buSzPct', 'indent', 'marL'):
                    pPr.remove(child)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        body_ph.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def fill_t3_slide_4_two_columns(slide, content_data):
    """
    3-Shablon Slayd 4 — Ikki ustunli (TITLE_AND_TWO_COLUMNS_1_1).
    idx=0: TITLE — sarlavha
    idx=1: SUBTITLE — o'ng ustun
    idx=2: SUBTITLE — chap ustun
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph  = find_placeholder_by_idx(slide, 0)
    right_ph  = find_placeholder_by_idx(slide, 1)
    left_ph   = find_placeholder_by_idx(slide, 2)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=28, min_font_pt=14, bold=True)

    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")
    if not col1_text:
        items = content_data.get("content", ["", ""])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""

    def write_col(ph, text):
        if not ph or not text:
            return
        tf = ph.text_frame
        txBody = tf._txBody
        for p_el in txBody.findall(qn('a:p')):
            txBody.remove(p_el)
        from lxml import etree as _etree
        new_p = _etree.SubElement(txBody, qn('a:p'))
        pPr = _etree.SubElement(new_p, qn('a:pPr'))
        _etree.SubElement(pPr, qn('a:buNone'))
        pPr.set('marL', '0')
        pPr.set('indent', '0')
        pPr.set('algn', 'l')
        new_r = _etree.SubElement(new_p, qn('a:r'))
        rPr = _etree.SubElement(new_r, qn('a:rPr'))
        rPr.set('lang', 'uz-UZ')
        rPr.set('dirty', '0')
        font_pt = calc_body_font_pt(len(text), base_pt=15, min_pt=10, max_pt=18)
        rPr.set('sz', str(font_pt * 100))
        t_el = _etree.SubElement(new_r, qn('a:t'))
        t_el.text = text
        tf.word_wrap = True

    write_col(left_ph, col1_text)
    write_col(right_ph, col2_text)


def fill_t3_slide_5_three_columns(slide, content_data):
    """
    3-Shablon Slayd 5 — Uch ustunli (CUSTOM_6).
    idx=0: TITLE — sarlavha
    idx=1: SUBTITLE — chap ustun
    idx=2: SUBTITLE — o'rta ustun
    idx=3: SUBTITLE — o'ng ustun
    """
    from pptx.oxml.ns import qn

    title_ph = find_placeholder_by_idx(slide, 0)
    col1_ph  = find_placeholder_by_idx(slide, 1)
    col2_ph  = find_placeholder_by_idx(slide, 2)
    col3_ph  = find_placeholder_by_idx(slide, 3)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=28, min_font_pt=14, bold=True)

    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")
    col3_text = content_data.get("col3", "")
    if not col1_text:
        items = content_data.get("content", ["", "", ""])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""
        col3_text = items[2] if len(items) > 2 else ""
    # Fallback: agar col3 bo'sh bo'lsa
    if not col3_text and col2_text:
        words = col2_text.split('. ')
        half = len(words) // 2
        if half > 0:
            col3_text = '. '.join(words[half:]).strip()
            col2_text = '. '.join(words[:half]).strip()
            if col3_text and not col3_text.endswith('.'):
                col3_text += '.'
    if not col3_text and col1_text:
        col3_text = col1_text

    def write_col(ph, text):
        if not ph or not text:
            return
        tf = ph.text_frame
        txBody = tf._txBody
        for p_el in txBody.findall(qn('a:p')):
            txBody.remove(p_el)
        from lxml import etree as _etree
        new_p = _etree.SubElement(txBody, qn('a:p'))
        pPr = _etree.SubElement(new_p, qn('a:pPr'))
        _etree.SubElement(pPr, qn('a:buNone'))
        pPr.set('marL', '0')
        pPr.set('indent', '0')
        pPr.set('algn', 'l')
        new_r = _etree.SubElement(new_p, qn('a:r'))
        rPr = _etree.SubElement(new_r, qn('a:rPr'))
        rPr.set('lang', 'uz-UZ')
        rPr.set('dirty', '0')
        font_pt = calc_body_font_pt(len(text), base_pt=13, min_pt=9, max_pt=16)
        rPr.set('sz', str(font_pt * 100))
        t_el = _etree.SubElement(new_r, qn('a:t'))
        t_el.text = text
        tf.word_wrap = True

    write_col(col1_ph, col1_text)
    write_col(col2_ph, col2_text)
    write_col(col3_ph, col3_text)


def fill_t3_slide_6_image_left(slide, content_data, image_query):
    """
    3-Shablon Slayd 6 — Chapda rasm, o'ngda matn (CUSTOM).
    idx=0: TITLE — sarlavha (yuqorida, keng)
    idx=1: SUBTITLE — matn (o'ngda)
    idx=2: PICTURE — rasm placeholder (chapda, katta)
    """
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=28, min_font_pt=14, bold=True)

    if body_ph:
        items = content_data.get("content", [])
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True

        total_chars = sum(len(str(i)) for i in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=17, min_pt=12, max_pt=21)

        for idx_p, item in enumerate(items):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            for child in list(pPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('buChar', 'buAutoNum', 'buNone', 'buClr', 'buFont', 'buSzPct', 'indent', 'marL'):
                    pPr.remove(child)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_pt)
        body_ph.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Rasm placeholder (idx=2) ga rasm qo'yish
    if image_query and os.path.isfile(image_query):
        place_image_in_placeholder(slide, 2, image_query)
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "nature"))
        img_path = fetch_image(query)
        place_image_in_placeholder(slide, 2, img_path)


def fill_t3_slide_7_quote(slide, content_data, image_query=None):
    """
    3-Shablon Slayd 7 — Chapda matn, o'ngda rasm (CUSTOM_4_1).
    idx=0: TITLE — sarlavha (chapda, yuqorida)
    idx=1: SUBTITLE — katta matn bloki (chapda)
    O'ng tomonda GROUP shape ustiga rasm qo'yiladi.
    """
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=26, min_font_pt=14, bold=True)

    if body_ph:
        items = content_data.get("content", [])
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True

        total_chars = sum(len(str(i)) for i in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=17, min_pt=12, max_pt=21)

        for idx_p, item in enumerate(items):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            for child in list(pPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('buChar', 'buAutoNum', 'buNone', 'buClr', 'buFont', 'buSzPct', 'indent', 'marL'):
                    pPr.remove(child)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_pt)
        body_ph.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # O'ng tomonga rasm qo'yish (GROUP shape o'rniga)
    from pptx.util import Inches
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "technology"))
        final_img_path = fetch_image(query)
    if final_img_path:
        # GROUP shape pozitsiyasi: left=4.46", top=1.24", width=4.39", height=2.93"
        left   = Inches(4.46)
        top    = Inches(1.24)
        width  = Inches(4.39)
        height = Inches(2.93)
        # GROUP shape ni o'chirish
        for shape in list(slide.shapes):
            if shape.shape_type == 6:  # GROUP
                sp = shape._element
                sp.getparent().remove(sp)
                break
        slide.shapes.add_picture(final_img_path, left, top, width, height)
        if os.path.isfile(final_img_path):
            os.remove(final_img_path)


def fill_t3_slide_8_conclusion(slide, conclusion_data):
    """
    3-Shablon Slayd 8 — Xulosa / Thanks (CUSTOM_11).
    idx=0: TITLE — sarlavha (chapda, yuqorida)
    TEXT_BOX shape — qo'shimcha matn (agar mavjud bo'lsa)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    title_ph = find_placeholder_by_idx(slide, 0)

    if title_ph:
        title_text = conclusion_data.get("title", "Xulosa")
        auto_shrink_text(title_ph, title_text, base_font_pt=32, min_font_pt=16, bold=False)

    # TEXT_BOX shape ga xulosa matnini yozish (agar mavjud bo'lsa)
    for shape in slide.shapes:
        try:
            pf = shape.placeholder_format
            if pf:
                continue
        except Exception:
            pass
        if shape.has_text_frame and shape.shape_type == 17:  # TEXT_BOX
            items = conclusion_data.get("content", [])
            if items:
                from pptx.dml.color import RGBColor
                tf = shape.text_frame
                tf.clear()
                tf.word_wrap = True
                total_chars = sum(len(str(it)) for it in items)
                font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)
                for i, item in enumerate(items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = str(item)
                    run.font.size = Pt(font_pt)
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            break


# ═══════════════════════════════════════════════════════════════
# 3-SHABLON ASOSIY FUNKSIYA
# ═══════════════════════════════════════════════════════════════

def generate_template_3_presentation(prs, topic, requested_slide_count, language,
                                      name_surname="", plan=None,
                                      content_data_list=None, user_images=None):
    """
    3-shablon (3.pptx) asosida taqdimot yaratadi.
    1-shablon va 2-shablon bilan bir xil 2-bosqichli tizim.
    """
    logging.info(f"[T3] Taqdimot yaratilmoqda: mavzu='{topic}', slaydlar={requested_slide_count}, til={language}")

    # ── 1. Shablon tuzilmasini qurish ──
    total_content_slides = build_slide_structure_3(prs, requested_slide_count)

    # ── 2. Kontent ma'lumotlarini tayyorlash ──

    # Reja
    if plan is None or not isinstance(plan, dict) or not plan.get("content"):
        plan = {"title": "Reja", "content": ["Kirish", "Asosiy qism", "Xulosa"]}

    # Kontent slaydlari
    if content_data_list is None:
        slide_type_map = {
            0: None,            # slayd 3: bir ustunli
            1: "two_columns",   # slayd 4: ikki ustunli
            2: "three_columns", # slayd 5: uch ustunli
            3: None,            # slayd 6: rasm chapda
            4: None,            # slayd 7: matn chapda
        }
        content_data_list = []
        for i in range(total_content_slides):
            stype = slide_type_map.get(i % 5, None)
            data = generate_slide_content(topic, i + 3, len(prs.slides), language, slide_type=stype)
            if not data:
                if stype == "three_columns":
                    data = {"title": f"{topic} — {i+1}", "col1": "Birinchi ustun", "col2": "Ikkinchi ustun", "col3": "Uchinchi ustun", "image_query": topic}
                elif stype == "two_columns":
                    data = {"title": f"{topic} — {i+1}", "col1": "Birinchi ustun", "col2": "Ikkinchi ustun", "image_query": topic}
                else:
                    data = {"title": f"{topic} — {i+1}", "content": ["Asosiy ma'lumot", "Qo'shimcha tafsilotlar"], "image_query": topic}
            content_data_list.append(data)
    else:
        while len(content_data_list) < total_content_slides:
            content_data_list.append({"title": topic, "content": ["Ma'lumot"], "image_query": topic})

    # Xulosa
    conclusion = generate_slide_content(topic, len(prs.slides), len(prs.slides), language, is_conclusion=True)
    if not conclusion:
        conclusion = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}

    # ── 3. Slaydlarni to'ldirish ──

    # Slayd 1: Muqova
    fill_t3_slide_1_cover(prs.slides[0], topic, name_surname)

    # Slayd 2: Reja
    fill_t3_slide_2_plan(prs.slides[1], plan)

    # Kontent slaydlar (3-7 takrorlanadi)
    # 3-shablon slayd turlari (0-indexed kontent pozitsiyasi):
    # 0 → slayd 3: ONE_COLUMN (bir ustunli matn)
    # 1 → slayd 4: TWO_COLUMNS (ikki ustunli)
    # 2 → slayd 5: THREE_COLUMNS (uch ustunli)
    # 3 → slayd 6: IMAGE_LEFT (chapda rasm)
    # 4 → slayd 7: QUOTE (chapda matn, o'ngda dekorativ)

    user_img_idx = 0
    for i in range(total_content_slides):
        slide_index = i + 2
        slide = prs.slides[slide_index]
        data = content_data_list[i]
        image_query = data.get("image_query", topic)
        slide_type = i % 5

        if slide_type == 0:
            fill_t3_slide_3_one_column(slide, data)
        elif slide_type == 1:
            fill_t3_slide_4_two_columns(slide, data)
        elif slide_type == 2:
            fill_t3_slide_5_three_columns(slide, data)
        elif slide_type == 3:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t3_slide_6_image_left(slide, data, img_path if img_path else image_query)
            else:
                fill_t3_slide_6_image_left(slide, data, image_query)
        elif slide_type == 4:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t3_slide_7_quote(slide, data, img_path if img_path else image_query)
            else:
                fill_t3_slide_7_quote(slide, data, image_query)

        logging.info(f"  [T3] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")

    # Oxirgi slayd: Xulosa
    fill_t3_slide_8_conclusion(prs.slides[-1], conclusion)

    # ── 4. Faylni xotiraga saqlash ──
    prs_bytes = BytesIO()
    prs.save(prs_bytes)
    prs_bytes.seek(0)

    logging.info(f"[T3] Taqdimot tayyor: {len(prs.slides)} ta slayd")
    return prs_bytes


# ═══════════════════════════════════════════════════════════════
# 4-SHABLON YORDAMCHI FUNKSIYALAR
# ═══════════════════════════════════════════════════════════════

# 4-shablon uchun slayd turi xaritasi
SLIDE_TYPE_NAMES_T4 = {
    0: "one_column",       # 3-slayd: bir ustunli matn
    1: "three_columns",    # 4-slayd: uch ustunli
    2: "two_columns",      # 5-slayd: ikki ustunli
    3: "image_left",       # 6-slayd: chapda matn, o'ngda rasm
    4: "three_textboxes",  # 7-slayd: uch TEXT_BOX
}

# 4-shablon kontent slayd indekslari (template da 3-7 slaydlar)
CONTENT_SLIDE_TEMPLATE_INDICES_4 = [2, 3, 4, 5, 6]


def build_slide_structure_4(prs, requested_content_count):
    """
    4-shablon uchun slayd tuzilmasini quradi.
    3-shablon kabi duplicate_slide va move_slide ishlatadi.
    1-slayd: muqova, 2-slayd: reja, 3-7: kontent (5 tur), 8: xulosa
    """
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5

    logging.info(f"[T4] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")

    extra_sets_needed = full_repeats - 1

    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_4:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T4] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")

    # Xulosa slaydini (index 7) oxiriga ko'chirish
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)

    logging.info(f"[T4] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def fill_t4_slide_1_cover(slide, topic, name_surname):
    """
    4-Shablon Slayd 1 — Muqova (TITLE layout).
    idx=0 CENTER_TITLE: katta sarlavha (chapda)
    idx=1 SUBTITLE: ism-familiya (pastda)
    """
    title_ph = find_placeholder_by_idx(slide, 0)
    sub_ph   = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, topic.upper(), base_font_pt=54, min_font_pt=28, bold=True)

    if sub_ph:
        tf = sub_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name_surname or ""
        run.font.size = Pt(16)


def fill_t4_slide_2_plan(slide, plan_data):
    """
    4-Shablon Slayd 2 — Reja (BLANK_1_1_1_1_1_1 layout).
    idx=0 TITLE: "REJA" — o'zgartirmaslik
    idx=1..4 SUBTITLE: reja bandlari (4 ta)
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    items = plan_data.get("content", [])

    for i in range(1, 5):
        ph = find_placeholder_by_idx(slide, i)
        if ph and (i - 1) < len(items):
            text = str(items[i - 1])
            # Raqam prefiksini olib tashlash (agar bor bo'lsa)
            import re
            text = re.sub(r'^\d+[\.\)]\s*', '', text).strip()
            tf = ph.text_frame
            txBody = tf._txBody
            for p_el in txBody.findall(qn('a:p')):
                txBody.remove(p_el)
            new_p = etree.SubElement(txBody, qn('a:p'))
            pPr = etree.SubElement(new_p, qn('a:pPr'))
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            pPr.set('algn', 'l')
            new_r = etree.SubElement(new_p, qn('a:r'))
            rPr = etree.SubElement(new_r, qn('a:rPr'))
            rPr.set('lang', 'uz-UZ')
            rPr.set('dirty', '0')
            rPr.set('sz', '1600')
            t_el = etree.SubElement(new_r, qn('a:t'))
            t_el.text = text
            tf.word_wrap = True
        elif ph:
            # Bo'sh qoldirilsin
            tf = ph.text_frame
            txBody = tf._txBody
            for p_el in txBody.findall(qn('a:p')):
                txBody.remove(p_el)
            etree.SubElement(txBody, qn('a:p'))


def fill_t4_slide_3_one_column(slide, content_data):
    """
    4-Shablon Slayd 3 — Bir ustunli matn (CUSTOM_4 layout).
    idx=0 TITLE: sarlavha (chapda)
    idx=1 SUBTITLE: asosiy matn bloki (chapda)
    """
    from pptx.enum.text import MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=24, min_font_pt=14, bold=True)

    if body_ph:
        items = content_data.get("content", [])
        total_chars = sum(len(str(i)) for i in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=15, min_pt=10, max_pt=18)

        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True

        for idx_p, item in enumerate(items):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            for child in list(pPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('buChar', 'buAutoNum', 'buNone', 'buClr', 'buFont', 'buSzPct', 'indent', 'marL'):
                    pPr.remove(child)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_pt)
        body_ph.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def fill_t4_slide_4_three_columns(slide, content_data):
    """
    4-Shablon Slayd 4 — Uch ustunli (CUSTOM_6_1_1 layout).
    idx=0 TITLE: sarlavha
    idx=1 SUBTITLE: chap ustun
    idx=2 SUBTITLE: o'rta ustun
    idx=3 SUBTITLE: o'ng ustun
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph = find_placeholder_by_idx(slide, 0)
    col1_ph  = find_placeholder_by_idx(slide, 1)
    col2_ph  = find_placeholder_by_idx(slide, 2)
    col3_ph  = find_placeholder_by_idx(slide, 3)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=24, min_font_pt=14, bold=True)

    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")
    col3_text = content_data.get("col3", "")
    if not col1_text:
        items = content_data.get("content", ["", "", ""])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""
        col3_text = items[2] if len(items) > 2 else ""
    if not col3_text and col2_text:
        words = col2_text.split('. ')
        half = len(words) // 2
        if half > 0:
            col3_text = '. '.join(words[half:]).strip()
            col2_text = '. '.join(words[:half]).strip()

    def write_col(ph, text):
        if not ph or not text:
            return
        tf = ph.text_frame
        txBody = tf._txBody
        for p_el in txBody.findall(qn('a:p')):
            txBody.remove(p_el)
        new_p = etree.SubElement(txBody, qn('a:p'))
        pPr = etree.SubElement(new_p, qn('a:pPr'))
        etree.SubElement(pPr, qn('a:buNone'))
        pPr.set('marL', '0')
        pPr.set('indent', '0')
        pPr.set('algn', 'l')
        new_r = etree.SubElement(new_p, qn('a:r'))
        rPr = etree.SubElement(new_r, qn('a:rPr'))
        rPr.set('lang', 'uz-UZ')
        rPr.set('dirty', '0')
        font_pt = calc_body_font_pt(len(text), base_pt=13, min_pt=9, max_pt=16)
        rPr.set('sz', str(font_pt * 100))
        t_el = etree.SubElement(new_r, qn('a:t'))
        t_el.text = text
        tf.word_wrap = True

    write_col(col1_ph, col1_text)
    write_col(col2_ph, col2_text)
    write_col(col3_ph, col3_text)


def fill_t4_slide_5_two_columns(slide, content_data):
    """
    4-Shablon Slayd 5 — Ikki ustunli (TITLE_AND_TWO_COLUMNS_1_1 layout).
    idx=0 TITLE: sarlavha
    idx=1 SUBTITLE: o'ng ustun
    idx=2 SUBTITLE: chap ustun
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    title_ph  = find_placeholder_by_idx(slide, 0)
    right_ph  = find_placeholder_by_idx(slide, 1)
    left_ph   = find_placeholder_by_idx(slide, 2)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=24, min_font_pt=14, bold=True)

    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")
    if not col1_text:
        items = content_data.get("content", ["", ""])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""

    def write_col(ph, text):
        if not ph or not text:
            return
        tf = ph.text_frame
        txBody = tf._txBody
        for p_el in txBody.findall(qn('a:p')):
            txBody.remove(p_el)
        new_p = etree.SubElement(txBody, qn('a:p'))
        pPr = etree.SubElement(new_p, qn('a:pPr'))
        etree.SubElement(pPr, qn('a:buNone'))
        pPr.set('marL', '0')
        pPr.set('indent', '0')
        pPr.set('algn', 'l')
        new_r = etree.SubElement(new_p, qn('a:r'))
        rPr = etree.SubElement(new_r, qn('a:rPr'))
        rPr.set('lang', 'uz-UZ')
        rPr.set('dirty', '0')
        font_pt = calc_body_font_pt(len(text), base_pt=15, min_pt=10, max_pt=18)
        rPr.set('sz', str(font_pt * 100))
        t_el = etree.SubElement(new_r, qn('a:t'))
        t_el.text = text
        tf.word_wrap = True

    write_col(left_ph, col1_text)
    write_col(right_ph, col2_text)


def fill_t4_slide_6_image_left(slide, content_data, image_query):
    """
    4-Shablon Slayd 6 — Chapda matn, o'ngda rasm (CUSTOM_4 layout).
    idx=0 TITLE: sarlavha
    idx=1 SUBTITLE: matn bloki (chapda)
    O'ngda katta gradient shape ustiga rasm qo'yiladi: (5.0",1.53") 5.5"x2.47"
    """
    from pptx.enum.text import MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree
    from pptx.util import Inches

    title_ph = find_placeholder_by_idx(slide, 0)
    body_ph  = find_placeholder_by_idx(slide, 1)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=24, min_font_pt=14, bold=True)

    if body_ph:
        items = content_data.get("content", [])
        total_chars = sum(len(str(i)) for i in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=15, min_pt=10, max_pt=18)

        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True

        for idx_p, item in enumerate(items):
            if idx_p == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            for child in list(pPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('buChar', 'buAutoNum', 'buNone', 'buClr', 'buFont', 'buSzPct', 'indent', 'marL'):
                    pPr.remove(child)
            etree.SubElement(pPr, qn('a:buNone'))
            pPr.set('marL', '0')
            pPr.set('indent', '0')
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_pt)
        body_ph.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # O'ngdagi gradient shape ustiga rasm qo'yish
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "technology"))
        final_img_path = fetch_image(query)
    if final_img_path:
        left   = Inches(5.0)
        top    = Inches(1.53)
        width  = Inches(4.77)
        height = Inches(2.47)
        slide.shapes.add_picture(final_img_path, left, top, width, height)
        if os.path.isfile(final_img_path):
            os.remove(final_img_path)


def fill_t4_slide_7_three_textboxes(slide, content_data):
    """
    4-Shablon Slayd 7 — Uch TEXT_BOX (TITLE_ONLY layout).
    idx=0 TITLE: sarlavha
    TEXT_BOX 1: (0.79",1.44") 3.8"x1.36" — chap yuqori
    TEXT_BOX 2: (5.41",2.63") 3.8"x1.36" — o'ng o'rta
    TEXT_BOX 3: (0.79",4.0") 3.8"x1.2" — chap pastki
    """
    from pptx.util import Pt

    title_ph = find_placeholder_by_idx(slide, 0)

    if title_ph:
        auto_shrink_text(title_ph, content_data.get("title", ""), base_font_pt=24, min_font_pt=14, bold=True)

    # 3 ta TEXT_BOX topish (shape_type == 17)
    textboxes = []
    for shape in slide.shapes:
        if shape.shape_type == 17 and shape.has_text_frame:  # TEXT_BOX
            textboxes.append(shape)

    items = content_data.get("content", [])
    # col1/col2/col3 dan ham olish
    if not items:
        items = [
            content_data.get("col1", ""),
            content_data.get("col2", ""),
            content_data.get("col3", ""),
        ]
    items = [str(i) for i in items if i]

    for i, tb in enumerate(textboxes[:3]):
        if i < len(items):
            text = items[i]
            font_pt = calc_body_font_pt(len(text), base_pt=14, min_pt=9, max_pt=16)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_pt)


def fill_t4_slide_8_conclusion(slide, conclusion_data):
    """
    4-Shablon Slayd 8 — Xulosa (TITLE_ONLY layout).
    idx=0 TITLE: sarlavha ("Thanks!" o'rniga xulosa sarlavhasi)
    idx=4294967295 BODY: katta matn maydoni
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    title_ph = find_placeholder_by_idx(slide, 0)
    if title_ph:
        title_text = conclusion_data.get("title", "Xulosa")
        auto_shrink_text(title_ph, title_text, base_font_pt=32, min_font_pt=16, bold=True)

    # BODY placeholder (idx=4294967295) ga xulosa matnini yozish
    for shape in slide.shapes:
        try:
            pf = shape.placeholder_format
            if pf and pf.idx == 4294967295:
                items = conclusion_data.get("content", [])
                if items:
                    tf = shape.text_frame
                    tf.clear()
                    tf.word_wrap = True
                    total_chars = sum(len(str(it)) for it in items)
                    font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)
                    for i, item in enumerate(items):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = str(item)
                        run.font.size = Pt(font_pt)
                break
        except Exception:
            pass


# ═══════════════════════════════════════════════════════════════
# 4-SHABLON ASOSIY FUNKSIYA
# ═══════════════════════════════════════════════════════════════

def generate_template_4_presentation(prs, topic, requested_slide_count, language,
                                      name_surname="", plan=None,
                                      content_data_list=None, user_images=None):
    """
    4-shablon (4.pptx) asosida taqdimot yaratadi.
    """
    logging.info(f"[T4] Taqdimot yaratilmoqda: mavzu='{topic}', slaydlar={requested_slide_count}, til={language}")

    # ── 1. Shablon tuzilmasini qurish ──
    total_content_slides = build_slide_structure_4(prs, requested_slide_count)

    # ── 2. Kontent ma'lumotlarini tayyorlash ──
    if plan is None or not isinstance(plan, dict) or not plan.get("content"):
        plan = {"title": "Reja", "content": ["Kirish", "Asosiy qism", "Xulosa"]}

    if content_data_list is None:
        content_data_list = []
        for i in range(total_content_slides):
            stype_name = SLIDE_TYPE_NAMES_T4.get(i % 5, "one_column")
            if stype_name == "three_columns":
                stype = "three_columns"
            elif stype_name == "two_columns":
                stype = "two_columns"
            elif stype_name == "image_left":
                stype = "image_left"
            elif stype_name == "three_textboxes":
                stype = "three_columns"  # GPT uchun three_columns formatini ishlatamiz
            else:
                stype = None
            data = generate_slide_content(topic, i + 3, len(prs.slides), language, slide_type=stype)
            if not data:
                data = {"title": f"{topic} — {i+1}", "content": ["Asosiy ma'lumot"], "image_query": topic}
            content_data_list.append(data)
    else:
        while len(content_data_list) < total_content_slides:
            content_data_list.append({"title": topic, "content": ["Ma'lumot"], "image_query": topic})

    conclusion = generate_slide_content(topic, len(prs.slides), len(prs.slides), language, is_conclusion=True)
    if not conclusion:
        conclusion = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}

    # ── 3. Slaydlarni to'ldirish ──
    fill_t4_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_t4_slide_2_plan(prs.slides[1], plan)

    user_img_idx = 0
    for i in range(total_content_slides):
        slide_index = i + 2
        slide = prs.slides[slide_index]
        data = content_data_list[i]
        image_query = data.get("image_query", topic)
        slide_type = i % 5

        if slide_type == 0:
            fill_t4_slide_3_one_column(slide, data)
        elif slide_type == 1:
            fill_t4_slide_4_three_columns(slide, data)
        elif slide_type == 2:
            fill_t4_slide_5_two_columns(slide, data)
        elif slide_type == 3:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t4_slide_6_image_left(slide, data, img_path if img_path else image_query)
            else:
                fill_t4_slide_6_image_left(slide, data, image_query)
        elif slide_type == 4:
            fill_t4_slide_7_three_textboxes(slide, data)

        logging.info(f"  [T4] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")

    fill_t4_slide_8_conclusion(prs.slides[-1], conclusion)

    # ── 4. Faylni xotiraga saqlash ──
    prs_bytes = BytesIO()
    prs.save(prs_bytes)
    prs_bytes.seek(0)

    logging.info(f"[T4] Taqdimot tayyor: {len(prs.slides)} ta slayd")
    return prs_bytes


# ═══════════════════════════════════════════════════════════════
# 5-SHABLON YORDAMCHI FUNKSIYALAR (YANGI 5.pptx ga mos)
# ═══════════════════════════════════════════════════════════════
# 5-shablon uchun slayd turi xaritasi (yangi tuzilma)
SLIDE_TYPE_NAMES_T5 = {
    0: "image_left",     # 3-slayd: rasm o'ngda, matn chapda
    1: "two_columns",    # 4-slayd: ikki ustunli (idx=2, idx=4)
    2: "two_staggered",  # 5-slayd: ikki offset blok
    3: "image_left",     # 6-slayd: rasm o'ngda, kichik sarlavha
    4: "two_columns",    # 7-slayd: ikki kvadrat blok
}
# 5-shablon kontent slayd indekslari (template da 3-7 slaydlar = index 2-6)
CONTENT_SLIDE_TEMPLATE_INDICES_5 = [2, 3, 4, 5, 6]


def build_slide_structure_5(prs, requested_content_count):
    """
    5-shablon uchun slayd tuzilmasini quradi.
    duplicate_slide va move_slide ishlatadi.
    """
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5
    logging.info(f"[T5] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_5:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T5] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    # Xulosa slaydini (index 7) oxiriga ko'chirish
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T5] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def fill_t5_slide_1_cover(slide, topic, name_surname):
    """
    5-Shablon Slayd 1 — Muqova.
    idx=0: Katta sarlavha (o'ngda) — KATTA HARFDA
    idx=1: Subtitle — ism-familiya
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not hasattr(shape, 'placeholder_format') or not shape.placeholder_format:
            continue
        idx = shape.placeholder_format.idx
        tf = shape.text_frame
        tf.clear()
        if idx == 0:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = topic.upper()
        elif idx == 1:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = name_surname


def fill_t5_slide_2_plan(slide, plan_data):
    """
    5-Shablon Slayd 2 — Reja.
    idx=0: 'REJA' — o'zgartirmaslik
    idx=1: Reja bandlari ro'yxati — chapga tekislangan, raqamli, bulletsiz
    """
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    titles = plan_data.get("content", plan_data.get("titles", []))
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not hasattr(shape, 'placeholder_format') or not shape.placeholder_format:
            continue
        idx = shape.placeholder_format.idx
        if idx == 0:
            continue  # "REJA" — o'zgartirmaslik
        if idx == 1:
            tf = shape.text_frame
            tf.word_wrap = True
            total_chars = sum(len(t) for t in titles)
            font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=22)
            # Remove all existing paragraphs from txBody
            txBody = tf._txBody
            for p_elem in txBody.findall(f'{{{ns_a}}}p'):
                txBody.remove(p_elem)
            # Add clean numbered paragraphs with buNone (no bullet)
            for title in titles:
                # title already has "1. " prefix from generate_plan_with_titles
                safe_title = title.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                p_xml = (
                    f'<a:p xmlns:a="{ns_a}">' 
                    f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                    f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                    f'<a:solidFill><a:srgbClr val="1A1A2E"/></a:solidFill></a:rPr>'
                    f'<a:t>{safe_title}</a:t></a:r></a:p>'
                )
                p_elem = etree.fromstring(p_xml)
                txBody.append(p_elem)


def fill_t5_slide_3_image_right(slide, content_data, image_query):
    """
    5-Shablon Slayd 3 — Rasm o'ngda, matn chapda.
    idx=0: Sarlavha (yuqorida)
    idx=1: Rasm (o'ngda, Pixabay)
    idx=2: Matn (chapda)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if not items:
        items = content_data.get("col1", []) + content_data.get("col2", [])

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not hasattr(shape, 'placeholder_format') or not shape.placeholder_format:
            continue
        idx = shape.placeholder_format.idx
        tf = shape.text_frame
        tf.clear()
        if idx == 0:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
        elif idx == 2:
            from lxml import etree
            ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            total_chars = sum(len(s) for s in items)
            font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=24)
            txBody = tf._txBody
            for p_elem in txBody.findall(f'{{{ns_a}}}p'):
                txBody.remove(p_elem)
            for item in items:
                safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                p_xml = (
                    f'<a:p xmlns:a="{ns_a}">'
                    f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                    f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                    f'<a:solidFill><a:srgbClr val="1A1A2E"/></a:solidFill></a:rPr>'
                    f'<a:t>{safe_item}</a:t></a:r></a:p>'
                )
                p_elem = etree.fromstring(p_xml)
                txBody.append(p_elem)
            tf.word_wrap = True

    # Rasm qo'yish (idx=1)
    if image_query and os.path.isfile(image_query):
        place_image_in_placeholder(slide, 1, image_query)
    else:
        img_path = fetch_image(image_query)
        if img_path:
            place_image_in_placeholder(slide, 1, img_path)


def fill_t5_slide_4_two_columns(slide, content_data):
    """
    5-Shablon Slayd 4 — Ikki ustunli.
    idx=0: Sarlavha
    idx=2: Chap ustun
    idx=4: O'ng ustun
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    col1 = content_data.get("col1", [])
    col2 = content_data.get("col2", [])
    # string bo'lsa, jumlalarga bo'lib list ga aylantir
    if isinstance(col1, str):
        col1 = [s.strip() for s in col1.replace('. ', '.\n').split('\n') if s.strip()]
    if isinstance(col2, str):
        col2 = [s.strip() for s in col2.replace('. ', '.\n').split('\n') if s.strip()]
    if not col1 and not col2:
        items = content_data.get("content", [])
        mid = len(items) // 2
        col1 = items[:mid]
        col2 = items[mid:]
    def write_col(tf, items, align):
        from lxml import etree
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        align_val = 'r' if align == PP_ALIGN.RIGHT else 'l'
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=12, max_pt=20)
        # Remove all existing paragraphs from txBody
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        # Add clean paragraphs with buNone (no bullet)
        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="{align_val}"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="1A1A2E"/></a:solidFill></a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)
        tf.word_wrap = True
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not hasattr(shape, 'placeholder_format') or not shape.placeholder_format:
            continue
        idx = shape.placeholder_format.idx
        tf = shape.text_frame
        if idx == 0:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
        elif idx == 2:
            write_col(tf, col1, PP_ALIGN.RIGHT)
        elif idx == 4:
            write_col(tf, col2, PP_ALIGN.LEFT)
def fill_t5_slide_5_two_staggered(slide, content_data):
    """
    5-Shablon Slayd 5 — Ikki offset blok.
    idx=0: Sarlavha
    idx=1: Yuqori blok (chapda)
    idx=2: Pastki blok (o'ngda offset)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if not items:
        items = content_data.get("col1", []) + content_data.get("col2", [])
    mid = len(items) // 2 if len(items) > 1 else 1
    block1 = items[:mid]
    block2 = items[mid:]

    def write_block(tf, items):
        from lxml import etree
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=24)
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="1A1A2E"/></a:solidFill></a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)
        tf.word_wrap = True
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not hasattr(shape, 'placeholder_format') or not shape.placeholder_format:
            continue
        idx = shape.placeholder_format.idx
        tf = shape.text_frame
        if idx == 0:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
        elif idx == 1:
            write_block(tf, block1)
        elif idx == 2:
            write_block(tf, block2)


def fill_t5_slide_6_image_right2(slide, content_data, image_query):
    """
    5-Shablon Slayd 6 — Rasm o'ngda, kichik sarlavha chapda.
    idx=0: Sarlavha (chapda, kichik)
    idx=1: Rasm (o'ngda, Pixabay)
    idx=2: Matn (chapda)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if not items:
        items = content_data.get("col1", []) + content_data.get("col2", [])

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not hasattr(shape, 'placeholder_format') or not shape.placeholder_format:
            continue
        idx = shape.placeholder_format.idx
        tf = shape.text_frame
        tf.clear()
        if idx == 0:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
        elif idx == 2:
            tf.word_wrap = True
            total_chars = sum(len(s) for s in items)
            font_pt = calc_body_font_pt(total_chars, base_pt=17, min_pt=12, max_pt=21)
            for j, item in enumerate(items):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = item
                run.font.size = Pt(font_pt)
                run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Rasm qo'yish (idx=1)
    if image_query and os.path.isfile(image_query):
        place_image_in_placeholder(slide, 1, image_query)
    else:
        img_path = fetch_image(image_query)
        if img_path:
            place_image_in_placeholder(slide, 1, img_path)


def fill_t5_slide_7_two_blocks(slide, content_data):
    """
    5-Shablon Slayd 7 — Ikki kvadrat blok.
    idx=0: Sarlavha
    idx=1: Chap blok
    idx=2: O'ng blok
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    col1 = content_data.get("col1", [])
    col2 = content_data.get("col2", [])
    # string bo'lsa, jumlalarga bo'lib list ga aylantir
    if isinstance(col1, str):
        col1 = [s.strip() for s in col1.replace('. ', '.\n').split('\n') if s.strip()]
    if isinstance(col2, str):
        col2 = [s.strip() for s in col2.replace('. ', '.\n').split('\n') if s.strip()]
    if not col1 and not col2:
        items = content_data.get("content", [])
        mid = len(items) // 2
        col1 = items[:mid]
        col2 = items[mid:]
    def write_block(tf, items):
        from lxml import etree
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=15, min_pt=11, max_pt=18)
        # Remove all existing paragraphs from txBody
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        # Add clean paragraphs with buNone (no bullet)
        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="1A1A2E"/></a:solidFill></a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)
        tf.word_wrap = True
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not hasattr(shape, 'placeholder_format') or not shape.placeholder_format:
            continue
        idx = shape.placeholder_format.idx
        tf = shape.text_frame
        if idx == 0:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
        elif idx == 1:
            write_block(tf, col1)
        elif idx == 2:
            write_block(tf, col2)
def fill_t5_slide_8_conclusion(slide, conclusion_data):
    """
    5-Shablon Slayd 8 — Xulosa.
    idx=0: 'THANKS!' — o'zgartirmaslik
    idx=1: Xulosa matni
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    items = conclusion_data.get("content", [])
    if not items:
        items = conclusion_data.get("col1", []) + conclusion_data.get("col2", [])

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not hasattr(shape, 'placeholder_format') or not shape.placeholder_format:
            continue
        idx = shape.placeholder_format.idx
        if idx == 0:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "XULOSA"
        if idx == 1:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            total_chars = sum(len(s) for s in items)
            font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)
            for j, item in enumerate(items):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = item
                run.font.size = Pt(font_pt)
                run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)


def generate_template_5_presentation(prs, topic, requested_slide_count, language,
                                      name_surname, plan, content_data_list, user_images=None):
    """
    5-shablon asosida to'liq prezentatsiya yaratadi.
    """
    import io
    total_content_slides = build_slide_structure_5(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}

    # Slayd 1 — Muqova
    fill_t5_slide_1_cover(prs.slides[0], topic, name_surname)

    # Slayd 2 — Reja
    fill_t5_slide_2_plan(prs.slides[1], plan_dict)

    # Kontent slaydlari (3-dan boshlab)
    user_img_idx = 0
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        image_query = data.get("image_query", topic)
        slide_type = i % 5

        if slide_type == 0:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t5_slide_3_image_right(slide, data, img_path if img_path else image_query)
            else:
                fill_t5_slide_3_image_right(slide, data, image_query)
        elif slide_type == 1:
            fill_t5_slide_4_two_columns(slide, data)
        elif slide_type == 2:
            fill_t5_slide_5_two_staggered(slide, data)
        elif slide_type == 3:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t5_slide_6_image_right2(slide, data, img_path if img_path else image_query)
            else:
                fill_t5_slide_6_image_right2(slide, data, image_query)
        elif slide_type == 4:
            fill_t5_slide_7_two_blocks(slide, data)
        logging.info(f"  [T5] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")

    # Xulosa slayd — alohida GPT so'rovi bilan
    conclusion_slide = prs.slides[-1]
    conclusion_data = generate_slide_content(topic, requested_slide_count, requested_slide_count, language, is_conclusion=True)
    if not conclusion_data:
        conclusion_data = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}
    fill_t5_slide_8_conclusion(conclusion_slide, conclusion_data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ─────────────────────────────────────────────────────────────
# 6-SHABLON
# ─────────────────────────────────────────────────────────────
SLIDE_TYPE_NAMES_T6 = {
    0: "one_column",      # 3-slayd: sarlavha + matn
    1: "one_column",      # 4-slayd: sarlavha + tavsif
    2: "three_columns",   # 5-slayd: 3 ta ustun (col1, col2, col3)
    3: "two_columns",     # 6-slayd: 2 ustun
    4: "one_column",      # 7-slayd: sarlavha + matn
}
CONTENT_SLIDE_TEMPLATE_INDICES_6 = [2, 3, 4, 5, 6]

def build_slide_structure_6(prs, requested_content_count):
    """
    6-shablon uchun slayd tuzilmasini quradi.
    duplicate_slide va move_slide ishlatadi.
    """
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5
    logging.info(f"[T6] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_6:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T6] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    
    # Xulosa slaydini (index 7) oxiriga ko'chirish
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T6] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides

def fill_t6_slide_1_cover(slide, topic, name_surname):
    """
    6-Shablon Slayd 1 — Muqova.
    Shape[1]: Sarlavha (F3F0DF rang)
    Shape[2]: Subtitle (CDD6E2 rang)
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = topic.upper()
        run.font.size = Pt(52.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xF3, 0xF0, 0xDF)
        
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name_surname
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0xCD, 0xD6, 0xE2)

def fill_t6_slide_2_plan(slide, plan_data):
    """
    6-Shablon Slayd 2 — Reja.
    Shape[1]: "Reja:"
    Shape[2]: Reja bandlari
    """
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    titles = plan_data.get("content", plan_data.get("titles", []))
    
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.word_wrap = True
        total_chars = sum(len(t) for t in titles)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=22)
        
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
            
        for title in titles:
            safe_title = title.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">' 
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="334155"/></a:solidFill></a:rPr>'
                f'<a:t>{safe_title}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)

def fill_t6_slide_3_text(slide, content_data, image_query=None):
    """
    6-Shablon Slayd 3.
    Shape[1], Shape[2]: Rasm placeholder (o'ng tomonda)
    Shape[3]: Matn (chap tomonda)
    Shape[4]: Sarlavha
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        tf = slide.shapes[4].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        tf = slide.shapes[3].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=24)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

    # Mavzu bo'yicha rasm qo'shish (Shape[1] va Shape[2] o'rniga)
    # Shape[1]: left=6.93", top=2.43", size=5.78"x4.17" => 17.60x10.59cm
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "presentation"))
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(17.60)
            top  = Cm(6.17)
            width = Cm(14.69)
            height = Cm(10.59)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T6] Slayd 3 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T6] Slayd 3 rasm xatolik: {e}")

def fill_t6_slide_4_text(slide, content_data):
    """
    6-Shablon Slayd 4.
    Shape[1]: Sarlavha
    Shape[2]: Matn
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(33.75)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=24)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

def fill_t6_slide_5_three_steps(slide, content_data):
    """
    6-Shablon Slayd 5.
    Shape[1]: 1-ustun matn (left=1.98cm)
    Shape[2]: 2-ustun matn (left=12.08cm)
    Shape[3]: 3-ustun matn (left=22.97cm)
    Shape[4]: Sarlavha
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    
    # col1/col2/col3 formatini qo'llab-quvvatlash
    col1 = content_data.get("col1", "")
    col2 = content_data.get("col2", "")
    col3 = content_data.get("col3", "")
    
    if col1 or col2 or col3:
        # three_columns formatida kelgan
        items = [col1, col2, col3]
    else:
        # content massivi formatida kelgan - 3 ta elementga to'ldirish
        raw_items = content_data.get("content", [])
        items = []
        for idx in range(3):
            if idx < len(raw_items):
                items.append(raw_items[idx])
            elif raw_items:
                # Mavjud elementlarni takrorlash o'rniga bo'sh qoldirmaslik uchun
                # oxirgi elementni ishlatamiz
                items.append(raw_items[-1])
            else:
                items.append("")
    
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        tf = slide.shapes[4].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    for i in range(3):
        shape_idx = i + 1
        if len(slide.shapes) > shape_idx and slide.shapes[shape_idx].has_text_frame:
            text = items[i] if i < len(items) else ""
            tf = slide.shapes[shape_idx].text_frame
            tf.clear()
            tf.word_wrap = True
            font_pt = calc_body_font_pt(len(text), base_pt=14, min_pt=10, max_pt=18)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

def fill_t6_slide_6_two_cols(slide, content_data):
    """
    6-Shablon Slayd 6.
    Shape[1]: 1-ustun matn (left=2.25cm, chap tomonda)
    Shape[2]: 2-ustun matn (left=18.52cm, o'ng tomonda)
    Shape[3]: Sarlavha
    Ustunlar gorizontal joylashgan, chapga tekislangan.
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    
    # col1/col2 string yoki list bo'lishi mumkin
    col1_raw = content_data.get("col1", "")
    col2_raw = content_data.get("col2", "")
    
    if col1_raw or col2_raw:
        # two_columns formatida kelgan
        col1_text = col1_raw if isinstance(col1_raw, str) else "\n".join(col1_raw)
        col2_text = col2_raw if isinstance(col2_raw, str) else "\n".join(col2_raw)
    else:
        # content massivi formatida kelgan - ikki qismga bo'lish
        items = content_data.get("content", [])
        mid = max(1, len(items) // 2)
        col1_list = items[:mid]
        col2_list = items[mid:]
        col1_text = "\n".join(col1_list)
        col2_text = "\n".join(col2_list)
        
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        tf = slide.shapes[3].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    for shape_idx, col_text in [(1, col1_text), (2, col2_text)]:
        if len(slide.shapes) > shape_idx and slide.shapes[shape_idx].has_text_frame:
            tf = slide.shapes[shape_idx].text_frame
            tf.clear()
            tf.word_wrap = True
            font_pt = calc_body_font_pt(len(col_text), base_pt=14, min_pt=10, max_pt=18)
            # Har bir satr uchun alohida paragraf - chapga tekislangan
            lines = col_text.split("\n") if col_text else [""]
            for j, line in enumerate(lines):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = line
                run.font.size = Pt(font_pt)
                run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

def fill_t6_slide_7_text(slide, content_data, image_query=None):
    """
    6-Shablon Slayd 7.
    Shape[0]: Fon rasm (to'liq)
    Shape[1]: Sarlavha (chap tomonda)
    Shape[2]: Chiziq
    Shape[3]: Matn (chap tomonda)
    Shape[4]: Rasm placeholder (o'ng tomonda)
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        tf = slide.shapes[3].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

    # Mavzu bo'yicha rasm qo'shish (Shape[4] o'rniga - o'ng tomonda)
    # Shape[4]: left=6.67", top=0.00", size=6.67"x7.50" => 16.93x19.05cm
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "presentation"))
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(16.93)
            top  = Cm(0.00)
            width = Cm(16.93)
            height = Cm(19.05)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T6] Slayd 7 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T6] Slayd 7 rasm xatolik: {e}")

def fill_t6_slide_8_conclusion(slide, conclusion_data):
    """
    6-Shablon Slayd 8 — Xulosa.
    Shape[1]: "XULOSA" (katta matn, pos=1.29cm, size=31.06x4.10cm)
    Shape[2]: Xulosa matni (pos=13.12cm, size=7.62cm -> 24cm ga kengaytiriladi)
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    items = conclusion_data.get("content", [])
    if not items:
        items = conclusion_data.get("col1", []) + conclusion_data.get("col2", [])
        
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "XULOSA"
        run.font.size = Pt(60)
        run.font.color.rgb = RGBColor(0xCD, 0xD6, 0xE2)
        
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        # Asosiy matn blokini 24 sm eniga kengaytirish va markazga joylashtirish
        shape2 = slide.shapes[2]
        shape2.width = Cm(24)
        # Slayd kengligi 33.87cm, markazda joylashish: left = (33.87 - 24) / 2 ≈ 4.94cm
        shape2.left = Cm(4.94)
        tf = shape2.text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=24)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0xCD, 0xD6, 0xE2)

def generate_template_6_presentation(prs, topic, requested_slide_count, language,
                                      name_surname, plan, content_data_list, user_images=None):
    """
    6-shablon asosida to'liq prezentatsiya yaratadi.
    """
    import io
    
    total_content_slides = build_slide_structure_6(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}
    
    # Slayd 1 — Muqova
    fill_t6_slide_1_cover(prs.slides[0], topic, name_surname)
    
    # Slayd 2 — Reja
    fill_t6_slide_2_plan(prs.slides[1], plan_dict)
    
    # Kontent slaydlari (3-dan boshlab)
    user_img_idx = 0
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
            
        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        image_query = data.get("image_query", topic)
        
        slide_type = i % 5
        if slide_type == 0:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t6_slide_3_text(slide, data, img_path if img_path else image_query)
            else:
                fill_t6_slide_3_text(slide, data)
        elif slide_type == 1:
            fill_t6_slide_4_text(slide, data)
        elif slide_type == 2:
            fill_t6_slide_5_three_steps(slide, data)
        elif slide_type == 3:
            fill_t6_slide_6_two_cols(slide, data)
        elif slide_type == 4:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t6_slide_7_text(slide, data, img_path if img_path else image_query)
            else:
                fill_t6_slide_7_text(slide, data)
            
        logging.info(f"  [T6] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
        
    # Xulosa slayd
    conclusion_slide = prs.slides[-1]
    conclusion_data = generate_slide_content(topic, requested_slide_count, requested_slide_count, language, is_conclusion=True)
    if not conclusion_data:
        conclusion_data = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}
    fill_t6_slide_8_conclusion(conclusion_slide, conclusion_data)
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════
# 7-SHABLON FUNKSIYALARI (7.pptx)
# ═══════════════════════════════════════════════════════════════
# Slayd tuzilmasi:
#   Slayd 1: Muqova — Shape[0]: sarlavha (o'ng), Shape[1]: muallif
#   Slayd 2: Reja — Shape[0]: sarlavha, Shape[1]: reja bandlari
#   Slayd 3: Rasm chap + matn o'ng — Shape[0]: sarlavha, Shape[1]: rasm, Shape[2]: qo'shimcha matn
#   Slayd 4: Katta matn — Shape[0]: sarlavha, Shape[1]: matn, Shape[2]: dekorativ group
#   Slayd 5: Ikki ustun — Shape[0]: sarlavha, Shape[1]: chap ustun, Shape[2]: o'ng ustun
#   Slayd 6: Sarlavha o'ngda + 2 matn — Shape[0]: sarlavha (o'ng), Shape[1]: yuqori matn, Shape[2]: pastki matn
#   Slayd 7: Rasm o'ng + matn chap — Shape[0]: sarlavha, Shape[1]: matn, Shape[2]: rasm
#   Slayd 8: Xulosa — Shape[0]: sarlavha, Shape[1]: matn

SLIDE_TYPE_NAMES_T7 = {
    0: "image_left",    # 3-slayd: rasm chap, matn o'ng
    1: "one_column",    # 4-slayd: katta matn bloki
    2: "two_columns",   # 5-slayd: ikki ustun
    3: "two_columns",   # 6-slayd: sarlavha o'ngda + 2 matn bloki
    4: "image_right",   # 7-slayd: rasm o'ng, matn chap
}
CONTENT_SLIDE_TEMPLATE_INDICES_7 = [2, 3, 4, 5, 6]


def build_slide_structure_7(prs, requested_content_count):
    """
    7-shablon uchun slayd tuzilmasini quradi.
    duplicate_slide va move_slide ishlatadi.
    """
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5
    logging.info(f"[T7] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_7:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T7] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")

    # Xulosa slaydini (index 7) oxiriga ko'chirish
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T7] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def fill_t7_slide_1_cover(slide, topic, name_surname):
    """
    7-Shablon Slayd 1 — Muqova.
    Shape[0]: Sarlavha (o'ng tomonda, katta)
    Shape[1]: Muallif ismi (o'ng pastda)
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = topic.upper()
        total_chars = len(topic)
        font_pt = calc_body_font_pt(total_chars, base_pt=52, min_pt=28, max_pt=60)
        run.font.size = Pt(font_pt)
        run.font.bold = True

    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = name_surname if name_surname and name_surname.strip() else ""
        run.font.size = Pt(18)


def fill_t7_slide_2_plan(slide, plan_data):
    """
    7-Shablon Slayd 2 — Reja.
    Shape[0]: Sarlavha ("Reja")
    Shape[1]: Reja bandlari
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    titles = plan_data.get("content", plan_data.get("titles", []))

    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Reja"
        run.font.size = Pt(36)
        run.font.bold = True

    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        import re
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(t) for t in titles)
        font_pt = calc_body_font_pt(total_chars, base_pt=20, min_pt=13, max_pt=24)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for i, title in enumerate(titles):
            # GPT allaqachon "1. Band" formatida qaytarishi mumkin - boshidagi raqamni olib tashlaymiz
            clean_title = re.sub(r'^\d+[\.)\-]\s*', '', title.strip())
            safe_title = clean_title.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{i+1}. {safe_title}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)


def fill_t7_slide_3_image_left(slide, content_data, image_query=None):
    """
    7-Shablon Slayd 3 — Rasm chap, matn o'ng.
    Shape[0]: Sarlavha (o'ng tomonda yuqori)
    Shape[1]: Rasm placeholder (chap tomonda, to'liq balandlik)
    Shape[2]: Qo'shimcha matn (o'ng pastda)
    Rasm: Shape[1] o'rniga joylashtiriladi
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    title = content_data.get("title", "")
    items = content_data.get("content", [])

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True

    # Shape[2]: Qo'shimcha matn (o'ng pastda)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)

    # Rasm: Shape[1] o'rniga (chap tomonda, to'liq balandlik)
    # Shape[1]: left=0.00cm, top=-0.03cm, size=16.09x19.11cm
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "presentation"))
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(0.00)
            top = Cm(0.00)
            width = Cm(16.09)
            height = Cm(19.05)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T7] Slayd 3 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T7] Slayd 3 rasm xatolik: {e}")


def fill_t7_slide_4_text(slide, content_data):
    """
    7-Shablon Slayd 4 — Katta matn bloki.
    Shape[0]: Sarlavha
    Shape[1]: Asosiy matn (bullet list)
    Shape[2]: Dekorativ group (o'zgartirilmaydi)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(36)
        run.font.bold = True

    # Shape[1]: Asosiy matn
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=22)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)


def fill_t7_slide_5_two_cols(slide, content_data):
    """
    7-Shablon Slayd 5 — Ikki ustun.
    Shape[0]: Sarlavha
    Shape[1]: Chap ustun matn
    Shape[2]: O'ng ustun matn
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")

    # Agar col1/col2 bo'lmasa, content dan olamiz
    if not col1_text and not col2_text:
        items = content_data.get("content", [])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True

    total_chars = len(col1_text) + len(col2_text)
    font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)

    # Shape[1]: Chap ustun
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        safe_col1 = col1_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        p_xml = (
            f'<a:p xmlns:a="{ns_a}">'
            f'<a:pPr algn="l"><a:buNone/></a:pPr>'
            f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
            f'</a:rPr>'
            f'<a:t>{safe_col1}</a:t></a:r></a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

    # Shape[2]: O'ng ustun
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.word_wrap = True
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        safe_col2 = col2_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        p_xml = (
            f'<a:p xmlns:a="{ns_a}">'
            f'<a:pPr algn="l"><a:buNone/></a:pPr>'
            f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
            f'</a:rPr>'
            f'<a:t>{safe_col2}</a:t></a:r></a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)


def fill_t7_slide_6_special(slide, content_data):
    """
    7-Shablon Slayd 6 — Sarlavha o'ngda, 2 matn bloki chapda.
    Shape[0]: Sarlavha (o'ng pastda)
    Shape[1]: Yuqori matn bloki (chap)
    Shape[2]: Pastki matn bloki (chap)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    col1_text = content_data.get("col1", "")
    col2_text = content_data.get("col2", "")

    if not col1_text and not col2_text:
        items = content_data.get("content", [])
        col1_text = items[0] if len(items) > 0 else ""
        col2_text = items[1] if len(items) > 1 else ""

    # Shape[0]: Sarlavha (o'ng pastda)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(28)
        run.font.bold = True

    total_chars = len(col1_text) + len(col2_text)
    font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)

    # Shape[1]: Yuqori matn
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        safe_col1 = col1_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        p_xml = (
            f'<a:p xmlns:a="{ns_a}">'
            f'<a:pPr algn="l"><a:buNone/></a:pPr>'
            f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
            f'</a:rPr>'
            f'<a:t>{safe_col1}</a:t></a:r></a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

    # Shape[2]: Pastki matn
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.word_wrap = True
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        safe_col2 = col2_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        p_xml = (
            f'<a:p xmlns:a="{ns_a}">'
            f'<a:pPr algn="l"><a:buNone/></a:pPr>'
            f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
            f'</a:rPr>'
            f'<a:t>{safe_col2}</a:t></a:r></a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)


def fill_t7_slide_7_image_right(slide, content_data, image_query=None):
    """
    7-Shablon Slayd 7 — Rasm o'ng, matn chap.
    Shape[0]: Sarlavha (chap yuqori)
    Shape[1]: Matn (chap pastda)
    Shape[2]: Rasm placeholder (o'ng tomonda, to'liq balandlik)
    Rasm: Shape[2] o'rniga joylashtiriladi
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True

    # Shape[1]: Matn
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)

    # Rasm: Shape[2] o'rniga (o'ng tomonda, to'liq balandlik)
    # Shape[2]: left=16.93cm, top=0.00cm, size=17.00x19.05cm
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", content_data.get("title", "presentation"))
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(16.93)
            top = Cm(0.00)
            width = Cm(17.00)
            height = Cm(19.05)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T7] Slayd 7 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T7] Slayd 7 rasm xatolik: {e}")


def fill_t7_slide_8_conclusion(slide, conclusion_data):
    """
    7-Shablon Slayd 8 — Xulosa.
    Shape[0]: Sarlavha ("Xulosa" yoki "Thanks!")
    Shape[1]: Qo'shimcha matn
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = conclusion_data.get("title", "Xulosa")
    items = conclusion_data.get("content", [])

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(48)
        run.font.bold = True

    # Shape[1]: Qo'shimcha matn - 2 sm yuqoriga ko'tarish
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        from pptx.util import Cm
        shape1 = slide.shapes[1]
        # top: 11.22cm -> 9.22cm (2 sm yuqoriga)
        shape1.top = Cm(9.22)
        tf = shape1.text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=22)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)


def generate_template_7_presentation(prs, topic, requested_slide_count, language,
                                      name_surname, plan, content_data_list, user_images=None):
    """
    7-shablon asosida to'liq prezentatsiya yaratadi.
    """
    import io

    total_content_slides = build_slide_structure_7(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}

    # Slayd 1 — Muqova
    fill_t7_slide_1_cover(prs.slides[0], topic, name_surname)

    # Slayd 2 — Reja
    fill_t7_slide_2_plan(prs.slides[1], plan_dict)

    # Kontent slaydlari (3-dan boshlab)
    user_img_idx = 0
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break

        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        image_query = data.get("image_query", topic)

        slide_type = i % 5
        if slide_type == 0:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t7_slide_3_image_left(slide, data, img_path if img_path else image_query)
            else:
                fill_t7_slide_3_image_left(slide, data)
        elif slide_type == 1:
            fill_t7_slide_4_text(slide, data)
        elif slide_type == 2:
            fill_t7_slide_5_two_cols(slide, data)
        elif slide_type == 3:
            fill_t7_slide_6_special(slide, data)
        elif slide_type == 4:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t7_slide_7_image_right(slide, data, img_path if img_path else image_query)
            else:
                fill_t7_slide_7_image_right(slide, data)

        logging.info(f"  [T7] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")

    # Xulosa slayd — GPT dan xulosa matni so'rash
    conclusion_slide = prs.slides[-1]
    conclusion_data = generate_slide_content(topic, len(prs.slides), len(prs.slides), language, is_conclusion=True)
    if not conclusion_data:
        conclusion_data = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}
    fill_t7_slide_8_conclusion(conclusion_slide, conclusion_data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════
# 8-SHABLON — "Biznes" uslubi
# ═══════════════════════════════════════════════════════════════

SLIDE_TYPE_NAMES_T8 = {
    0: "image_right",    # Slayd 3: sarlavha + matn chap + rasm o'ng
    1: "one_column",     # Slayd 4: sarlavha katta chap + matn o'ng
    2: "image_right",    # Slayd 5: sarlavha + matn past + rasm o'ng
    3: "two_columns",    # Slayd 6: sarlavha + 2 teng ustun
    4: "two_columns",    # Slayd 7: sarlavha + tor chap + keng o'ng
}

CONTENT_SLIDE_TEMPLATE_INDICES_T8 = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7


def build_slide_structure_8(prs, requested_content_count):
    """
    8-shablon uchun slayd tuzilmasini yaratadi.
    Muqova(1) + Reja(1) + Kontent(N) + Xulosa(1) = jami
    5 ta kontent shablon slayd bor (index 2-6), to'liq to'plamlar sifatida takrorlanadi.
    """
    n_templates = len(CONTENT_SLIDE_TEMPLATE_INDICES_T8)  # 5
    full_repeats = max(1, round(requested_content_count / n_templates))
    total_content_slides = full_repeats * n_templates
    logging.info(f"[T8] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")

    conclusion_current_index = 7  # 0-indexed: slayd 8

    # Kerakli to'plamlar soniga qarab nusxalash
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_T8:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T8] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")

    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T8] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def fill_t8_slide_1_cover(slide, topic, name_surname):
    """
    8-Shablon Slayd 1 — Muqova.
    Shape[0]: Fon rasm (PICTURE placeholder)
    Shape[1]: Sarlavha (TITLE)
    Shape[2]: Ism-familiya (TEXT_BOX)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    # Shape[1]: Sarlavha
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = topic
        total_chars = len(topic)
        font_pt = calc_body_font_pt(total_chars, base_pt=48, min_pt=28, max_pt=60)
        run.font.size = Pt(font_pt)
        run.font.bold = True

    # Shape[2]: Ism-familiya (TEXT_BOX)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = name_surname if name_surname and name_surname.strip() else ""
        run.font.size = Pt(18)


def fill_t8_slide_2_plan(slide, plan_data):
    """
    8-Shablon Slayd 2 — Reja.
    Shape[0]: Sarlavha ("Reja")
    Shape[1]: Reja bandlari
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    titles = plan_data.get("content", plan_data.get("titles", []))

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Reja"
        run.font.size = Pt(32)
        run.font.bold = True

    # Shape[1]: Reja bandlari
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        import re
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(t) for t in titles)
        font_pt = calc_body_font_pt(total_chars, base_pt=22, min_pt=14, max_pt=28)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for i, title in enumerate(titles):
            # GPT allaqachon "1. Band" formatida qaytarishi mumkin - boshidagi raqamni olib tashlaymiz
            clean_title = re.sub(r'^\d+[\.)\-]\s*', '', title.strip())
            safe_title = clean_title.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{i+1}. {safe_title}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)


def fill_t8_slide_3_image_right(slide, content_data, image_query=None):
    """
    8-Shablon Slayd 3 — Sarlavha + Matn chap + Rasm o'ng.
    Shape[0]: Sarlavha (TITLE) - chapda yuqori
    Shape[1]: Matn (OBJECT idx=11) - chapda pastda
    Shape[2]: Rasm (PICTURE idx=10) - o'ngda
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True

    # Shape[1]: Matn (chapda pastda)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)

    # Rasm: Shape[2] o'rniga (o'ng tomonda)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(18.0)
            top = Cm(2.56)
            width = Cm(15.87)
            height = Cm(13.97)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T8] Slayd 3 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T8] Slayd 3 rasm xatolik: {e}")


def fill_t8_slide_4_title_left(slide, content_data):
    """
    8-Shablon Slayd 4 — Sarlavha katta chapda + Matn o'ngda.
    Shape[0]: Sarlavha (TITLE) - chapda katta
    Shape[1]: Matn (OBJECT idx=11) - o'ngda
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[0]: Sarlavha (chapda, katta)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        total_chars = len(title)
        font_pt = calc_body_font_pt(total_chars, base_pt=40, min_pt=24, max_pt=52)
        run.font.size = Pt(font_pt)
        run.font.bold = True

    # Shape[1]: Matn (o'ngda)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)


def fill_t8_slide_5_image_right2(slide, content_data, image_query=None):
    """
    8-Shablon Slayd 5 — Sarlavha + Matn pastda + Rasm o'ng.
    Shape[0]: Sarlavha (TITLE) - chapda yuqori
    Shape[1]: Matn (OBJECT idx=11) - chapda pastda
    Shape[2]: Rasm (PICTURE idx=10) - o'ngda
    Shape[3]: Slide number
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[0]: Sarlavha — 1 sm yuqoriga ko'tarish (2.50 -> 1.50)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        slide.shapes[0].top = Cm(1.50)
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True

    # Shape[1]: Matn (chapda pastda) — 2 sm yuqoriga ko'tarish (10.63 -> 8.63)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        slide.shapes[1].top = Cm(8.63)
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)

    # Rasm: Shape[2] o'rniga (o'ng tomonda)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(18.0)
            top = Cm(2.56)
            width = Cm(15.87)
            height = Cm(13.97)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T8] Slayd 5 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T8] Slayd 5 rasm xatolik: {e}")


def fill_t8_slide_6_two_cols(slide, content_data):
    """
    8-Shablon Slayd 6 — Sarlavha + 2 teng ustun.
    Shape[0]: Sarlavha (TITLE) - yuqori
    Shape[1]: Chap ustun (OBJECT idx=10)
    Shape[2]: O'ng ustun (OBJECT idx=11)
    Shape[3]: Slide number
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    col1 = content_data.get("col1", "")
    col2 = content_data.get("col2", "")
    if isinstance(col1, list):
        col1 = " ".join(col1)
    if isinstance(col2, list):
        col2 = " ".join(col2)
    # Fallback: content dan olamiz
    if not col1 and not col2:
        items = content_data.get("content", [])
        if isinstance(items, list) and len(items) >= 2:
            mid = len(items) // 2
            col1 = " ".join(items[:mid])
            col2 = " ".join(items[mid:])
        elif isinstance(items, str):
            col1 = items
            col2 = ""

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True

    # Shape[1]: Chap ustun
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = len(col1)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = col1
        run.font.size = Pt(font_pt)

    # Shape[2]: O'ng ustun
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = len(col2)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = col2
        run.font.size = Pt(font_pt)


def fill_t8_slide_7_narrow_wide(slide, content_data):
    """
    8-Shablon Slayd 7 — Sarlavha + Tor chap (kalit so'zlar) + Keng o'ng (tavsif).
    Shape[0]: Sarlavha (TITLE) - yuqori
    Shape[1]: Tor chap ustun (OBJECT idx=10) - kalit so'zlar
    Shape[2]: Keng o'ng ustun (OBJECT idx=11) - tavsif
    Shape[3]: Slide number
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    col1 = content_data.get("col1", "")
    col2 = content_data.get("col2", "")
    if isinstance(col1, list):
        col1 = "\n".join(col1)
    if isinstance(col2, list):
        col2 = "\n".join(col2)
    # Fallback: content dan olamiz
    if not col1 and not col2:
        items = content_data.get("content", [])
        if isinstance(items, list) and len(items) >= 2:
            # Birinchi yarmini kalit so'z sifatida, ikkinchisini tavsif sifatida
            mid = max(1, len(items) // 3)
            col1 = "\n".join(items[:mid])
            col2 = "\n".join(items[mid:])
        elif isinstance(items, str):
            col1 = items
            col2 = ""

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True

    # Shape[1]: Tor chap ustun (kalit so'zlar)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        lines = col1.split("\n") if col1 else []
        total_chars = len(col1)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=13, max_pt=22)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for line in lines:
            safe_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" b="1" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{safe_line}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)

    # Shape[2]: Keng o'ng ustun (tavsif)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.word_wrap = True
        lines = col2.split("\n") if col2 else []
        total_chars = len(col2)
        font_pt = calc_body_font_pt(total_chars, base_pt=16, min_pt=11, max_pt=20)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for line in lines:
            safe_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{safe_line}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)


def fill_t8_slide_8_conclusion(slide, conclusion_data):
    """
    8-Shablon Slayd 8 — Xulosa.
    Shape[0]: Sarlavha (TITLE)
    Shape[1]: Xulosa matn (OBJECT idx=10)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = conclusion_data.get("title", "Xulosa")
    items = conclusion_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(48)
        run.font.bold = True

    # Shape[1]: Xulosa matn
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=20, min_pt=13, max_pt=24)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'</a:rPr>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)


def generate_template_8_presentation(prs, topic, requested_slide_count, language,
                                      name_surname, plan, content_data_list, user_images=None):
    """
    8-shablon asosida to'liq prezentatsiya yaratadi.
    """
    import io

    total_content_slides = build_slide_structure_8(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}

    # Slayd 1 — Muqova
    fill_t8_slide_1_cover(prs.slides[0], topic, name_surname)

    # Slayd 2 — Reja
    fill_t8_slide_2_plan(prs.slides[1], plan_dict)

    # Kontent slaydlari (3-dan boshlab)
    user_img_idx = 0
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break

        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        image_query = data.get("image_query", topic)

        slide_type = i % 5
        if slide_type == 0:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t8_slide_3_image_right(slide, data, img_path if img_path else image_query)
            else:
                fill_t8_slide_3_image_right(slide, data)
        elif slide_type == 1:
            fill_t8_slide_4_title_left(slide, data)
        elif slide_type == 2:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                fill_t8_slide_5_image_right2(slide, data, img_path if img_path else image_query)
            else:
                fill_t8_slide_5_image_right2(slide, data)
        elif slide_type == 3:
            fill_t8_slide_6_two_cols(slide, data)
        elif slide_type == 4:
            fill_t8_slide_7_narrow_wide(slide, data)

        logging.info(f"  [T8] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")

    # Xulosa slayd — GPT dan xulosa matni so'rash
    conclusion_slide = prs.slides[-1]
    conclusion_data = generate_slide_content(topic, len(prs.slides), len(prs.slides), language, is_conclusion=True)
    if not conclusion_data:
        conclusion_data = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}
    fill_t8_slide_8_conclusion(conclusion_slide, conclusion_data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════
# TEMPLATE 9 — FUTURISTIC
# Tuzilma:
#   Slayd 1: Muqova (TITLE + SUBTITLE)
#   Slayd 2: Reja (TITLE + BODY)
#   Slayd 3: 3 ustun (3x SUBTITLE placeholder)
#   Slayd 4: Sarlavha + Rasm chap + Matn o'ng (TITLE + BODY + PICTURE)
#   Slayd 5: Sarlavha + Matn o'ng + 2 rasm (TITLE + SUBTITLE + 2xPICTURE)
#   Slayd 6: Sarlavha + Rasm o'ng + Matn chap (TITLE + SUBTITLE + PICTURE)
#   Slayd 7: Sarlavha + Matn o'ng + Rasm chap (TITLE + SUBTITLE + PICTURE)
#   Slayd 8: Xulosa (TEXT_BOX)
# ═══════════════════════════════════════════════════════════════

SLIDE_TYPE_NAMES_T9 = {
    0: "three_columns",  # Slayd 3: 3 ustun matn
    1: "image_left",     # Slayd 4: sarlavha + rasm chap + matn o'ng
    2: "two_images",     # Slayd 5: sarlavha + matn o'ng + 2 rasm
    3: "image_right",    # Slayd 6: sarlavha + rasm o'ng + matn chap
    4: "image_left2",    # Slayd 7: sarlavha + matn o'ng + rasm chap
}

CONTENT_SLIDE_TEMPLATE_INDICES_T9 = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7


def build_slide_structure_9(prs, requested_content_count):
    """
    9-shablon uchun slayd tuzilmasini yaratadi.
    Muqova(1) + Reja(1) + Kontent(N) + Xulosa(1) = jami
    5 ta kontent shablon slayd bor (index 2-6), to'liq to'plamlar sifatida takrorlanadi.
    """
    n_templates = len(CONTENT_SLIDE_TEMPLATE_INDICES_T9)  # 5
    full_repeats = max(1, round(requested_content_count / n_templates))
    total_content_slides = full_repeats * n_templates
    logging.info(f"[T9] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")

    conclusion_current_index = 7  # 0-indexed: slayd 8

    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_T9:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T9] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")

    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T9] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def fill_t9_slide_1_cover(slide, topic, name_surname):
    """
    9-Shablon Slayd 1 — Muqova.
    Shape[0]: Sarlavha (CENTER_TITLE idx=0)
    Shape[1]: Ism-familiya / Subtitle (SUBTITLE idx=1)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    # Shape[0]: Sarlavha (CENTER_TITLE)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = topic.upper() if topic else "TAQDIMOT"
        font_pt = calc_body_font_pt(len(topic), base_pt=32, min_pt=18, max_pt=40)
        run.font.size = Pt(font_pt)
        run.font.bold = True

    # Shape[1]: Ism-familiya (SUBTITLE)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name_surname if name_surname else "Taqdimot"
        run.font.size = Pt(16)


def fill_t9_slide_2_plan(slide, plan_data):
    """
    9-Shablon Slayd 2 — Reja.
    Shape[0]: Sarlavha (TITLE idx=0)
    Shape[1]: Reja matni (BODY idx=1)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = plan_data.get("title", "Reja") if isinstance(plan_data, dict) else "Reja"
    items = plan_data.get("content", []) if isinstance(plan_data, dict) else []

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title.upper()
        run.font.size = Pt(28)
        run.font.bold = True

    # Shape[1]: Reja ro'yxati
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(str(s)) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=22)

        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for idx, item in enumerate(items):
            # Agar item allaqachon "1. ..." yoki "1.1 ..." formatida bo'lsa, raqamni olib tashlaymiz
            import re as _re
            clean_item = _re.sub(r'^\d+\.\s*\d*\.?\s*', '', str(item)).strip()
            safe_item = clean_item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" b="1" dirty="0"/>'
                f'<a:t>{idx+1}. {safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))


def fill_t9_slide_3_three_cols(slide, content_data):
    """
    9-Shablon Slayd 3 — 3 ustun.
    Shape[0]: O'rta ustun (SUBTITLE idx=1)
    Shape[1]: O'ng ustun (SUBTITLE idx=3)
    Shape[2]: Chap ustun (SUBTITLE idx=5)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")

    # GPT col1/col2/col3 formatida qaytarsa — to'g'ridan-to'g'ri ishlatamiz
    col1_raw = content_data.get("col1", "")
    col2_raw = content_data.get("col2", "")
    col3_raw = content_data.get("col3", "")

    if col1_raw or col2_raw or col3_raw:
        # col1/col2/col3 mavjud — ularni list sifatida ishlatamiz
        def split_to_sentences(text):
            """Matnni jumlalarga bo'lish"""
            import re
            if not text:
                return []
            # Jumlalarga bo'lish
            sentences = re.split(r'(?<=[.!?])\s+', str(text).strip())
            return [s.strip() for s in sentences if s.strip()]

        col1_items = split_to_sentences(col1_raw) or [str(col1_raw)]
        col2_items = split_to_sentences(col2_raw) or [str(col2_raw)]
        col3_items = split_to_sentences(col3_raw) or [str(col3_raw)]
    else:
        # content list formatida — teng 3 ga bo'lamiz
        items = content_data.get("content", [])
        if isinstance(items, str):
            items = [items]
        n = len(items)
        if n == 0:
            # Matn yo'q — sarlavhani 3 ga bo'lib yozamiz
            blks = split_text_into_blocks(title if title else "—", 3)
            items = [blks[0] or "—", blks[1] or "—", blks[2] or "—"]
            n = 3
        elif n == 1:
            blks = split_text_into_blocks(str(items[0]), 3)
            items = [blks[0] or str(items[0]), blks[1] or str(items[0]), blks[2] or str(items[0])]
            n = 3
        elif n == 2:
            items = [items[0], items[1], items[0]]
            n = 3
        base = n // 3
        rem = n % 3
        sizes = [base + (1 if i < rem else 0) for i in range(3)]
        col1_items = items[:sizes[0]]
        col2_items = items[sizes[0]:sizes[0]+sizes[1]]
        col3_items = items[sizes[0]+sizes[1]:]

    def write_col(shape, col_items, header=""):
        if shape is None or not shape.has_text_frame:
            return
        tf = shape.text_frame
        tf.word_wrap = True
        all_items = ([header] if header else []) + list(col_items)
        total_chars = sum(len(s) for s in all_items) if all_items else 10
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)

        txBody = tf._txBody
        # lstStyle tozalash (shablon indent ni olib tashlash)
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        if header:
            safe_h = header.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100+200)}" b="1" dirty="0"/>'
                f'<a:t>{safe_h}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

        for item in col_items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Sarlavha TextBox qo'shish (yuqori qismga, 3 ustundan yuqorida)
    if title:
        from pptx.util import Cm, Pt as _Pt
        from pptx.dml.color import RGBColor
        txBox = slide.shapes.add_textbox(Cm(2.91), Cm(0.5), Cm(19.58), Cm(2.5))
        tf_title = txBox.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.alignment = PP_ALIGN.LEFT
        run_title = p_title.add_run()
        run_title.text = title
        run_title.font.size = _Pt(28)
        run_title.font.bold = True
        run_title.font.color.rgb = RGBColor(0xFF, 0xA5, 0x00)  # To'q sariq rang (shablon rangiga mos)

    # Shablon: Shape[2]=chap(idx=5), Shape[0]=o'rta(idx=1), Shape[1]=o'ng(idx=3)
    write_col(slide.shapes[2] if len(slide.shapes) > 2 else None, col1_items)
    write_col(slide.shapes[0] if len(slide.shapes) > 0 else None, col2_items)
    write_col(slide.shapes[1] if len(slide.shapes) > 1 else None, col3_items)


def fill_t9_slide_4_image_left(slide, content_data, image_query=None):
    """
    9-Shablon Slayd 4 — Sarlavha + Rasm chap + Matn o'ng.
    Shape[0]: Sarlavha (TITLE idx=0) — yuqori chap
    Shape[1]: Matn (BODY idx=4294967295) — o'ng
    Shape[2]: Rasm (PICTURE) — chap
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(28)
        run.font.bold = True

    # Shape[1]: Matn (o'ng) - indent=0, marL=0 bilan tekis
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)

        txBody = tf._txBody
        # lstStyle ni ham tozalash (shablon indent ni olib tashlash)
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Shape[2]: Rasm (chap tomonda)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(2.48)
            top = Cm(4.60)
            width = Cm(13.32)
            height = Cm(8.24)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T9] Slayd 4 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T9] Slayd 4 rasm xatolik: {e}")


def fill_t9_slide_5_two_images(slide, content_data, image_query=None):
    """
    9-Shablon Slayd 5 — Sarlavha + Matn o'ng + 2 rasm.
    Shape[0]: Sarlavha (TITLE idx=0)
    Shape[1]: Matn (SUBTITLE idx=1) — o'ng
    Shape[2]: Katta rasm (PICTURE) — o'rta
    Shape[3]: Kichik rasm (PICTURE) — chap past
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(28)
        run.font.bold = True

    # Shape[1]: Matn (o'ng) - indent=0, marL=0 bilan tekis
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)

        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Rasm 1 (katta, o'rta — Shape[2])
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(7.91)
            top = Cm(3.41)
            width = Cm(5.88)
            height = Cm(9.73)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T9] Slayd 5 katta rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T9] Slayd 5 katta rasm xatolik: {e}")

    # Rasm 2 (kichik, chap past — Shape[3])
    img2_path = fetch_image(query if 'query' in dir() else title)
    if img2_path:
        try:
            left2 = Cm(3.40)
            top2 = Cm(6.87)
            width2 = Cm(3.07)
            height2 = Cm(5.47)
            slide.shapes.add_picture(img2_path, left2, top2, width2, height2)
            if os.path.isfile(img2_path):
                os.remove(img2_path)
            logging.info(f"[T9] Slayd 5 kichik rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T9] Slayd 5 kichik rasm xatolik: {e}")


def fill_t9_slide_6_image_right(slide, content_data, image_query=None):
    """
    9-Shablon Slayd 6 — Sarlavha + Rasm o'ng + Matn chap.
    Shape[0]: Rasm (PICTURE) — o'ng
    Shape[1]: Sarlavha (TITLE idx=0)
    Shape[2]: Matn (SUBTITLE idx=1) — chap
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[1]: Sarlavha
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(28)
        run.font.bold = True

    # Shape[2]: Matn (chap) - sarlavhadan pastga tushiriladi, indent=0
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        from pptx.util import Cm as _Cm
        slide.shapes[2].top = _Cm(4.0)
        slide.shapes[2].height = _Cm(8.5)
        tf = slide.shapes[2].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)

        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Shape[0]: Rasm (o'ng tomonda)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(14.01)
            top = Cm(3.68)
            width = Cm(5.48)
            height = Cm(7.22)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T9] Slayd 6 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T9] Slayd 6 rasm xatolik: {e}")


def fill_t9_slide_7_image_left2(slide, content_data, image_query=None):
    """
    9-Shablon Slayd 7 — Sarlavha + Matn o'ng + Rasm chap.
    Shape[0]: Matn (SUBTITLE idx=1) — o'ng
    Shape[1]: Sarlavha (TITLE idx=0)
    Shape[2]: Rasm (PICTURE) — chap
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[1]: Sarlavha
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(28)
        run.font.bold = True

    # Shape[0]: Matn (o'ng) - chapga tekislangan, indent=0
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)

        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Shape[2]: Rasm (chap)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(3.39)
            top = Cm(4.69)
            width = Cm(9.22)
            height = Cm(5.59)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T9] Slayd 7 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T9] Slayd 7 rasm xatolik: {e}")


def fill_t9_slide_8_conclusion(slide, content_data):
    """
    9-Shablon Oxirgi slayd — har doim "E'TIBORINGIZ UCHUN RAHMAT!" bilan tugaydi.
    Shape[0]: TEXT_BOX — xulosa matni
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    # Oxirgi slayd har doim shu matn bilan tugaydi
    text = "E'TIBORINGIZ UCHUN RAHMAT!"

    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(32)
        run.font.bold = True


def generate_template_9_presentation(prs, topic, requested_slide_count, language,
                                      name_surname, plan, content_data_list, user_images=None):
    """
    9-shablon asosida to'liq prezentatsiya yaratadi.
    """
    import io

    total_content_slides = build_slide_structure_9(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}

    # Slayd 1 — Muqova
    fill_t9_slide_1_cover(prs.slides[0], topic, name_surname)

    # Slayd 2 — Reja
    fill_t9_slide_2_plan(prs.slides[1], plan_dict)

    # Kontent slaydlari (3-dan boshlab)
    user_img_idx = 0
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break

        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        image_query = data.get("image_query", topic)

        slide_type = i % 5
        has_image = slide_type in [1, 2, 3, 4]  # Rasm ishlatadigan slayd turlari

        if has_image:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                img_arg = img_path if img_path else image_query
            else:
                img_arg = image_query
        else:
            img_arg = None

        if slide_type == 0:
            fill_t9_slide_3_three_cols(slide, data)
        elif slide_type == 1:
            fill_t9_slide_4_image_left(slide, data, img_arg)
        elif slide_type == 2:
            fill_t9_slide_5_two_images(slide, data, img_arg)
        elif slide_type == 3:
            fill_t9_slide_6_image_right(slide, data, img_arg)
        elif slide_type == 4:
            fill_t9_slide_7_image_left2(slide, data, img_arg)

        logging.info(f"  [T9] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")

    # Xulosa slayd
    conclusion_slide = prs.slides[-1]
    conclusion_data = generate_slide_content(topic, len(prs.slides), len(prs.slides), language, is_conclusion=True)
    if not conclusion_data:
        conclusion_data = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}
    fill_t9_slide_8_conclusion(conclusion_slide, conclusion_data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════
# 10-SHABLON FUNKSIYALARI (10.pptx)
# ═══════════════════════════════════════════════════════════════
# Slayd tuzilmasi:
#   Slayd 1: TITLE         — Muqova (sarlavha + subtitle)
#   Slayd 2: TITLE_AND_BODY — Reja
#   Slayd 3: CUSTOM_10     — Sarlavha + Matn (o'ng) + Rasm (chap, to'liq)
#   Slayd 4: CUSTOM_2      — Sarlavha (chap) + 2 matn bloki (o'ng yuqori/quyi)
#   Slayd 5: CUSTOM_17     — Sarlavha + Rasm (chap) + Matn (o'ng)
#   Slayd 6: CUSTOM_18     — Sarlavha + Kichik rasm (chap) + Katta matn (o'ng)
#   Slayd 7: CUSTOM_19     — Sarlavha + Matn (chap) + Rasm (o'ng)
#   Slayd 8: Title and text 7 — Xulosa: "E'TIBORINGIZ UCHUN RAHMAT!"
# ═══════════════════════════════════════════════════════════════

SLIDE_TYPE_NAMES_T10 = {
    0: "image_left",    # Slayd 3: sarlavha + matn o'ng + rasm chap (to'liq)
    1: "two_columns",   # Slayd 4: sarlavha chap + 2 matn bloki o'ng
    2: "image_left",    # Slayd 5: sarlavha + rasm chap + matn o'ng
    3: "image_left",    # Slayd 6: sarlavha + kichik rasm chap + katta matn o'ng
    4: "image_right",   # Slayd 7: sarlavha + matn chap + rasm o'ng
}

CONTENT_SLIDE_TEMPLATE_INDICES_T10 = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7


def build_slide_structure_10(prs, requested_content_count):
    """
    10-shablon uchun slayd tuzilmasini yaratadi.
    Muqova(1) + Reja(1) + Kontent(N) + Xulosa(1) = jami
    5 ta kontent shablon slayd bor (index 2-6), to'liq to'plamlar sifatida takrorlanadi.
    """
    n_templates = len(CONTENT_SLIDE_TEMPLATE_INDICES_T10)  # 5
    full_repeats = max(1, round(requested_content_count / n_templates))
    total_content_slides = full_repeats * n_templates
    logging.info(f"[T10] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")

    conclusion_current_index = 7  # 0-indexed: slayd 8

    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_T10:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T10] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")

    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T10] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def fill_t10_slide_1_cover(slide, topic, name_surname):
    """
    10-Shablon Slayd 1 — Muqova.
    Shape[0]: Sarlavha (TITLE idx=0)
    Shape[1]: Subtitle (SUBTITLE idx=1)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = topic.upper() if topic else "TAQDIMOT"
        font_pt = calc_body_font_pt(len(topic), base_pt=36, min_pt=20, max_pt=44)
        run.font.size = Pt(font_pt)
        run.font.bold = True

    # Shape[1]: Ism-familiya
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = name_surname if name_surname else ""
        run.font.size = Pt(16)


def fill_t10_slide_2_plan(slide, plan_data):
    """
    10-Shablon Slayd 2 — Reja.
    Shape[0]: Sarlavha (CENTER_TITLE idx=0)
    Shape[1]: Reja matni (SUBTITLE idx=1)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = plan_data.get("title", "Reja") if isinstance(plan_data, dict) else "Reja"
    items = plan_data.get("content", []) if isinstance(plan_data, dict) else []

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "REJA"
        run.font.size = Pt(28)
        run.font.bold = True

    # Shape[1]: Reja ro'yxati
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(str(s)) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=22)

        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        import re as _re
        for idx, item in enumerate(items):
            clean_item = _re.sub(r'^\d+\.\s*\d*\.?\s*', '', str(item)).strip()
            safe_item = clean_item.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" b="1" dirty="0"/>'
                f'<a:t>{idx+1}. {safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))


def fill_t10_slide_3_image_full_left(slide, content_data, image_query=None):
    """
    10-Shablon Slayd 3 (CUSTOM_10) — Sarlavha + Matn (o'ng) + Rasm (chap, to'liq balandlik).
    Shape[0]: Sarlavha (CENTER_TITLE idx=0) — o'ng yuqori
    Shape[1]: Matn (SUBTITLE idx=1) — o'ng
    Shape[2]: Rasm (PICTURE) — chap, to'liq balandlik
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[0]: Sarlavha (o'ng yuqori)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(26)
        run.font.bold = True

    # Shape[1]: Matn (o'ng)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)

        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Shape[2]: Rasm (chap, to'liq balandlik: left=-0.67cm, top=0.49cm, w=14.50cm, h=13.17cm)
    # Eski PICTURE shape larni o'chirish
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
    pic_shapes = [s for s in slide.shapes if s.shape_type == _MSO.PICTURE]
    for ps in pic_shapes:
        sp = ps._element
        sp.getparent().remove(sp)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            # Rasm o'ng chegarasi = matn chap chegarasi bilan bir chiziqda, 1 sm chapga surilgan
            # left = 1.8002 - 1.0 = 0.8002 cm
            left = Cm(0.8002)
            top = Cm(0.4943)
            width = Cm(8.0)
            height = Cm(13.17)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T10] Slayd 3 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T10] Slayd 3 rasm xatolik: {e}")


def fill_t10_slide_4_two_cols(slide, content_data, image_query=None):
    """
    10-Shablon Slayd 4 (CUSTOM_2) — Sarlavha (chap markazda) + 2 matn bloki (o'ng yuqori/quyi).
    Shape[0]: Sarlavha (CENTER_TITLE idx=0) — chap, vertikal markazda
    Shape[1]: Yuqori matn bloki (SUBTITLE idx=1) — o'ng yuqori
    Shape[2]: Quyi matn bloki (SUBTITLE idx=2) — o'ng quyi
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    col1 = content_data.get("col1", "")
    col2 = content_data.get("col2", "")
    # col1/col2 yo'q bo'lsa content dan olamiz
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    if not col1 and items:
        col1 = items[0] if len(items) > 0 else ""
    if not col2 and items:
        col2 = items[1] if len(items) > 1 else items[0] if items else ""

    # Shape[0]: Sarlavha (chap, vertikal markazda)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(22)
        run.font.bold = True

    def write_text_block(shape, text):
        if shape is None or not shape.has_text_frame:
            return
        tf = shape.text_frame
        tf.word_wrap = True
        font_pt = calc_body_font_pt(len(str(text)), base_pt=14, min_pt=10, max_pt=16)
        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        # Jumlalarga bo'lish
        import re as _re
        sentences = _re.split(r'(?<=[.!?])\s+', str(text).strip())
        for sentence in sentences:
            if not sentence.strip():
                continue
            safe = sentence.strip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Shape[1]: Yuqori matn bloki
    if len(slide.shapes) > 1:
        write_text_block(slide.shapes[1], col1)

    # Shape[2]: Quyi matn bloki
    if len(slide.shapes) > 2:
        write_text_block(slide.shapes[2], col2)


def fill_t10_slide_5_image_left(slide, content_data, image_query=None):
    """
    10-Shablon Slayd 5 (CUSTOM_17) — Sarlavha + Rasm (chap) + Matn (o'ng).
    Shape[0]: Rasm (PICTURE) — chap: left=3.27cm, top=2.80cm, w=6.92cm, h=9.19cm
    Shape[1]: Sarlavha (CENTER_TITLE idx=0) — yuqori
    Shape[2]: Matn (TEXT_BOX) — o'ng: left=12.96cm, top=3.93cm, w=10.30cm, h=9.10cm
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[1]: Sarlavha
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(26)
        run.font.bold = True

    # Shape[2]: Matn (o'ng)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)

        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Shape[0]: Rasm (chap: left=1.27cm, top=2.80cm, w=10.30cm, h=9.19cm - landscape)
    # Eski PICTURE shape larni o'chirish
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
    pic_shapes = [s for s in slide.shapes if s.shape_type == _MSO.PICTURE]
    for ps in pic_shapes:
        sp = ps._element
        sp.getparent().remove(sp)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(3.2663)
            top = Cm(2.7973)
            width = Cm(6.9250)
            height = Cm(9.1867)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T10] Slayd 5 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T10] Slayd 5 rasm xatolik: {e}")


def fill_t10_slide_6_small_image_left(slide, content_data, image_query=None):
    """
    10-Shablon Slayd 6 (CUSTOM_18) — Sarlavha + Kichik rasm (chap) + Katta matn (o'ng).
    Shape[0]: Sarlavha (CENTER_TITLE idx=0) — yuqori
    Shape[1]: Kichik rasm (PICTURE) — chap: left=2.66cm, top=3.98cm, w=3.61cm, h=7.36cm
    Shape[2]: Katta matn (TEXT_BOX) — o'ng: left=8.46cm, top=3.34cm, w=13.28cm, h=9.28cm
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(26)
        run.font.bold = True

    # Shape[2]: Katta matn (o'ng)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)

        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Shape[1]: Rasm (chap: left=1.27cm, top=3.34cm, w=6.50cm, h=9.28cm - landscape uchun kengaytirilgan)
    # Eski PICTURE shape larni o'chirish
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
    pic_shapes = [s for s in slide.shapes if s.shape_type == _MSO.PICTURE]
    for ps in pic_shapes:
        sp = ps._element
        sp.getparent().remove(sp)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(2.6635)
            top = Cm(3.9813)
            width = Cm(3.6117)
            height = Cm(7.3583)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T10] Slayd 6 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T10] Slayd 6 rasm xatolik: {e}")


def fill_t10_slide_7_image_right(slide, content_data, image_query=None):
    """
    10-Shablon Slayd 7 (CUSTOM_19) — Sarlavha + Matn (o'ng) + Rasm (chap).
    Shape[0]: Sarlavha (CENTER_TITLE idx=0) — yuqori
    Shape[1]: Matn (TEXT_BOX) — o'ng: left=12.70cm, top=3.39cm, w=9.88cm, h=9.22cm
    Shape[2]: Rasm (PICTURE) — chap: left=1.79cm, top=3.90cm, w=8.77cm, h=5.46cm
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]

    # Shape[0]: Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(26)
        run.font.bold = True

    # Shape[1]: Matn (o'ng tomonda - shablonda left=12.70)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)

        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)

        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))

    # Shape[2]: Rasm (chap tomonda: left=1.79cm, top=3.39cm, w=9.88cm, h=9.22cm)
    # Eski PICTURE shape larni o'chirish
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
    pic_shapes = [s for s in slide.shapes if s.shape_type == _MSO.PICTURE]
    for ps in pic_shapes:
        sp = ps._element
        sp.getparent().remove(sp)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(1.7895)
            top = Cm(3.9003)
            width = Cm(8.7675)
            height = Cm(5.4633)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T10] Slayd 7 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T10] Slayd 7 rasm xatolik: {e}")


def fill_t10_slide_8_conclusion(slide, content_data):
    """
    10-Shablon Oxirgi slayd — har doim "E'TIBORINGIZ UCHUN RAHMAT!" bilan tugaydi.
    Shape[0]: CENTER_TITLE idx=0
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    text = "E'TIBORINGIZ UCHUN RAHMAT!"

    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(32)
        run.font.bold = True


def generate_template_10_presentation(prs, topic, requested_slide_count, language,
                                       name_surname, plan, content_data_list, user_images=None):
    """
    10-shablon asosida to'liq prezentatsiya yaratadi.
    """
    import io

    total_content_slides = build_slide_structure_10(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}

    # Slayd 1 — Muqova
    fill_t10_slide_1_cover(prs.slides[0], topic, name_surname)

    # Slayd 2 — Reja
    fill_t10_slide_2_plan(prs.slides[1], plan_dict)

    # Kontent slaydlari (3-dan boshlab)
    user_img_idx = 0
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break

        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        image_query = data.get("image_query", topic)

        slide_type = i % 5
        has_image = slide_type in [0, 1, 2, 3, 4]  # Barcha slayd turlari rasm ishlatadi

        if has_image:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                img_arg = img_path if img_path else image_query
            else:
                img_arg = image_query
        else:
            img_arg = None

        if slide_type == 0:
            fill_t10_slide_3_image_full_left(slide, data, img_arg)
        elif slide_type == 1:
            fill_t10_slide_4_two_cols(slide, data, img_arg)
        elif slide_type == 2:
            fill_t10_slide_5_image_left(slide, data, img_arg)
        elif slide_type == 3:
            fill_t10_slide_6_small_image_left(slide, data, img_arg)
        elif slide_type == 4:
            fill_t10_slide_7_image_right(slide, data, img_arg)

        logging.info(f"  [T10] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")

    # Xulosa slayd
    conclusion_slide = prs.slides[-1]
    fill_t10_slide_8_conclusion(conclusion_slide, {})

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════════════
# 11-SHABLON (Technology Consulting) FUNKSIYALARI
# ═══════════════════════════════════════════════════════════════════════

SLIDE_TYPE_NAMES_T11 = {
    0: "image_left",    # Slayd 4: sarlavha + rasm chap + 2 matn bloki o'ng
    1: "image_right",   # Slayd 5: sarlavha + quote chap + rasm o'ng
    2: "two_columns",   # Slayd 6: sarlavha + matn o'ng (dekorativ rasm fonda)
    3: "four_blocks",   # Slayd 7: sarlavha + 4 ta matn bloki (infografika)
    4: "image_left",    # Slayd 3: sarlavha + 2 ta matn bloki (reja uslubi)
}

CONTENT_SLIDE_TEMPLATE_INDICES_T11 = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7


def build_slide_structure_11(prs, requested_content_count):
    """
    11-shablon uchun slayd tuzilmasini yaratadi.
    Muqova(1) + Reja(1) + Kontent(N) + Xulosa(1) = jami
    5 ta kontent shablon slayd bor (index 2-6), to'liq to'plamlar sifatida takrorlanadi.
    """
    n_templates = len(CONTENT_SLIDE_TEMPLATE_INDICES_T11)  # 5
    full_repeats = max(1, round(requested_content_count / n_templates))
    total_content_slides = full_repeats * n_templates
    logging.info(f"[T11] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    conclusion_current_index = 7  # 0-indexed: slayd 8
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_T11:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T11] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T11] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def fill_t11_slide_1_cover(slide, topic, name_surname):
    """
    11-Shablon Slayd 1 — Muqova.
    Shape[0]: Sarlavha (CENTER_TITLE idx=0) — o'ng yuqori
    Shape[1]: Subtitle (SUBTITLE idx=1) — o'ng pastki
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = topic.upper() if topic else "TAQDIMOT"
        font_pt = calc_body_font_pt(len(topic), base_pt=32, min_pt=18, max_pt=40)
        run.font.size = Pt(font_pt)
        run.font.bold = True
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = name_surname if name_surname else ""
        run.font.size = Pt(14)
def fill_t11_slide_2_plan(slide, plan_data):
    """
    11-Shablon Slayd 2 — Reja.
    Shape[0]: Sarlavha (TITLE idx=0)
    Shape[1]: Matn (BODY idx=1)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = plan_data.get("title", "REJA")
    items = plan_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title.upper() if title else "REJA"
        run.font.size = Pt(28)
        run.font.bold = True
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        import re as _re
        for idx, item in enumerate(items, 1):
            # item ichida allaqachon raqam bor bo'lsa (masalan "1.1. ...") — tozalash
            item_str = str(item).strip()
            # "1.1." yoki "1." kabi boshlanishni olib tashlash
            item_str = _re.sub(r'^\d+[\d\.]*\.?\s*', '', item_str).strip()
            safe_item = item_str.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="342900" indent="-342900"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="1600" b="1" dirty="0"/>'
                f'<a:t>{idx}. {safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))
def fill_t11_slide_3_two_text(slide, content_data, image_query=None):
    """
    11-Shablon Slayd 3 (CUSTOM) — Sarlavha + 2 ta matn bloki (o'ng tekislash).
    Shape[0]: Sarlavha (TITLE idx=0): left=1.0583, top=0.8667, w=23.0928, h=1.6000
    Shape[1]: Matn blok 1 (SUBTITLE idx=2): left=4.6228, top=2.1464, w=9.8944, h=4.0725
    Shape[2]: Matn blok 2 (SUBTITLE idx=8): left=2.1476, top=7.8317, w=9.7393, h=4.7411
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    # Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(24)
        run.font.bold = True
    # Matn blok 1 (yuqori) va 2 (pastki) — ikkalasi ham chap tekislash
    half = max(1, len(items) // 2)
    items1 = items[:half]
    items2 = items[half:]
    for shape_idx, item_list in [(1, items1), (2, items2)]:
        if len(slide.shapes) > shape_idx and slide.shapes[shape_idx].has_text_frame:
            tf = slide.shapes[shape_idx].text_frame
            tf.word_wrap = True
            total_chars = sum(len(s) for s in item_list)
            font_pt = calc_body_font_pt(total_chars, base_pt=13, min_pt=10, max_pt=16)
            txBody = tf._txBody
            for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
                lst.clear()
            for p_elem in txBody.findall(f'{{{ns_a}}}p'):
                txBody.remove(p_elem)
            for item in item_list:
                safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                p_xml = (
                    f'<a:p xmlns:a="{ns_a}">'
                    f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                    f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                    f'<a:t>{safe_item}</a:t></a:r></a:p>'
                )
                txBody.append(etree.fromstring(p_xml))


def fill_t11_slide_4_image_left(slide, content_data, image_query=None):
    """
    11-Shablon Slayd 4 (TITLE_ONLY) — Sarlavha + Rasm chap + 2 matn bloki o'ng.
    Shape[0]: PICTURE: left=3.1502, top=3.8740, w=11.4259, h=7.7496
    Shape[1]: Sarlavha (TITLE idx=0): left=2.7672, top=0.9398, w=20.4439, h=1.5908
    Shape[2]: Matn o'ng pastki: left=15.1097, top=8.6434, w=7.6450, h=3.5975
    Shape[3]: Matn o'ng yuqori: left=15.1097, top=3.8740, w=7.6450, h=3.4261
    Shape[4]: Izoh pastki: left=3.4925, top=11.3438, w=10.5500, h=1.4617
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    # Sarlavha (Shape[1])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(24)
        run.font.bold = True
    # Matn o'ng yuqori (Shape[3])
    half = max(1, len(items) // 2)
    items1 = items[:half]
    items2 = items[half:]
    for shape_idx, item_list in [(3, items1), (2, items2)]:
        if len(slide.shapes) > shape_idx and slide.shapes[shape_idx].has_text_frame:
            tf = slide.shapes[shape_idx].text_frame
            tf.word_wrap = True
            total_chars = sum(len(s) for s in item_list)
            font_pt = calc_body_font_pt(total_chars, base_pt=12, min_pt=9, max_pt=15)
            txBody = tf._txBody
            for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
                lst.clear()
            for p_elem in txBody.findall(f'{{{ns_a}}}p'):
                txBody.remove(p_elem)
            for item in item_list:
                safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                p_xml = (
                    f'<a:p xmlns:a="{ns_a}">'
                    f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                    f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                    f'<a:t>{safe_item}</a:t></a:r></a:p>'
                )
                txBody.append(etree.fromstring(p_xml))
    # Izoh (Shape[4]) - bo'sh qilish
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        tf = slide.shapes[4].text_frame
        tf.clear()
    # Eski PICTURE shape larni o'chirish va yangi rasm qo'shish
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
    pic_shapes = [s for s in slide.shapes if s.shape_type == _MSO.PICTURE]
    for ps in pic_shapes:
        sp = ps._element
        sp.getparent().remove(sp)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(3.1502)
            top = Cm(3.8740)
            width = Cm(11.4259)
            height = Cm(7.7496)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T11] Slayd 4 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T11] Slayd 4 rasm xatolik: {e}")


def fill_t11_slide_5_quote_image(slide, content_data, image_query=None):
    """
    11-Shablon Slayd 5 (ONE_COLUMN_TEXT) — Sarlavha + Quote chap + Rasm o'ng.
    Shape[0]: Sarlavha (TITLE idx=0): left=2.9861, top=1.9812, w=10.4369, h=1.1175
    Shape[1]: Quote/Matn (SUBTITLE idx=1): left=2.2140, top=4.8941, w=10.2588, h=6.5668
    Shape[2]: PICTURE: left=13.1851, top=1.6425, w=10.4369, h=10.4648
    """
    from pptx.util import Pt, Cm
    from pptx.enum.text import PP_ALIGN
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    # col1/col2 formatidan ham matn olish (two_columns formatida kelsa)
    if not items:
        col1 = content_data.get("col1", "")
        col2 = content_data.get("col2", "")
        if col1 or col2:
            items = [x for x in [col1, col2] if x]
    # Sarlavha (Shape[0])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(22)
        run.font.bold = True
    # Matn/Quote (Shape[1]) - placeholder idx=1
    # Barcha placeholder larni ko'rib chiqib, idx=1 ni topish
    target_shape = None
    for s in slide.shapes:
        if s.is_placeholder and s.placeholder_format.idx == 1 and s.has_text_frame:
            target_shape = s
            break
    if target_shape is None and len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        target_shape = slide.shapes[1]
    if target_shape is not None:
        tf = target_shape.text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)
        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))
    # Eski PICTURE shape larni o'chirish va yangi rasm qo'shish (Shape[2])
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
    pic_shapes = [s for s in slide.shapes if s.shape_type == _MSO.PICTURE]
    for ps in pic_shapes:
        sp = ps._element
        sp.getparent().remove(sp)
    if image_query and os.path.isfile(image_query):
        final_img_path = image_query
    else:
        query = image_query or content_data.get("image_query", title)
        final_img_path = fetch_image(query)
    if final_img_path:
        try:
            left = Cm(13.1851)
            top = Cm(1.6425)
            width = Cm(10.4369)
            height = Cm(10.4648)
            slide.shapes.add_picture(final_img_path, left, top, width, height)
            if os.path.isfile(final_img_path):
                os.remove(final_img_path)
            logging.info(f"[T11] Slayd 5 rasm joylashtirildi.")
        except Exception as e:
            logging.error(f"[T11] Slayd 5 rasm xatolik: {e}")


def fill_t11_slide_6_title_body(slide, content_data, image_query=None):
    """
    11-Shablon Slayd 6 (CUSTOM_10) — Sarlavha + Matn o'ng (dekorativ rasm fonda).
    Shape[0]: Sarlavha (TITLE idx=0): left=3.5306, top=0.9398, w=18.3642, h=1.5242
    Shape[1]: Matn (SUBTITLE idx=1): left=13.0302, top=4.1426, w=9.8808, h=7.3783
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    # col1/col2 formatidan ham matn olish (two_columns formatida kelsa)
    if not items:
        col1 = content_data.get("col1", "")
        col2 = content_data.get("col2", "")
        if col1 or col2:
            items = [x for x in [col1, col2] if x]
    # Sarlavha (Shape[0])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(26)
        run.font.bold = True
    # Matn (Shape[1])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)
        txBody = tf._txBody
        for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
            lst.clear()
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
        for item in items:
            safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">'
                f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                f'<a:t>{safe_item}</a:t></a:r></a:p>'
            )
            txBody.append(etree.fromstring(p_xml))


def fill_t11_slide_7_four_blocks(slide, content_data, image_query=None):
    """
    11-Shablon Slayd 7 (CUSTOM_15) — Sarlavha + 2 ta matn bloki (chap va o'ng).
    Faqat 2 ta matn bloki ishlatiladi: chap va o'ng.
    Bo'sh 2 ta blok (pastki) o'chiriladi, matn joylashgan 2 ta blok 3sm pastga siljiydi.
    Shape[0]: Sarlavha (TITLE idx=0): left=1.0372, top=1.4637, w=23.5797, h=1.5908
    Shape[1]: Blok 1 chap (SUBTITLE idx=2): left=2.2733, top=4.5665 -> 7.5665, w=6.3651, h=3.1115
    Shape[2]: Blok 2 chap pastki (SUBTITLE idx=3): O'CHIRILADI
    Shape[3]: Blok 3 o'ng (SUBTITLE idx=5): left=17.0245, top=4.5665 -> 7.5665, w=6.3651, h=3.1115
    Shape[4]: Blok 4 o'ng pastki (SUBTITLE idx=7): O'CHIRILADI
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    # Sarlavha (Shape[0])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(26)
        run.font.bold = True
    from pptx.util import Cm, Emu
    # Pastki 2 ta blokni (Shape[2] idx=3 va Shape[4] idx=7) o'chirish
    # Placeholder idx=3 va idx=7 ni topib o'chirish
    shapes_to_remove = []
    for s in slide.shapes:
        if s.is_placeholder and s.placeholder_format.idx in [3, 7]:
            shapes_to_remove.append(s)
    for s in shapes_to_remove:
        sp = s._element
        sp.getparent().remove(sp)
    # Yuqori 2 ta blokni (idx=2 va idx=5) 3sm pastga siljitish
    # va matn yozish
    # Chap blok (idx=2) va o'ng blok (idx=5)
    shift_emu = int(3.0 * 914400 / 2.54)  # 3sm EMU ga
    left_items = items[:max(1, len(items)//2)] if items else []
    right_items = items[max(1, len(items)//2):] if len(items) > 1 else []
    for s in slide.shapes:
        if not s.is_placeholder:
            continue
        ph_idx = s.placeholder_format.idx
        if ph_idx == 2:  # Chap yuqori blok
            # 3sm pastga siljitish
            s.top = s.top + shift_emu
            # Matn yozish
            if s.has_text_frame:
                tf = s.text_frame
                tf.word_wrap = True
                total_chars = sum(len(x) for x in left_items)
                font_pt = calc_body_font_pt(total_chars, base_pt=13, min_pt=10, max_pt=16)
                txBody = tf._txBody
                for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
                    lst.clear()
                for p_elem in txBody.findall(f'{{{ns_a}}}p'):
                    txBody.remove(p_elem)
                for item in left_items:
                    safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    p_xml = (
                        f'<a:p xmlns:a="{ns_a}">'
                        f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                        f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                        f'<a:t>{safe_item}</a:t></a:r></a:p>'
                    )
                    txBody.append(etree.fromstring(p_xml))
        elif ph_idx == 5:  # O'ng yuqori blok
            # 3sm pastga siljitish
            s.top = s.top + shift_emu
            # Matn yozish
            if s.has_text_frame:
                tf = s.text_frame
                tf.word_wrap = True
                total_chars = sum(len(x) for x in right_items)
                font_pt = calc_body_font_pt(total_chars, base_pt=13, min_pt=10, max_pt=16)
                txBody = tf._txBody
                for lst in txBody.findall(f'{{{ns_a}}}lstStyle'):
                    lst.clear()
                for p_elem in txBody.findall(f'{{{ns_a}}}p'):
                    txBody.remove(p_elem)
                for item in right_items:
                    safe_item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    p_xml = (
                        f'<a:p xmlns:a="{ns_a}">'
                        f'<a:pPr algn="l" marL="0" indent="0"><a:buNone/></a:pPr>'
                        f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0"/>'
                        f'<a:t>{safe_item}</a:t></a:r></a:p>'
                    )
                    txBody.append(etree.fromstring(p_xml))


def fill_t11_slide_8_conclusion(slide, content_data):
    """
    11-Shablon Oxirgi slayd — "E'TIBORINGIZ UCHUN RAHMAT!" bilan tugaydi.
    Shape[0]: Sarlavha (TITLE idx=0): left=5.6698, top=5.0333, w=13.7417, h=3.8463
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        tf = slide.shapes[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "E'TIBORINGIZ UCHUN RAHMAT!"
        run.font.size = Pt(32)
        run.font.bold = True


def generate_template_11_presentation(prs, topic, requested_slide_count, language,
                                       name_surname, plan, content_data_list, user_images=None):
    """
    11-shablon asosida to'liq prezentatsiya yaratadi.
    """
    import io
    total_content_slides = build_slide_structure_11(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}
    # Slayd 1 — Muqova
    fill_t11_slide_1_cover(prs.slides[0], topic, name_surname)
    # Slayd 2 — Reja
    fill_t11_slide_2_plan(prs.slides[1], plan_dict)
    # Kontent slaydlari (3-dan boshlab)
    user_img_idx = 0
    # Slayd turlari: 0=slayd3(two_text), 1=slayd4(image_left), 2=slayd5(quote+image),
    #                3=slayd6(title+body), 4=slayd7(four_blocks)
    IMAGE_SLIDE_TYPES = [1, 2]  # Faqat slayd4 va slayd5 da rasm bor
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        image_query = data.get("image_query", topic)
        slide_type = i % 5
        has_image = slide_type in IMAGE_SLIDE_TYPES
        if has_image:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                img_arg = img_path if img_path else image_query
            else:
                img_arg = image_query
        else:
            img_arg = None
        if slide_type == 0:
            fill_t11_slide_3_two_text(slide, data, img_arg)
        elif slide_type == 1:
            fill_t11_slide_4_image_left(slide, data, img_arg)
        elif slide_type == 2:
            fill_t11_slide_5_quote_image(slide, data, img_arg)
        elif slide_type == 3:
            fill_t11_slide_6_title_body(slide, data, img_arg)
        elif slide_type == 4:
            fill_t11_slide_7_four_blocks(slide, data, img_arg)
        logging.info(f"  [T11] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    # Xulosa slayd
    conclusion_slide = prs.slides[-1]
    fill_t11_slide_8_conclusion(conclusion_slide, {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ============================================================
# 12-SHABLON (AESTHETIC / BOTANICAL) — QAYTA YOZILGAN
# ============================================================

SLIDE_TYPE_NAMES_T12 = {
    0: "one_column",      # Slayd 3: sarlavha + bitta katta matn bloki (chap), dekor o'ng
    1: "image_left",      # Slayd 4: sarlavha + rasm chap (freeform), matn o'ng
    2: "two_columns",     # Slayd 5: sarlavha + 2 ta ustun matni
    3: "three_columns",   # Slayd 6: sarlavha + 3 ta ustun matni
    4: "image_right",     # Slayd 7: sarlavha + matn chap, rasm o'ng (freeform)
}
CONTENT_SLIDE_TEMPLATE_INDICES_T12 = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7


def build_slide_structure_12(prs, requested_content_count):
    """
    12-shablon uchun slayd tuzilmasini yaratadi.
    Muqova(1) + Reja(1) + Kontent(N) + Xulosa(1) = jami
    5 ta kontent shablon slayd bor (index 2-6), to'liq to'plamlar sifatida takrorlanadi.
    """
    n_templates = len(CONTENT_SLIDE_TEMPLATE_INDICES_T12)  # 5
    full_repeats = max(1, round(requested_content_count / n_templates))
    total_content_slides = full_repeats * n_templates
    logging.info(f"[T12] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    conclusion_current_index = 7  # 0-indexed: slayd 8
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_T12:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T12] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T12] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def _t12_clear_and_write(txBody, ns_a, paragraphs_data):
    """
    txBody dagi barcha paragraflarni o'chirib, yangi paragraflar yozadi.
    paragraphs_data: list of dict {algn, marL, indent, spcPts, runs: [{sz, b, text}]}
    """
    from lxml import etree
    import re
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        # pPr atributlari
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"' 
        if indent:
            pPr_attrs += f' indent="{indent}"' 
        # spcBef
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        # runs
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            text = run.get('text', '')
            safe_text = str(text).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b}" dirty="0">' 
                f'<a:solidFill><a:srgbClr val="000000"/></a:solidFill>' 
                f'</a:rPr><a:t>{safe_text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="{ns_a}">' 
            f'<a:pPr {pPr_attrs}><a:buNone/>{spcBef_xml}</a:pPr>' 
            f'{runs_xml}</a:p>'
        )
        txBody.append(etree.fromstring(p_xml))


def fill_t12_slide_1_cover(slide, topic, name_surname):
    """
    12-Shablon Slayd 1 — Muqova.
    Shape[0] 'TextBox 9': sarlavha — sz=6000, b=0, algn=ctr, spcPts=17205
    Shape[1] 'TextBox 11': ism-familiya — sz=2251, b=1, algn=ctr
    """
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    # Sarlavha (Shape[0])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        title_text = (topic or "TAQDIMOT").upper()
        _t12_clear_and_write(slide.shapes[0].text_frame._txBody, ns_a, [
            {'algn': 'ctr', 'spcPts': 17205, 'runs': [{'sz': 6000, 'b': 0, 'text': title_text}]}
        ])
    # Ism-familiya (Shape[1])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        name_text = name_surname or ""
        _t12_clear_and_write(slide.shapes[1].text_frame._txBody, ns_a, [
            {'algn': 'ctr', 'runs': [{'sz': 2251, 'b': 1, 'text': name_text}]}
        ])


def fill_t12_slide_2_plan(slide, plan_dict):
    """
    12-Shablon Slayd 2 — Reja.
    Shape[0] 'TextBox 7': 'REJA' sarlavhasi — sz=8135, b=0, algn=ctr (o'zgarmaydi)
    Shape[1] 'TextBox 8': reja ro'yxati — sz=2771, b=1, algn=l, marL=514350, indent=-514350
    """
    import re
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    # Shape[1] — reja punktlari
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        items = []
        # plan_dict = {"title": "Reja", "content": ["1. ...", "2. ..."]}
        if isinstance(plan_dict, dict):
            content = plan_dict.get("content", [])
            if isinstance(content, list) and content:
                items = [str(x) for x in content]
            elif isinstance(content, str) and content:
                items = [content]
            else:
                # content yo'q bo'lsa, boshqa keylardan olish
                for k, v in plan_dict.items():
                    if k == "title":
                        continue
                    if isinstance(v, list):
                        items.extend([str(x) for x in v])
                    elif isinstance(v, str) and v:
                        items.append(v)
        elif isinstance(plan_dict, list):
            items = [str(x) for x in plan_dict]
        # Raqamlarni tozalash - faqat boshidagi raqamlarni olib tashlash
        clean_items = []
        for item in items:
            cleaned = re.sub(r'^\s*[\d]+[\.\'\:\)]+\s*', '', str(item)).strip()
            if cleaned:
                clean_items.append(cleaned)
        paragraphs = []
        for idx, item in enumerate(clean_items):
            paragraphs.append({
                'algn': 'l',
                'marL': 514350,
                'indent': -514350,
                'runs': [{'sz': 2771, 'b': 1, 'text': f'{idx+1}. {item}'}]
            })
        if paragraphs:
            _t12_clear_and_write(slide.shapes[1].text_frame._txBody, ns_a, paragraphs)


def fill_t12_slide_3_text_left(slide, content_data, image_query=None):
    """
    12-Shablon Slayd 3 — Sarlavha + Matn chap, Dekor o'ng.
    Shape[0] 'Freeform 6': dekor (o'zgarmaydi)
    Shape[1] 'TextBox 7': asosiy matn — sz=2000, b=1, algn=l
    Shape[2] 'TextBox 8': sarlavha — sz=4800, b=0, algn=l
    """
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    if not items:
        col1 = content_data.get("col1", "")
        col2 = content_data.get("col2", "")
        items = [x for x in [col1, col2] if x]
    # Sarlavha (Shape[2]) — sz=4800, b=0, algn=l
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t12_clear_and_write(slide.shapes[2].text_frame._txBody, ns_a, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'text': title.upper()}]}
        ])
    # Asosiy matn (Shape[1]) — sz=2000, b=1, algn=l
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paragraphs = []
        for item in items:
            paragraphs.append({'algn': 'l', 'runs': [{'sz': 2000, 'b': 1, 'text': str(item)}]})
        if paragraphs:
            _t12_clear_and_write(slide.shapes[1].text_frame._txBody, ns_a, paragraphs)


def fill_t12_slide_4_image_left(slide, content_data, image_query=None):
    """
    12-Shablon Slayd 4 — Rasm chap (freeform), Sarlavha + Matn o'ng.
    Shape[0] 'TextBox 4': asosiy matn — sz=2400, b=1, algn=l
    Shape[1] 'TextBox 7': sarlavha — sz=4800, b=0, algn=l
    Shape[2] 'Freeform 7': rasm joyi (freeform o'chirib rasm qo'yiladi)
    """
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    if not items:
        col1 = content_data.get("col1", "")
        col2 = content_data.get("col2", "")
        items = [x for x in [col1, col2] if x]
    # Sarlavha (Shape[1]) — sz=4800, b=0, algn=l
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t12_clear_and_write(slide.shapes[1].text_frame._txBody, ns_a, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'text': title.upper()}]}
        ])
    # Asosiy matn (Shape[0]) — sz=2400, b=1, algn=l
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        paragraphs = []
        for item in items:
            paragraphs.append({'algn': 'l', 'runs': [{'sz': 2400, 'b': 1, 'text': str(item)}]})
        if paragraphs:
            _t12_clear_and_write(slide.shapes[0].text_frame._txBody, ns_a, paragraphs)
    # Rasm (Shape[2] — Freeform 7 o'chirib, rasm qo'yish)
    if image_query and len(slide.shapes) > 2:
        try:
            img_path = None
            if isinstance(image_query, str) and image_query.endswith(('.jpg', '.jpeg', '.png', '.webp')):
                # Fayl yo'li berilgan
                img_path = image_query
            elif isinstance(image_query, str):
                # Qidiruv so'zi berilgan - fetch_image fayl yo'li qaytaradi
                img_path = fetch_image(image_query)
            if img_path and os.path.exists(img_path):
                freeform_shape = slide.shapes[2]
                left = freeform_shape.left
                top = freeform_shape.top
                width = freeform_shape.width
                height = freeform_shape.height
                sp = freeform_shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(img_path, left, top, width, height)
        except Exception as e:
            logging.warning(f"[T12] Slayd 4 rasm qo'shishda xatolik: {e}")


def fill_t12_slide_5_two_columns(slide, content_data, image_query=None):
    """
    12-Shablon Slayd 5 — Sarlavha + 2 ta ustun matni.
    Shape[0] 'TextBox 4': sarlavha — sz=7200, b=0, algn=ctr
    Shape[1] 'TextBox 9': chap ustun — sz=2567, b=1, algn=ctr
    Shape[2] 'TextBox 10': o'ng ustun — sz=2567, b=1, algn=ctr
    """
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    col1 = content_data.get("col1", "")
    col2 = content_data.get("col2", "")
    if not col1 and not col2:
        items = content_data.get("content", [])
        if isinstance(items, list) and len(items) >= 2:
            col1, col2 = items[0], items[1]
        elif isinstance(items, list) and len(items) == 1:
            col1 = items[0]
        elif isinstance(items, str):
            col1 = items
    # Sarlavha (Shape[0]) — sz=7200, b=0, algn=ctr
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t12_clear_and_write(slide.shapes[0].text_frame._txBody, ns_a, [
            {'algn': 'ctr', 'runs': [{'sz': 7200, 'b': 0, 'text': title.upper()}]}
        ])
    # Chap ustun (Shape[1]) — sz=2567, b=1, algn=ctr
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t12_clear_and_write(slide.shapes[1].text_frame._txBody, ns_a, [
            {'algn': 'ctr', 'runs': [{'sz': 2567, 'b': 1, 'text': str(col1)}]}
        ])
    # O'ng ustun (Shape[2]) — sz=2567, b=1, algn=ctr
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t12_clear_and_write(slide.shapes[2].text_frame._txBody, ns_a, [
            {'algn': 'ctr', 'runs': [{'sz': 2567, 'b': 1, 'text': str(col2)}]}
        ])


def fill_t12_slide_6_three_columns(slide, content_data, image_query=None):
    """
    12-Shablon Slayd 6 — Sarlavha + 3 ta ustun matni.
    Shape[0] 'TextBox 3': sarlavha — sz=7200, b=0, algn=ctr
    Shape[1] 'TextBox 7': chap ustun — sz=2899, b=1, algn=ctr
    Shape[2] 'TextBox 8': markaz ustun — sz=2899, b=1, algn=ctr
    Shape[3] 'TextBox 9': o'ng ustun — sz=2899, b=1, algn=ctr
    """
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    col1 = content_data.get("col1", "")
    col2 = content_data.get("col2", "")
    col3 = content_data.get("col3", "")
    if not col1 and not col2 and not col3:
        items = content_data.get("content", [])
        if isinstance(items, list):
            if len(items) >= 3:
                col1, col2, col3 = items[0], items[1], items[2]
            elif len(items) == 2:
                col1, col2 = items[0], items[1]
            elif len(items) == 1:
                col1 = items[0]
        elif isinstance(items, str):
            col1 = items
    # Sarlavha (Shape[0]) — sz=7200, b=0, algn=ctr
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t12_clear_and_write(slide.shapes[0].text_frame._txBody, ns_a, [
            {'algn': 'ctr', 'runs': [{'sz': 7200, 'b': 0, 'text': title.upper()}]}
        ])
    # Chap ustun (Shape[1]) — sz=2899, b=1, algn=ctr
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t12_clear_and_write(slide.shapes[1].text_frame._txBody, ns_a, [
            {'algn': 'ctr', 'runs': [{'sz': 2899, 'b': 1, 'text': str(col1)}]}
        ])
    # Markaz ustun (Shape[2]) — sz=2899, b=1, algn=ctr
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t12_clear_and_write(slide.shapes[2].text_frame._txBody, ns_a, [
            {'algn': 'ctr', 'runs': [{'sz': 2899, 'b': 1, 'text': str(col2)}]}
        ])
    # O'ng ustun (Shape[3]) — sz=2899, b=1, algn=ctr
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t12_clear_and_write(slide.shapes[3].text_frame._txBody, ns_a, [
            {'algn': 'ctr', 'runs': [{'sz': 2899, 'b': 1, 'text': str(col3)}]}
        ])


def fill_t12_slide_7_image_right(slide, content_data, image_query=None):
    """
    12-Shablon Slayd 7 — Sarlavha + Matn chap, Rasm o'ng (freeform).
    Shape[0] 'TextBox 8': asosiy matn — sz=2899, b=1, algn=l
    Shape[1] 'TextBox 13': sarlavha — sz=7200, b=0, algn=l
    Shape[2] 'Freeform 10': rasm joyi (freeform o'chirib rasm qo'yiladi)
    """
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    if isinstance(items, str):
        items = [items]
    if not items:
        col1 = content_data.get("col1", "")
        col2 = content_data.get("col2", "")
        items = [x for x in [col1, col2] if x]
    # Sarlavha (Shape[1]) — sz=7200, b=0, algn=l
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t12_clear_and_write(slide.shapes[1].text_frame._txBody, ns_a, [
            {'algn': 'l', 'runs': [{'sz': 7200, 'b': 0, 'text': title.upper()}]}
        ])
    # Asosiy matn (Shape[0]) — sz=2899, b=1, algn=l
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        paragraphs = []
        for item in items:
            paragraphs.append({'algn': 'l', 'runs': [{'sz': 2899, 'b': 1, 'text': str(item)}]})
        if paragraphs:
            _t12_clear_and_write(slide.shapes[0].text_frame._txBody, ns_a, paragraphs)
    # Rasm (Shape[2] — Freeform 10 o'chirib, rasm qo'yish)
    if image_query and len(slide.shapes) > 2:
        try:
            img_path = None
            if isinstance(image_query, str) and image_query.endswith(('.jpg', '.jpeg', '.png', '.webp')):
                img_path = image_query
            elif isinstance(image_query, str):
                img_path = fetch_image(image_query)
            if img_path and os.path.exists(img_path):
                freeform_shape = slide.shapes[2]
                left = freeform_shape.left
                top = freeform_shape.top
                width = freeform_shape.width
                height = freeform_shape.height
                sp = freeform_shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(img_path, left, top, width, height)
        except Exception as e:
            logging.warning(f"[T12] Slayd 7 rasm qo'shishda xatolik: {e}")


def fill_t12_slide_8_conclusion(slide, content_data):
    """
    12-Shablon Slayd 8 — Xulosa.
    Shape[0] 'TextBox 9': xulosa matni — sz=8000, b=0, algn=ctr, spcPts=17205
    Shablon matni saqlanadi (o'zgarmaydi).
    """
    pass  # Xulosa slaydida shablon matni saqlanadi


def generate_template_12_presentation(prs, topic, requested_slide_count, language,
                                       name_surname, plan, content_data_list, user_images=None):
    """
    12-shablon asosida to'liq prezentatsiya yaratadi.
    """
    import io
    total_content_slides = build_slide_structure_12(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}
    # Slayd 1 — Muqova
    fill_t12_slide_1_cover(prs.slides[0], topic, name_surname)
    # Slayd 2 — Reja
    fill_t12_slide_2_plan(prs.slides[1], plan_dict)
    # Kontent slaydlari (3-dan boshlab)
    user_img_idx = 0
    # Rasm ishlatadigan slayd turlari: 1=slayd4(image_left), 4=slayd7(image_right)
    IMAGE_SLIDE_TYPES = [1, 4]
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        image_query = data.get("image_query", topic)
        slide_type = i % 5
        has_image = slide_type in IMAGE_SLIDE_TYPES
        if has_image:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                img_arg = img_path if img_path else image_query
            else:
                img_arg = image_query
        else:
            img_arg = None
        if slide_type == 0:
            fill_t12_slide_3_text_left(slide, data, img_arg)
        elif slide_type == 1:
            fill_t12_slide_4_image_left(slide, data, img_arg)
        elif slide_type == 2:
            fill_t12_slide_5_two_columns(slide, data, img_arg)
        elif slide_type == 3:
            fill_t12_slide_6_three_columns(slide, data, img_arg)
        elif slide_type == 4:
            fill_t12_slide_7_image_right(slide, data, img_arg)
        logging.info(f"  [T12] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    # Xulosa slayd
    conclusion_slide = prs.slides[-1]
    fill_t12_slide_8_conclusion(conclusion_slide, {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ============================================================
# 13-SHABLON (Pastel Botanical) FUNKSIYALARI
# ============================================================

SLIDE_TYPE_NAMES_T13 = {
    0: "one_column",       # Slayd 3: rasm chap (freeform), sarlavha + matn o'ng
    1: "image_right",      # Slayd 4: sarlavha chap yuqori, rasm o'ng, matn chap pastda
    2: "title_body",       # Slayd 5: sarlavha + bitta qisqa body matn
    3: "image_left_text",  # Slayd 6: rasm chap (freeform), sarlavha + matn o'ng
    4: "two_text_blocks",  # Slayd 7: sarlavha chap, matn chap, izoh matn o'ng, rasm o'ng
}
CONTENT_SLIDE_TEMPLATE_INDICES_T13 = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7


def build_slide_structure_13(prs, requested_content_count):
    n_templates = len(CONTENT_SLIDE_TEMPLATE_INDICES_T13)
    full_repeats = max(1, round(requested_content_count / n_templates))
    total_content_slides = full_repeats * n_templates
    logging.info(f"[T13] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    conclusion_current_index = 7
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_T13:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T13] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T13] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def _t13_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            color = run.get('color', '404040')
            text = run.get('text', '')
            safe_text = str(text).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{safe_text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="{ns_a}">'
            f'<a:pPr {pPr_attrs}><a:buNone/>{spcBef_xml}</a:pPr>'
            f'{runs_xml}</a:p>'
        )
        txBody.append(etree.fromstring(p_xml))


def _t13_replace_blip(slide, shape_idx, img_arg):
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        blip = el.find('.//a:blip', {'a': ns_a})
        if blip is not None:
            part = slide.part
            _, img_rId = part.get_or_add_image_part(img_path)
            blip.set(f'{{{ns_r}}}embed', img_rId)
    except Exception as e:
        logging.warning(f"[T13] Rasm almashtirish xatoligi: {e}")


def fill_t13_slide_1_cover(slide, topic, name_surname):
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t13_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6000, 'b': 0, 'color': '4a6741', 'text': topic or "TAQDIMOT"}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t13_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '404040', 'text': name_surname or ""}]}
        ])


def fill_t13_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        content = plan_dict.get("content", [])
        if isinstance(content, str):
            content = [content]
        if not content:
            content = plan_dict.get("items", [])
        if not content:
            content = [plan_dict.get("title", "")]
        clean_items = []
        for item in content:
            if isinstance(item, str):
                item_clean = re.sub(r'^\s*\d+[\.\)]\s*', '', item.strip())
                if item_clean:
                    clean_items.append(item_clean)
        if not clean_items:
            clean_items = ["Reja mavjud emas"]
        paragraphs = []
        for idx, item in enumerate(clean_items):
            paragraphs.append({
                'algn': 'l',
                'marL': 816609,
                'indent': -514350,
                'runs': [{'sz': 2800, 'b': 0, 'color': '404040', 'text': f"{idx + 1}.  {item}"}]
            })
        _t13_clear_and_write(slide.shapes[1].text_frame._txBody, paragraphs)


def _t13_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text


def fill_t13_slide_3_image_left_text_right(slide, data, img_arg=None):
    """Slayd 3: Shape[0]=Freeform(rasm), Shape[1]=sarlavha, Shape[2]=matn"""
    if not isinstance(data, dict):
        data = {}
    if img_arg and len(slide.shapes) > 0:
        _t13_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t13_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': '4a6741', 'text': data.get("title", "")}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t13_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2499, 'b': 0, 'color': '404040', 'text': _t13_get_body_text(data)}]}
        ])


def fill_t13_slide_4_image_right_text_left(slide, data, img_arg=None):
    """Slayd 4: Shape[0]=sarlavha, Shape[1]=Group(rasm), Shape[2]=matn"""
    if not isinstance(data, dict):
        data = {}
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t13_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4000, 'b': 0, 'color': '4a6741', 'text': data.get("title", "")}]}
        ])
    if img_arg and len(slide.shapes) > 1:
        _t13_replace_blip(slide, 1, img_arg)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t13_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2574, 'b': 0, 'color': '404040', 'text': _t13_get_body_text(data)}]}
        ])


def fill_t13_slide_5_title_body(slide, data, img_arg=None):
    """Slayd 5: Shape[0]=body matn, Shape[1]=sarlavha"""
    if not isinstance(data, dict):
        data = {}
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t13_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': '404040', 'text': _t13_get_body_text(data)}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t13_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': '4a6741', 'text': data.get("title", "")}]}
        ])


def fill_t13_slide_6_image_left_text_right(slide, data, img_arg=None):
    """Slayd 6: Shape[0]=sarlavha, Shape[1]=Freeform(rasm), Shape[2]=matn"""
    if not isinstance(data, dict):
        data = {}
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t13_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': '4a6741', 'text': data.get("title", "")}]}
        ])
    if img_arg and len(slide.shapes) > 1:
        _t13_replace_blip(slide, 1, img_arg)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t13_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2275, 'b': 0, 'color': '404040', 'text': _t13_get_body_text(data)}]}
        ])


def fill_t13_slide_7_two_text_image(slide, data, img_arg=None):
    """Slayd 7: Shape[0]=Picture(rasm), Shape[1]=sarlavha, Shape[2]=matn chap, Shape[3]=izoh o'ng"""
    if not isinstance(data, dict):
        data = {}
    content = data.get("content", [])
    if isinstance(content, list) and len(content) >= 2:
        main_text = str(content[0]) if content[0] else ""
        caption_text = str(content[1]) if content[1] else ""
    elif isinstance(content, list) and len(content) == 1:
        main_text = str(content[0])
        caption_text = ""
    elif isinstance(content, str):
        main_text = content
        caption_text = ""
    else:
        main_text = data.get("col1", "") or data.get("text", "")
        caption_text = data.get("col2", "")
    if img_arg and len(slide.shapes) > 0:
        _t13_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t13_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': '4a6741', 'text': data.get("title", "")}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t13_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2650, 'b': 0, 'color': '404040', 'text': main_text}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t13_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1850, 'b': 0, 'color': '404040', 'text': caption_text}]}
        ])


def fill_t13_slide_8_conclusion(slide, data):
    pass  # Xulosa slayd shablondagi holicha qoladi


def generate_template_13_presentation(prs, topic, requested_slide_count, language,
                                       name_surname, plan, content_data_list, user_images=None):
    import io
    total_content_slides = build_slide_structure_13(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}
    fill_t13_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_t13_slide_2_plan(prs.slides[1], plan_dict)
    user_img_idx = 0
    IMAGE_SLIDE_TYPES = [0, 1, 3, 4]
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        image_query = data.get("image_query", topic)
        slide_type = i % 5
        has_image = slide_type in IMAGE_SLIDE_TYPES
        if has_image:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                img_arg = img_path if img_path else image_query
            else:
                img_arg = image_query
        else:
            img_arg = None
        if slide_type == 0:
            fill_t13_slide_3_image_left_text_right(slide, data, img_arg)
        elif slide_type == 1:
            fill_t13_slide_4_image_right_text_left(slide, data, img_arg)
        elif slide_type == 2:
            fill_t13_slide_5_title_body(slide, data, img_arg)
        elif slide_type == 3:
            fill_t13_slide_6_image_left_text_right(slide, data, img_arg)
        elif slide_type == 4:
            fill_t13_slide_7_two_text_image(slide, data, img_arg)
        logging.info(f"  [T13] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t13_slide_8_conclusion(prs.slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ==================== 14-SHABLON (Thesis Defense) ====================

SLIDE_TYPE_NAMES_T14 = {
    0: "image_right",      # Slayd 3: sarlavha + matn chap, freeform rasm o'ng
    1: "image_left",       # Slayd 4: freeform rasm chap, sarlavha + matn o'ng
    2: "image_left_list",  # Slayd 5: sarlavha tepada, freeform rasm chap, ro'yxat matn o'ng
    3: "title_body",       # Slayd 6: sarlavha + matn chap, rasm/diagramma o'ng (rasm qo'yiladi)
    4: "image_left",       # Slayd 7: freeform rasm chap, sarlavha + matn o'ng
}
CONTENT_SLIDE_TEMPLATE_INDICES_T14 = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7

def build_slide_structure_14(prs, requested_content_count):
    n_templates = len(CONTENT_SLIDE_TEMPLATE_INDICES_T14)
    full_repeats = max(1, round(requested_content_count / n_templates))
    total_content_slides = full_repeats * n_templates
    logging.info(f"[T14] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    conclusion_current_index = 7
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_T14:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T14] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T14] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides

def _t14_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            color = run.get('color', '000000')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t14_replace_blip(slide, shape_idx, img_arg):
    """Freeform yoki Picture ichidagi blip ni yangi rasm bilan almashtirish"""
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        blip = el.find('.//a:blip', {'a': ns_a})
        if blip is not None:
            part = slide.part
            _, img_rId = part.get_or_add_image_part(img_path)
            blip.set(f'{{{ns_r}}}embed', img_rId)
    except Exception as e:
        logging.warning(f"[T14] Rasm almashtirish xatoligi: {e}")

def _t14_replace_picture(slide, shape_idx, img_arg):
    """Oddiy Picture shape ni yangi rasm bilan almashtirish"""
    try:
        import os
        from pptx.util import Emu
        shape = slide.shapes[shape_idx]
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        sp = shape._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(img_path, left, top, width, height)
    except Exception as e:
        logging.warning(f"[T14] Picture almashtirish xatoligi: {e}")

def _t14_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t14_slide_1_cover(slide, topic, name_surname):
    """Slayd 1: Shape[0]=katta sarlavha, Shape[1]=ism-familiya"""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t14_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6000, 'b': 1, 'color': '000000', 'text': (topic or "TAQDIMOT").upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t14_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3872, 'b': 0, 'color': '000000', 'text': name_surname or ''}]}
        ])

def fill_t14_slide_2_plan(slide, plan_dict):
    """Slayd 2: Shape[0]=sarlavha 'Reja', Shape[1]=reja matn"""
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    # Sarlavha
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t14_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 8000, 'b': 0, 'color': 'ffffff', 'text': 'Reja'}]}
        ])
    # Reja matn
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        content = plan_dict.get("content", [])
        if isinstance(content, str):
            content = [content]
        if not content:
            content = plan_dict.get("items", [])
        if not content:
            content = [plan_dict.get("title", "")]
        clean_items = []
        for item in content:
            if isinstance(item, str):
                item_clean = re.sub(r'^\s*\d+[\.\)]\s*', '', item.strip())
                if item_clean:
                    clean_items.append(item_clean)
        if not clean_items:
            clean_items = ["Reja mavjud emas"]
        paragraphs = []
        for idx, item in enumerate(clean_items):
            paragraphs.append({
                'algn': 'just',
                'marL': 816609,
                'indent': -514350,
                'runs': [{'sz': 2500, 'b': 0, 'color': '000000', 'text': f"{idx + 1}.  {item}"}]
            })
        _t14_clear_and_write(slide.shapes[1].text_frame._txBody, paragraphs)

def fill_t14_slide_3_text_left_image_right(slide, data, img_arg=None):
    """Slayd 3: Shape[0]=Freeform(rasm o'ng), Shape[1]=sarlavha, Shape[2]=matn"""
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t14_get_body_text(data)
    # Rasm (Freeform o'ng tomonda)
    if img_arg and len(slide.shapes) > 0:
        _t14_replace_blip(slide, 0, img_arg)
    # Sarlavha
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t14_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 0, 'color': '000000', 'text': title}]}
        ])
    # Matn
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t14_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 3500, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t14_slide_4_image_left_text_right(slide, data, img_arg=None):
    """Slayd 4: Shape[0]=Freeform(rasm chap), Shape[1]=matn, Shape[2]=sarlavha"""
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t14_get_body_text(data)
    # Rasm (Freeform chap tomonda)
    if img_arg and len(slide.shapes) > 0:
        _t14_replace_blip(slide, 0, img_arg)
    # Matn (o'ng tomonda)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t14_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 3500, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])
    # Sarlavha (o'ng yuqori)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t14_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 0, 'color': '1a3a5c', 'text': title}]}
        ])

def fill_t14_slide_5_image_left_list_right(slide, data, img_arg=None):
    """Slayd 5: Shape[0]=matn o'ng, Shape[1]=Freeform(rasm chap), Shape[2]=sarlavha tepada"""
    import re
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    content = data.get("content", [])
    if isinstance(content, str):
        content = [content]
    if not content:
        content = [_t14_get_body_text(data)]
    # Rasm (Freeform chap tomonda)
    if img_arg and len(slide.shapes) > 1:
        _t14_replace_blip(slide, 1, img_arg)
    # Sarlavha (tepada)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t14_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 0, 'color': '000000', 'text': title}]}
        ])
    # Matn (o'ng tomonda) - raqamlangan
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        clean_items = []
        for item in content:
            if isinstance(item, str):
                item_clean = re.sub(r'^\s*\d+[\.\)]\s*', '', item.strip())
                if item_clean:
                    clean_items.append(item_clean)
        if not clean_items:
            clean_items = [_t14_get_body_text(data)]
        paragraphs = []
        for idx, item in enumerate(clean_items):
            paragraphs.append({
                'algn': 'just',
                'runs': [{'sz': 3000, 'b': 0, 'color': '000000', 'text': f"{idx + 1}.{item}"}]
            })
        _t14_clear_and_write(slide.shapes[0].text_frame._txBody, paragraphs)

def fill_t14_slide_6_text_left_image_right(slide, data, img_arg=None):
    """Slayd 6: Shape[0]=Picture(rasm o'ng), Shape[1]=sarlavha, Shape[2]=matn"""
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t14_get_body_text(data)
    # Rasm (Picture o'ng tomonda)
    if img_arg and len(slide.shapes) > 0:
        _t14_replace_picture(slide, 0, img_arg)
    # Sarlavha
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t14_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 0, 'color': '1a3a5c', 'text': title}]}
        ])
    # Matn
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t14_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 3000, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t14_slide_7_image_left_text_right(slide, data, img_arg=None):
    """Slayd 7: Shape[0]=Freeform(rasm chap), Shape[1]=matn, Shape[2]=sarlavha"""
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t14_get_body_text(data)
    # Rasm (Freeform chap tomonda)
    if img_arg and len(slide.shapes) > 0:
        _t14_replace_blip(slide, 0, img_arg)
    # Matn (o'ng tomonda)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t14_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 3000, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])
    # Sarlavha (tepada)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t14_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 0, 'color': '1a3a5c', 'text': title}]}
        ])

def fill_t14_slide_8_conclusion(slide, data):
    """Slayd 8: Yakuniy slayd - shablondagi matn saqlanadi"""
    pass

def generate_template_14_presentation(prs, topic, requested_slide_count, language,
                                       name_surname, plan, content_data_list, user_images=None):
    import io
    total_content_slides = build_slide_structure_14(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}
    fill_t14_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_t14_slide_2_plan(prs.slides[1], plan_dict)
    user_img_idx = 0
    IMAGE_SLIDE_TYPES = [0, 1, 2, 3, 4]  # Barcha slaydlarda rasm bor
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        image_query = data.get("image_query", topic)
        slide_type = i % 5
        has_image = slide_type in IMAGE_SLIDE_TYPES
        if has_image:
            if user_images and user_img_idx < len(user_images):
                img_path = save_user_image_to_tmp(user_images[user_img_idx])
                user_img_idx += 1
                img_arg = img_path if img_path else image_query
            else:
                img_arg = image_query
        else:
            img_arg = None
        if slide_type == 0:
            fill_t14_slide_3_text_left_image_right(slide, data, img_arg)
        elif slide_type == 1:
            fill_t14_slide_4_image_left_text_right(slide, data, img_arg)
        elif slide_type == 2:
            fill_t14_slide_5_image_left_list_right(slide, data, img_arg)
        elif slide_type == 3:
            fill_t14_slide_6_text_left_image_right(slide, data, img_arg)
        elif slide_type == 4:
            fill_t14_slide_7_image_left_text_right(slide, data, img_arg)
        logging.info(f"  [T14] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t14_slide_8_conclusion(prs.slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ============================================================
# 15-SHABLON (Green Forest / Nature)
# ============================================================
SLIDE_TYPE_NAMES_T15 = {
    0: "image_left_text_right",
    1: "text_left_image_right",
    2: "title_text_image_right",
    3: "title_text_image_right2",
    4: "image_left_text_right2",
}
CONTENT_SLIDE_TEMPLATE_INDICES_T15 = [2, 3, 4, 5, 6]

def build_slide_structure_15(prs, requested_content_count):
    n_templates = len(CONTENT_SLIDE_TEMPLATE_INDICES_T15)
    full_repeats = max(1, round(requested_content_count / n_templates))
    total_content_slides = full_repeats * n_templates
    logging.info(f"[T15] Kontent slaydlari: {requested_content_count} so'raldi, {full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    conclusion_current_index = 7
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_T15:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T15] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T15] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides

def _t15_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            color = run.get('color', '000000')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t15_replace_blip(slide, shape_idx, img_arg):
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        blip = el.find('.//a:blip', {'a': ns_a})
        if blip is not None:
            part = slide.part
            _, img_rId = part.get_or_add_image_part(img_path)
            blip.set(f'{{{ns_r}}}embed', img_rId)
    except Exception as e:
        logging.warning(f"[T15] Rasm almashtirish xatoligi: {e}")

def _t15_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t15_slide_1_cover(slide, topic, name_surname):
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t15_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 1, 'color': 'FFFFFF', 'text': topic}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t15_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 1, 'color': '000000', 'text': name_surname}]}
        ])

def fill_t15_slide_2_plan(slide, plan_dict):
    import re
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t15_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 8741, 'b': 1, 'color': 'FFFFFF', 'text': 'Reja'}]}
        ])
    items = []
    if isinstance(plan_dict, dict):
        raw = plan_dict.get("content", plan_dict.get("items", []))
        if isinstance(raw, list):
            items = raw
        elif isinstance(raw, str):
            items = [raw]
    if not items:
        items = ["Kirish", "Asosiy qism", "Xulosa"]
    paragraphs = []
    for idx, item in enumerate(items):
        clean = re.sub(r'^\d+[\.\)]\s*', '', str(item)).strip()
        paragraphs.append({
            'algn': 'l',
            'runs': [{'sz': 3047, 'b': 1, 'color': 'FFFFFF', 'text': f"{idx+1}. {clean}"}]
        })
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t15_clear_and_write(slide.shapes[1].text_frame._txBody, paragraphs)

def fill_t15_slide_3_image_left_text_right(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t15_get_body_text(data)
    if img_arg:
        _t15_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t15_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 5400, 'b': 1, 'color': '749835', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t15_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1813, 'b': 0, 'color': '14190D', 'text': body_text}]}
        ])

def fill_t15_slide_4_text_left_image_right(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t15_get_body_text(data)
    if img_arg:
        _t15_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t15_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': '749835', 'text': title}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t15_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1813, 'b': 0, 'color': '14190D', 'text': body_text}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t15_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1813, 'b': 0, 'color': '14190D', 'text': ''}]}
        ])

def fill_t15_slide_5_title_text_image_right(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t15_get_body_text(data)
    if img_arg:
        _t15_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t15_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': 'FFFFFF', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t15_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2924, 'b': 0, 'color': 'FFFFFF', 'text': body_text}]}
        ])

def fill_t15_slide_6_title_text_image_right2(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t15_get_body_text(data)
    if img_arg:
        _t15_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t15_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 6000, 'b': 1, 'color': '749835', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t15_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2924, 'b': 0, 'color': 'FFFFFF', 'text': body_text}]}
        ])

def fill_t15_slide_7_image_left_text_right2(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t15_get_body_text(data)
    if img_arg:
        _t15_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t15_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 1, 'color': '749835', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t15_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2132, 'b': 0, 'color': '14190D', 'text': body_text}]}
        ])

def fill_t15_slide_8_conclusion(slide, data):
    pass

def generate_template_15_presentation(prs, topic, requested_slide_count, language,
                                       name_surname, plan, content_data_list, user_images=None):
    import io
    total_content_slides = build_slide_structure_15(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}
    fill_t15_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_t15_slide_2_plan(prs.slides[1], plan_dict)
    user_img_idx = 0
    IMAGE_SLIDE_TYPES = [0, 1, 2, 3, 4]
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        image_query = data.get("image_query", topic)
        slide_type = i % 5
        if user_images and user_img_idx < len(user_images):
            img_path = save_user_image_to_tmp(user_images[user_img_idx])
            user_img_idx += 1
            img_arg = img_path if img_path else image_query
        else:
            img_arg = image_query
        if slide_type == 0:
            fill_t15_slide_3_image_left_text_right(slide, data, img_arg)
        elif slide_type == 1:
            fill_t15_slide_4_text_left_image_right(slide, data, img_arg)
        elif slide_type == 2:
            fill_t15_slide_5_title_text_image_right(slide, data, img_arg)
        elif slide_type == 3:
            fill_t15_slide_6_title_text_image_right2(slide, data, img_arg)
        elif slide_type == 4:
            fill_t15_slide_7_image_left_text_right2(slide, data, img_arg)
        logging.info(f"  [T15] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t15_slide_8_conclusion(prs.slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ============================================================
# 16-SHABLON (Medical Blue) funksiyalari
# ============================================================

SLIDE_TYPE_NAMES_T16 = {
    "cover": "Muqova",
    "plan": "Reja",
    "text_left_image_right": "Matn chap, rasm o'ng",
    "image_left_two_texts": "Sarlavha, matn chap, rasm o'ng",
    "image_left_text_right_two": "Rasm chap, ikki matn o'ng",
    "text_left_image_right_quote": "Matn chap, rasm o'ng, quote",
    "image_left_colored_text": "Rasm chap, rangli matn o'ng",
    "conclusion": "Xulosa",
}

def build_slide_structure_16(prs, requested_content_count):
    content_slide_types = [
        "text_left_image_right",
        "image_left_two_texts",
        "image_left_text_right_two",
        "text_left_image_right_quote",
        "image_left_colored_text",
    ]
    structure = ["cover", "plan"]
    for i in range(requested_content_count):
        structure.append(content_slide_types[i % len(content_slide_types)])
    structure.append("conclusion")
    return structure

def _t16_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 1500)
            b = run.get('b', 0)
            color = run.get('color', '19437A')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t16_replace_blip(slide, shape_idx, img_arg):
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        blip = el.find('.//a:blip', {'a': ns_a})
        if blip is not None:
            part = slide.part
            _, img_rId = part.get_or_add_image_part(img_path)
            blip.set(f'{{{ns_r}}}embed', img_rId)
        else:
            for child in el.iter():
                if child.tag == f'{{{ns_a}}}blip':
                    part = slide.part
                    _, img_rId = part.get_or_add_image_part(img_path)
                    child.set(f'{{{ns_r}}}embed', img_rId)
                    break
    except Exception as e:
        logging.warning(f"[T16] Rasm almashtirish xatoligi: {e}")

def _t16_replace_group_blip(slide, group_shape_idx, img_arg):
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[group_shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        for child in el.iter():
            if child.tag == f'{{{ns_a}}}blip':
                part = slide.part
                _, img_rId = part.get_or_add_image_part(img_path)
                child.set(f'{{{ns_r}}}embed', img_rId)
                break
    except Exception as e:
        logging.warning(f"[T16] Group rasm almashtirish xatoligi: {e}")

def _t16_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t16_slide_1_cover(slide, topic, name_surname):
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t16_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 7200, 'b': 0, 'color': '19437A', 'text': topic}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t16_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '19437A', 'text': name_surname}]}
        ])

def fill_t16_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t16_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 5499, 'b': 0, 'color': '19437A', 'text': plan_title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'marL': 342900,
                'indent': -342900,
                'spcPts': 150,
                'runs': [{'sz': 3200, 'b': 0, 'color': '19437A', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t16_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t16_slide_3_text_left_image_right(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t16_get_body_text(data)
    if img_arg:
        _t16_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t16_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 5499, 'b': 0, 'color': '19437A', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t16_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '19437A', 'text': body_text}]}
        ])

def fill_t16_slide_4_image_left_two_texts(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t16_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t16_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 5000, 'b': 0, 'color': '19437A', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t16_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '19437A', 'text': body_text}]}
        ])
    if img_arg:
        _t16_replace_group_blip(slide, 2, img_arg)
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t16_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '19437A', 'text': body_text}]}
        ])

def fill_t16_slide_5_image_left_text_right_two(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t16_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if img_arg:
        _t16_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t16_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 5499, 'b': 0, 'color': '19437A', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t16_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '19437A', 'text': text1}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t16_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '19437A', 'text': text2}]}
        ])

def fill_t16_slide_6_text_left_image_right_quote(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t16_get_body_text(data)
    if img_arg:
        _t16_replace_group_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t16_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '19437A', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t16_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '19437A', 'text': body_text}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        words = body_text.split()
        quote_text = " ".join(words[-15:]) if len(words) > 15 else body_text
        _t16_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 2400, 'b': 1, 'color': '19437A', 'text': quote_text}]}
        ])

def fill_t16_slide_7_image_left_colored_text(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t16_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if img_arg:
        _t16_replace_group_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t16_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '19437A', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t16_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '19437A', 'text': text2}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t16_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': 'FFFFFF', 'text': text1}]}
        ])

def fill_t16_slide_8_conclusion(slide, data):
    pass

def generate_template_16_presentation(prs, topic, requested_slide_count, language,
                                       name_surname, plan=None, content_data_list=None,
                                       user_images=None):
    import io
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T16] Shablon slaydlari yetarli emas")
        return None
    fill_t16_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t16_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t16_slide_3_text_left_image_right,
        fill_t16_slide_4_image_left_two_texts,
        fill_t16_slide_5_image_left_text_right_two,
        fill_t16_slide_6_text_left_image_right_quote,
        fill_t16_slide_7_image_left_colored_text,
    ]
    user_img_idx = 0
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        image_query = data.get("image_query", topic)
        slide_type = i % len(content_slide_funcs)
        if user_images and user_img_idx < len(user_images):
            img_path = save_user_image_to_tmp(user_images[user_img_idx])
            user_img_idx += 1
            img_arg = img_path if img_path else image_query
        else:
            img_arg = image_query
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T16] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t16_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# 17-SHABLON funksiyalari
# ============================================================

SLIDE_TYPE_NAMES_T17 = {
    "cover": "Muqova",
    "plan": "Reja",
    "text_left_image_right": "Matn chap, rasm o'ng",
    "image_left_two_texts": "Sarlavha, rasm chap, ikki matn o'ng",
    "image_left_text_right_two": "Sarlavha, rasm chap, matn o'ng",
    "text_left_image_right_quote": "Sarlavha markazda, matn chap, rasm o'ng",
    "image_left_colored_text": "Sarlavha chap, rasm chap, ikki matn o'ng",
    "conclusion": "Xulosa",
}

def build_slide_structure_17(prs, requested_content_count):
    content_slide_types = [
        "text_left_image_right",
        "image_left_two_texts",
        "image_left_text_right_two",
        "text_left_image_right_quote",
        "image_left_colored_text",
    ]
    structure = ["cover", "plan"]
    for i in range(requested_content_count):
        structure.append(content_slide_types[i % len(content_slide_types)])
    structure.append("conclusion")
    return structure

def _t17_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2200)
            b = run.get('b', 0)
            color = run.get('color', '222223')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t17_replace_blip(slide, shape_idx, img_arg):
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        blip = el.find('.//a:blip', {'a': ns_a})
        if blip is not None:
            part = slide.part
            _, img_rId = part.get_or_add_image_part(img_path)
            blip.set(f'{{{ns_r}}}embed', img_rId)
        else:
            for child in el.iter():
                if child.tag == f'{{{ns_a}}}blip':
                    part = slide.part
                    _, img_rId = part.get_or_add_image_part(img_path)
                    child.set(f'{{{ns_r}}}embed', img_rId)
                    break
    except Exception as e:
        logging.warning(f"[T17] Rasm almashtirish xatoligi: {e}")

def _t17_replace_group_blip(slide, group_shape_idx, img_arg):
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[group_shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        for child in el.iter():
            if child.tag == f'{{{ns_a}}}blip':
                part = slide.part
                _, img_rId = part.get_or_add_image_part(img_path)
                child.set(f'{{{ns_r}}}embed', img_rId)
                break
    except Exception as e:
        logging.warning(f"[T17] Group rasm almashtirish xatoligi: {e}")

def _t17_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t17_slide_1_cover(slide, topic, name_surname):
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t17_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 5400, 'b': 0, 'color': '23374D', 'text': topic}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t17_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '30466B', 'text': name_surname}]}
        ])

def fill_t17_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t17_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 10650, 'b': 0, 'color': '23374D', 'text': plan_title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'marL': 342900,
                'indent': -342900,
                'spcPts': 150,
                'runs': [{'sz': 2800, 'b': 0, 'color': '23374D', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t17_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t17_slide_3_text_left_image_right(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t17_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if img_arg:
        _t17_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t17_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '222223', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t17_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '222223', 'text': text2}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t17_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '23374D', 'text': title}]}
        ])

def fill_t17_slide_4_image_left_two_texts(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t17_get_body_text(data)
    # Matn yetarli bo'lishi uchun 3 marta takrorlash
    if len(body_text.split()) < 40:
        body_text = (body_text + " " + body_text + " " + body_text).strip()
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if img_arg:
        # Ikkala rasm uchun har xil query ishlatish
        img_arg2 = data.get("image_query2", "") or (img_arg + " medical treatment" if isinstance(img_arg, str) and not os.path.exists(img_arg) else img_arg)
        _t17_replace_group_blip(slide, 0, img_arg)
        _t17_replace_group_blip(slide, 1, img_arg2)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t17_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '222223', 'text': text1}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t17_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '222223', 'text': text2}]}
        ])
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        _t17_clear_and_write(slide.shapes[4].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '23374D', 'text': title}]}
        ])

def fill_t17_slide_5_image_left_text_right_two(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t17_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t17_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '23374D', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t17_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '222223', 'text': body_text}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t17_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '222223', 'text': body_text}]}
        ])

def fill_t17_slide_6_text_left_image_right_quote(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t17_get_body_text(data)
    # Matn yetarli bo'lishi uchun 3 marta takrorlash
    if len(body_text.split()) < 50:
        body_text = (body_text + " " + body_text + " " + body_text).strip()
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t17_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 5400, 'b': 0, 'color': '23374D', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t17_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '222223', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t17_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '222223', 'text': text2}]}
        ])

def fill_t17_slide_7_image_left_colored_text(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t17_get_body_text(data)
    if img_arg:
        # Ikkala rasm uchun har xil query ishlatish
        img_arg2 = data.get("image_query2", "") or (img_arg + " healthcare facility" if isinstance(img_arg, str) and not os.path.exists(img_arg) else img_arg)
        _t17_replace_blip(slide, 0, img_arg)
        _t17_replace_group_blip(slide, 1, img_arg2)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t17_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '23374D', 'text': body_text}]}
        ])
    # 2-matn bloki (shapes[3]) - sarlavha emas, oddiy matn bilan to'ldirish
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        words = body_text.split()
        extra_text = " ".join(words[len(words)//2:]) if len(words) > 10 else body_text
        _t17_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '222223', 'text': extra_text}]}
        ])
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        _t17_clear_and_write(slide.shapes[4].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 6000, 'b': 0, 'color': '23374D', 'text': title}]}
        ])

def fill_t17_slide_8_conclusion(slide, data):
    pass

def generate_template_17_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T17] Shablon slaydlari yetarli emas")
        return None
    fill_t17_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t17_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t17_slide_3_text_left_image_right,
        fill_t17_slide_4_image_left_two_texts,
        fill_t17_slide_5_image_left_text_right_two,
        fill_t17_slide_6_text_left_image_right_quote,
        fill_t17_slide_7_image_left_colored_text,
    ]
    user_img_idx = 0
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        image_query = data.get("image_query", topic)
        slide_type = i % len(content_slide_funcs)
        if user_images and user_img_idx < len(user_images):
            img_path = save_user_image_to_tmp(user_images[user_img_idx])
            user_img_idx += 1
            img_arg = img_path if img_path else image_query
        else:
            img_arg = image_query
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T17] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t17_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
# 18-SHABLON funksiyalari
# ============================================================

SLIDE_TYPE_NAMES_T18 = {
    "cover": "Muqova",
    "plan": "Reja",
    "image_left_two_texts": "Rasm chap, ikki matn o'ng",
    "text_center_two_texts": "Sarlavha markazda, ikki matn",
    "image_right_text_left": "Sarlavha va matn chap, rasm o'ng",
    "text_center_two_texts_shape": "Sarlavha markazda, ikki matn, shakl chapda",
    "image_right_text_left_colored": "Sarlavha va matn chap, rasm o'ng (to'q fon)",
    "conclusion": "Xulosa",
}

def build_slide_structure_18(prs, requested_content_count):
    content_slide_types = [
        "image_left_two_texts",
        "text_center_two_texts",
        "image_right_text_left",
        "text_center_two_texts_shape",
        "image_right_text_left_colored",
    ]
    structure = ["cover", "plan"]
    for i in range(requested_content_count):
        structure.append(content_slide_types[i % len(content_slide_types)])
    structure.append("conclusion")
    return structure

def _t18_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            color = run.get('color', '101010')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t18_replace_blip(slide, shape_idx, img_arg):
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        blip = el.find('.//a:blip', {'a': ns_a})
        if blip is not None:
            part = slide.part
            _, img_rId = part.get_or_add_image_part(img_path)
            blip.set(f'{{{ns_r}}}embed', img_rId)
        else:
            for child in el.iter():
                if child.tag == f'{{{ns_a}}}blip':
                    part = slide.part
                    _, img_rId = part.get_or_add_image_part(img_path)
                    child.set(f'{{{ns_r}}}embed', img_rId)
                    break
    except Exception as e:
        logging.warning(f"[T18] Rasm almashtirish xatoligi: {e}")

def _t18_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t18_slide_1_cover(slide, topic, name_surname):
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        # Sarlavha matn blokini 4sm pastga tushirish (4sm = 4/2.54*914400 = 1440000 EMU)
        slide.shapes[0].top = slide.shapes[0].top + 1440000
        _t18_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 6000, 'b': 1, 'color': '000000', 'text': topic}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t18_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': '000000', 'text': name_surname}]}
        ])

def fill_t18_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t18_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 7340, 'b': 1, 'color': '000000', 'text': plan_title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'marL': 342900,
                'indent': -342900,
                'spcPts': 150,
                'runs': [{'sz': 3200, 'b': 0, 'color': '101010', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t18_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t18_slide_3_image_left_two_texts(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t18_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if img_arg:
        _t18_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t18_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 4400, 'b': 1, 'color': '000000', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t18_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '101010', 'text': text1}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t18_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 2000, 'b': 0, 'color': '101010', 'text': text2}]}
        ])

def fill_t18_slide_4_text_center_two_texts(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t18_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t18_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 1, 'color': '101010', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t18_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '101010', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t18_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '101010', 'text': text2}]}
        ])

def fill_t18_slide_5_image_right_text_left(slide, data, img_arg=None):
    import os
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t18_get_body_text(data)
    if img_arg:
        # shapes[0] = Freeform (ikonka rasm) - mavzuga doir haqiqiy rasm bilan almashtirish
        # Freeform o'lchamida add_picture qo'shamiz
        try:
            freeform_sh = slide.shapes[0]
            left = freeform_sh.left
            top = freeform_sh.top
            width = freeform_sh.width
            height = freeform_sh.height
            if isinstance(img_arg, str) and os.path.exists(img_arg):
                img_path = img_arg
            else:
                img_path = fetch_image(img_arg)
            if img_path and os.path.exists(img_path):
                slide.shapes.add_picture(img_path, left, top, width, height)
        except Exception as e:
            logging.warning(f"[T18] Slayd5 rasm qo'shish xatoligi: {e}")
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t18_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 1, 'color': '101010', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t18_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '101010', 'text': body_text}]}
        ])

def fill_t18_slide_6_text_center_two_texts_shape(slide, data, img_arg=None):
    import os
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t18_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t18_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 0, 'color': '000000', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t18_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '101010', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t18_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '101010', 'text': text2}]}
        ])
    # shapes[3] = Freeform (solid fill) - rasm qo'shish (add_picture o'sha joydagi o'lchamda)
    if img_arg and len(slide.shapes) > 3:
        try:
            freeform_sh = slide.shapes[3]
            left = freeform_sh.left
            top = freeform_sh.top
            width = freeform_sh.width
            height = freeform_sh.height
            if isinstance(img_arg, str) and os.path.exists(img_arg):
                img_path = img_arg
            else:
                img_path = fetch_image(img_arg)
            if img_path and os.path.exists(img_path):
                slide.shapes.add_picture(img_path, left, top, width, height)
        except Exception as e:
            logging.warning(f"[T18] Slayd6 rasm qo'shish xatoligi: {e}")

def fill_t18_slide_7_image_right_text_left_colored(slide, data, img_arg=None):
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t18_get_body_text(data)
    if img_arg:
        # shapes[0] = Picture - rasm almashtirish va 2sm o'ngga siljitish
        _t18_replace_blip(slide, 0, img_arg)
        # Picture ni 2sm o'ngga siljitish (2sm = 2/2.54*914400 = 720000 EMU)
        if len(slide.shapes) > 0:
            slide.shapes[0].left = slide.shapes[0].left + 720000
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t18_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4000, 'b': 1, 'color': 'F4F6FC', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t18_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2090, 'b': 0, 'color': 'FFFFFF', 'text': body_text}]}
        ])

def fill_t18_slide_8_conclusion(slide, data):
    pass

def generate_template_18_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T18] Shablon slaydlari yetarli emas")
        return None
    fill_t18_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t18_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t18_slide_3_image_left_two_texts,
        fill_t18_slide_4_text_center_two_texts,
        fill_t18_slide_5_image_right_text_left,
        fill_t18_slide_6_text_center_two_texts_shape,
        fill_t18_slide_7_image_right_text_left_colored,
    ]
    user_img_idx = 0
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        image_query = data.get("image_query", topic)
        slide_type = i % len(content_slide_funcs)
        if user_images and user_img_idx < len(user_images):
            img_path = save_user_image_to_tmp(user_images[user_img_idx])
            user_img_idx += 1
            img_arg = img_path if img_path else image_query
        else:
            img_arg = image_query
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T18] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t18_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ============================================================
# 19-SHABLON (Cyber Security / Dark theme)
# ============================================================

SLIDE_TYPE_NAMES_T19 = {
    "cover": "Muqova",
    "plan": "Reja",
    "image_left_text_right": "Rasm chap, matn o'ng",
    "image_left_group_text_right": "Rasm chap (group), matn o'ng",
    "two_images_two_texts": "Ikki rasm, ikki matn",
    "image_right_text_left": "Rasm o'ng, matn chap",
    "three_images_text_right": "Uch rasm, matn o'ng",
    "conclusion": "Xulosa",
}

def _t19_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            color = run.get('color', 'FFFFFF')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t19_replace_blip(slide, shape_idx, img_arg):
    """Freeform yoki Picture ichidagi blip ni yangi rasm bilan almashtirish"""
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        blip = el.find('.//a:blip', {'a': ns_a})
        if blip is not None:
            part = slide.part
            _, img_rId = part.get_or_add_image_part(img_path)
            blip.set(f'{{{ns_r}}}embed', img_rId)
    except Exception as e:
        logging.warning(f"[T19] Rasm almashtirish xatoligi (idx={shape_idx}): {e}")

def _t19_replace_group_blip(slide, group_idx, img_arg):
    """Group ichidagi birinchi Freeform blipini almashtirish"""
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        group = slide.shapes[group_idx]
        el = group._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        blip = el.find('.//a:blip', {'a': ns_a})
        if blip is not None:
            part = slide.part
            _, img_rId = part.get_or_add_image_part(img_path)
            blip.set(f'{{{ns_r}}}embed', img_rId)
    except Exception as e:
        logging.warning(f"[T19] Group rasm almashtirish xatoligi (idx={group_idx}): {e}")

def _t19_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t19_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (60pt, oq, center)
    # shapes[1] = ism (20pt, oq, center)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t19_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6000, 'b': 1, 'color': 'FFFFFF', 'text': topic}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t19_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': name_surname}]}
        ])

def fill_t19_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    # shapes[0] = sarlavha (80pt, oq, center)
    # shapes[1] = reja matn (24pt, oq, justify) - ko'p paragraf
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t19_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 8000, 'b': 1, 'color': 'FFFFFF', 'text': plan_title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'marL': 342900,
                'indent': -342900,
                'spcPts': 200,
                'runs': [{'sz': 2400, 'b': 0, 'color': 'FFFFFF', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t19_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t19_slide_3_image_left_text_right(slide, data, img_arg=None):
    # shapes[0] = sarlavha (48pt, oq, left)
    # shapes[1] = Freeform rasm (blip bor)
    # shapes[2] = matn (20pt, oq, justify)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t19_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t19_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': 'FFFFFF', 'text': title}]}
        ])
    if img_arg:
        _t19_replace_blip(slide, 1, img_arg)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t19_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': body_text}]}
        ])

def fill_t19_slide_4_image_left_group_text_right(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt, oq, left)
    # shapes[1] = matn (20pt, oq, justify)
    # shapes[2] = Group (Freeform rasm bor)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t19_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t19_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': 'FFFFFF', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t19_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': body_text}]}
        ])
    if img_arg:
        _t19_replace_group_blip(slide, 2, img_arg)

def fill_t19_slide_5_two_images_two_texts(slide, data, img_arg=None):
    # shapes[0] = sarlavha (48pt, oq, left)
    # shapes[1] = Freeform rasm 1 (blip bor)
    # shapes[2] = Freeform rasm 2 (blip bor)
    # shapes[3] = matn 1 (20pt, oq, justify)
    # shapes[4] = matn 2 (20pt, oq, justify)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t19_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t19_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': 'FFFFFF', 'text': title}]}
        ])
    if img_arg:
        img_arg2 = img_arg + " detail" if isinstance(img_arg, str) and not __import__('os').path.exists(img_arg) else img_arg
        _t19_replace_blip(slide, 1, img_arg)
        _t19_replace_blip(slide, 2, img_arg2)
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t19_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text1}]}
        ])
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        _t19_clear_and_write(slide.shapes[4].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text2}]}
        ])

def fill_t19_slide_6_image_right_text_left(slide, data, img_arg=None):
    # shapes[0] = sarlavha (40pt, oq, left)
    # shapes[1] = Freeform rasm (blip bor)
    # shapes[2] = matn (20pt, oq, justify)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t19_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t19_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4000, 'b': 1, 'color': 'FFFFFF', 'text': title}]}
        ])
    if img_arg:
        _t19_replace_blip(slide, 1, img_arg)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t19_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': body_text}]}
        ])

def fill_t19_slide_7_three_images_text_right(slide, data, img_arg=None):
    # shapes[0] = sarlavha (40pt, oq)
    # shapes[1] = matn (20pt, oq, justify)
    # shapes[2] = Freeform rasm 1 (blip bor)
    # shapes[3] = Freeform rasm 2 (blip bor)
    # shapes[4] = Freeform rasm 3 (blip bor)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t19_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t19_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4000, 'b': 1, 'color': 'FFFFFF', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t19_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': body_text}]}
        ])
    if img_arg:
        import os
        img_arg2 = img_arg + " close-up" if isinstance(img_arg, str) and not os.path.exists(img_arg) else img_arg
        img_arg3 = img_arg + " background" if isinstance(img_arg, str) and not os.path.exists(img_arg) else img_arg
        _t19_replace_blip(slide, 2, img_arg)
        _t19_replace_blip(slide, 3, img_arg2)
        _t19_replace_blip(slide, 4, img_arg3)

def fill_t19_slide_8_conclusion(slide, data):
    pass

def generate_template_19_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T19] Shablon slaydlari yetarli emas")
        return None
    fill_t19_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t19_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t19_slide_3_image_left_text_right,
        fill_t19_slide_4_image_left_group_text_right,
        fill_t19_slide_5_two_images_two_texts,
        fill_t19_slide_6_image_right_text_left,
        fill_t19_slide_7_three_images_text_right,
    ]
    user_img_idx = 0
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        image_query = data.get("image_query", topic)
        slide_type = i % len(content_slide_funcs)
        if user_images and user_img_idx < len(user_images):
            img_path = save_user_image_to_tmp(user_images[user_img_idx])
            user_img_idx += 1
            img_arg = img_path if img_path else image_query
        else:
            img_arg = image_query
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T19] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t19_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ============================================================
# 20-SHABLON (Group Project - qizil/ko'k rang sxemasi)
# ============================================================

SLIDE_TYPE_NAMES_T20 = {
    "cover": "Muqova",
    "plan": "Reja",
    "text_center": "Sarlavha va matn markazda",
    "three_text_columns": "Uch ustun matn",
    "two_text_blocks_title_left": "Sarlavha chap, ikki matn bloki",
    "image_right_text_left": "Rasm o'ng, matn chap",
    "three_circle_images_text": "Uch doira rasm, matn pastda",
    "conclusion": "Xulosa",
}

def _t20_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            color = run.get('color', '466FB8')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t20_replace_blip(slide, shape_idx, img_arg):
    """Picture yoki Freeform ichidagi blipni yangi rasm bilan almashtirish"""
    try:
        import os
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        shape = slide.shapes[shape_idx]
        el = shape._element
        if isinstance(img_arg, str) and os.path.exists(img_arg):
            img_path = img_arg
        else:
            img_path = fetch_image(img_arg)
        if not img_path or not os.path.exists(img_path):
            return
        blip = el.find('.//a:blip', {'a': ns_a})
        if blip is not None:
            part = slide.part
            _, img_rId = part.get_or_add_image_part(img_path)
            blip.set(f'{{{ns_r}}}embed', img_rId)
    except Exception as e:
        logging.warning(f"[T20] Rasm almashtirish xatoligi (idx={shape_idx}): {e}")

def _t20_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t20_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (66pt, qizil FF4036)
    # shapes[1] = ism (32pt, ko'k 018AD0)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t20_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 6600, 'b': 0, 'color': 'FF4036', 'text': topic}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t20_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3200, 'b': 1, 'color': '018AD0', 'text': name_surname}]}
        ])

def fill_t20_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    # shapes[0] = matn bloki (28pt, ko'k 466FB8)
    # shapes[1] = sarlavha (75pt, ko'k 466FB8, center)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t20_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 7500, 'b': 0, 'color': '466FB8', 'text': plan_title}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'j',
                'marL': 342900,
                'indent': -342900,
                'spcPts': 150,
                'runs': [{'sz': 2800, 'b': 0, 'color': '466FB8', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t20_clear_and_write(slide.shapes[0].text_frame._txBody, paras)

def fill_t20_slide_3_text_center(slide, data, img_arg=None):
    # shapes[0] = matn (28pt, ko'k 466FB8, justify)
    # shapes[1] = sarlavha (48pt, ko'k 466FB8, center)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t20_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t20_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 0, 'color': '466FB8', 'text': title}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t20_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2800, 'b': 0, 'color': '466FB8', 'text': body_text}]}
        ])

def fill_t20_slide_4_three_text_columns(slide, data, img_arg=None):
    # shapes[0] = matn 1 (27pt, ko'k 466FB8, justify)
    # shapes[1] = matn 2 (27pt, ko'k 466FB8, justify)
    # shapes[2] = matn 3 (27pt, ko'k 466FB8, justify)
    # shapes[3] = sarlavha (48pt, ko'k 466FB8, center)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t20_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 3)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    text3 = blocks[2] if len(blocks) > 2 else ""
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t20_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 0, 'color': '466FB8', 'text': title}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t20_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2700, 'b': 0, 'color': '466FB8', 'text': text1}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t20_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2700, 'b': 0, 'color': '466FB8', 'text': text2}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t20_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2700, 'b': 0, 'color': '466FB8', 'text': text3}]}
        ])

def fill_t20_slide_5_two_text_blocks_title_left(slide, data, img_arg=None):
    # shapes[0] = matn 1 (28pt, ko'k 466FB8, justify) - o'ng yuqori
    # shapes[1] = sarlavha (44pt, ko'k 466FB8, center) - chap
    # shapes[2] = matn 2 (28pt, ko'k 466FB8, justify) - o'ng pastki
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t20_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t20_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 0, 'color': '466FB8', 'text': title}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t20_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2800, 'b': 0, 'color': '466FB8', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t20_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2800, 'b': 0, 'color': '466FB8', 'text': text2}]}
        ])

def fill_t20_slide_6_image_right_text_left(slide, data, img_arg=None):
    # shapes[0] = Picture rasm (o'ng tomonda)
    # shapes[1] = matn (28pt, ko'k 466FB8, justify)
    # shapes[2] = sarlavha (40pt, ko'k 466FB8, center)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t20_get_body_text(data)
    if img_arg:
        _t20_replace_blip(slide, 0, img_arg)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t20_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4000, 'b': 0, 'color': '466FB8', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t20_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'j', 'runs': [{'sz': 2800, 'b': 0, 'color': '466FB8', 'text': body_text}]}
        ])

def fill_t20_slide_7_three_circle_images_text(slide, data, img_arg=None):
    # shapes[0] = Freeform rasm 1 (doira)
    # shapes[1] = Freeform rasm 2 (doira)
    # shapes[2] = Freeform rasm 3 (doira)
    # shapes[3] = matn (28pt, ko'k 466FB8, center)
    # shapes[4] = sarlavha (44pt, ko'k 466FB8, center)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t20_get_body_text(data)
    if img_arg:
        import os
        img_arg2 = img_arg + " close-up" if isinstance(img_arg, str) and not os.path.exists(img_arg) else img_arg
        img_arg3 = img_arg + " detail" if isinstance(img_arg, str) and not os.path.exists(img_arg) else img_arg
        _t20_replace_blip(slide, 0, img_arg)
        _t20_replace_blip(slide, 1, img_arg2)
        _t20_replace_blip(slide, 2, img_arg3)
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        _t20_clear_and_write(slide.shapes[4].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 0, 'color': '466FB8', 'text': title}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t20_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 2800, 'b': 0, 'color': '466FB8', 'text': body_text}]}
        ])

def fill_t20_slide_8_conclusion(slide, data):
    pass

def generate_template_20_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T20] Shablon slaydlari yetarli emas")
        return None
    fill_t20_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t20_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t20_slide_3_text_center,
        fill_t20_slide_4_three_text_columns,
        fill_t20_slide_5_two_text_blocks_title_left,
        fill_t20_slide_6_image_right_text_left,
        fill_t20_slide_7_three_circle_images_text,
    ]
    user_img_idx = 0
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        image_query = data.get("image_query", topic)
        slide_type = i % len(content_slide_funcs)
        if user_images and user_img_idx < len(user_images):
            img_path = save_user_image_to_tmp(user_images[user_img_idx])
            user_img_idx += 1
            img_arg = img_path if img_path else image_query
        else:
            img_arg = image_query
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T20] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t20_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ============================================================
# 21-SHABLON (Elements of Poetry - to'q ko'k/oq, rasm yo'q)
# ============================================================

SLIDE_TYPE_NAMES_T21 = {
    "cover": "Muqova",
    "plan": "Reja",
    "title_right_text_right": "Sarlavha o'ng, matn o'ng pastda",
    "title_left_text_left": "Sarlavha chap, matn chap pastda",
    "title_right_text_right_2": "Sarlavha o'ng, matn o'ng (2)",
    "title_left_text_left_2": "Sarlavha chap, matn chap (2)",
    "title_right_text_right_3": "Sarlavha o'ng, matn o'ng (3)",
    "conclusion": "Xulosa",
}

def _t21_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2800)
            b = run.get('b', 0)
            color = run.get('color', '1C2143')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t21_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t21_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (48pt, to'q ko'k 1C2143, left)
    # shapes[1] = ism/tavsif (20pt, to'q ko'k 1C2143)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t21_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': '1C2143', 'text': topic}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t21_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '1C2143', 'text': name_surname}]}
        ])

def fill_t21_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    # shapes[0] = sarlavha (70pt, to'q ko'k 1C2143, center)
    # shapes[1] = reja matn (28pt, to'q ko'k 1C2143, left)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t21_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 7000, 'b': 1, 'color': '1C2143', 'text': plan_title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'marL': 342900,
                'indent': -342900,
                'spcPts': 150,
                'runs': [{'sz': 2800, 'b': 0, 'color': '1C2143', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t21_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t21_slide_3_title_right_text_right(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt, to'q ko'k 1C2143, right)
    # shapes[1] = matn (28pt, to'q ko'k 1C2143, left) - o'ng pastda
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t21_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t21_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 4400, 'b': 1, 'color': '1C2143', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t21_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': '1C2143', 'text': body_text}]}
        ])

def fill_t21_slide_4_title_left_text_left(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt, to'q ko'k 1C2143, left)
    # shapes[1] = matn (28pt, to'q ko'k 1C2143, left) - chap pastda
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t21_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t21_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 1, 'color': '1C2143', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t21_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': '1C2143', 'text': body_text}]}
        ])

def fill_t21_slide_5_title_right_text_right_2(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt, to'q ko'k 1C2143, right)
    # shapes[1] = matn (28pt, to'q ko'k 1C2143, left) - o'ng pastda
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t21_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t21_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 4400, 'b': 1, 'color': '1C2143', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t21_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': '1C2143', 'text': body_text}]}
        ])

def fill_t21_slide_6_title_left_text_left_2(slide, data, img_arg=None):
    # shapes[0] = sarlavha (48pt, to'q ko'k 1C2143, left)
    # shapes[1] = matn (28pt, to'q ko'k 1C2143, left) - chap pastda
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t21_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t21_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': '1C2143', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t21_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': '1C2143', 'text': body_text}]}
        ])

def fill_t21_slide_7_title_right_text_right_3(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt, to'q ko'k 1C2143, left)
    # shapes[1] = matn (28pt, to'q ko'k 1C2143, left) - o'ng pastda
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t21_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t21_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 1, 'color': '1C2143', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t21_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': '1C2143', 'text': body_text}]}
        ])

def fill_t21_slide_8_conclusion(slide, data):
    pass

def generate_template_21_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T21] Shablon slaydlari yetarli emas")
        return None
    fill_t21_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t21_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t21_slide_3_title_right_text_right,
        fill_t21_slide_4_title_left_text_left,
        fill_t21_slide_5_title_right_text_right_2,
        fill_t21_slide_6_title_left_text_left_2,
        fill_t21_slide_7_title_right_text_right_3,
    ]
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        slide_type = i % len(content_slide_funcs)
        content_slide_funcs[slide_type](slide, data, None)
        logging.info(f"  [T21] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t21_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ============================================================

# ============================================================
# 22-SHABLON (Informative Presentation v2 - terakota C15943, yangi layout)
# ============================================================

SLIDE_TYPE_NAMES_T22 = {
    "cover": "Muqova",
    "plan": "Reja",
    "title_left_two_text_rows_dark": "Sarlavha chap, ikki qator matn (to'q rang)",
    "two_text_columns_center_title": "Ikki ustun matn, sarlavha markazda",
    "title_left_two_text_rows": "Sarlavha chap, ikki qator matn",
    "title_left_text_left_freeform": "Sarlavha chap, matn chap, dekor",
    "two_text_columns_left_title": "Ikki ustun matn, sarlavha chap",
    "conclusion": "Xulosa",
}

def _t22_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2800)
            b = run.get('b', 0)
            color = run.get('color', 'C15943')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t22_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t22_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (48pt, C15943, bold, left)
    # shapes[1] = ism (24pt, C15943, left)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t22_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': 'C15943', 'text': topic.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t22_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': 'C15943', 'text': f"Presented by {name_surname}"}]}
        ])

def fill_t22_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    # shapes[0] = matn (28pt, C15943, left)
    # shapes[1] = sarlavha (79pt, C15943, bold, left)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t22_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 7900, 'b': 1, 'color': 'C15943', 'text': plan_title.upper()}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'marL': 342900,
                'indent': -342900,
                'spcPts': 150,
                'runs': [{'sz': 2800, 'b': 0, 'color': 'C15943', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t22_clear_and_write(slide.shapes[0].text_frame._txBody, paras)

def fill_t22_slide_3_title_left_two_text_rows_dark(slide, data, img_arg=None):
    # shapes[0] = matn 1 (24pt, 571A0E to'q jigarrang, left) - yuqori
    # shapes[1] = sarlavha (44pt, C15943, bold, left)
    # shapes[2] = matn 2 (24pt, 571A0E to'q jigarrang, left) - pastki
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t22_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t22_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 1, 'color': 'C15943', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t22_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '571A0E', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t22_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '571A0E', 'text': text2}]}
        ])

def fill_t22_slide_4_two_text_columns_center_title(slide, data, img_arg=None):
    # shapes[0] = matn 1 (26pt, C15943, left) - chap ustun
    # shapes[1] = sarlavha (44pt, C15943, bold, center)
    # shapes[2] = matn 2 (26pt, C15943, left) - o'ng ustun
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t22_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t22_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 1, 'color': 'C15943', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t22_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2600, 'b': 0, 'color': 'C15943', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t22_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2600, 'b': 0, 'color': 'C15943', 'text': text2}]}
        ])

def fill_t22_slide_5_title_left_two_text_rows(slide, data, img_arg=None):
    # shapes[0] = matn 1 (24pt, C15943, left) - yuqori
    # shapes[1] = sarlavha (44pt, C15943, bold, left)
    # shapes[2] = matn 2 (24pt, C15943, left) - pastki
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t22_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t22_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 1, 'color': 'C15943', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t22_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': 'C15943', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t22_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': 'C15943', 'text': text2}]}
        ])

def fill_t22_slide_6_title_left_text_left_freeform(slide, data, img_arg=None):
    # shapes[0] = Freeform (rasm joyi - solid fill)
    # shapes[1] = matn (28pt, C15943, left)
    # shapes[2] = sarlavha (40pt, C15943, bold, left)
    import os
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t22_get_body_text(data)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t22_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4000, 'b': 1, 'color': 'C15943', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t22_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': 'C15943', 'text': body_text}]}
        ])
    # Freeform (shapes[0]) o'rniga rasm qo'shish
    if img_arg and len(slide.shapes) > 0:
        try:
            freeform = slide.shapes[0]
            left = freeform.left
            top = freeform.top
            width = freeform.width
            height = freeform.height
            if isinstance(img_arg, str) and os.path.exists(img_arg):
                img_path = img_arg
            else:
                img_path = fetch_image(img_arg)
            if img_path and os.path.exists(img_path):
                slide.shapes.add_picture(img_path, left, top, width, height)
        except Exception as e:
            logging.warning(f"[T22] Slayd6 rasm qo'shish xatoligi: {e}")

def fill_t22_slide_7_two_text_columns_left_title(slide, data, img_arg=None):
    # shapes[0] = matn 1 (28pt, C15943, left) - chap
    # shapes[1] = sarlavha (44pt, C15943, bold, left)
    # shapes[2] = matn 2 (28pt, C15943, left) - o'ng
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t22_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t22_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 1, 'color': 'C15943', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t22_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': 'C15943', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t22_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': 'C15943', 'text': text2}]}
        ])

def fill_t22_slide_8_conclusion(slide, data):
    pass

def generate_template_22_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T22] Shablon slaydlari yetarli emas")
        return None
    fill_t22_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t22_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t22_slide_3_title_left_two_text_rows_dark,
        fill_t22_slide_4_two_text_columns_center_title,
        fill_t22_slide_5_title_left_two_text_rows,
        fill_t22_slide_6_title_left_text_left_freeform,
        fill_t22_slide_7_two_text_columns_left_title,
    ]
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        slide_type = i % len(content_slide_funcs)
        # Slayd 6 (idx 3) uchun rasm
        img_arg = None
        if slide_type == 3:
            if user_images and i < len(user_images):
                img_arg = user_images[i]
            else:
                img_query = data.get("title", topic) if isinstance(data, dict) else topic
                img_arg = fetch_image(img_query)
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T22] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t22_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# 23-SHABLON (World of Sports - oq fon, to'q 202020 matn, har slaydda Freeform rasm)
# ============================================================

SLIDE_TYPE_NAMES_T23 = {
    "cover": "Muqova",
    "plan": "Reja (rasm chap, matn o'ng)",
    "title_bottom_text_right_img_top": "Rasm yuqori, sarlavha va matn pastda",
    "title_top_text_right_img_bottom": "Rasm pastda, sarlavha va matn yuqorida",
    "title_left_text_left_img_right": "Rasm o'ng, sarlavha va matn chap",
    "title_bottom_text_right_img_topleft": "Rasm chap yuqori, sarlavha va matn pastda",
    "title_right_text_right_img_left": "Rasm chap, sarlavha va matn o'ng",
    "conclusion": "Xulosa",
}

def _t23_replace_blip(shape, img_path):
    """Freeform ichidagi blip (rasm) ni yangi rasm bilan almashtirish - get_or_add_image_part usuli"""
    try:
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        blips = shape._element.findall('.//a:blip', {'a': ns_a})
        if not blips:
            return False
        blip = blips[0]
        rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
        if not rId:
            return False
        slide_part = shape.part
        # get_or_add_image_part - eng ishonchli usul
        img_part, new_rId = slide_part.get_or_add_image_part(img_path)
        blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)
        return True
    except Exception as e:
        logging.warning(f"[T23] blip almashtirish xatoligi: {e}")
    return False

def _t23_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2400)
            b = run.get('b', 0)
            color = run.get('color', '202020')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t23_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t23_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = Freeform (rasm, o'ng tomonida)
    # shapes[1] = sarlavha (54pt, 202020, left)
    # shapes[2] = tagline (26pt, 919191, left)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t23_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 5400, 'b': 0, 'color': '202020', 'text': topic.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t23_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2600, 'b': 0, 'color': '919191', 'text': name_surname}]}
        ])

def fill_t23_slide_2_plan(slide, plan_dict, img_arg=None):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    # shapes[0] = Freeform (rasm, chap pastda)
    # shapes[1] = sarlavha (98pt, 202020, center)
    # shapes[2] = reja matn (24pt, 202020, justify)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t23_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 9800, 'b': 0, 'color': '202020', 'text': plan_title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'marL': 342900,
                'indent': -342900,
                'spcPts': 150,
                'runs': [{'sz': 2400, 'b': 0, 'color': '202020', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t23_clear_and_write(slide.shapes[2].text_frame._txBody, paras)
    if img_arg and len(slide.shapes) > 0:
        try:
            img_path = img_arg if (isinstance(img_arg, str) and __import__('os').path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t23_replace_blip(slide.shapes[0], img_path)
        except Exception as e:
            logging.warning(f"[T23] Slayd2 rasm: {e}")

def fill_t23_slide_3_img_top_title_text(slide, data, img_arg=None):
    # shapes[0] = Freeform (rasm, yuqori to'la kenglik)
    # shapes[1] = sarlavha (54pt, 202020, left) - pastda
    # shapes[2] = matn (19pt, 202020, justify) - o'ng pastda
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t23_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t23_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 5400, 'b': 0, 'color': '202020', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t23_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1900, 'b': 0, 'color': '202020', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            img_path = img_arg if (isinstance(img_arg, str) and __import__('os').path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t23_replace_blip(slide.shapes[0], img_path)
        except Exception as e:
            logging.warning(f"[T23] Slayd3 rasm: {e}")

def fill_t23_slide_4_img_bottom_title_text(slide, data, img_arg=None):
    # shapes[0] = Freeform (rasm, pastda to'la kenglik)
    # shapes[1] = sarlavha (48pt, 202020, left) - yuqorida
    # shapes[2] = matn (19pt, 202020, justify) - o'ng yuqorida
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t23_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t23_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '202020', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t23_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1900, 'b': 0, 'color': '202020', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            img_path = img_arg if (isinstance(img_arg, str) and __import__('os').path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t23_replace_blip(slide.shapes[0], img_path)
        except Exception as e:
            logging.warning(f"[T23] Slayd4 rasm: {e}")

def fill_t23_slide_5_img_right_title_text(slide, data, img_arg=None):
    # shapes[0] = Freeform (rasm, o'ng tomonida)
    # shapes[1] = sarlavha (44pt, 919191, left) - chap yuqori
    # shapes[2] = matn (19pt, 202020, justify) - chap pastda
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t23_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t23_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': '919191', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t23_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1900, 'b': 0, 'color': '202020', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            img_path = img_arg if (isinstance(img_arg, str) and __import__('os').path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t23_replace_blip(slide.shapes[0], img_path)
        except Exception as e:
            logging.warning(f"[T23] Slayd5 rasm: {e}")

def fill_t23_slide_6_img_topleft_title_bottom(slide, data, img_arg=None):
    # shapes[0] = Freeform (rasm, chap yuqori)
    # shapes[1] = sarlavha (48pt, 202020, left) - pastda
    # shapes[2] = matn (19pt, 202020, justify) - o'ng
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t23_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t23_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '202020', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t23_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1900, 'b': 0, 'color': '202020', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            img_path = img_arg if (isinstance(img_arg, str) and __import__('os').path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t23_replace_blip(slide.shapes[0], img_path)
        except Exception as e:
            logging.warning(f"[T23] Slayd6 rasm: {e}")

def fill_t23_slide_7_img_left_title_text(slide, data, img_arg=None):
    # shapes[0] = Freeform (rasm, chap tomonida)
    # shapes[1] = sarlavha (44pt, 919191, left) - o'ng yuqori
    # shapes[2] = matn (19pt, 202020, justify) - o'ng pastda
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t23_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t23_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': '919191', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t23_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1900, 'b': 0, 'color': '202020', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            img_path = img_arg if (isinstance(img_arg, str) and __import__('os').path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t23_replace_blip(slide.shapes[0], img_path)
        except Exception as e:
            logging.warning(f"[T23] Slayd7 rasm: {e}")

def fill_t23_slide_8_conclusion(slide, data):
    pass

def generate_template_23_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T23] Shablon slaydlari yetarli emas")
        return None
    fill_t23_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    content_slide_funcs = [
        fill_t23_slide_3_img_top_title_text,
        fill_t23_slide_4_img_bottom_title_text,
        fill_t23_slide_5_img_right_title_text,
        fill_t23_slide_6_img_topleft_title_bottom,
        fill_t23_slide_7_img_left_title_text,
    ]
    # Reja slaydiga rasm
    plan_img = None
    if user_images and len(user_images) > 0:
        plan_img = user_images[0]
    elif content_data_list and len(content_data_list) > 0:
        q = content_data_list[0].get("title", topic) if isinstance(content_data_list[0], dict) else topic
        plan_img = fetch_image(q + " sport")
    fill_t23_slide_2_plan(slides[1], plan_dict, plan_img)
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        slide_type = i % len(content_slide_funcs)
        # Rasm
        img_arg = None
        if user_images and i < len(user_images):
            img_arg = user_images[i]
        else:
            img_query = data.get("title", topic) if isinstance(data, dict) else topic
            img_arg = fetch_image(img_query + " sport")
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T23] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t23_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ============================================================
# 24-SHABLON (Sports Dark - to'q fon, oq FFFFFF matn, Freeform blip rasmlar)
# ============================================================

SLIDE_TYPE_NAMES_T24 = {
    "cover": "Muqova",
    "plan": "Reja (oq matn, to'q fon)",
    "img_right_title_text": "Rasm o'ng, sarlavha va matn chap",
    "title_center_two_text_cols": "Sarlavha markazda, ikki ustun matn",
    "img_right_title_two_text_bottom": "Rasm o'ng yuqori, sarlavha va ikki matn pastda",
    "img_right_title_two_text_bottom2": "Rasm o'ng, sarlavha chap, ikki matn pastda",
    "title_left_three_text": "Sarlavha chap, uch matn bloki",
    "conclusion": "Xulosa",
}

def _t24_replace_blip(shape, img_path):
    """Freeform ichidagi blip (rasm) ni yangi rasm bilan almashtirish"""
    try:
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        blips = shape._element.findall('.//a:blip', {'a': ns_a})
        if not blips:
            return False
        blip = blips[0]
        rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
        if not rId:
            return False
        slide_part = shape.part
        img_part, new_rId = slide_part.get_or_add_image_part(img_path)
        blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)
        return True
    except Exception as e:
        logging.warning(f"[T24] blip almashtirish xatoligi: {e}")
    return False

def _t24_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            color = run.get('color', 'FFFFFF')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t24_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("col1", "") or data.get("text", "")
    return body_text

def fill_t24_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (80pt, FFFFFF, left)
    # shapes[1] = tagline/url (20pt, FFFFFF, left) - pastda
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t24_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 8000, 'b': 0, 'color': 'FFFFFF', 'text': topic.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t24_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': name_surname}]}
        ])

def fill_t24_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    # shapes[0] = sarlavha (123pt, FFFFFF, center)
    # shapes[1] = reja matn (20pt, FFFFFF, justify)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t24_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 12300, 'b': 0, 'color': 'FFFFFF', 'text': plan_title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'marL': 342900,
                'indent': -342900,
                'spcPts': 150,
                'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t24_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t24_slide_3_img_right_title_text(slide, data, img_arg=None):
    # shapes[0] = Freeform (rasm, o'ng tomonida)
    # shapes[1] = sarlavha (48pt, FFFFFF, left)
    # shapes[2] = matn (20pt, FFFFFF, justify)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t24_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t24_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': 'FFFFFF', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t24_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t24_replace_blip(slide.shapes[0], img_path)
        except Exception as e:
            logging.warning(f"[T24] Slayd3 rasm: {e}")

def fill_t24_slide_4_title_center_two_text_cols(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt, FFFFFF, center)
    # shapes[1] = matn chap (20pt, FFFFFF, justify)
    # shapes[2] = matn o'ng (20pt, FFFFFF, justify)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t24_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t24_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 0, 'color': 'FFFFFF', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t24_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t24_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text2}]}
        ])

def fill_t24_slide_5_img_right_title_two_text_bottom(slide, data, img_arg=None):
    # shapes[0] = Freeform (rasm, o'ng yuqori)
    # shapes[1] = sarlavha (44pt, FFFFFF, right)
    # shapes[2] = matn chap pastda (20pt, FFFFFF, justify)
    # shapes[3] = matn o'ng pastda (20pt, FFFFFF, justify)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t24_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t24_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 4400, 'b': 0, 'color': 'FFFFFF', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t24_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text1}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t24_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text2}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t24_replace_blip(slide.shapes[0], img_path)
        except Exception as e:
            logging.warning(f"[T24] Slayd5 rasm: {e}")

def fill_t24_slide_6_img_right_title_two_text_bottom2(slide, data, img_arg=None):
    # shapes[0] = Freeform (rasm, o'ng yuqori)
    # shapes[1] = sarlavha (44pt, FFFFFF, left)
    # shapes[2] = matn chap pastda (20pt, FFFFFF, justify)
    # shapes[3] = matn o'ng pastda (20pt, FFFFFF, justify)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t24_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t24_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': 'FFFFFF', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t24_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text1}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t24_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text2}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t24_replace_blip(slide.shapes[0], img_path)
        except Exception as e:
            logging.warning(f"[T24] Slayd6 rasm: {e}")

def fill_t24_slide_7_title_left_three_text(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt, FFFFFF, left)
    # shapes[1] = matn chap pastda (20pt, FFFFFF, justify)
    # shapes[2] = matn o'ng yuqori (20pt, FFFFFF, justify)
    # shapes[3] = matn o'ng pastda (20pt, FFFFFF, justify)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t24_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 3)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    text3 = blocks[2] if len(blocks) > 2 else ""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t24_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': 'FFFFFF', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t24_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t24_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text2}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t24_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2000, 'b': 0, 'color': 'FFFFFF', 'text': text3}]}
        ])

def fill_t24_slide_8_conclusion(slide, data):
    pass

def generate_template_24_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T24] Shablon slaydlari yetarli emas")
        return None
    fill_t24_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t24_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t24_slide_3_img_right_title_text,
        fill_t24_slide_4_title_center_two_text_cols,
        fill_t24_slide_5_img_right_title_two_text_bottom,
        fill_t24_slide_6_img_right_title_two_text_bottom2,
        fill_t24_slide_7_title_left_three_text,
    ]
    # Rasm kerak bo'lgan slayd turlari (0-indexed content_slide_funcs)
    img_slide_types = {0, 2, 3}
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        slide_type = i % len(content_slide_funcs)
        img_arg = None
        if slide_type in img_slide_types:
            if user_images and i < len(user_images):
                img_arg = user_images[i]
            else:
                img_query = data.get("title", topic) if isinstance(data, dict) else topic
                img_arg = fetch_image(img_query + " sport")
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T24] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t24_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ============================================================
# 25-SHABLON (Sports Blue Dark - to'q ko'k 01002A fon, oq FFFFFF matn, Picture rasmlar)
# ============================================================

SLIDE_TYPE_NAMES_T25 = {
    "cover": "Muqova",
    "plan": "Reja",
    "img_right_title_text": "Rasm o'ng, sarlavha va matn chap",
    "img_left_title_two_text": "Rasm chap, sarlavha va ikki matn o'ng",
    "img_top_right_title_two_text": "Rasm yuqori o'ng, sarlavha va ikki matn pastda",
    "two_img_title_three_text": "Ikki rasm, sarlavha va uch matn",
    "img_left_title_text": "Rasm chap, sarlavha va matn o'ng",
    "conclusion": "Xulosa",
}

def _t25_replace_picture(slide, shape_index, img_path):
    """Picture shape ichidagi rasmni yangi rasm bilan almashtirish"""
    try:
        import os
        from lxml import etree
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        shape = slide.shapes[shape_index]
        blips = shape._element.findall('.//a:blip', {'a': ns_a})
        if not blips:
            return False
        blip = blips[0]
        slide_part = shape.part
        img_part, new_rId = slide_part.get_or_add_image_part(img_path)
        blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)
        return True
    except Exception as e:
        import logging
        logging.warning(f"[T25] Picture almashtirish xatoligi: {e}")
    return False

def _t25_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        marL = para.get('marL', 0)
        indent = para.get('indent', 0)
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        pPr_attrs = f'algn="{algn}"'
        if marL:
            pPr_attrs += f' marL="{marL}"'
        if indent:
            pPr_attrs += f' indent="{indent}"'
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            color = run.get('color', 'FFFFFF')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr {pPr_attrs}>{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t25_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("text", "")
    return body_text

def _t25_split_into_blocks(body_text, n):
    """Matnni n ta alohida gapga bo'lish - har blok alohida gap bo'lishi kerak"""
    import re
    # Avval gaplarga bo'lish
    sentences = re.split(r'(?<=[.!?])\s+', body_text.strip())
    sentences = [s.strip() for s in sentences if s.strip()]
    if len(sentences) >= n:
        # Gaplarni n ta blokga teng taqsimlash
        blocks = []
        per_block = max(1, len(sentences) // n)
        for i in range(n):
            start = i * per_block
            end = start + per_block if i < n - 1 else len(sentences)
            block = " ".join(sentences[start:end])
            if block and not block.endswith('.'):
                block += '.'
            blocks.append(block)
        return blocks
    else:
        # Gaplar yetarli emas, so'zlarga bo'lish
        words = body_text.split()
        blocks = []
        per_block = max(1, len(words) // n)
        for i in range(n):
            start = i * per_block
            end = start + per_block if i < n - 1 else len(words)
            block = " ".join(words[start:end])
            if block and not block.endswith('.'):
                block += '.'
            # Bosh harf bilan boshlash
            if block:
                block = block[0].upper() + block[1:]
            blocks.append(block)
        return blocks

def fill_t25_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (60pt, FFFFFF, center)
    # shapes[1] = tagline (18pt, 5DDFE6, center)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t25_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6000, 'b': 0, 'color': 'FFFFFF', 'text': topic.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t25_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 1800, 'b': 0, 'color': '5DDFE6', 'text': name_surname}]}
        ])

def fill_t25_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    # shapes[0] = sarlavha (80pt, 01002A, left)
    # shapes[1] = reja matn (32pt, FFFFFF, left)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t25_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 8000, 'b': 0, 'color': '01002A', 'text': plan_title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'spcPts': 150,
                'runs': [{'sz': 3200, 'b': 0, 'color': 'FFFFFF', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t25_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t25_slide_3_img_right_title_text(slide, data, img_arg=None):
    # shapes[0] = Picture (rasm, o'ng)
    # shapes[1] = sarlavha (48pt, 01002A, left)
    # shapes[2] = matn (22pt, FFFFFF, left)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t25_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t25_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '01002A', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t25_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': 'FFFFFF', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t25_replace_picture(slide, 0, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T25] Slayd3 rasm: {e}")

def fill_t25_slide_4_img_left_title_two_text(slide, data, img_arg=None):
    # shapes[0] = Picture (rasm, chap)
    # shapes[1] = sarlavha (44pt, 01002A, left)
    # shapes[2] = matn1 (22pt, 1F2020, left) - yuqori o'ng
    # shapes[3] = matn2 (22pt, 1F2020, left) - pastki o'ng
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t25_get_body_text(data)
    blocks = _t25_split_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t25_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': '01002A', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t25_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '1F2020', 'text': text1}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t25_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '1F2020', 'text': text2}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t25_replace_picture(slide, 0, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T25] Slayd4 rasm: {e}")

def fill_t25_slide_5_img_top_right_title_two_text(slide, data, img_arg=None):
    # shapes[0] = Picture (rasm, yuqori o'ng)
    # shapes[1] = sarlavha (44pt, 01002A, left)
    # shapes[2] = matn1 (22pt, FFFFFF, left) - pastki chap
    # shapes[3] = matn2 (22pt, FFFFFF, left) - pastki o'ng
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t25_get_body_text(data)
    blocks = _t25_split_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t25_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': '01002A', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t25_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': 'FFFFFF', 'text': text1}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t25_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': 'FFFFFF', 'text': text2}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t25_replace_picture(slide, 0, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T25] Slayd5 rasm: {e}")

def fill_t25_slide_6_two_img_title_three_text(slide, data, img_arg=None, img_arg2=None):
    # shapes[0] = Picture1 (rasm, markazda)
    # shapes[1] = Picture2 (rasm, o'ng yuqori)
    # shapes[2] = sarlavha (60pt, 01002A, left)
    # shapes[3] = matn1 (22pt, 1F2020, left) - o'ng pastda
    # shapes[4] = matn2 (22pt, 1F2020, left) - pastda
    # shapes[5] = matn3 (22pt, 1F2020, left) - chap
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t25_get_body_text(data)
    blocks = _t25_split_into_blocks(body_text, 3)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    text3 = blocks[2] if len(blocks) > 2 else ""
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t25_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 6000, 'b': 0, 'color': '01002A', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t25_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '1F2020', 'text': text1}]}
        ])
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        _t25_clear_and_write(slide.shapes[4].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '1F2020', 'text': text2}]}
        ])
    if len(slide.shapes) > 5 and slide.shapes[5].has_text_frame:
        _t25_clear_and_write(slide.shapes[5].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '1F2020', 'text': text3}]}
        ])
    # Rasmlarni almashtirish
    if img_arg and len(slide.shapes) > 0:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t25_replace_picture(slide, 0, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T25] Slayd6 rasm1: {e}")
    if img_arg2 and len(slide.shapes) > 1:
        try:
            import os
            img_path2 = img_arg2 if (isinstance(img_arg2, str) and os.path.exists(img_arg2)) else fetch_image(img_arg2 + " sport action")
            if img_path2:
                _t25_replace_picture(slide, 1, img_path2)
        except Exception as e:
            import logging
            logging.warning(f"[T25] Slayd6 rasm2: {e}")

def fill_t25_slide_7_img_left_title_text(slide, data, img_arg=None):
    # shapes[0] = Picture (rasm, chap)
    # shapes[1] = sarlavha (48pt, 01002A, left)
    # shapes[2] = matn (28pt, 1F2020, left)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t25_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t25_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '01002A', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t25_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': '1F2020', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 0:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t25_replace_picture(slide, 0, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T25] Slayd7 rasm: {e}")

def fill_t25_slide_8_conclusion(slide, data):
    pass

def generate_template_25_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    import logging
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T25] Shablon slaydlari yetarli emas")
        return None
    fill_t25_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t25_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t25_slide_3_img_right_title_text,
        fill_t25_slide_4_img_left_title_two_text,
        fill_t25_slide_5_img_top_right_title_two_text,
        fill_t25_slide_6_two_img_title_three_text,
        fill_t25_slide_7_img_left_title_text,
    ]
    # Rasm kerak bo'lgan slayd turlari (0-indexed)
    img_slide_types = {0, 1, 2, 3, 4}
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        slide_type = i % len(content_slide_funcs)
        img_arg = None
        img_arg2 = None
        if slide_type in img_slide_types:
            if user_images and i < len(user_images):
                img_arg = user_images[i]
            else:
                img_query = data.get("title", topic) if isinstance(data, dict) else topic
                img_arg = fetch_image(img_query + " sports")
            # Slayd6 uchun ikkinchi rasm
            if slide_type == 3:
                img_query2 = data.get("title", topic) if isinstance(data, dict) else topic
                img_arg2 = fetch_image(img_query2 + " athlete")
        if slide_type == 3:
            content_slide_funcs[slide_type](slide, data, img_arg, img_arg2)
        else:
            content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T25] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t25_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ============================================================
# 26-SHABLON (Science Presentation - to'q ko'k 032C54 fon, oq EFECE0 matn, Freeform blip rasmlar)
# Slayd tuzilishi:
# 1: Muqova - sarlavha 60pt EFECE0 + ism 24pt bold EFECE0
# 2: Reja - sarlavha 84pt EFECE0 center + reja matn 34pt 032C54 bold
# 3: Rasm o'ng (Freeform blip), sarlavha center, matn chap 28pt - shapes: [matn, freeform, sarlavha]
# 4: Rasm o'ng (Freeform blip), sarlavha center, matn chap 28pt - shapes: [matn, sarlavha, freeform]
# 5: Rasm chap (Freeform blip), sarlavha center, ikki matn o'ng 25pt - shapes: [sarlavha, matn1, matn2, freeform]
# 6: Rasm chap (Freeform blip), sarlavha center, matn o'ng 28pt - shapes: [sarlavha, freeform, matn]
# 7: Rasm pastki o'ng (Freeform blip), sarlavha center, uch matn 25pt - shapes: [sarlavha, matn_pastki_chap, matn_yuqori_chap, matn_yuqori_ong, freeform]
# 8: Xulosa - "E'TIBORINGIZ UCHUN RAHMAT!" 88pt EFECE0 center
# ============================================================

SLIDE_TYPE_NAMES_T26 = {
    "cover": "Muqova",
    "plan": "Reja",
    "img_right_title_text": "Rasm o'ng, sarlavha va matn chap (1)",
    "img_right_title_text2": "Rasm o'ng, sarlavha va matn chap (2)",
    "img_left_title_two_text": "Rasm chap, sarlavha va ikki matn o'ng",
    "img_left_title_text": "Rasm chap, sarlavha va matn o'ng",
    "img_bottom_right_title_three_text": "Rasm pastki o'ng, sarlavha va uch matn",
    "conclusion": "Xulosa",
}

def _t26_clear_and_write(txBody, paragraphs_data):
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_elem in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p_elem)
    for para in paragraphs_data:
        algn = para.get('algn', 'l')
        spcPts = para.get('spcPts', None)
        runs = para.get('runs', [])
        spcBef_xml = ''
        if spcPts:
            spcBef_xml = f'<a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2000)
            b = run.get('b', 0)
            color = run.get('color', 'EFECE0')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            b_val = '1' if b else '0'
            runs_xml += (
                f'<a:r><a:rPr lang="uz-UZ" sz="{sz}" b="{b_val}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'</a:rPr><a:t>{text}</a:t></a:r>'
            )
        p_xml = (
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr algn="{algn}">{spcBef_xml}</a:pPr>'
            f'{runs_xml}'
            f'</a:p>'
        )
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)

def _t26_replace_blip(slide, shape_index, img_path):
    """Freeform blip shapeni topib, uni o'chirib o'rniga add_picture bilan rasm qo'yish.
    shape_index berilsa shu indeksdan, aks holda slayddagi birinchi blip topiladi."""
    import logging
    try:
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        # Freeform blip shapeni topish
        target_shape = None
        if shape_index is not None and shape_index < len(slide.shapes):
            candidate = slide.shapes[shape_index]
            blips = candidate._element.findall(f'.//{{{ns_a}}}blip')
            if blips:
                target_shape = candidate
        # Topilmasa, slayddagi barcha shapelarda blip qidirish
        if target_shape is None:
            for shape in slide.shapes:
                blips = shape._element.findall(f'.//{{{ns_a}}}blip')
                if blips:
                    target_shape = shape
                    break
        if target_shape is None:
            logging.warning(f"[T26] Slaydda blip topilmadi")
            return False
        # Freeform o'lchamlarini saqlash
        left = target_shape.left
        top = target_shape.top
        width = target_shape.width
        height = target_shape.height
        # Freeform ni o'chirish
        sp = target_shape._element
        sp.getparent().remove(sp)
        # O'rniga rasm qo'yish
        slide.shapes.add_picture(img_path, left, top, width, height)
        return True
    except Exception as e:
        logging.warning(f"[T26] Blip almashtirish xatoligi: {e}")
    return False

def _t26_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("text", "")
    return body_text

def fill_t26_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (60pt, EFECE0)
    # shapes[1] = ism (24pt bold, EFECE0)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t26_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 6000, 'b': 0, 'color': 'EFECE0', 'text': topic.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t26_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2400, 'b': 1, 'color': 'EFECE0', 'text': name_surname}]}
        ])

def fill_t26_slide_2_plan(slide, plan_dict):
    import re
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    plan_content = plan_dict.get("content", [])
    if isinstance(plan_content, list):
        items = plan_content
    else:
        items = [str(plan_content)]
    clean_items = []
    for item in items:
        item_str = str(item).strip()
        item_str = re.sub(r'^\d+[\.\)]\s*', '', item_str)
        clean_items.append(item_str)
    # shapes[0] = sarlavha (84pt, EFECE0, center)
    # shapes[1] = reja matn (34pt bold, 032C54)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t26_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 8400, 'b': 0, 'color': 'EFECE0', 'text': plan_title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            paras.append({
                'algn': 'l',
                'spcPts': 200,
                'runs': [{'sz': 3400, 'b': 1, 'color': '032C54', 'text': f"{idx}. {item}"}]
            })
        if paras:
            _t26_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t26_slide_3_img_right_title_text(slide, data, img_arg=None):
    # shapes[0] = matn (28pt, EFECE0, justify) - chap
    # shapes[1] = Freeform blip (rasm, o'ng)
    # shapes[2] = sarlavha (48pt, EFECE0, center)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t26_get_body_text(data)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t26_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 0, 'color': 'EFECE0', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t26_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': 'EFECE0', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 1:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t26_replace_blip(slide, 1, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T26] Slayd3 rasm: {e}")

def fill_t26_slide_4_img_right_title_text2(slide, data, img_arg=None):
    # shapes[0] = matn (28pt, EFECE0, justify) - chap
    # shapes[1] = sarlavha (40pt, EFECE0, center)
    # shapes[2] = Freeform blip (rasm, o'ng)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t26_get_body_text(data)
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t26_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4000, 'b': 0, 'color': 'EFECE0', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t26_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': 'EFECE0', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 2:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t26_replace_blip(slide, 2, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T26] Slayd4 rasm: {e}")

def fill_t26_slide_5_img_left_title_two_text(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt, EFECE0, center)
    # shapes[1] = matn1 (25pt, EFECE0, justify) - o'ng yuqori
    # shapes[2] = matn2 (25pt, EFECE0, justify) - o'ng pastki
    # shapes[3] = Freeform blip (rasm, chap)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t26_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t26_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 0, 'color': 'EFECE0', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t26_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2500, 'b': 0, 'color': 'EFECE0', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t26_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2500, 'b': 0, 'color': 'EFECE0', 'text': text2}]}
        ])
    if img_arg and len(slide.shapes) > 3:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t26_replace_blip(slide, 3, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T26] Slayd5 rasm: {e}")

def fill_t26_slide_6_img_left_title_text(slide, data, img_arg=None):
    # shapes[0] = sarlavha (48pt, EFECE0, center)
    # shapes[1] = Freeform blip (rasm, chap)
    # shapes[2] = matn (28pt, EFECE0, justify) - o'ng
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t26_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t26_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 0, 'color': 'EFECE0', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t26_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 0, 'color': 'EFECE0', 'text': body_text}]}
        ])
    if img_arg and len(slide.shapes) > 1:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t26_replace_blip(slide, 1, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T26] Slayd6 rasm: {e}")

def fill_t26_slide_7_img_bottom_right_title_three_text(slide, data, img_arg=None):
    # shapes[0] = sarlavha (48pt, EFECE0, center)
    # shapes[1] = matn_pastki_chap (25pt, EFECE0, justify)
    # shapes[2] = matn_yuqori_chap (25pt, EFECE0, justify)
    # shapes[3] = matn_yuqori_ong (25pt, EFECE0, justify)
    # shapes[4] = Freeform blip (rasm, pastki o'ng)
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t26_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 3)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    text3 = blocks[2] if len(blocks) > 2 else ""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t26_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 0, 'color': 'EFECE0', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t26_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2500, 'b': 0, 'color': 'EFECE0', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t26_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2500, 'b': 0, 'color': 'EFECE0', 'text': text2}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t26_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2500, 'b': 0, 'color': 'EFECE0', 'text': text3}]}
        ])
    if img_arg and len(slide.shapes) > 4:
        try:
            import os
            img_path = img_arg if (isinstance(img_arg, str) and os.path.exists(img_arg)) else fetch_image(img_arg)
            if img_path:
                _t26_replace_blip(slide, 4, img_path)
        except Exception as e:
            import logging
            logging.warning(f"[T26] Slayd7 rasm: {e}")

def fill_t26_slide_8_conclusion(slide, data):
    pass

def generate_template_26_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    import logging
    slides = prs.slides
    if len(slides) < 2:
        logging.error("[T26] Shablon slaydlari yetarli emas")
        return None
    fill_t26_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t26_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t26_slide_3_img_right_title_text,
        fill_t26_slide_4_img_right_title_text2,
        fill_t26_slide_5_img_left_title_two_text,
        fill_t26_slide_6_img_left_title_text,
        fill_t26_slide_7_img_bottom_right_title_three_text,
    ]
    # Barcha kontent slaydlarda Freeform blip rasm bor - har doim rasm olinadi
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        slide_type = i % len(content_slide_funcs)
        # Har doim rasm olish - barcha kontent slaydlarda blip rasm bor
        if user_images and i < len(user_images):
            img_arg = user_images[i]
        else:
            img_query = data.get("title", topic) if isinstance(data, dict) else topic
            img_arg = fetch_image(img_query + " science")
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T26] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t26_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

SLIDE_TYPE_NAMES_T27 = {
    "cover": "Muqova",
    "plan": "Reja",
    "img_left_title_text": "Rasm chap, sarlavha va matn o'ng",
    "img_right_title_two_text": "Rasm o'ng, sarlavha va ikki matn",
    "img_left_title_two_text": "Rasm chap, sarlavha va uch matn",
    "title_two_text": "Sarlavha markazda, ikki ustun matn",
    "title_text": "Sarlavha markazda, katta matn",
    "conclusion": "Xulosa",
}

# 27-SHABLON (Minimalist Blue) funksiyalari
# Fon: oq, Matn: to'q ko'k 01324F
# 8 slayd: muqova, reja, 5 kontent, xulosa
# Slayd 3,4,5 da Freeform blip rasm bor

def _t27_clear_and_write(txBody, paras_data):
    from lxml import etree
    for child in list(txBody):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'p':
            txBody.remove(child)
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for pd in paras_data:
        algn = pd.get('algn', 'l')
        spcPts = pd.get('spcPts', 0)
        runs = pd.get('runs', [])
        pPr_xml = f'<a:pPr xmlns:a="{ns}" algn="{algn}"'
        if spcPts:
            pPr_xml += f'><a:spcBef><a:spcPts val="{spcPts}"/></a:spcBef></a:pPr>'
        else:
            pPr_xml += '/>'
        runs_xml = ''
        for run in runs:
            sz = run.get('sz', 2400)
            b = run.get('b', 0)
            color = run.get('color', '01324F')
            text = run.get('text', '')
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
            runs_xml += f'''<a:r xmlns:a="{ns}">
  <a:rPr lang="uz-UZ" sz="{sz}" b="{b}" dirty="0">
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:latin typeface="+mj-lt"/>
  </a:rPr>
  <a:t>{text}</a:t>
</a:r>'''
        p_xml = f'<a:p xmlns:a="{ns}">{pPr_xml}{runs_xml}</a:p>'
        p_elem = etree.fromstring(p_xml)
        txBody.append(p_elem)


def _t27_replace_freeform_with_picture(slide, img_path):
    """Freeform blip shapeni topib, uni o'chirib o'rniga add_picture bilan rasm qo'yish."""
    import logging
    import os
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    if not img_path or not os.path.exists(img_path):
        logging.warning(f"[T27] Rasm fayli topilmadi: {img_path}")
        return False
    try:
        target_shape = None
        for shape in slide.shapes:
            blips = shape._element.findall(f'.//{{{ns_a}}}blip')
            if blips:
                target_shape = shape
                break
        if target_shape is None:
            # blip yo'q bo'lsa, birinchi Freeform ni topamiz
            for shape in slide.shapes:
                if shape.shape_type == 5:  # FREEFORM
                    target_shape = shape
                    break
        if target_shape is None:
            logging.warning("[T27] Slaydda Freeform topilmadi")
            return False
        left = target_shape.left
        top = target_shape.top
        width = target_shape.width
        height = target_shape.height
        sp = target_shape._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(img_path, left, top, width, height)
        logging.info(f"[T27] Freeform o'chirildi, rasm qo'yildi: {img_path}")
        return True
    except Exception as e:
        logging.warning(f"[T27] Freeform almashtirish xatoligi: {e}")
    return False


def _t27_get_body_text(data):
    content = data.get("content", [])
    if isinstance(content, list):
        body_text = " ".join(str(c) for c in content if c)
    else:
        body_text = str(content) if content else ""
    if not body_text:
        body_text = data.get("text", "")
    return body_text


def fill_t27_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (66pt, 86AADF, center)
    # shapes[1] = ism (27pt, 01324F)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t27_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6600, 'b': 0, 'color': '86AADF', 'text': topic.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t27_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2700, 'b': 0, 'color': '01324F', 'text': name_surname}]}
        ])


def fill_t27_slide_2_plan(slide, plan_dict):
    # shapes[0] = "Reja" sarlavha (113pt, 01324F, center)
    # shapes[1] = reja elementlari (27pt, 01324F)
    if not isinstance(plan_dict, dict):
        plan_dict = {}
    plan_title = plan_dict.get("title", "Reja")
    items = plan_dict.get("content", [])
    if isinstance(items, str):
        items = [items]
    clean_items = [str(i).strip() for i in items if i]
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t27_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 11300, 'b': 0, 'color': '01324F', 'text': plan_title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        paras = []
        for idx, item in enumerate(clean_items[:7], 1):
            # AI allaqachon "1. Mavzu" formatida berishi mumkin, shuning uchun raqamni olib tashlaymiz
            import re
            clean_item = re.sub(r'^\d+[\.\.\s]+', '', item).strip()
            paras.append({
                'algn': 'l',
                'spcPts': 150,
                'runs': [{'sz': 2700, 'b': 0, 'color': '01324F', 'text': f"{idx}. {clean_item}"}]
            })
        if paras:
            _t27_clear_and_write(slide.shapes[1].text_frame._txBody, paras)


def fill_t27_slide_3_img_left_title_text(slide, data, img_arg=None):
    # shapes[0] = Freeform blip (rasm, chap) L=1" T=3" W=4" H=5"
    # shapes[1] = matn (28pt, 01324F, justify) o'ng
    # shapes[2] = sarlavha (48pt, 01324F, left) yuqori
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t27_get_body_text(data)
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t27_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 0, 'color': '01324F', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t27_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2800, 'b': 0, 'color': '01324F', 'text': body_text}]}
        ])
    if img_arg:
        import os
        import logging
        try:
            if isinstance(img_arg, (bytes, bytearray)):
                img_path = save_user_image_to_tmp(img_arg)
            elif isinstance(img_arg, str) and os.path.exists(img_arg):
                img_path = img_arg
            else:
                img_path = fetch_image(img_arg if isinstance(img_arg, str) else title or "science")
            if img_path and os.path.exists(img_path):
                _t27_replace_freeform_with_picture(slide, img_path)
            else:
                logging.warning(f"[T27] Slayd3 rasm fayli yo'q: {img_path}")
        except Exception as e:
            logging.warning(f"[T27] Slayd3 rasm: {e}")


def fill_t27_slide_4_img_right_title_two_text(slide, data, img_arg=None):
    # shapes[0] = Freeform blip (rasm, o'ng) L=14" T=0" W=2" H=4"
    # shapes[1] = sarlavha (48pt, 01324F, center) yuqori
    # shapes[2] = matn1 (24pt, 01324F, justify) chap pastki
    # shapes[3] = matn2 (24pt, 01324F, justify) o'ng pastki
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t27_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t27_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 0, 'color': '01324F', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t27_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2400, 'b': 0, 'color': '01324F', 'text': text1}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t27_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2400, 'b': 0, 'color': '01324F', 'text': text2}]}
        ])
    if img_arg:
        import os
        import logging
        try:
            if isinstance(img_arg, (bytes, bytearray)):
                img_path = save_user_image_to_tmp(img_arg)
            elif isinstance(img_arg, str) and os.path.exists(img_arg):
                img_path = img_arg
            else:
                img_path = fetch_image(img_arg if isinstance(img_arg, str) else title or "science")
            if img_path and os.path.exists(img_path):
                _t27_replace_freeform_with_picture(slide, img_path)
            else:
                logging.warning(f"[T27] Slayd4 rasm fayli yo'q: {img_path}")
        except Exception as e:
            logging.warning(f"[T27] Slayd4 rasm: {e}")


def fill_t27_slide_5_img_left_title_two_text(slide, data, img_arg=None):
    # shapes[0] = Freeform blip (rasm, chap) L=1" T=0" W=5" H=6"
    # shapes[1] = sarlavha (44pt, 01324F, left) pastki chap
    # shapes[2] = matn1 (24pt, 01324F, justify) pastki chap
    # shapes[3] = matn2 (22pt, 01324F, justify) o'ng yuqori
    # shapes[4] = matn3 (22pt, 01324F, justify) o'ng pastki
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t27_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 3)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    text3 = blocks[2] if len(blocks) > 2 else ""
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t27_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 0, 'color': '01324F', 'text': title}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t27_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2400, 'b': 0, 'color': '01324F', 'text': text1}]}
        ])
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        _t27_clear_and_write(slide.shapes[3].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2200, 'b': 0, 'color': '01324F', 'text': text2}]}
        ])
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        _t27_clear_and_write(slide.shapes[4].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2200, 'b': 0, 'color': '01324F', 'text': text3}]}
        ])
    if img_arg:
        import os
        import logging
        try:
            if isinstance(img_arg, (bytes, bytearray)):
                img_path = save_user_image_to_tmp(img_arg)
            elif isinstance(img_arg, str) and os.path.exists(img_arg):
                img_path = img_arg
            else:
                img_path = fetch_image(img_arg if isinstance(img_arg, str) else title or "science")
            if img_path and os.path.exists(img_path):
                _t27_replace_freeform_with_picture(slide, img_path)
            else:
                logging.warning(f"[T27] Slayd5 rasm fayli yo'q: {img_path}")
        except Exception as e:
            logging.warning(f"[T27] Slayd5 rasm: {e}")


def fill_t27_slide_6_title_two_text(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt, 01324F, center) yuqori
    # shapes[1] = matn1 (22pt, 01324F, justify) chap
    # shapes[2] = matn2 (22pt, 01324F, justify) o'ng
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t27_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t27_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 0, 'color': '01324F', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t27_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2200, 'b': 0, 'color': '01324F', 'text': text1}]}
        ])
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t27_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2200, 'b': 0, 'color': '01324F', 'text': text2}]}
        ])


def fill_t27_slide_7_title_text(slide, data, img_arg=None):
    # shapes[0] = sarlavha (54pt, 01324F, center) yuqori
    # shapes[1] = matn (28pt, 01324F, justify) markazda
    if not isinstance(data, dict):
        data = {}
    title = data.get("title", "")
    body_text = _t27_get_body_text(data)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t27_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 5400, 'b': 0, 'color': '01324F', 'text': title}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t27_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 2800, 'b': 0, 'color': '01324F', 'text': body_text}]}
        ])


def fill_t27_slide_8_conclusion(slide, data):
    # shapes[0] = "E'TIBORINGIZ UCHUN RAHMAT!" (96pt, 01324F, center)
    pass


CONTENT_SLIDE_TEMPLATE_INDICES_27 = [2, 3, 4, 5, 6]  # 27-shablondagi 3-7 slaydlar


def build_slide_structure_27(prs, requested_content_count):
    """27-shablon uchun slayd tuzilmasini quradi. 3-7 slaydlar takrorlanadi, 8-slayd oxirida."""
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5
    logging.info(f"[T27] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_27:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T27] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T27] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides


def generate_template_27_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    import logging
    import os
    if len(prs.slides) < 2:
        logging.error("[T27] Shablon slaydlari yetarli emas")
        return None
    # ── 1. Slayd tuzilmasini qurish (requested_slide_count bo'yicha ko'paytirish) ──
    total_content_slides = build_slide_structure_27(prs, requested_slide_count)
    # content_data_list ni to'ldirish yoki kengaytirish
    if not content_data_list:
        content_data_list = []
    while len(content_data_list) < total_content_slides:
        idx = len(content_data_list)
        content_data_list.append({"title": f"{topic} — {idx + 1}", "content": [topic]})
    slides = prs.slides
    fill_t27_slide_1_cover(slides[0], topic, name_surname)
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t27_slide_2_plan(slides[1], plan_dict)
    content_slide_funcs = [
        fill_t27_slide_3_img_left_title_text,
        fill_t27_slide_4_img_right_title_two_text,
        fill_t27_slide_5_img_left_title_two_text,
        fill_t27_slide_6_title_two_text,
        fill_t27_slide_7_title_text,
    ]
    # Slayd 3, 4, 5 da rasm bor (index 0, 1, 2)
    img_slide_types = {0, 1, 2}
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
        slide_type = i % len(content_slide_funcs)
        img_arg = None
        if slide_type in img_slide_types:
            if user_images and i < len(user_images):
                raw = user_images[i]
                if isinstance(raw, (bytes, bytearray)):
                    img_arg = save_user_image_to_tmp(raw)
                elif isinstance(raw, str) and os.path.exists(raw):
                    img_arg = raw
                else:
                    img_query = data.get("title", topic) if isinstance(data, dict) else topic
                    img_arg = fetch_image(img_query) or fetch_image(topic)
            else:
                img_query = data.get("title", topic) if isinstance(data, dict) else topic
                img_arg = fetch_image(img_query) or fetch_image(topic)
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T27] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
    fill_t27_slide_8_conclusion(slides[-1], {})
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ============================================================
# 28-SHABLON (Zamonaviy PPT Dizayni 2.0)
# Fon: to'q ko'k, Matn: Oq (F8FAFC), Moviy (CBD5E1, 818CF8, 22D3EE)
# 8 slayd: muqova, reja, 5 kontent, xulosa
# Slayd 4,5,6 da rasm bor (PICTURE shapes)

CONTENT_SLIDE_TEMPLATE_INDICES_28 = [2, 3, 4, 5, 6]

def build_slide_structure_28(prs, requested_content_count):
    """28-shablon uchun slayd tuzilmasini quradi. 3-7 slaydlar takrorlanadi, 8-slayd oxirida."""
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5
    logging.info(f"[T28] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_28:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T28] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T28] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides

def _t28_clear_and_write(txBody, paras_data):
    from lxml import etree
    for child in list(txBody):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'p':
            child.getparent().remove(child)
    for p_data in paras_data:
        p_elem = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}p')
        pPr = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        algn = p_data.get('algn', 'l')
        if algn != 'l':
            pPr.set('algn', algn)
        p_elem.append(pPr)
        for run_data in p_data.get('runs', []):
            r_elem = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}r')
            rPr = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            rPr.set('lang', 'en-US')
            if run_data.get('sz'):
                rPr.set('sz', str(run_data['sz']))
            if run_data.get('b') is not None:
                rPr.set('b', str(run_data['b']))
            if run_data.get('color'):
                solidFill = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                srgbClr = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                srgbClr.set('val', run_data['color'])
                solidFill.append(srgbClr)
                rPr.append(solidFill)
            r_elem.append(rPr)
            t_elem = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}t')
            t_elem.text = run_data.get('text', '')
            r_elem.append(t_elem)
            p_elem.append(r_elem)
        txBody.append(p_elem)

def _t28_get_body_text(data):
    if isinstance(data, dict):
        content = data.get("content", [])
        if isinstance(content, list):
            return " ".join(str(x) for x in content)
        return str(content)
    return str(data)

def _t28_replace_picture(slide, shape_idx, img_path):
    """Slayddagi rasmni (PICTURE) yangisi bilan almashtiradi."""
    import os
    import logging
    if not img_path or not os.path.exists(img_path):
        return False
    try:
        if shape_idx >= len(slide.shapes):
            return False
        target_shape = slide.shapes[shape_idx]
        left = target_shape.left
        top = target_shape.top
        width = target_shape.width
        height = target_shape.height
        
        # Eski rasmni o'chirish
        sp = target_shape._element
        sp.getparent().remove(sp)
        
        # Yangi rasmni qo'shish
        slide.shapes.add_picture(img_path, left, top, width, height)
        logging.info(f"[T28] Rasm almashtirildi: {img_path}")
        return True
    except Exception as e:
        logging.error(f"[T28] Rasm almashtirish xatosi: {e}")
        return False

def fill_t28_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (54pt bold F8FAFC center)
    # shapes[1] = subtitle (16pt CBD5E1 center)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t28_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 5400, 'b': 1, 'color': 'F8FAFC', 'text': str(topic).upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        sub = name_surname if name_surname else "Zamonaviy taqdimot"
        _t28_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 1600, 'b': 0, 'color': 'CBD5E1', 'text': str(sub)}]}
        ])

def fill_t28_slide_2_plan(slide, plan_data):
    # shapes[0] = sarlavha (54pt bold F8FAFC center)
    # shapes[1] = reja_matn (22pt 818CF8)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t28_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 5400, 'b': 1, 'color': 'F8FAFC', 'text': "REJA"}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        content = plan_data.get("content", []) if isinstance(plan_data, dict) else []
        paras = []
        import re as _re
        for idx, item in enumerate(content):
            # item da allaqachon raqam bo'lishi mumkin (masalan "1. Mavzu") - tozalash
            clean_item = _re.sub(r'^\d+[.)\s]+', '', str(item).strip())
            paras.append({'algn': 'l', 'runs': [{'sz': 2200, 'b': 0, 'color': '818CF8', 'text': f"{idx+1}. {clean_item}"}]})
        _t28_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t28_slide_3_two_col(slide, data, img_arg=None):
    # shapes[2] = sarlavha (32pt bold F8FAFC left)
    # shapes[0] = matn_chap (16pt CBD5E1)
    # shapes[1] = matn_ong (16pt CBD5E1)
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t28_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        _t28_clear_and_write(slide.shapes[2].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3200, 'b': 1, 'color': 'F8FAFC', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t28_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1600, 'b': 0, 'color': 'CBD5E1', 'text': text1}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t28_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1600, 'b': 0, 'color': 'CBD5E1', 'text': text2}]}
        ])

def fill_t28_slide_4_img_left(slide, data, img_arg=None):
    # shapes[0] = PICTURE (chap), shapes[1] = matn, shapes[2] = sarlavha
    # MUHIM: avval rasmni almashtirish, keyin matnlarni yozish
    # (rasm almashtirish shapes tartibini o'zgartiradi)
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t28_get_body_text(data)
    # Avval matnlarni yozish (shapes tartibini bilgan holda)
    # shapes[2] = sarlavha, shapes[1] = matn
    title_shape = None
    text_shape = None
    pic_shape_idx = None
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            pic_shape_idx = i
        elif hasattr(s, 'has_text_frame') and s.has_text_frame:
            txt = s.text_frame.text.strip()
            if s.left < 7000000:  # chap tomonda
                if title_shape is None or s.top < title_shape.top:
                    if s.width > 4000000:  # keng shape = sarlavha
                        title_shape = s
                    else:
                        text_shape = s
            else:  # o'ng tomonda
                if title_shape is None:
                    title_shape = s
                else:
                    text_shape = s
    # Oddiy indeks usuli (ishonchli)
    if pic_shape_idx is not None and img_arg:
        _t28_replace_picture(slide, pic_shape_idx, img_arg)
    # Rasm almashtirilgandan keyin shapes yangilandi - TOP koordinatasi bo'yicha topish
    # Sarlavha = eng tepada turgan matn shape (min top)
    # Matn = qolgan matn shape
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top, i))
    text_shapes.sort(key=lambda x: x[0])  # top bo'yicha saralash
    sarlavha_idx = text_shapes[0][1] if len(text_shapes) > 0 else None
    matn_idx = text_shapes[1][1] if len(text_shapes) > 1 else None
    if sarlavha_idx is not None and slide.shapes[sarlavha_idx].has_text_frame:
        _t28_clear_and_write(slide.shapes[sarlavha_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3200, 'b': 1, 'color': 'F8FAFC', 'text': title.upper()}]}
        ])
    if matn_idx is not None and slide.shapes[matn_idx].has_text_frame:
        _t28_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1600, 'b': 0, 'color': 'CBD5E1', 'text': body_text}]}
        ])

def fill_t28_slide_5_img_right(slide, data, img_arg=None):
    # shapes[0] = matn (chap), shapes[1] = PICTURE (ong), shapes[2] = sarlavha
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t28_get_body_text(data)
    # Avval rasm indeksini topish
    pic_shape_idx = None
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            pic_shape_idx = i
            break
    if pic_shape_idx is not None and img_arg:
        _t28_replace_picture(slide, pic_shape_idx, img_arg)
    # TOP koordinatasi bo'yicha sarlavha va matn topish
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top, i))
    text_shapes.sort(key=lambda x: x[0])
    sarlavha_idx = text_shapes[0][1] if len(text_shapes) > 0 else None
    matn_idx = text_shapes[1][1] if len(text_shapes) > 1 else None
    if sarlavha_idx is not None and slide.shapes[sarlavha_idx].has_text_frame:
        _t28_clear_and_write(slide.shapes[sarlavha_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3200, 'b': 1, 'color': 'F8FAFC', 'text': title.upper()}]}
        ])
    if matn_idx is not None and slide.shapes[matn_idx].has_text_frame:
        _t28_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1600, 'b': 0, 'color': 'CBD5E1', 'text': body_text}]}
        ])

def fill_t28_slide_6_img_right_large(slide, data, img_arg=None):
    # shapes[0] = sarlavha, shapes[1] = matn, shapes[2] = PICTURE (ong katta)
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t28_get_body_text(data)
    # Avval rasm indeksini topish
    pic_shape_idx = None
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            pic_shape_idx = i
            break
    if pic_shape_idx is not None and img_arg:
        _t28_replace_picture(slide, pic_shape_idx, img_arg)
    # Matn shapelari
    sarlavha_idx = None
    matn_idx = None
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            if sarlavha_idx is None:
                sarlavha_idx = i
            else:
                matn_idx = i
    if sarlavha_idx is not None and slide.shapes[sarlavha_idx].has_text_frame:
        _t28_clear_and_write(slide.shapes[sarlavha_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2800, 'b': 1, 'color': 'F8FAFC', 'text': title.upper()}]}
        ])
    if matn_idx is not None and slide.shapes[matn_idx].has_text_frame:
        _t28_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1600, 'b': 0, 'color': 'CBD5E1', 'text': body_text}]}
        ])

def fill_t28_slide_7_quote(slide, data, img_arg=None):
    # shapes[1] = muallif (24pt bold 22D3EE center)
    # shapes[0] = iqtibos (20pt F8FAFC center)
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t28_get_body_text(data)
    
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t28_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 2400, 'b': 1, 'color': '22D3EE', 'text': title}]}
        ])
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t28_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 2000, 'b': 0, 'color': 'F8FAFC', 'text': body_text}]}
        ])

def fill_t28_slide_8_conclusion(slide, data):
    # shapes[0] = xulosa_matn (60pt bold 22D3EE center)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t28_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6000, 'b': 1, 'color': '22D3EE', 'text': "E'TIBORINGIZ UCHUN RAHMAT!"}]}
        ])

def generate_template_28_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    import logging
    import os
    if len(prs.slides) < 2:
        logging.error("[T28] Shablon slaydlari yetarli emas")
        return None
        
    total_content_slides = build_slide_structure_28(prs, requested_slide_count)
    
    if not content_data_list:
        content_data_list = []
    while len(content_data_list) < total_content_slides:
        idx = len(content_data_list)
        content_data_list.append({"title": f"{topic} — {idx + 1}", "content": [topic]})
        
    slides = prs.slides
    fill_t28_slide_1_cover(slides[0], topic, name_surname)
    
    plan_dict = plan if isinstance(plan, dict) else {}
    if not plan_dict and content_data_list:
        titles = [d.get("title", "") for d in content_data_list if isinstance(d, dict)]
        plan_dict = {"title": "Reja", "content": titles}
    fill_t28_slide_2_plan(slides[1], plan_dict)
    
    content_slide_funcs = [
        fill_t28_slide_3_two_col,
        fill_t28_slide_4_img_left,
        fill_t28_slide_5_img_right,
        fill_t28_slide_6_img_right_large,
        fill_t28_slide_7_quote,
    ]
    
    # Slayd 4, 5, 6 da rasm bor (index 1, 2, 3)
    img_slide_types = {1, 2, 3}
    
    img_counter = 0  # Rasm slaydlari uchun alohida hisoblagich
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
            
        slide_type = i % len(content_slide_funcs)
        img_arg = None
        
        if slide_type in img_slide_types:
            if user_images and img_counter < len(user_images):
                raw = user_images[img_counter]
                if isinstance(raw, (bytes, bytearray)):
                    img_arg = save_user_image_to_tmp(raw)
                elif isinstance(raw, str) and os.path.exists(raw):
                    img_arg = raw
                else:
                    img_query = data.get("title", topic) if isinstance(data, dict) else topic
                    img_arg = fetch_image(img_query) or fetch_image(topic)
            else:
                img_query = data.get("title", topic) if isinstance(data, dict) else topic
                img_arg = fetch_image(img_query) or fetch_image(topic)
            img_counter += 1  # Har bir rasm slaydida hisoblagichni oshirish
                
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T28] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
        
    fill_t28_slide_8_conclusion(slides[-1], {})
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

SLIDE_TYPE_NAMES_T28 = {
    "cover": "Muqova",
    "plan": "Reja",
    "two_col": "Ikki ustunli matn",
    "img_left": "Rasm chap, sarlavha va matn o'ng",
    "img_right": "Rasm o'ng, sarlavha va matn chap",
    "img_right_large": "Katta rasm o'ng, sarlavha va matn chap",
    "quote": "Iqtibos va muallif",
    "conclusion": "Xulosa",
}
def _t29_clear_and_write(txBody, paras_data):
    from lxml import etree
    for child in list(txBody):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'p':
            txBody.remove(child)
    
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_data in paras_data:
        p = etree.SubElement(txBody, f"{{{ns_a}}}p")
        pPr = etree.SubElement(p, f"{{{ns_a}}}pPr")
        if 'algn' in p_data:
            pPr.set('algn', p_data['algn'])
        # Bullet/numbering ni o'chirish - har doim buNone qo'shish
        etree.SubElement(pPr, f"{{{ns_a}}}buNone")
        for run_data in p_data.get('runs', []):
            r = etree.SubElement(p, f"{{{ns_a}}}r")
            rPr = etree.SubElement(r, f"{{{ns_a}}}rPr")
            if 'sz' in run_data:
                rPr.set('sz', str(run_data['sz']))
            if 'b' in run_data:
                rPr.set('b', str(run_data['b']))
            if 'color' in run_data:
                solidFill = etree.SubElement(rPr, f"{{{ns_a}}}solidFill")
                srgbClr = etree.SubElement(solidFill, f"{{{ns_a}}}srgbClr")
                srgbClr.set('val', run_data['color'])
            t = etree.SubElement(r, f"{{{ns_a}}}t")
            t.text = run_data.get('text', '')

def _t29_get_body_text(data):
    if isinstance(data, dict):
        content = data.get("content", [])
        if isinstance(content, list):
            return " ".join([str(c) for c in content])
        return str(content)
    return str(data)

def _t29_replace_picture(slide, shape_idx, img_path):
    import os
    import logging
    if not img_path or not os.path.exists(img_path):
        return False
    try:
        if shape_idx >= len(slide.shapes):
            return False
        target_shape = slide.shapes[shape_idx]
        left = target_shape.left
        top = target_shape.top
        width = target_shape.width
        height = target_shape.height
        
        sp = target_shape._element
        sp.getparent().remove(sp)
        
        slide.shapes.add_picture(img_path, left, top, width, height)
        logging.info(f"[T29] Rasm almashtirildi: {img_path}")
        return True
    except Exception as e:
        logging.error(f"[T29] Rasm almashtirish xatosi: {e}")
        return False

def fill_t29_slide_1_cover(slide, topic, name_surname):
    # shapes[0] = sarlavha (54pt bold FFFFFF left)
    # shapes[1] = subtitle (18pt bold FFFFFF left)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t29_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 5400, 'b': 1, 'color': 'FFFFFF', 'text': str(topic).upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        sub = f"Muallif: {name_surname}" if name_surname else "Zamonaviy taqdimot"
        _t29_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 1, 'color': 'FFFFFF', 'text': str(sub)}]}
        ])

def fill_t29_slide_2_plan(slide, plan_data):
    # shapes[0] = sarlavha (54pt bold 000000 center)
    # shapes[1] = reja matn (24pt 000000 left)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t29_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 5400, 'b': 1, 'color': '000000', 'text': "REJA"}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        content = plan_data.get("content", []) if isinstance(plan_data, dict) else []
        paras = []
        import re as _re
        for idx, item in enumerate(content):
            clean_item = _re.sub(r'^\d+[.)\s]+', '', str(item).strip())
            paras.append({'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '000000', 'text': f"{idx+1}. {clean_item}"}]})
        _t29_clear_and_write(slide.shapes[1].text_frame._txBody, paras)

def fill_t29_slide_3_title_text(slide, data, img_arg=None):
    # shapes[0] = sarlavha (44pt bold 000000 left)
    # shapes[1] = matn (20pt 000000 left)
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t29_get_body_text(data)
    
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t29_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        _t29_clear_and_write(slide.shapes[1].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t29_slide_4_two_col(slide, data, img_arg=None):
    # Slayd 4: chap blok = SARLAVHA, o'ng blok = ASOSIY MATN
    # left koordinatasi bo'yicha: eng chapda = sarlavha, o'ngda = matn
    body_text = _t29_get_body_text(data)
    title = data.get("title", "") if isinstance(data, dict) else ""
    
    # Matn shapelari ni left koordinatasi bo'yicha saralash
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.left or 0, i))
    text_shapes.sort(key=lambda x: x[0])  # chapdan o'ngga
    
    if len(text_shapes) >= 2:
        # Chap blok = sarlavha
        left_idx = text_shapes[0][1]
        # O'ng blok = asosiy matn
        right_idx = text_shapes[1][1]
        _t29_clear_and_write(slide.shapes[left_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
        _t29_clear_and_write(slide.shapes[right_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])
    elif len(text_shapes) == 1:
        idx = text_shapes[0][1]
        _t29_clear_and_write(slide.shapes[idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])

def fill_t29_slide_5_img_left(slide, data, img_arg=None):
    # shapes[0] = sarlavha (36pt bold 000000 left)
    # shapes[1] = matn (18pt 000000 left)
    # shapes[2] = PICTURE (chap)
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t29_get_body_text(data)
    
    pic_shape_idx = None
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            pic_shape_idx = i
            break
    if pic_shape_idx is not None and img_arg:
        _t29_replace_picture(slide, pic_shape_idx, img_arg)
        
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top, i))
    text_shapes.sort(key=lambda x: x[0])
    sarlavha_idx = text_shapes[0][1] if len(text_shapes) > 0 else None
    matn_idx = text_shapes[1][1] if len(text_shapes) > 1 else None
    
    if sarlavha_idx is not None and slide.shapes[sarlavha_idx].has_text_frame:
        _t29_clear_and_write(slide.shapes[sarlavha_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
    if matn_idx is not None and slide.shapes[matn_idx].has_text_frame:
        _t29_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t29_slide_6_img_right(slide, data, img_arg=None):
    # shapes: sarlavha + matn + PICTURE (ong) yoki sarlavha + 2 matn
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t29_get_body_text(data)
    
    pic_shape_idx = None
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            pic_shape_idx = i
            break
    if pic_shape_idx is not None and img_arg:
        _t29_replace_picture(slide, pic_shape_idx, img_arg)
        
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        if blips:
            continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else body_text
    text2 = blocks[1] if len(blocks) > 1 else ""
    
    if len(text_shapes) >= 3:
        # sarlavha + chap matn + o'ng matn
        sarlavha_idx = text_shapes[0][1]
        left_idx = text_shapes[1][1]
        right_idx = text_shapes[2][1]
        _t29_clear_and_write(slide.shapes[sarlavha_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
        _t29_clear_and_write(slide.shapes[left_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': text1}]}
        ])
        _t29_clear_and_write(slide.shapes[right_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': text2}]}
        ])
    elif len(text_shapes) == 2:
        # sarlavha + matn
        sarlavha_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t29_clear_and_write(slide.shapes[sarlavha_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
        _t29_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])
    elif len(text_shapes) == 1:
        sarlavha_idx = text_shapes[0][1]
        _t29_clear_and_write(slide.shapes[sarlavha_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])

def fill_t29_slide_7_conclusion(slide, data):
    # shapes[0] = xulosa_matn (60pt bold FFFFFF center)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t29_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6000, 'b': 1, 'color': 'FFFFFF', 'text': "E'TIBORINGIZ UCHUN RAHMAT!"}]}
        ])

def build_slide_structure_29(prs, requested_slide_count):
    import logging
    from utils import duplicate_slide, move_slide
    if len(prs.slides) < 8:
        logging.error("[T29] Shablon slaydlari yetarli emas")
        return
    
    content_slides_count = requested_slide_count
    
    # 29-shablonda kontent slaydlar: 2, 3, 4, 5, 6 (index)
    template_indices = [2, 3, 4, 5, 6]
    base_count = len(template_indices)
    
    if content_slides_count <= base_count:
        return
    
    sets_needed = (content_slides_count - 1) // base_count
    logging.info(f"[T29] Kontent slaydlari: {content_slides_count} so'raldi, {sets_needed} marta takrorlanadi")
    
    for s_idx in range(sets_needed):
        for idx in template_indices:
            new_slide = duplicate_slide(prs, idx)
            # Yangi slaydni xulosadan oldinga qo'yish
            move_slide(prs, len(prs.slides) - 1, len(prs.slides) - 2)
        logging.info(f"  [T29] {s_idx+1}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
        
    logging.info(f"[T29] Yakuniy tuzilma: {len(prs.slides)} ta slayd")

def generate_template_29_presentation(prs, topic, requested_slide_count, language,
                                       name_surname="", plan=None, content_data_list=None,
                                       user_images=None):
    import io
    import logging
    import os
    if len(prs.slides) < 8:
        logging.error("[T29] Shablon slaydlari yetarli emas")
        return None
        
    build_slide_structure_29(prs, requested_slide_count)
    
    slides = prs.slides
    fill_t29_slide_1_cover(slides[0], topic, name_surname)
    
    plan_dict = plan if isinstance(plan, dict) else {"title": "Reja", "content": []}
    fill_t29_slide_2_plan(slides[1], plan_dict)
    
    content_slide_funcs = [
        fill_t29_slide_3_title_text,
        fill_t29_slide_4_two_col,
        fill_t29_slide_5_img_left,
        fill_t29_slide_6_img_right
    ]
    
    # Slayd 5 (index 4) va 7 (index 6) da rasm bor -> slide_type 2 va 3
    img_slide_types = {2, 3}
    
    img_counter = 0
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(slides) - 1:
            break
        slide = slides[slide_index]
        if not isinstance(data, dict):
            data = {"title": str(data)[:80] if data else "", "content": [str(data)] if data else []}
            
        slide_type = i % len(content_slide_funcs)
        img_arg = None
        
        if slide_type in img_slide_types:
            if user_images and img_counter < len(user_images):
                raw = user_images[img_counter]
                if isinstance(raw, (bytes, bytearray)):
                    img_arg = save_user_image_to_tmp(raw)
                elif isinstance(raw, str) and os.path.exists(raw):
                    img_arg = raw
                else:
                    img_query = data.get("title", topic) if isinstance(data, dict) else topic
                    img_arg = fetch_image(img_query) or fetch_image(topic)
            else:
                img_query = data.get("title", topic) if isinstance(data, dict) else topic
                img_arg = fetch_image(img_query) or fetch_image(topic)
            img_counter += 1
                
        content_slide_funcs[slide_type](slide, data, img_arg)
        logging.info(f"  [T29] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
        
    fill_t29_slide_7_conclusion(slides[-1], {})
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

SLIDE_TYPE_NAMES_T29 = {
    "cover": "Muqova",
    "plan": "Reja",
    "title_text": "Sarlavha va matn",
    "two_col": "Ikki ustunli matn",
    "img_left": "Rasm chap, sarlavha va matn o'ng",
    "img_right": "Rasm o'ng, sarlavha va matn chap",
    "conclusion": "Xulosa",
}
# === 30-shablon ===
SLIDE_TYPE_NAMES_T30 = {
    0: "two_col",
    1: "title_text",
    2: "img_right",
    3: "three_col",
    4: "img_left"
}

def _t30_clear_and_write(txBody, paras_data):
    from lxml import etree
    for child in list(txBody):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'p':
            txBody.remove(child)
    
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for p_data in paras_data:
        p = etree.SubElement(txBody, f"{{{ns_a}}}p")
        pPr = etree.SubElement(p, f"{{{ns_a}}}pPr")
        if 'algn' in p_data:
            pPr.set('algn', p_data['algn'])
        etree.SubElement(pPr, f"{{{ns_a}}}buNone")
        for run_data in p_data.get('runs', []):
            r = etree.SubElement(p, f"{{{ns_a}}}r")
            rPr = etree.SubElement(r, f"{{{ns_a}}}rPr")
            if 'sz' in run_data:
                rPr.set('sz', str(run_data['sz']))
            if 'b' in run_data:
                rPr.set('b', str(run_data['b']))
            if 'color' in run_data:
                solidFill = etree.SubElement(rPr, f"{{{ns_a}}}solidFill")
                srgbClr = etree.SubElement(solidFill, f"{{{ns_a}}}srgbClr")
                srgbClr.set('val', run_data['color'])
            t = etree.SubElement(r, f"{{{ns_a}}}t")
            t.text = run_data.get('text', '')

def _t30_get_body_text(data):
    if isinstance(data, dict):
        c = data.get("content", [])
        if isinstance(c, list):
            return "\n".join(str(x) for x in c)
        return str(c)
    return str(data)

def _t30_replace_picture(slide, shape_idx, img_arg):
    import os
    from pptx.util import Inches
    if not img_arg:
        return
    if isinstance(img_arg, str) and not os.path.exists(img_arg):
        return
    old_shape = slide.shapes[shape_idx]
    left, top = old_shape.left, old_shape.top
    width, height = old_shape.width, old_shape.height
    try:
        spTree = slide.shapes._spTree
        spTree.remove(old_shape._element)
        slide.shapes.add_picture(img_arg, left, top, width, height)
    except Exception as e:
        import logging
        logging.error(f"Error replacing picture in T30: {e}")

def _t30_find_pic_placeholder(slide):
    """Slayddagi picture placeholder yoki blip shape indeksini topish."""
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for i, s in enumerate(slide.shapes):
        # Avval blip bor shapeni tekshirish
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips:
            return i
        # Keyin picture placeholder ni tekshirish
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic':
            return i
    return None

def fill_t30_slide_1_cover(slide, title, name_surname):
    # [0] sarlavha (44pt bold 000000 center) top=0.83
    # [1] muallif (24pt 000000 center) top=6.11
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        author_idx = text_shapes[1][1]
        _t30_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4400, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
        _t30_clear_and_write(slide.shapes[author_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 2400, 'b': 0, 'color': '000000', 'text': f"Muallif: {name_surname}"}]}
        ])

def fill_t30_slide_2_plan(slide, plan_data):
    # [0] reja matni (24pt 000000 left) top=2.19
    # [1] "REJA" (36pt bold 000000 left) top=1.12
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t30_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': "REJA"}]}
        ])
        
        content = plan_data.get("content", []) if isinstance(plan_data, dict) else []
        paras = []
        import re as _re
        for idx, item in enumerate(content):
            clean_item = _re.sub(r'^\d+[.)\s]+', '', str(item).strip())
            paras.append({'algn': 'l', 'runs': [{'sz': 2400, 'b': 0, 'color': '000000', 'text': f"{idx+1}. {clean_item}"}]})
        _t30_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, paras)

def fill_t30_slide_3_two_col(slide, data, img_arg=None):
    # [0] sarlavha (36pt bold 000000 left) top=0.88
    # [1] chap matn (18pt 000000 left) top=3.36
    # [2] o'ng matn (18pt 000000 left) top=3.36
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t30_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 2)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, s.left or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    if len(text_shapes) >= 3:
        title_idx = text_shapes[0][2]
        matn_shapes = sorted(text_shapes[1:3], key=lambda x: x[1]) # chapdan o'ngga
        left_idx = matn_shapes[0][2]
        right_idx = matn_shapes[1][2]
        
        _t30_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
        _t30_clear_and_write(slide.shapes[left_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': text1}]}
        ])
        _t30_clear_and_write(slide.shapes[right_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': text2}]}
        ])

def fill_t30_slide_4_title_text(slide, data, img_arg=None):
    # [0] sarlavha (36pt bold 000000 left) top=0.80
    # [1] matn (20pt 000000 left) top=2.53
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t30_get_body_text(data)
    
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t30_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
        _t30_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t30_slide_5_img_right(slide, data, img_arg=None):
    # [0] sarlavha (36pt bold 000000 left) top=0.57
    # [1] matn (18pt 000000 left) top=2.04
    # [2] PICTURE (o'ng) top=2.04
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t30_get_body_text(data)
    
    pic_shape_idx = _t30_find_pic_placeholder(slide)
    if pic_shape_idx is not None and img_arg:
        _t30_replace_picture(slide, pic_shape_idx, img_arg)
        
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips:
            continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic':
            continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    # Barcha matn shapelari ni tozalash (placeholder larni ham)
    for _, idx in text_shapes:
        _t30_clear_and_write(slide.shapes[idx].text_frame._txBody, [])
    
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t30_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
        _t30_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t30_slide_6_three_col(slide, data, img_arg=None):
    # [0] sarlavha (36pt bold 000000 left) top=0.98
    # [1] chap matn (16pt 000000 left) top=2.52
    # [2] o'rta matn (16pt 000000 left) top=2.52
    # [3] o'ng matn (16pt 000000 left) top=2.52
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t30_get_body_text(data)
    blocks = split_text_into_blocks(body_text, 3)
    text1 = blocks[0] if len(blocks) > 0 else ""
    text2 = blocks[1] if len(blocks) > 1 else ""
    text3 = blocks[2] if len(blocks) > 2 else ""
    
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, s.left or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    if len(text_shapes) >= 4:
        title_idx = text_shapes[0][2]
        matn_shapes = sorted(text_shapes[1:4], key=lambda x: x[1]) # chapdan o'ngga
        left_idx = matn_shapes[0][2]
        mid_idx = matn_shapes[1][2]
        right_idx = matn_shapes[2][2]
        
        _t30_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
        _t30_clear_and_write(slide.shapes[left_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1600, 'b': 0, 'color': '000000', 'text': text1}]}
        ])
        _t30_clear_and_write(slide.shapes[mid_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1600, 'b': 0, 'color': '000000', 'text': text2}]}
        ])
        _t30_clear_and_write(slide.shapes[right_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1600, 'b': 0, 'color': '000000', 'text': text3}]}
        ])

def fill_t30_slide_7_img_left(slide, data, img_arg=None):
    # [0] sarlavha (36pt bold 000000 left) top=0.70
    # [1] matn (18pt 000000 left) top=2.16
    # [2] PICTURE (chap) top=2.16
    # [3] page number (ignore)
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t30_get_body_text(data)
    
    pic_shape_idx = _t30_find_pic_placeholder(slide)
    if pic_shape_idx is not None and img_arg:
        _t30_replace_picture(slide, pic_shape_idx, img_arg)
        
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips:
            continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic':
            continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            # page number shape larni o'tkazib yuborish (kichik yoki pastda)
            if (s.top or 0) > 6000000:
                continue
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    # Barcha matn shapelari ni tozalash (placeholder larni ham)
    for _, idx in text_shapes:
        _t30_clear_and_write(slide.shapes[idx].text_frame._txBody, [])
    
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t30_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 3600, 'b': 1, 'color': '000000', 'text': title.upper()}]}
        ])
        _t30_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t30_slide_8_conclusion(slide, data):
    # [0] xulosa_matn (60pt bold 000000 center)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t30_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6000, 'b': 1, 'color': '000000', 'text': "E'TIBORINGIZ UCHUN RAHMAT!"}]}
        ])

def build_slide_structure_30(prs, requested_slide_count):
    import logging
    from utils import duplicate_slide, move_slide
    if len(prs.slides) < 8:
        logging.error("T30 shablonida kamida 8 ta slayd bo'lishi kerak.")
        return

    content_slide_indices = [2, 3, 4, 5, 6]
    needed = requested_slide_count - len(content_slide_indices)
    
    if needed > 0:
        for i in range(needed):
            src_idx = content_slide_indices[i % len(content_slide_indices)]
            new_slide = duplicate_slide(prs, src_idx)
            move_slide(prs, len(prs.slides) - 1, 7 + i)

def generate_template_30_presentation(prs, topic, requested_slide_count, language, name_surname, plan, content_data_list, user_images=None):
    import os
    from utils import save_user_image_to_tmp, fetch_image
    
    build_slide_structure_30(prs, requested_slide_count)
    
    fill_t30_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_t30_slide_2_plan(prs.slides[1], plan)
    
    content_slide_funcs = [
        fill_t30_slide_3_two_col,
        fill_t30_slide_4_title_text,
        fill_t30_slide_5_img_right,
        fill_t30_slide_6_three_col,
        fill_t30_slide_7_img_left
    ]
    
    img_slide_types = {2, 4}  # img_right va img_left
    
    slides = list(prs.slides)
    slide_index = 2
    img_counter = 0
    
    for i, data in enumerate(content_data_list):
        if slide_index >= len(slides) - 1:
            break
            
        slide = slides[slide_index]
        slide_type = i % len(content_slide_funcs)
        func = content_slide_funcs[slide_type]
        
        img_arg = None
        if slide_type in img_slide_types:
            if user_images and img_counter < len(user_images):
                raw = user_images[img_counter]
                if isinstance(raw, (bytes, bytearray)):
                    img_arg = save_user_image_to_tmp(raw)
                elif isinstance(raw, str) and os.path.exists(raw):
                    img_arg = raw
            if not img_arg:
                img_query = data.get("image_query", "")
                img_arg = fetch_image(img_query) or fetch_image(topic)
            img_counter += 1
            func(slide, data, img_arg)
        else:
            func(slide, data, None)
            
        slide_index += 1
        
    fill_t30_slide_8_conclusion(slides[-1], {})
    
    from io import BytesIO
    out = BytesIO()
    prs.save(out)
    return out.getvalue()
def _t31_clear_and_write(txBody, paras_data):
    """paras_data ichidagi 'color' qiymati:
       - 'bg1', 'tx1', 'accent1' kabi string => schemeClr
       - 'FFFFFF', 'FE5A5B' kabi 6 xonali hex => srgbClr
    """
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SCHEME_COLORS = {'bg1','bg2','dk1','dk2','lt1','lt2','tx1','tx2',
                     'accent1','accent2','accent3','accent4','accent5','accent6',
                     'hlink','folHlink'}
    for p in list(txBody):
        if p.tag == f'{{{ns_a}}}p':
            txBody.remove(p)
    for p_data in paras_data:
        p_elem = etree.SubElement(txBody, f'{{{ns_a}}}p')
        pPr = etree.SubElement(p_elem, f'{{{ns_a}}}pPr')
        algn = p_data.get('algn', 'l')
        if algn:
            pPr.set('algn', algn)
        etree.SubElement(pPr, f'{{{ns_a}}}buNone')
        for run in p_data.get('runs', []):
            r_elem = etree.SubElement(p_elem, f'{{{ns_a}}}r')
            rPr = etree.SubElement(r_elem, f'{{{ns_a}}}rPr')
            rPr.set('lang', 'en-US')
            rPr.set('dirty', '0')
            if 'sz' in run:
                rPr.set('sz', str(run['sz']))
            if run.get('b'):
                rPr.set('b', '1')
            else:
                rPr.set('b', '0')
            if 'color' in run:
                solidFill = etree.SubElement(rPr, f'{{{ns_a}}}solidFill')
                color_val = run['color']
                if color_val in SCHEME_COLORS:
                    schemeClr = etree.SubElement(solidFill, f'{{{ns_a}}}schemeClr')
                    schemeClr.set('val', color_val)
                else:
                    srgbClr = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
                    srgbClr.set('val', color_val.upper().lstrip('#'))
            t_elem = etree.SubElement(r_elem, f'{{{ns_a}}}t')
            t_elem.text = run.get('text', '')

def _t31_get_body_text(data):
    if isinstance(data, dict):
        # col1/col2 formatini ham handle qilish
        if "col1" in data or "col2" in data:
            parts = []
            for k in ["col1", "col2", "col3"]:
                v = data.get(k, "")
                if v:
                    parts.append(str(v))
            return "\n".join(parts)
        c = data.get("content", [])
        if isinstance(c, list):
            return "\n".join(str(x) for x in c)
        return str(c)
    return str(data)

def _t31_replace_picture(slide, shape_idx, img_arg):
    import os
    from pptx.util import Inches
    if not img_arg:
        return
    if isinstance(img_arg, str) and not os.path.exists(img_arg):
        return
    old_shape = slide.shapes[shape_idx]
    left, top = old_shape.left, old_shape.top
    width, height = old_shape.width, old_shape.height
    try:
        spTree = slide.shapes._spTree
        spTree.remove(old_shape._element)
        slide.shapes.add_picture(img_arg, left, top, width, height)
    except Exception as e:
        import logging
        logging.error(f"Error replacing picture in T31: {e}")

def _t31_find_pic_placeholder(slide):
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips:
            return i
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic':
            return i
    return None

def fill_t31_slide_1_cover(slide, title, name_surname):
    # [0] TEXT 48pt top=0.69 left=0.7 (Sarlavha) — QORA fon, OQ matn
    # [1] TEXT 18pt top=6.5 left=0.17 (Muallif) — QORA fon, OQ matn
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        author_idx = text_shapes[1][1]
        _t31_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': 'bg1', 'text': title.upper()}]}
        ])
        _t31_clear_and_write(slide.shapes[author_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': 'bg1', 'text': name_surname}]}
        ])

def fill_t31_slide_2_plan(slide, plan_data):
    # [0] TEXT BOLD 48pt top=1.33 left=3.6 (REJA)
    # [1] TEXT 18pt top=2.66 left=2.29 (Matn)
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t31_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4800, 'b': 1, 'color': 'bg1', 'text': "REJA"}]}
        ])
        
        items = plan_data.get("content", []) if isinstance(plan_data, dict) else []
        if not items:
            items = ["Kirish", "Asosiy qism", "Xulosa"]
            
        import re
        paras = []
        for idx, item in enumerate(items):
            clean_item = re.sub(r'^\d+[.)]\s*', '', str(item)).strip()
            paras.append({'algn': 'l', 'runs': [{'sz': 1867, 'b': 0, 'color': 'bg1', 'text': f"{idx+1}. {clean_item}"}]})
        _t31_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, paras)

def fill_t31_slide_3_two_col(slide, data, img_arg=None):
    # [0] TEXT 54pt top=0.5 left=0 (Sarlavha)
    # [1] TEXT 12pt top=1.93 left=1.1 (Chap matn)
    # [2] TEXT 12pt top=1.92 left=6.17 (O'ng matn)
    title = data.get("title", "") if isinstance(data, dict) else ""
    if isinstance(data, dict) and ("col1" in data or "col2" in data):
        text1 = str(data.get("col1", ""))
        text2 = str(data.get("col2", ""))
    else:
        body_text = _t31_get_body_text(data)
        from utils import split_text_into_blocks
        blocks = split_text_into_blocks(body_text, 2)
        text1 = blocks[0] if len(blocks) > 0 else ""
        text2 = blocks[1] if len(blocks) > 1 else ""
    
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, s.left or 0, i))
    
    # Eng tepada sarlavha
    text_shapes.sort(key=lambda x: x[0])
    title_idx = text_shapes[0][2]
    
    # Qolgan ikkitasi chap va o'ng
    others = text_shapes[1:]
    others.sort(key=lambda x: x[1])
    
    _t31_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
        {'algn': 'ctr', 'runs': [{'sz': 5400, 'b': 1, 'color': 'bg1', 'text': title.upper()}]}
    ])
    if len(others) >= 2:
        left_idx = others[0][2]
        right_idx = others[1][2]
        _t31_clear_and_write(slide.shapes[left_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1200, 'b': 0, 'color': 'bg1', 'text': text1}]}
        ])
        _t31_clear_and_write(slide.shapes[right_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1200, 'b': 0, 'color': 'bg1', 'text': text2}]}
        ])

def fill_t31_slide_4_two_col(slide, data, img_arg=None):
    # [0] TEXT 48pt top=0.17 left=0.35 (Sarlavha)
    # [1] TEXT 14pt top=2.1 left=1.0 (Chap matn)
    # [2] TEXT 14pt top=2.1 left=7.38 (O'ng matn)
    title = data.get("title", "") if isinstance(data, dict) else ""
    if isinstance(data, dict) and ("col1" in data or "col2" in data):
        text1 = str(data.get("col1", ""))
        text2 = str(data.get("col2", ""))
    else:
        body_text = _t31_get_body_text(data)
        from utils import split_text_into_blocks
        blocks = split_text_into_blocks(body_text, 2)
        text1 = blocks[0] if len(blocks) > 0 else ""
        text2 = blocks[1] if len(blocks) > 1 else ""
    
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, s.left or 0, i))
    
    text_shapes.sort(key=lambda x: x[0])
    title_idx = text_shapes[0][2]
    
    others = text_shapes[1:]
    others.sort(key=lambda x: x[1])
    
    _t31_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
        {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': 'tx1', 'text': title.upper()}]}
    ])
    if len(others) >= 2:
        left_idx = others[0][2]
        right_idx = others[1][2]
        _t31_clear_and_write(slide.shapes[left_idx].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1400, 'b': 0, 'color': 'tx1', 'text': text1}]}
        ])
        _t31_clear_and_write(slide.shapes[right_idx].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1400, 'b': 0, 'color': 'tx1', 'text': text2}]}
        ])

def fill_t31_slide_5_img_left(slide, data, img_arg=None):
    # [0] TEXT 44pt top=0.17 left=0.35 (Sarlavha)
    # [1] TEXT 14pt top=1.69 left=7.14 (O'ng matn)
    # [2] TEXT top=2.13 left=1.34 (bo'sh placeholder)
    # [3] PICTURE top=1.52 left=1.34 w=5.24 h=4.46 (Chap rasm)
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t31_get_body_text(data)
    
    pic_shape_idx = _t31_find_pic_placeholder(slide)
    if pic_shape_idx is not None and img_arg:
        _t31_replace_picture(slide, pic_shape_idx, img_arg)
        
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips: continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic': continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
            
    text_shapes.sort(key=lambda x: x[0])
    for _, idx in text_shapes:
        _t31_clear_and_write(slide.shapes[idx].text_frame._txBody, [])
        
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t31_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4400, 'b': 1, 'color': 'bg1', 'text': title.upper()}]}
        ])
        _t31_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1400, 'b': 0, 'color': 'bg1', 'text': body_text}]}
        ])

def fill_t31_slide_6_three_col(slide, data, img_arg=None):
    # [0] TEXT 48pt top=0.16 left=0.35 (Sarlavha)
    # [1] TEXT 14pt top=2.74 left=1.14 (Matn 1)
    # [2] TEXT 14pt top=2.74 left=5.22 (Matn 2)
    # [3] TEXT 14pt top=2.75 left=9.38 (Matn 3)
    title = data.get("title", "") if isinstance(data, dict) else ""
    # col1/col2/col3 formatini to'g'ridan olish
    if isinstance(data, dict) and ("col1" in data or "col2" in data):
        text1 = str(data.get("col1", ""))
        text2 = str(data.get("col2", ""))
        text3 = str(data.get("col3", ""))
    else:
        body_text = _t31_get_body_text(data)
        from utils import split_text_into_blocks
        blocks = split_text_into_blocks(body_text, 3)
        text1 = blocks[0] if len(blocks) > 0 else ""
        text2 = blocks[1] if len(blocks) > 1 else ""
        text3 = blocks[2] if len(blocks) > 2 else ""
    
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, s.left or 0, i))
            
    text_shapes.sort(key=lambda x: x[0])
    title_idx = text_shapes[0][2]
    
    others = text_shapes[1:]
    others.sort(key=lambda x: x[1])
    
    _t31_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
        {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': 'tx1', 'text': title.upper()}]}
    ])
    
    for idx, txt in zip([x[2] for x in others], [text1, text2, text3]):
        _t31_clear_and_write(slide.shapes[idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 1400, 'b': 0, 'color': 'tx1', 'text': txt}]}
        ])

def fill_t31_slide_7_img_right(slide, data, img_arg=None):
    # [0] TEXT top=2.74 left=5.28 (bo'sh)
    # [1] TEXT top=2.55 left=5.18 (pic placeholder)
    # [2] TEXT BOLD 28pt top=0.58 left=1.09 (Sarlavha)
    # [3] TEXT 14pt top=2.35 left=1.08 (Chap matn)
    title = data.get("title", "") if isinstance(data, dict) else ""
    body_text = _t31_get_body_text(data)
    
    pic_shape_idx = _t31_find_pic_placeholder(slide)
    if pic_shape_idx is not None and img_arg:
        _t31_replace_picture(slide, pic_shape_idx, img_arg)
        
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips: continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic': continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
            
    text_shapes.sort(key=lambda x: x[0])
    for _, idx in text_shapes:
        _t31_clear_and_write(slide.shapes[idx].text_frame._txBody, [])
        
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        # Sarlavha: accent1 (tema qizil rangi), algn=r — asl shablon kabi
        _t31_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'r', 'runs': [
                {'sz': 2800, 'b': 1, 'color': 'accent1', 'text': title},
                {'sz': 2800, 'b': 1, 'color': '5D6268', 'text': ''}
            ]}
        ])
        # Matn: tx1 (qora) rang, algn=r
        _t31_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 1400, 'b': 0, 'color': 'tx1', 'text': body_text}]}
        ])

def fill_t31_slide_8_conclusion(slide, data):
    # [0] TEXT top=2.65 left=0 (E'TIBORINGIZ UCHUN RAHMAT!)
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t31_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6000, 'b': 1, 'color': 'tx1', 'text': "E'TIBORINGIZ UCHUN RAHMAT!"}]}
        ])

def build_slide_structure_31(prs, requested_slide_count):
    import logging
    from utils import duplicate_slide, move_slide
    if len(prs.slides) < 8:
        logging.error("T31 shablonida kamida 8 ta slayd bo'lishi kerak.")
        return

    content_slide_indices = [2, 3, 4, 5, 6]
    needed = requested_slide_count - len(content_slide_indices)
    
    if needed > 0:
        base_indices = content_slide_indices.copy()
        for i in range(needed):
            src_idx = base_indices[i % len(base_indices)]
            duplicate_slide(prs, src_idx)
            
        last_content_idx = 6 + needed
        move_slide(prs, len(prs.slides)-1, last_content_idx + 1)

SLIDE_TYPE_NAMES_T31 = {
    0: "two_columns",
    1: "two_columns",
    2: "image_left",
    3: "three_columns",
    4: "image_right"
}

def generate_template_31_presentation(prs, topic, requested_slide_count, language, name_surname, plan, content_data_list, user_images=None):
    import logging
    from utils import fetch_image, save_user_image_to_tmp
    import os
    
    build_slide_structure_31(prs, requested_slide_count)
    
    fill_t31_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_t31_slide_2_plan(prs.slides[1], plan)
    
    content_slide_funcs = [
        fill_t31_slide_3_two_col,
        fill_t31_slide_4_two_col,
        fill_t31_slide_5_img_left,
        fill_t31_slide_6_three_col,
        fill_t31_slide_7_img_right
    ]
    
    img_counter = 0
    img_slide_types = {2, 4}  # img_left va img_right
    
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
            
        slide = prs.slides[slide_index]
        slide_type = i % len(content_slide_funcs)
        func = content_slide_funcs[slide_type]
        
        img_arg = None
        if slide_type in img_slide_types:
            img_query = data.get("image_query", "")
            if user_images and img_counter < len(user_images):
                raw = user_images[img_counter]
                if isinstance(raw, (bytes, bytearray)):
                    img_arg = save_user_image_to_tmp(raw)
                elif isinstance(raw, str) and os.path.exists(raw):
                    img_arg = raw
                else:
                    img_arg = fetch_image(img_query) or fetch_image(topic)
                img_counter += 1
            else:
                img_arg = fetch_image(img_query) or fetch_image(topic)
                
        func(slide, data, img_arg)
        
    fill_t31_slide_8_conclusion(prs.slides[-1], {})
    
    import io
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ============================================================
# 32-SHABLON FUNKSIYALARI
# ============================================================

SLIDE_TYPE_NAMES_T32 = {
    0: 'image_left',    # Slayd 3 — kulrang fon, rasm chap, matn o'ng
    1: 'image_right',   # Slayd 4 — oq fon, rasm o'ng, sarlavha+matn chap
    2: 'image_right',   # Slayd 5 — qora fon, rasm o'ng, sarlavha+matn chap
    3: 'image_left',    # Slayd 6 — kulrang fon, rasm chap, matn o'ng
    4: 'three_columns', # Slayd 7 — oq fon, 3 ustun + rasm yuqori-o'ng
}

def _t32_clear_and_write(txBody, paras_data):
    """paras_data ichidagi 'color':
       - 'bg1','tx1','accent1' kabi => schemeClr
       - 'F3F3F3','000000' kabi 6-xonali hex => srgbClr
    """
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SCHEME_COLORS = {'bg1','bg2','dk1','dk2','lt1','lt2','tx1','tx2',
                     'accent1','accent2','accent3','accent4','accent5','accent6',
                     'hlink','folHlink'}
    for p in list(txBody):
        if p.tag == f'{{{ns_a}}}p':
            txBody.remove(p)
    for p_data in paras_data:
        p_elem = etree.SubElement(txBody, f'{{{ns_a}}}p')
        pPr = etree.SubElement(p_elem, f'{{{ns_a}}}pPr')
        algn = p_data.get('algn', 'l')
        if algn:
            pPr.set('algn', algn)
        etree.SubElement(pPr, f'{{{ns_a}}}buNone')
        for run in p_data.get('runs', []):
            r_elem = etree.SubElement(p_elem, f'{{{ns_a}}}r')
            rPr = etree.SubElement(r_elem, f'{{{ns_a}}}rPr')
            rPr.set('lang', 'en-US')
            rPr.set('dirty', '0')
            if 'sz' in run:
                rPr.set('sz', str(run['sz']))
            if run.get('b'):
                rPr.set('b', '1')
            else:
                rPr.set('b', '0')
            if 'color' in run:
                solidFill = etree.SubElement(rPr, f'{{{ns_a}}}solidFill')
                color_val = run['color']
                if color_val in SCHEME_COLORS:
                    schemeClr = etree.SubElement(solidFill, f'{{{ns_a}}}schemeClr')
                    schemeClr.set('val', color_val)
                else:
                    srgbClr = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
                    srgbClr.set('val', color_val.upper().lstrip('#'))
            t_elem = etree.SubElement(r_elem, f'{{{ns_a}}}t')
            t_elem.text = run.get('text', '')

def _t32_get_body_text(data):
    if isinstance(data, dict):
        if 'col1' in data or 'col2' in data:
            parts = []
            for k in ['col1', 'col2', 'col3']:
                v = data.get(k, '')
                if v:
                    parts.append(str(v))
            return '\n'.join(parts)
        c = data.get('content', [])
        if isinstance(c, list):
            return '\n'.join(str(x) for x in c)
        return str(c)
    return str(data)

def _t32_find_pic_shape(slide):
    """Slayddagi rasm yoki pic-placeholder indeksini qaytaradi."""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips:
            return i
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic':
            return i
    return None

def _t32_replace_picture(slide, shape_idx, img_arg):
    import os
    from pptx.util import Inches
    if not img_arg:
        return
    if isinstance(img_arg, str) and not os.path.exists(img_arg):
        return
    old_shape = slide.shapes[shape_idx]
    left, top = old_shape.left, old_shape.top
    width, height = old_shape.width, old_shape.height
    try:
        spTree = slide.shapes._spTree
        spTree.remove(old_shape._element)
        slide.shapes.add_picture(img_arg, left, top, width, height)
    except Exception as e:
        import logging
        logging.error(f"T32 rasm almashtirish xatosi: {e}")

def fill_t32_slide_1_cover(slide, title, name_surname):
    # [0] Sarlavha: top=1.84 left=1.49, sz=6600, srgb:000000, algn=ctr
    # [1] Subtitr: top=9.71 left=3.33, sz=3600, srgb:000000, algn=ctr
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        sub_idx = text_shapes[1][1]
        _t32_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6600, 'b': 1, 'color': '000000', 'text': title}]}
        ])
        _t32_clear_and_write(slide.shapes[sub_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 3600, 'b': 0, 'color': '000000', 'text': name_surname}]}
        ])

def fill_t32_slide_2_plan(slide, plan_data):
    # [0] Sarlavha: top=1.69 left=3.15, sz=7200, srgb:F3F3F3, algn=ctr
    # [1] Ro'yxat: top=4.52 left=4.58, sz=2800, srgb:F3F3F3
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))
    text_shapes.sort(key=lambda x: x[0])
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t32_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 7200, 'b': 1, 'color': 'F3F3F3', 'text': 'REJA'}]}
        ])
        items = plan_data.get('content', []) if isinstance(plan_data, dict) else []
        if not items:
            items = ['Kirish', 'Asosiy qism', 'Xulosa']
        import re
        paras = []
        for idx, item in enumerate(items):
            clean = re.sub(r'^\d+[.)]\s*', '', str(item)).strip()
            paras.append({'algn': 'l', 'runs': [{'sz': 2800, 'b': 1, 'color': 'F3F3F3', 'text': f'{idx+1}.  {clean}'}]})
        _t32_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, paras)

def fill_t32_slide_3_img_left(slide, data, img_arg=None):
    # Kulrang fon | Rasm CHAP | Sarlavha + Matn O'NG
    # [0] Sarlavha: top=0.71 left=4.08, sz=4800, srgb:F3F3F3, algn=r
    # [1] PICTURE chap
    # [2] Matn: top=3.46 left=10.42, sz=1800, scheme:bg1
    title = data.get('title', '') if isinstance(data, dict) else ''
    body_text = _t32_get_body_text(data)

    pic_idx = _t32_find_pic_shape(slide)
    if pic_idx is not None and img_arg:
        _t32_replace_picture(slide, pic_idx, img_arg)

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips: continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic': continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, s.left or 0, i))

    text_shapes.sort(key=lambda x: x[0])
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][2]
        matn_idx = text_shapes[1][2]
        _t32_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 4800, 'b': 1, 'color': 'F3F3F3', 'text': title}]}
        ])
        _t32_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1800, 'b': 0, 'color': 'bg1', 'text': body_text}]}
        ])

def fill_t32_slide_4_img_right(slide, data, img_arg=None):
    # Oq fon | Rasm O'NG | Sarlavha + Matn CHAP
    # [0] PICTURE o'ng
    # [1] Matn: top=4.03 left=1.92, sz=2800, srgb:000000, algn=ctr
    # [2] Sarlavha: top=0.71 left=1.08, sz=4800, NO_COLOR (qora)
    title = data.get('title', '') if isinstance(data, dict) else ''
    body_text = _t32_get_body_text(data)

    pic_idx = _t32_find_pic_shape(slide)
    if pic_idx is not None and img_arg:
        _t32_replace_picture(slide, pic_idx, img_arg)

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips: continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic': continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, s.left or 0, i))

    # Sarlavha eng yuqorida (top kichik), matn pastda
    text_shapes.sort(key=lambda x: x[0])
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][2]
        matn_idx = text_shapes[1][2]
        _t32_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': '000000', 'text': title}]}
        ])
        _t32_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 2800, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t32_slide_5_img_right_dark(slide, data, img_arg=None):
    # Qora fon | Rasm O'NG | Sarlavha + Matn CHAP
    # [0] Matn: top=4.29 left=1.13, sz=1800, srgb:F3F3F3
    # [1] PICTURE o'ng
    # [2] Sarlavha: top=1.04 left=1.12, sz=4800, srgb:F3F3F3
    title = data.get('title', '') if isinstance(data, dict) else ''
    body_text = _t32_get_body_text(data)

    pic_idx = _t32_find_pic_shape(slide)
    if pic_idx is not None and img_arg:
        _t32_replace_picture(slide, pic_idx, img_arg)

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips: continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic': continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))

    text_shapes.sort(key=lambda x: x[0])
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t32_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': 'F3F3F3', 'text': title}]}
        ])
        _t32_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1800, 'b': 0, 'color': 'F3F3F3', 'text': body_text}]}
        ])

def fill_t32_slide_6_img_left2(slide, data, img_arg=None):
    # Kulrang fon | Rasm CHAP | Sarlavha + Matn O'NG
    # [0] PICTURE chap
    # [1] Matn: top=3.79 left=9.33, sz=1800, srgb:F3F3F3
    # [2] Sarlavha: top=1.28 left=7.71, sz=4800, srgb:F3F3F3
    title = data.get('title', '') if isinstance(data, dict) else ''
    body_text = _t32_get_body_text(data)

    pic_idx = _t32_find_pic_shape(slide)
    if pic_idx is not None and img_arg:
        _t32_replace_picture(slide, pic_idx, img_arg)

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips: continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic': continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, i))

    text_shapes.sort(key=lambda x: x[0])
    if len(text_shapes) >= 2:
        title_idx = text_shapes[0][1]
        matn_idx = text_shapes[1][1]
        _t32_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'r', 'runs': [{'sz': 4800, 'b': 1, 'color': 'F3F3F3', 'text': title}]}
        ])
        _t32_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'just', 'runs': [{'sz': 1800, 'b': 0, 'color': 'F3F3F3', 'text': body_text}]}
        ])

def fill_t32_slide_7_three_col(slide, data, img_arg=None):
    # Oq fon | Sarlavha + 3 ustun + Rasm yuqori-o'ng
    # [0] Sarlavha: top=1.10 left=1.12, sz=4800, srgb:000000
    # [1] Col1: top=5.88 left=1.12, sz=1800, srgb:000000
    # [2] Col2: top=5.88 left=7.60, sz=1800, srgb:000000
    # [3] Col3: top=5.88 left=14.07, sz=1800, srgb:000000
    # [4] PICTURE: top=0.59 left=12.65
    title = data.get('title', '') if isinstance(data, dict) else ''
    if isinstance(data, dict) and ('col1' in data or 'col2' in data):
        text1 = str(data.get('col1', ''))
        text2 = str(data.get('col2', ''))
        text3 = str(data.get('col3', ''))
    else:
        body_text = _t32_get_body_text(data)
        from utils import split_text_into_blocks
        blocks = split_text_into_blocks(body_text, 3)
        text1 = blocks[0] if len(blocks) > 0 else ''
        text2 = blocks[1] if len(blocks) > 1 else ''
        text3 = blocks[2] if len(blocks) > 2 else ''

    pic_idx = _t32_find_pic_shape(slide)
    if pic_idx is not None and img_arg:
        _t32_replace_picture(slide, pic_idx, img_arg)

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    text_shapes = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips: continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic': continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            text_shapes.append((s.top or 0, s.left or 0, i))

    text_shapes.sort(key=lambda x: x[0])
    title_idx = text_shapes[0][2]
    others = text_shapes[1:]
    others.sort(key=lambda x: x[1])

    _t32_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
        {'algn': 'l', 'runs': [{'sz': 4800, 'b': 1, 'color': '000000', 'text': title}]}
    ])
    for idx, txt in zip([x[2] for x in others], [text1, text2, text3]):
        _t32_clear_and_write(slide.shapes[idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': txt}]}
        ])

def fill_t32_slide_8_conclusion(slide, data):
    # [0] Markaziy matn: top=2.84 left=3.58, sz=7200, srgb:F3F3F3, algn=ctr
    if len(slide.shapes) > 0 and slide.shapes[0].has_text_frame:
        _t32_clear_and_write(slide.shapes[0].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 7200, 'b': 1, 'color': 'F3F3F3', 'text': "E'TIBORINGIZ UCHUN RAHMAT!"}]}
        ])

def build_slide_structure_32(prs, requested_slide_count):
    import logging
    from utils import duplicate_slide, move_slide
    if len(prs.slides) < 8:
        logging.error("T32 shablonida kamida 8 ta slayd bo'lishi kerak.")
        return

    content_slide_indices = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7
    needed = requested_slide_count - len(content_slide_indices)

    if needed > 0:
        for _ in range(needed):
            src = content_slide_indices[len(content_slide_indices) % 5]
            duplicate_slide(prs, src)
            content_slide_indices.append(len(prs.slides) - 1)
    elif needed < 0:
        remove_count = -needed
        indices_to_remove = sorted(content_slide_indices[len(content_slide_indices) - remove_count:], reverse=True)
        for idx in indices_to_remove:
            rId = prs.slides._sldIdLst[idx].get('r:id') or prs.slides._sldIdLst[idx].get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rId:
                prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]

def generate_template_32_presentation(prs, topic, requested_slide_count, language, name_surname, plan, content_data_list, user_images=None):
    import logging
    from utils import fetch_image, save_user_image_to_tmp
    import os
    build_slide_structure_32(prs, requested_slide_count)
    fill_t32_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_t32_slide_2_plan(prs.slides[1], plan)
    fill_funcs = [
        fill_t32_slide_3_img_left,
        fill_t32_slide_4_img_right,
        fill_t32_slide_5_img_right_dark,
        fill_t32_slide_6_img_left2,
        fill_t32_slide_7_three_col,
    ]
    img_counter = 0
    img_slide_types = {0, 1, 2, 3, 4}  # barcha content slaydlarda rasm bor
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        slide_type = i % len(fill_funcs)
        func = fill_funcs[slide_type]
        img_arg = None
        if slide_type in img_slide_types:
            img_query = data.get('image_query', '') if isinstance(data, dict) else ''
            if user_images and img_counter < len(user_images):
                raw = user_images[img_counter]
                if isinstance(raw, (bytes, bytearray)):
                    img_arg = save_user_image_to_tmp(raw)
                elif isinstance(raw, str) and os.path.exists(raw):
                    img_arg = raw
                else:
                    img_arg = fetch_image(img_query) or fetch_image(topic)
                img_counter += 1
            else:
                img_arg = fetch_image(img_query) or fetch_image(topic)
        func(slide, data, img_arg)
    fill_t32_slide_8_conclusion(prs.slides[-1], {})
    import io
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ============================================================
# 33-SHABLON FUNKSIYALARI
# ============================================================

SLIDE_TYPE_NAMES_T33 = {
    0: 'image_left',    # Slayd 3 — phone mockup, rasm chap
    1: 'image_right',   # Slayd 4 — laptop mockup, rasm o'ng
    2: 'image_right',   # Slayd 5 — photo, rasm o'ng
    3: 'two_columns',   # Slayd 6 — ikki ustun
    4: 'three_columns', # Slayd 7 — uch ustun
}

def _t33_clear_and_write(txBody, paras_data):
    """paras_data ichidagi 'color':
       - 'bg1','tx1','accent1' kabi => schemeClr
       - '000000','FFFFFF' kabi 6-xonali hex => srgbClr
       - None/'' => rang yozilmaydi (shablon default)
    """
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SCHEME_COLORS = {'bg1','bg2','dk1','dk2','lt1','lt2','tx1','tx2',
                     'accent1','accent2','accent3','accent4','accent5','accent6',
                     'hlink','folHlink'}
    # lstStyle ni to'liq olib tashlab, toza yangi lstStyle qo'yish
    # Bu barcha darajalar (lvl1pPr...lvl9pPr) dagi bullet, hanging indent,
    # hangingPunct va boshqa meros formatlashni butunlay yo'q qiladi
    old_lstStyle = txBody.find(f'{{{ns_a}}}lstStyle')
    if old_lstStyle is not None:
        txBody.remove(old_lstStyle)
    # Toza lstStyle — barcha 9 daraja uchun marL=0, indent=0, buNone
    new_lstStyle = etree.SubElement(txBody, f'{{{ns_a}}}lstStyle')
    for lvl in range(1, 10):
        lvl_pPr = etree.SubElement(new_lstStyle, f'{{{ns_a}}}lvl{lvl}pPr')
        lvl_pPr.set('marL', '0')
        lvl_pPr.set('indent', '0')
        lvl_pPr.set('algn', 'l')
        etree.SubElement(lvl_pPr, f'{{{ns_a}}}buNone')
    for p in list(txBody):
        if p.tag == f'{{{ns_a}}}p':
            txBody.remove(p)
    for p_data in paras_data:
        p_elem = etree.SubElement(txBody, f'{{{ns_a}}}p')
        pPr = etree.SubElement(p_elem, f'{{{ns_a}}}pPr')
        algn = p_data.get('algn', 'l')
        if algn:
            pPr.set('algn', algn)
        # Barcha qatorlar bir xil chapga tekislansin — indent va margin nolga tushiriladi
        pPr.set('marL', '0')
        pPr.set('marR', '0')
        pPr.set('indent', '0')
        etree.SubElement(pPr, f'{{{ns_a}}}buNone')
        for run in p_data.get('runs', []):
            r_elem = etree.SubElement(p_elem, f'{{{ns_a}}}r')
            rPr = etree.SubElement(r_elem, f'{{{ns_a}}}rPr')
            rPr.set('lang', 'en-US')
            rPr.set('dirty', '0')
            if 'sz' in run:
                rPr.set('sz', str(run['sz']))
            if run.get('b'):
                rPr.set('b', '1')
            else:
                rPr.set('b', '0')
            color_val = run.get('color', None)
            if color_val:
                solidFill = etree.SubElement(rPr, f'{{{ns_a}}}solidFill')
                if color_val in SCHEME_COLORS:
                    schemeClr = etree.SubElement(solidFill, f'{{{ns_a}}}schemeClr')
                    schemeClr.set('val', color_val)
                else:
                    srgbClr = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
                    srgbClr.set('val', color_val.upper().lstrip('#'))
            t_elem = etree.SubElement(r_elem, f'{{{ns_a}}}t')
            t_elem.text = run.get('text', '')

def _t33_get_body_text(data):
    if isinstance(data, dict):
        c = data.get('content', [])
        if isinstance(c, list):
            return '\n'.join(str(x) for x in c)
        return str(c)
    return str(data)

def _t33_find_pic_shape(slide):
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips:
            return i
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic':
            return i
    return None

def _t33_replace_picture(slide, shape_idx, img_arg):
    import os
    if not img_arg:
        return
    if isinstance(img_arg, str) and not os.path.exists(img_arg):
        return
    old_shape = slide.shapes[shape_idx]
    left, top = old_shape.left, old_shape.top
    width, height = old_shape.width, old_shape.height
    try:
        spTree = slide.shapes._spTree
        spTree.remove(old_shape._element)
        slide.shapes.add_picture(img_arg, left, top, width, height)
    except Exception as e:
        import logging
        logging.error(f"T33 rasm almashtirish xatosi: {e}")

def _t33_text_shapes(slide):
    """Rasm va pic-placeholder bo'lmagan text shape larni top bo'yicha saralangan holda qaytaradi."""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    result = []
    for i, s in enumerate(slide.shapes):
        blips = s._element.findall(f'.//{{{ns_a}}}blip')
        if blips: continue
        ph = s._element.find(f'.//{{{ns_p}}}ph')
        if ph is not None and ph.get('type', '') == 'pic': continue
        if hasattr(s, 'has_text_frame') and s.has_text_frame:
            result.append((s.top or 0, s.left or 0, i))
    return result

def fill_t33_slide_1_cover(slide, title, name_surname):
    # [0] Sarlavha: top=1.37 left=0.78, sz=5000, b=0, algn=ctr
    # [1] Subtitr: top=6.30 left=0.91, sz=2000, b=0, algn=ctr
    shapes = _t33_text_shapes(slide)
    shapes.sort(key=lambda x: x[0])
    if len(shapes) >= 2:
        title_idx = shapes[0][2]
        sub_idx = shapes[1][2]
        _t33_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 5000, 'b': 0, 'color': '000000', 'text': title}]}
        ])
        _t33_clear_and_write(slide.shapes[sub_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 2000, 'b': 0, 'color': '000000', 'text': name_surname}]}
        ])

def fill_t33_slide_2_plan(slide, plan_data):
    # [0] Sarlavha: top=0.96 left=2.88, sz=4000, b=0, algn=ctr
    # [1] Ro'yxat: top=2.98 left=1.45, sz=2000, b=0, algn=l
    shapes = _t33_text_shapes(slide)
    shapes.sort(key=lambda x: x[0])
    if len(shapes) >= 2:
        title_idx = shapes[0][2]
        matn_idx = shapes[1][2]
        _t33_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4000, 'b': 0, 'color': '000000', 'text': 'Reja'}]}
        ])
        items = plan_data.get('content', []) if isinstance(plan_data, dict) else []
        if not items:
            items = ['Kirish', 'Asosiy qism', 'Xulosa']
        import re
        paras = []
        for idx, item in enumerate(items):
            clean = re.sub(r'^\d+[.)]\s*', '', str(item)).strip()
            paras.append({'algn': 'l', 'runs': [{'sz': 2000, 'b': 0, 'color': '000000', 'text': f'{idx+1}.  {clean}'}]})
        _t33_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, paras)

def fill_t33_slide_3_img_left(slide, data, img_arg=None):
    # Rasm CHAP | Sarlavha ctr + Matn l O'NG
    # [0] PICTURE chap
    # [1] Sarlavha: top=0.64 left=4.49, sz=4000, b=1, algn=ctr
    # [2] Matn: top=2.74 left=5.28, sz=1800, b=0, algn=l
    title = data.get('title', '') if isinstance(data, dict) else ''
    body_text = _t33_get_body_text(data)

    pic_idx = _t33_find_pic_shape(slide)
    if pic_idx is not None:
        _t33_replace_picture(slide, pic_idx, img_arg)

    shapes = _t33_text_shapes(slide)
    shapes.sort(key=lambda x: x[0])
    if len(shapes) >= 2:
        title_idx = shapes[0][2]
        matn_idx = shapes[1][2]
        _t33_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 4000, 'b': 1, 'color': '000000', 'text': title}]}
        ])
        _t33_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t33_slide_4_img_right(slide, data, img_arg=None):
    # Rasm O'NG | Sarlavha l + Matn l CHAP
    # [0] PICTURE o'ng
    # [1] Sarlavha: top=1.06 left=0.41, sz=4000, b=1, algn=l
    # [2] Matn: top=2.78 left=0.78, sz=1800, b=0, algn=l
    title = data.get('title', '') if isinstance(data, dict) else ''
    body_text = _t33_get_body_text(data)

    pic_idx = _t33_find_pic_shape(slide)
    if pic_idx is not None:
        _t33_replace_picture(slide, pic_idx, img_arg)

    shapes = _t33_text_shapes(slide)
    shapes.sort(key=lambda x: x[0])
    if len(shapes) >= 2:
        title_idx = shapes[0][2]
        matn_idx = shapes[1][2]
        _t33_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 4000, 'b': 1, 'color': '000000', 'text': title}]}
        ])
        _t33_clear_and_write(slide.shapes[matn_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': body_text}]}
        ])

def fill_t33_slide_5_img_right2(slide, data, img_arg=None):
    # Rasm O'NG | Sarlavha l + Matn l CHAP (slayd 5 ham slayd 4 bilan bir xil tuzilma)
    fill_t33_slide_4_img_right(slide, data, img_arg)

def fill_t33_slide_6_two_col(slide, data, img_arg=None):
    # Sarlavha ctr | Col1 chap + Col2 o'ng
    # [0] Sarlavha: top=0.44 left=1.15, sz=4000, b=1, algn=ctr
    # [1] Col2 (o'ng): top=2.34 left=6.89, sz=2000, b=0, algn=l
    # [2] Col1 (chap): top=2.34 left=1.39, sz=2000, b=0, algn=l
    title = data.get('title', '') if isinstance(data, dict) else ''
    if isinstance(data, dict) and ('col1' in data or 'col2' in data):
        text1 = str(data.get('col1', ''))
        text2 = str(data.get('col2', ''))
    else:
        body_text = _t33_get_body_text(data)
        parts = body_text.split('\n', 1)
        text1 = parts[0] if parts else ''
        text2 = parts[1] if len(parts) > 1 else ''

    shapes = _t33_text_shapes(slide)
    shapes.sort(key=lambda x: x[0])
    title_idx = shapes[0][2]
    # Qolgan 2 ta shape ni left bo'yicha saralash
    others = sorted(shapes[1:], key=lambda x: x[1])

    _t33_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
        {'algn': 'ctr', 'runs': [{'sz': 4000, 'b': 1, 'color': '000000', 'text': title}]}
    ])
    if len(others) >= 2:
        col1_idx = others[0][2]
        col2_idx = others[1][2]
        _t33_clear_and_write(slide.shapes[col1_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': text1}]}
        ])
        _t33_clear_and_write(slide.shapes[col2_idx].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': text2}]}
        ])

def fill_t33_slide_7_three_col(slide, data, img_arg=None):
    # Sarlavha ctr | 3 ustun ctr
    # [0] Sarlavha: top=0.49 left=0.79, sz=4000, b=1, algn=ctr
    # [1] Col1: top=3.94 left=1.53, sz=2000, b=0, algn=ctr
    # [2] Col2: top=3.94 left=5.27, sz=2000, b=0, algn=ctr
    # [3] Col3: top=3.94 left=9.00, sz=2000, b=0, algn=ctr
    title = data.get('title', '') if isinstance(data, dict) else ''
    if isinstance(data, dict) and ('col1' in data or 'col2' in data):
        text1 = str(data.get('col1', ''))
        text2 = str(data.get('col2', ''))
        text3 = str(data.get('col3', ''))
    else:
        body_text = _t33_get_body_text(data)
        parts = body_text.split('\n')
        n = len(parts)
        third = max(1, n // 3)
        text1 = '\n'.join(parts[:third])
        text2 = '\n'.join(parts[third:2*third])
        text3 = '\n'.join(parts[2*third:])

    shapes = _t33_text_shapes(slide)
    shapes.sort(key=lambda x: x[0])
    title_idx = shapes[0][2]
    others = sorted(shapes[1:], key=lambda x: x[1])

    _t33_clear_and_write(slide.shapes[title_idx].text_frame._txBody, [
        {'algn': 'ctr', 'runs': [{'sz': 4000, 'b': 1, 'color': '000000', 'text': title}]}
    ])
    for idx_info, txt in zip(others, [text1, text2, text3]):
        _t33_clear_and_write(slide.shapes[idx_info[2]].text_frame._txBody, [
            {'algn': 'l', 'runs': [{'sz': 1800, 'b': 0, 'color': '000000', 'text': txt}]}
        ])

def fill_t33_slide_8_conclusion(slide, data):
    # [0] Markaziy matn: top=2.09 left=0.88, sz=6600, b=1, algn=ctr
    shapes = _t33_text_shapes(slide)
    if shapes:
        idx = shapes[0][2]
        _t33_clear_and_write(slide.shapes[idx].text_frame._txBody, [
            {'algn': 'ctr', 'runs': [{'sz': 6600, 'b': 1, 'color': '000000', 'text': "E'TIBORINGIZ UCHUN RAHMAT!"}]}
        ])

def build_slide_structure_33(prs, requested_slide_count):
    import logging
    from pptx.oxml.ns import qn
    import copy
    if len(prs.slides) < 8:
        logging.error("T33 shablonida kamida 8 ta slayd bo'lishi kerak.")
        return

    content_templates = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7
    needed = requested_slide_count - len(content_templates)

    if needed > 0:
        for extra in range(needed):
            src_idx = content_templates[extra % len(content_templates)]
            src_slide = prs.slides[src_idx]
            xml_str = src_slide._element.xml
            from lxml import etree
            new_el = etree.fromstring(xml_str)
            rId = prs.slides._sldIdLst[-1].get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            ) or prs.slides._sldIdLst[-1].get('r:id')
            # Oddiy duplicate — oxirgi content slaydni nusxalash
            from pptx.util import Inches
            slide_layout = src_slide.slide_layout
            new_slide = prs.slides.add_slide(slide_layout)
            new_slide._element.getparent().replace(new_slide._element, copy.deepcopy(src_slide._element))
    elif needed < 0:
        remove_count = -needed
        indices_to_remove = sorted(
            content_templates[len(content_templates) - remove_count:],
            reverse=True
        )
        for idx in indices_to_remove:
            sldIdLst = prs.slides._sldIdLst
            sld_elem = sldIdLst[idx]
            rId = sld_elem.get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            ) or sld_elem.get('r:id')
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass
            sldIdLst.remove(sld_elem)

def generate_template_33_presentation(prs, topic, requested_slide_count, language, name_surname, plan, content_data_list, user_images=None):
    import logging
    from utils import fetch_image, save_user_image_to_tmp
    import os
    build_slide_structure_33(prs, requested_slide_count)
    fill_t33_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_t33_slide_2_plan(prs.slides[1], plan)
    fill_funcs = [
        fill_t33_slide_3_img_left,
        fill_t33_slide_4_img_right,
        fill_t33_slide_5_img_right2,
        fill_t33_slide_6_two_col,
        fill_t33_slide_7_three_col,
    ]
    img_slide_types = {0, 1, 2}  # slayd 3,4,5 da rasm bor
    img_counter = 0
    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        slide_type = i % len(fill_funcs)
        func = fill_funcs[slide_type]
        img_arg = None
        if slide_type in img_slide_types:
            img_query = data.get('image_query', '') if isinstance(data, dict) else ''
            if user_images and img_counter < len(user_images):
                raw = user_images[img_counter]
                if isinstance(raw, (bytes, bytearray)):
                    img_arg = save_user_image_to_tmp(raw)
                elif isinstance(raw, str) and os.path.exists(raw):
                    img_arg = raw
                else:
                    img_arg = fetch_image(img_query) or fetch_image(topic)
                img_counter += 1
            else:
                img_arg = fetch_image(img_query) or fetch_image(topic)
        func(slide, data, img_arg)
    fill_t33_slide_8_conclusion(prs.slides[-1], {})
    import io
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ============================================================
# 34-SHABLON — Architect Minitheme (oq fon, ko'k sarlavha)
# ============================================================

def build_slide_structure_34(prs, requested_slide_count):
    import logging
    import copy
    if len(prs.slides) < 8:
        logging.error("T34 shablonida kamida 8 ta slayd bo'lishi kerak.")
        return

    content_templates = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7
    needed = requested_slide_count - len(content_templates)

    if needed > 0:
        for extra in range(needed):
            src_idx = content_templates[extra % len(content_templates)]
            src_slide = prs.slides[src_idx]
            slide_layout = src_slide.slide_layout
            new_slide = prs.slides.add_slide(slide_layout)
            new_slide._element.getparent().replace(new_slide._element, copy.deepcopy(src_slide._element))
    elif needed < 0:
        remove_count = -needed
        indices_to_remove = sorted(
            content_templates[len(content_templates) - remove_count:],
            reverse=True
        )
        for idx in indices_to_remove:
            sldIdLst = prs.slides._sldIdLst
            sld_elem = sldIdLst[idx]
            rId = sld_elem.get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            ) or sld_elem.get('r:id')
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass
            sldIdLst.remove(sld_elem)


SLIDE_TYPE_NAMES_T34 = {
    0: 'two_columns',   # Slayd 3 — ikki ustun (markazda dekor rasm)
    1: 'two_columns',   # Slayd 4 — ikki ustun
    2: 'image_left',    # Slayd 5 — rasm chap
    3: 'image_right',   # Slayd 6 — rasm o'ng
    4: 'two_columns',   # Slayd 7 — ikki ustun (icons yuqorida)
}

def _t34_clear_and_write(txBody, paras_data):
    """paras_data ichidagi 'color':
       - 'bg1','tx1','accent1','bg2' kabi => schemeClr
       - '0B3763','FFFFFF' kabi 6-xonali hex => srgbClr
       - None/'' => rang yozilmaydi (shablon default)
    """
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SCHEME_COLORS = {'bg1','bg2','dk1','dk2','lt1','lt2','tx1','tx2',
                     'accent1','accent2','accent3','accent4','accent5','accent6',
                     'hlink','folHlink'}
    # lstStyle ni to'liq olib tashlab, toza yangi lstStyle qo'yish
    old_lstStyle = txBody.find(f'{{{ns_a}}}lstStyle')
    if old_lstStyle is not None:
        txBody.remove(old_lstStyle)
    new_lstStyle = etree.SubElement(txBody, f'{{{ns_a}}}lstStyle')
    for lvl in range(1, 10):
        lvl_pPr = etree.SubElement(new_lstStyle, f'{{{ns_a}}}lvl{lvl}pPr')
        lvl_pPr.set('marL', '0')
        lvl_pPr.set('indent', '0')
        lvl_pPr.set('algn', 'l')
        etree.SubElement(lvl_pPr, f'{{{ns_a}}}buNone')

    for p in list(txBody):
        if p.tag == f'{{{ns_a}}}p':
            txBody.remove(p)

    for p_data in paras_data:
        p_elem = etree.SubElement(txBody, f'{{{ns_a}}}p')
        pPr = etree.SubElement(p_elem, f'{{{ns_a}}}pPr')
        algn = p_data.get('algn', 'l')
        pPr.set('algn', algn)
        pPr.set('marL', '0')
        pPr.set('marR', '0')
        pPr.set('indent', '0')
        etree.SubElement(pPr, f'{{{ns_a}}}buNone')

        r = etree.SubElement(p_elem, f'{{{ns_a}}}r')
        rPr = etree.SubElement(r, f'{{{ns_a}}}rPr')
        rPr.set('lang', 'uz-UZ')
        rPr.set('altLang', 'en-US')
        sz = p_data.get('sz', 1800)
        rPr.set('sz', str(sz))
        b = p_data.get('b', 0)
        rPr.set('b', str(b))
        rPr.set('dirty', '0')

        color = p_data.get('color', None)
        if color:
            solidFill = etree.SubElement(rPr, f'{{{ns_a}}}solidFill')
            if color in SCHEME_COLORS:
                clr = etree.SubElement(solidFill, f'{{{ns_a}}}schemeClr')
                clr.set('val', color)
            else:
                clr = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
                clr.set('val', color.lstrip('#'))

        t = etree.SubElement(r, f'{{{ns_a}}}t')
        t.text = p_data.get('text', '')


def build_slide_structure_34(prs, requested_slide_count):
    import logging
    import copy
    if len(prs.slides) < 8:
        logging.error("T34 shablonida kamida 8 ta slayd bo'lishi kerak.")
        return

    content_templates = [2, 3, 4, 5, 6]  # 0-indexed: slayd 3,4,5,6,7
    needed = requested_slide_count - len(content_templates)

    if needed > 0:
        for extra in range(needed):
            src_idx = content_templates[extra % len(content_templates)]
            src_slide = prs.slides[src_idx]
            slide_layout = src_slide.slide_layout
            new_slide = prs.slides.add_slide(slide_layout)
            new_slide._element.getparent().replace(new_slide._element, copy.deepcopy(src_slide._element))
    elif needed < 0:
        remove_count = -needed
        indices_to_remove = sorted(
            content_templates[len(content_templates) - remove_count:],
            reverse=True
        )
        for idx in indices_to_remove:
            sldIdLst = prs.slides._sldIdLst
            sld_elem = sldIdLst[idx]
            rId = sld_elem.get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            ) or sld_elem.get('r:id')
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass
            sldIdLst.remove(sld_elem)


SLIDE_TYPE_NAMES_T34 = {
    0: 'two_columns',   # Slayd 3 — ikki ustun (markazda dekor rasm)
    1: 'two_columns',   # Slayd 4 — ikki ustun
    2: 'image_left',    # Slayd 5 — rasm chap
    3: 'image_right',   # Slayd 6 — rasm o'ng
    4: 'two_columns',   # Slayd 7 — ikki ustun (icons yuqorida)
}

def _t34_clear_and_write(txBody, paras_data):
    """paras_data ichidagi 'color':
       - 'bg1','tx1','accent1','bg2' kabi => schemeClr
       - '0B3763','FFFFFF' kabi 6-xonali hex => srgbClr
       - None/'' => rang yozilmaydi (shablon default)
    """
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SCHEME_COLORS = {'bg1','bg2','dk1','dk2','lt1','lt2','tx1','tx2',
                     'accent1','accent2','accent3','accent4','accent5','accent6',
                     'hlink','folHlink'}
    # lstStyle ni to'liq olib tashlab, toza yangi lstStyle qo'yish
    old_lstStyle = txBody.find(f'{{{ns_a}}}lstStyle')
    if old_lstStyle is not None:
        txBody.remove(old_lstStyle)
    new_lstStyle = etree.SubElement(txBody, f'{{{ns_a}}}lstStyle')
    for lvl in range(1, 10):
        lvl_pPr = etree.SubElement(new_lstStyle, f'{{{ns_a}}}lvl{lvl}pPr')
        lvl_pPr.set('marL', '0')
        lvl_pPr.set('indent', '0')
        lvl_pPr.set('algn', 'l')
        etree.SubElement(lvl_pPr, f'{{{ns_a}}}buNone')

    for p in list(txBody):
        if p.tag == f'{{{ns_a}}}p':
            txBody.remove(p)

    for p_data in paras_data:
        p_elem = etree.SubElement(txBody, f'{{{ns_a}}}p')
        pPr = etree.SubElement(p_elem, f'{{{ns_a}}}pPr')
        algn = p_data.get('algn', 'l')
        pPr.set('algn', algn)
        pPr.set('marL', '0')
        pPr.set('marR', '0')
        pPr.set('indent', '0')
        etree.SubElement(pPr, f'{{{ns_a}}}buNone')

        r = etree.SubElement(p_elem, f'{{{ns_a}}}r')
        rPr = etree.SubElement(r, f'{{{ns_a}}}rPr')
        rPr.set('lang', 'uz-UZ')
        rPr.set('altLang', 'en-US')
        sz = p_data.get('sz', 1800)
        rPr.set('sz', str(sz))
        b = p_data.get('b', 0)
        rPr.set('b', str(b))
        rPr.set('dirty', '0')

        color = p_data.get('color', None)
        if color:
            solidFill = etree.SubElement(rPr, f'{{{ns_a}}}solidFill')
            if color in SCHEME_COLORS:
                clr = etree.SubElement(solidFill, f'{{{ns_a}}}schemeClr')
                clr.set('val', color)
            else:
                clr = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
                clr.set('val', color.lstrip('#'))

        t = etree.SubElement(r, f'{{{ns_a}}}t')
        t.text = p_data.get('text', '')


def fill_t34_slide_1_cover(slide, title, name_surname):
    """1-slayd: Muqova — sarlavha + muallif"""
    from pptx.util import Pt
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    # [0] sarlavha, [1] subtitle/muallif
    if len(shapes) >= 1:
        _t34_clear_and_write(shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 5400, 'b': 0, 'color': '0B3763', 'algn': 'ctr'}
        ])
    if len(shapes) >= 2:
        _t34_clear_and_write(shapes[1].text_frame._txBody, [
            {'text': name_surname, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'ctr'}
        ])


def fill_t34_slide_2_plan(slide, plan_items):
    """2-slayd: Reja — sarlavha + raqamlangan ro'yxat"""
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    # plan_items dict bo'lsa content ni ol, list bo'lsa to'g'ridan ishlat
    if isinstance(plan_items, dict):
        items = plan_items.get('content', [])
    elif isinstance(plan_items, list):
        items = plan_items
    else:
        items = []
    shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    # [0] sarlavha
    if len(shapes) >= 1:
        _t34_clear_and_write(shapes[0].text_frame._txBody, [
            {'text': 'REJA', 'sz': 4000, 'b': 0, 'color': '0B3763', 'algn': 'ctr'}
        ])
    # [1] body — raqamlangan ro'yxat
    if len(shapes) >= 2:
        import re
        paras = []
        for idx, item in enumerate(items, 1):
            item_str = str(item).strip()
            # Agar element allaqachon raqam bilan boshlangan bo'lsa ("1. ..."), tozalab faqat matnni ol
            clean = re.sub(r'^\d+[\.\.\)\:]\s*', '', item_str).strip()
            text = f"{idx}. {clean}" if clean else f"{idx}. {item_str}"
            paras.append({'text': text, 'sz': 2000, 'b': 0, 'color': '44546A', 'algn': 'l'})
        _t34_clear_and_write(shapes[1].text_frame._txBody, paras)


def fill_t34_slide_3_two_col(slide, title, col1_text, col2_text):
    """3-slayd: Ikki ustun (markazda dekor rasm saqlanadi)"""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    # [0] sarlavha, [1] chap body, [2] o'ng body
    if len(text_shapes) >= 1:
        _t34_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 4000, 'b': 0, 'color': '0B3763', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        paras = [{'text': p, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}
                 for p in col1_text.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': col1_text, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}]
        _t34_clear_and_write(text_shapes[1].text_frame._txBody, paras)
    if len(text_shapes) >= 3:
        paras = [{'text': p, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}
                 for p in col2_text.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': col2_text, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}]
        _t34_clear_and_write(text_shapes[2].text_frame._txBody, paras)


def fill_t34_slide_4_two_col(slide, title, col1_text, col2_text):
    """4-slayd: Ikki ustun"""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(text_shapes) >= 1:
        _t34_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 4000, 'b': 0, 'color': '0B3763', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        paras = [{'text': p, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}
                 for p in col1_text.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': col1_text, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}]
        _t34_clear_and_write(text_shapes[1].text_frame._txBody, paras)
    if len(text_shapes) >= 3:
        paras = [{'text': p, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}
                 for p in col2_text.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': col2_text, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}]
        _t34_clear_and_write(text_shapes[2].text_frame._txBody, paras)


def fill_t34_slide_5_img_left(slide, title, body_text, image_path=None):
    """5-slayd: Rasm chap, matn o'ng"""
    from pptx.util import Emu
    from PIL import Image as PILImage
    import os
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    pic_shapes = [s for s in slide.shapes
                  if s._element.findall(f'.//{{{ns_a}}}blip')]

    if len(text_shapes) >= 1:
        _t34_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 4000, 'b': 0, 'color': '0B3763', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        paras = [{'text': p, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}
                 for p in body_text.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': body_text, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}]
        _t34_clear_and_write(text_shapes[1].text_frame._txBody, paras)

    if image_path and os.path.exists(image_path) and pic_shapes:
        try:
            pic_shape = pic_shapes[0]
            left = pic_shape.left
            top = pic_shape.top
            width = pic_shape.width
            height = pic_shape.height
            sp = pic_shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception:
            pass


def fill_t34_slide_6_img_right(slide, title, body_text, image_path=None):
    """6-slayd: Matn chap, rasm o'ng"""
    from pptx.util import Emu
    import os
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    pic_shapes = [s for s in slide.shapes
                  if s._element.findall(f'.//{{{ns_a}}}blip')]

    if len(text_shapes) >= 1:
        _t34_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 4000, 'b': 0, 'color': '0B3763', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        paras = [{'text': p, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}
                 for p in body_text.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': body_text, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}]
        _t34_clear_and_write(text_shapes[1].text_frame._txBody, paras)

    if image_path and os.path.exists(image_path) and pic_shapes:
        try:
            pic_shape = pic_shapes[0]
            left = pic_shape.left
            top = pic_shape.top
            width = pic_shape.width
            height = pic_shape.height
            sp = pic_shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception:
            pass


def fill_t34_slide_7_two_col_icons(slide, title, col1_text, col2_text):
    """7-slayd: Ikki ustun (icons yuqorida — dekor, o'zgartirilmaydi)"""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(text_shapes) >= 1:
        _t34_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 4000, 'b': 0, 'color': '0B3763', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        paras = [{'text': p, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}
                 for p in col1_text.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': col1_text, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}]
        _t34_clear_and_write(text_shapes[1].text_frame._txBody, paras)
    if len(text_shapes) >= 3:
        paras = [{'text': p, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}
                 for p in col2_text.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': col2_text, 'sz': 1800, 'b': 0, 'color': '44546A', 'algn': 'l'}]
        _t34_clear_and_write(text_shapes[2].text_frame._txBody, paras)


def fill_t34_slide_8_outro(slide):
    """8-slayd: Xulosa — E'TIBORINGIZ UCHUN RAHMAT! (o'zgartirilmaydi)"""
    pass


def _t34_get_body_text(slide_data):
    """slide_data dan body matnini olish"""
    if isinstance(slide_data, dict):
        content = slide_data.get('content') or slide_data.get('col1') or ''
        if isinstance(content, list):
            return '\n'.join(str(c) for c in content)
        return str(content)
    return str(slide_data)


def _t34_get_col_text(slide_data, col_key):
    """slide_data dan ustun matnini olish"""
    if isinstance(slide_data, dict):
        val = slide_data.get(col_key, '')
        if isinstance(val, list):
            return '\n'.join(str(v) for v in val)
        return str(val)
    return ''


def generate_template_34_presentation(prs, topic, requested_slide_count, language,
                                       name_surname, plan, content_data_list,
                                       user_images=None):
    """34-shablon asosida taqdimot yaratish — 33-shablon uslubida"""
    import logging
    from utils import fetch_image, save_user_image_to_tmp
    import os

    build_slide_structure_34(prs, requested_slide_count)
    fill_t34_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_t34_slide_2_plan(prs.slides[1], plan)

    fill_funcs = [
        fill_t34_slide_3_two_col,        # 0: two_col_center
        fill_t34_slide_4_two_col,        # 1: two_col
        fill_t34_slide_5_img_left,       # 2: img_left
        fill_t34_slide_6_img_right,      # 3: img_right
        fill_t34_slide_7_two_col_icons,  # 4: two_col_icons
    ]
    img_slide_types = {2, 3}  # slayd 5 (img_left) va slayd 6 (img_right) da rasm bor
    img_counter = 0

    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        slide_type = i % len(fill_funcs)
        func = fill_funcs[slide_type]
        img_arg = None
        if slide_type in img_slide_types:
            img_query = data.get('image_query', '') if isinstance(data, dict) else ''
            if user_images and img_counter < len(user_images):
                raw = user_images[img_counter]
                if isinstance(raw, (bytes, bytearray)):
                    img_arg = save_user_image_to_tmp(raw)
                elif isinstance(raw, str) and os.path.exists(raw):
                    img_arg = raw
                else:
                    img_arg = fetch_image(img_query) or fetch_image(topic)
                img_counter += 1
            else:
                img_arg = fetch_image(img_query) or fetch_image(topic)

        # fill funksiyalarini chaqirish
        title = data.get('title', topic) if isinstance(data, dict) else topic
        if slide_type in {0, 1, 4}:  # two_col turlari
            col1 = _t34_get_col_text(data, 'col1')
            col2 = _t34_get_col_text(data, 'col2')
            if not col1 and not col2:
                col1 = _t34_get_body_text(data)
            func(slide, title, col1, col2)
        elif slide_type in {2, 3}:  # img turlari
            body = _t34_get_body_text(data)
            func(slide, title, body, img_arg)
        else:
            func(slide, data, img_arg)

    fill_t34_slide_8_outro(prs.slides[-1])

    import io
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ============================================================
# ODDIY1 SHABLON FUNKSIYALARI
# ============================================================

def _oddiy1_clear_and_write(txBody, paras_data):
    """
    txBody ichidagi barcha paragraflarni tozalab, yangi matn yozadi.
    paras_data: list of dict:
      - 'text': str
      - 'sz': int (half-points, e.g. 3200 = 32pt)
      - 'b': bool/int
      - 'color': str hex (e.g. '003366') yoki schemeClr nomi
      - 'algn': str ('l','r','ctr','just') yoki None
    """
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SCHEME_COLORS = {'bg1','bg2','dk1','dk2','lt1','lt2','tx1','tx2',
                     'accent1','accent2','accent3','accent4','accent5','accent6',
                     'hlink','folHlink'}

    # Barcha mavjud paragraflarni o'chirish
    for p in txBody.findall(f'{{{ns_a}}}p'):
        txBody.remove(p)
    # lstStyle ni tozalash — hanging indent va bullet muammolarini oldini olish
    lstStyle = txBody.find(f'{{{ns_a}}}lstStyle')
    if lstStyle is not None:
        for child in list(lstStyle):
            lstStyle.remove(child)
        # Barcha 9 daraja uchun toza pPr qo'yish
        for lvl in range(1, 10):
            lvl_elem = etree.SubElement(lstStyle, f'{{{ns_a}}}lvl{lvl}pPr')
            lvl_elem.set('marL', '0')
            lvl_elem.set('indent', '0')
            lvl_elem.set('algn', 'l')
            buNone = etree.SubElement(lvl_elem, f'{{{ns_a}}}buNone')

    for para_data in paras_data:
        text = para_data.get('text', '')
        sz = para_data.get('sz', 1800)
        b = para_data.get('b', 0)
        color = para_data.get('color', '334155')
        algn = para_data.get('algn', None)

        p_elem = etree.SubElement(txBody, f'{{{ns_a}}}p')
        pPr = etree.SubElement(p_elem, f'{{{ns_a}}}pPr')
        if algn:
            pPr.set('algn', algn)
        pPr.set('marL', '0')
        pPr.set('indent', '0')
        buNone = etree.SubElement(pPr, f'{{{ns_a}}}buNone')

        r_elem = etree.SubElement(p_elem, f'{{{ns_a}}}r')
        rPr = etree.SubElement(r_elem, f'{{{ns_a}}}rPr', attrib={'lang': 'uz-UZ', 'dirty': '0'})
        if sz:
            rPr.set('sz', str(sz))
        if b:
            rPr.set('b', '1')
        else:
            rPr.set('b', '0')
        # Rang
        solidFill = etree.SubElement(rPr, f'{{{ns_a}}}solidFill')
        if color in SCHEME_COLORS:
            clr_elem = etree.SubElement(solidFill, f'{{{ns_a}}}schemeClr')
            clr_elem.set('val', color)
        else:
            clr_elem = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
            clr_elem.set('val', color.lstrip('#'))
        t_elem = etree.SubElement(r_elem, f'{{{ns_a}}}t')
        t_elem.text = text


def fill_oddiy1_slide_1_cover(slide, title, name_surname):
    """1-slayd: Muqova — fon rasmi saqlanadi, sarlavha va muallif yoziladi"""
    # Shape[0]: sarlavha (katta oq, markazda), Shape[1]: muallif (kichik, pastda)
    shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(shapes) >= 1:
        _oddiy1_clear_and_write(shapes[0].text_frame._txBody, [
            {'text': title.upper() if title else '', 'sz': 5400, 'b': 1, 'color': 'FFFFFF', 'algn': 'ctr'}
        ])
    if len(shapes) >= 2:
        _oddiy1_clear_and_write(shapes[1].text_frame._txBody, [
            {'text': name_surname, 'sz': 1650, 'b': 1, 'color': 'CBD5E1', 'algn': 'ctr'}
        ])


def fill_oddiy1_slide_2_plan(slide, plan_items):
    """2-slayd: Reja — sarlavha + raqamlangan ro'yxat"""
    import re
    if isinstance(plan_items, dict):
        items = plan_items.get('content', [])
    elif isinstance(plan_items, list):
        items = plan_items
    else:
        items = []
    shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(shapes) >= 1:
        _oddiy1_clear_and_write(shapes[0].text_frame._txBody, [
            {'text': 'Reja:', 'sz': 4800, 'b': 1, 'color': '003366', 'algn': 'ctr'}
        ])
    if len(shapes) >= 2:
        paras = []
        for idx, item in enumerate(items, 1):
            item_str = str(item).strip()
            clean = re.sub(r'^\d+[\.\)\:]\s*', '', item_str).strip()
            text = f"{idx}. {clean}" if clean else f"{idx}. {item_str}"
            paras.append({'text': text, 'sz': 2800, 'b': 0, 'color': '64748B', 'algn': 'l'})
        _oddiy1_clear_and_write(shapes[1].text_frame._txBody, paras)


def fill_oddiy1_slide_3_two_blocks(slide, title, col1_text, col2_text):
    """3-slayd: Sarlavha + ikki alohida matn bloki (yuqori va pastda)"""
    # Shape[0]: sarlavha (katta ko'k, markazda)
    # Shape[1]: birinchi matn bloki (chapda yuqori)
    # Shape[2]: ikkinchi matn bloki (chapda pastda)
    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(text_shapes) >= 1:
        _oddiy1_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title.upper() if title else '', 'sz': 3200, 'b': 1, 'color': '003366', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        paras = []
        lines = col1_text.split('\n') if isinstance(col1_text, str) else [str(col1_text)]
        for i, line in enumerate(lines):
            if line.strip():
                paras.append({'text': line.strip(), 'sz': 1800, 'b': 1 if i == 0 else 0,
                               'color': '334155', 'algn': 'l'})
        if not paras:
            paras = [{'text': str(col1_text), 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}]
        _oddiy1_clear_and_write(text_shapes[1].text_frame._txBody, paras)
    if len(text_shapes) >= 3:
        paras = []
        lines = col2_text.split('\n') if isinstance(col2_text, str) else [str(col2_text)]
        for i, line in enumerate(lines):
            if line.strip():
                paras.append({'text': line.strip(), 'sz': 1800, 'b': 1 if i == 0 else 0,
                               'color': '334155', 'algn': 'l'})
        if not paras:
            paras = [{'text': str(col2_text), 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}]
        _oddiy1_clear_and_write(text_shapes[2].text_frame._txBody, paras)


def fill_oddiy1_slide_3_two_col(slide, title, col1_text, col2_text):
    """3-slayd alias — fill_oddiy1_slide_3_two_blocks ga yo'naltiradi"""
    fill_oddiy1_slide_3_two_blocks(slide, title, col1_text, col2_text)


def fill_oddiy1_slide_3_two_col_UNUSED(slide, title, col1_text, col2_text):
    """3-slayd: Ikki matn bloki (eski, ishlatilmaydi)"""
    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(text_shapes) >= 1:
        _oddiy1_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 3200, 'b': 1, 'color': '003366', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        paras = []
        lines = col1_text.split('\n') if col1_text else []
        for i, line in enumerate(lines):
            if line.strip():
                paras.append({'text': line.strip(), 'sz': 1800, 'b': 1 if i == 0 else 0, 'color': '334155', 'algn': 'l'})
        if not paras:
            paras = [{'text': col1_text, 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}]
        _oddiy1_clear_and_write(text_shapes[1].text_frame._txBody, paras)
    if len(text_shapes) >= 3:
        paras = []
        lines = col2_text.split('\n') if col2_text else []
        for i, line in enumerate(lines):
            if line.strip():
                paras.append({'text': line.strip(), 'sz': 1800, 'b': 1 if i == 0 else 0, 'color': '334155', 'algn': 'l'})
        if not paras:
            paras = [{'text': col2_text, 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}]
        _oddiy1_clear_and_write(text_shapes[2].text_frame._txBody, paras)


def fill_oddiy1_slide_4_single_body(slide, title, body_text, image_path=None):
    """4-slayd: Sarlavha + bitta matn bloki (rasm yo'q)"""
    # Shape[0]: sarlavha (katta ko'k, markazda)
    # Shape[1]: asosiy matn bloki
    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(text_shapes) >= 1:
        _oddiy1_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title.upper() if title else '', 'sz': 3200, 'b': 1, 'color': '003366', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        body = body_text if isinstance(body_text, str) else str(body_text)
        paras = [{'text': p.strip(), 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}
                 for p in body.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': body, 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}]
        _oddiy1_clear_and_write(text_shapes[1].text_frame._txBody, paras)


def fill_oddiy1_slide_4_img_right(slide, title, body_text, image_path=None):
    """4-slayd: Sarlavha + chap matn + o'ng rasm"""
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    import os

    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    pic_shapes = [s for s in slide.shapes if s.shape_type == 13]

    if len(text_shapes) >= 1:
        _oddiy1_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 3200, 'b': 1, 'color': '003366', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        paras = [{'text': p.strip(), 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}
                 for p in body_text.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': body_text, 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}]
        _oddiy1_clear_and_write(text_shapes[1].text_frame._txBody, paras)

    # Rasm almashtirish
    if image_path and os.path.exists(image_path) and pic_shapes:
        pic = pic_shapes[0]
        left, top, width, height = pic.left, pic.top, pic.width, pic.height
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(image_path, left, top, width, height)


def fill_oddiy1_slide_5_three_col(slide, title, col1_text, col2_text, col3_text):
    """5-slayd: Sarlavha + 3 ustun (kartalar)"""
    # Shape[0]: sarlavha, Shape[1]: 1-karta, Shape[2]: 2-karta, Shape[3]: 3-karta
    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(text_shapes) >= 1:
        _oddiy1_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title.upper() if title else '', 'sz': 3200, 'b': 1, 'color': '003366', 'algn': 'ctr'}
        ])
    for col_idx, col_text in enumerate([col1_text, col2_text, col3_text]):
        shape_idx = col_idx + 1
        if len(text_shapes) > shape_idx:
            col_str = col_text if isinstance(col_text, str) else str(col_text)
            paras = [{'text': p.strip(), 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'just'}
                     for p in col_str.split('\n') if p.strip()]
            if not paras:
                paras = [{'text': col_str, 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'just'}]
            _oddiy1_clear_and_write(text_shapes[shape_idx].text_frame._txBody, paras)


def fill_oddiy1_slide_6_two_col_bold(slide, title, col1_text, col2_text):
    """6-slayd: Sarlavha + 2 ustun (ustun sarlavhalari qalin ko'k, birinchi qator sarlavha)"""
    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(text_shapes) >= 1:
        _oddiy1_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 3200, 'b': 1, 'color': '003366', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        lines = col1_text.split('\n') if col1_text else []
        paras = []
        for i, line in enumerate(lines):
            if line.strip():
                paras.append({'text': line.strip(), 'sz': 1800, 'b': 1 if i == 0 else 0,
                               'color': '003366' if i == 0 else '334155', 'algn': 'just'})
        if not paras:
            paras = [{'text': col1_text, 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'just'}]
        _oddiy1_clear_and_write(text_shapes[1].text_frame._txBody, paras)
    if len(text_shapes) >= 3:
        lines = col2_text.split('\n') if col2_text else []
        paras = []
        for i, line in enumerate(lines):
            if line.strip():
                paras.append({'text': line.strip(), 'sz': 1800, 'b': 1 if i == 0 else 0,
                               'color': '003366' if i == 0 else '334155', 'algn': 'just'})
        if not paras:
            paras = [{'text': col2_text, 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'just'}]
        _oddiy1_clear_and_write(text_shapes[2].text_frame._txBody, paras)


def fill_oddiy1_slide_7_single_body(slide, title, body_text, image_path=None):
    """7-slayd: Sarlavha + bitta matn bloki (rasm yo'q)"""
    # Shape[0]: sarlavha, Shape[1]: asosiy matn bloki
    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    if len(text_shapes) >= 1:
        _oddiy1_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title.upper() if title else '', 'sz': 3200, 'b': 1, 'color': '003366', 'algn': 'ctr'}
        ])
    if len(text_shapes) >= 2:
        body = body_text if isinstance(body_text, str) else str(body_text)
        paras = [{'text': p.strip(), 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}
                 for p in body.split('\n') if p.strip()]
        if not paras:
            paras = [{'text': body, 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}]
        _oddiy1_clear_and_write(text_shapes[1].text_frame._txBody, paras)


def fill_oddiy1_slide_7_img_center(slide, title, body_text, image_path=None):
    """7-slayd: Sarlavha + markazda katta rasm (PLACEHOLDER tipida)"""
    import os
    from pptx.util import Inches, Emu

    text_shapes = [s for s in slide.shapes if hasattr(s, 'has_text_frame') and s.has_text_frame]
    pic_shapes = [s for s in slide.shapes if s.shape_type == 13]

    if len(text_shapes) >= 1:
        _oddiy1_clear_and_write(text_shapes[0].text_frame._txBody, [
            {'text': title, 'sz': 3200, 'b': 1, 'color': '003366', 'algn': 'ctr'}
        ])
    # 2-shape PLACEHOLDER — bu rasm uchun joy. Rasm qo'shish
    if image_path and os.path.exists(image_path):
        # Placeholder ni olib tashlash va rasm qo'yish
        if len(text_shapes) >= 2:
            ph = text_shapes[1]
            left, top, width, height = ph.left, ph.top, ph.width, ph.height
            sp = ph._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(image_path, left, top, width, height)
        elif pic_shapes:
            pic = pic_shapes[0]
            left, top, width, height = pic.left, pic.top, pic.width, pic.height
            sp = pic._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(image_path, left, top, width, height)
    else:
        # Rasm yo'q bo'lsa, placeholder ga matn yozish
        if len(text_shapes) >= 2:
            _oddiy1_clear_and_write(text_shapes[1].text_frame._txBody, [
                {'text': body_text or '', 'sz': 1800, 'b': 0, 'color': '334155', 'algn': 'l'}
            ])


def fill_oddiy1_slide_8_outro(slide):
    """8-slayd: Xulosa — o'zgartirmasdan qoldiriladi"""
    pass


def build_slide_structure_oddiy1(prs, requested_slide_count):
    """
    Shablon slaydlarini kerakli songa moslash.
    Slayd tuzilmasi:
      [0] cover
      [1] plan
      [2..N-2] content (two_col, img_right, three_col, two_col_bold, img_center navbatma-navbat)
      [N-1] outro
    """
    from pptx.util import Emu
    import copy
    from lxml import etree

    CONTENT_SLIDE_TYPES = [
        2,  # two_col (3-slayd)
        3,  # img_right (4-slayd)
        4,  # three_col (5-slayd)
        5,  # two_col_bold (6-slayd)
        6,  # img_center (7-slayd)
    ]
    content_count = max(1, requested_slide_count - 2)
    total_needed = content_count + 2  # cover + content + outro

    # Outro slaydini saqlash
    outro_slide_xml = copy.deepcopy(prs.slides[-1]._element)

    # Mavjud content slaydlarini o'chirish (cover va outro dan tashqari)
    slide_list = prs.slides._sldIdLst
    existing_slides = list(prs.slides)
    for slide in existing_slides[1:-1]:
        rId = prs.slides._sldIdLst.findall(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}sldId'
        )

    # Yangi content slaydlarini qo'shish
    template_slides = list(prs.slides)
    n_templates = len(CONTENT_SLIDE_TYPES)

    # Avval barcha content slaydlarini o'chirish
    slides_to_remove = list(prs.slides)[1:-1]
    for slide in slides_to_remove:
        rId_to_remove = None
        for rId, rel in prs.part.rels.items():
            if hasattr(rel, 'target_part') and rel.target_part == slide.part:
                rId_to_remove = rId
                break
        if rId_to_remove:
            prs.part.drop_rel(rId_to_remove)
        sldId_elem = None
        for elem in prs.slides._sldIdLst:
            if elem.get('r:id') == rId_to_remove:
                sldId_elem = elem
                break
        if sldId_elem is not None:
            prs.slides._sldIdLst.remove(sldId_elem)

    # Kerakli content slaydlarini template dan nusxa ko'chirib qo'shish
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    import copy

    for i in range(content_count):
        tmpl_idx = CONTENT_SLIDE_TYPES[i % n_templates]
        if tmpl_idx < len(template_slides):
            tmpl_slide = template_slides[tmpl_idx]
        else:
            tmpl_slide = template_slides[2]

        # Yangi slayd yaratish
        new_slide_part = copy.deepcopy(tmpl_slide.part)
        rId = prs.part.relate_to(new_slide_part,
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        sldId_elem = etree.SubElement(prs.slides._sldIdLst, f'{{{ns_p}}}sldId')
        max_id = max((int(e.get('id', 256)) for e in prs.slides._sldIdLst), default=256)
        sldId_elem.set('id', str(max_id + 1))
        sldId_elem.set(f'{{{ns_r}}}id', rId)


SLIDE_TYPE_NAMES_ODDIY1 = {
    0: 'two_columns',
    1: 'single_body',
    2: 'three_columns',
    3: 'two_columns',
    4: 'single_body',
}


def generate_template_oddiy1_presentation(prs, topic, requested_slide_count, language,
                                           name_surname, plan, content_data_list,
                                           user_images=None):
    """oddiy1-shablon asosida taqdimot yaratish"""
    import logging
    import os
    try:
        from utils import fetch_image, save_user_image_to_tmp
    except ImportError:
        from student_helper_bot_git.utils import fetch_image, save_user_image_to_tmp

    logger = logging.getLogger(__name__)

    fill_oddiy1_slide_1_cover(prs.slides[0], topic, name_surname)
    fill_oddiy1_slide_2_plan(prs.slides[1], plan)

    fill_funcs = [
        fill_oddiy1_slide_3_two_blocks,     # 0: two_columns
        fill_oddiy1_slide_4_single_body,    # 1: single_body
        fill_oddiy1_slide_5_three_col,      # 2: three_columns
        fill_oddiy1_slide_6_two_col_bold,   # 3: two_columns
        fill_oddiy1_slide_7_single_body,    # 4: single_body
    ]

    for i, data in enumerate(content_data_list):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
        slide = prs.slides[slide_index]
        slide_type = i % len(fill_funcs)
        func = fill_funcs[slide_type]

        # Data dan matn olish
        if isinstance(data, dict):
            title = data.get('title', topic)
            col1 = data.get('col1', data.get('content', ''))
            col2 = data.get('col2', '')
            col3 = data.get('col3', '')
            body = data.get('content', col1)
            # col1/col2/col3 list bo'lsa string ga o'tkazish
            if isinstance(col1, list): col1 = '\n'.join(str(x) for x in col1)
            if isinstance(col2, list): col2 = '\n'.join(str(x) for x in col2)
            if isinstance(col3, list): col3 = '\n'.join(str(x) for x in col3)
            if isinstance(body, list): body = '\n'.join(str(x) for x in body)
        else:
            title = topic
            col1 = str(data)
            col2 = ''
            col3 = ''
            body = str(data)

        try:
            if slide_type == 0:   # two_columns
                func(slide, title, col1, col2)
            elif slide_type == 1: # single_body
                func(slide, title, body)
            elif slide_type == 2: # three_columns
                func(slide, title, col1, col2, col3)
            elif slide_type == 3: # two_columns
                func(slide, title, col1, col2)
            elif slide_type == 4: # single_body
                func(slide, title, body)
        except Exception as e:
            logger.warning(f"oddiy1 slide {slide_index} fill xatolik: {e}")

    fill_oddiy1_slide_8_outro(prs.slides[-1])

    from io import BytesIO
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ============================================================
# ODDIY2 SHABLON (template_num=36) — Nature/Green style
# Slaydlar: 1-muqova, 2-reja, 3-quote_img, 4-single_body,
#           5-two_col, 6-two_col_dark, 7-two_col_dark2, 8-xulosa
# ============================================================

SLIDE_TYPE_NAMES_ODDIY2 = {
    0: 'two_columns',
    1: 'two_columns',
    2: 'single_body',
    3: 'two_columns',
    4: 'two_columns',
    5: 'two_columns',
}

def _t_oddiy2_clear_and_write(shape, text, sz=None, bold=None, algn=None, color=None):
    """oddiy2 shablon uchun matn yozish — lstStyle ni tozalab, toza matn yozish."""
    from lxml import etree
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    tf = shape.text_frame
    txBody = tf._txBody
    nsA = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # lstStyle ni tozalash
    lstStyle = txBody.find(f'{{{nsA}}}lstStyle')
    if lstStyle is not None:
        for child in list(lstStyle):
            lstStyle.remove(child)
        # Toza lvl1pPr qo'shish
        lvl1pPr = etree.SubElement(lstStyle, f'{{{nsA}}}lvl1pPr')
        lvl1pPr.set('marL', '0')
        lvl1pPr.set('indent', '0')
        if algn:
            lvl1pPr.set('algn', algn)
        buNone = etree.SubElement(lvl1pPr, f'{{{nsA}}}buNone')

    # Barcha paragraflarni tozalash
    for p in txBody.findall(f'{{{nsA}}}p'):
        txBody.remove(p)

    # Matnni satrlarga bo'lib yozish
    if isinstance(text, list):
        text = '\n'.join(str(t) for t in text)
    elif not isinstance(text, str):
        text = str(text) if text else ''
    lines = text.split('\n') if text else ['']
    for line in lines:
        p = etree.SubElement(txBody, f'{{{nsA}}}p')
        pPr = etree.SubElement(p, f'{{{nsA}}}pPr')
        pPr.set('marL', '0')
        pPr.set('indent', '0')
        if algn:
            pPr.set('algn', algn)
        buNone = etree.SubElement(pPr, f'{{{nsA}}}buNone')

        if line.strip():
            r = etree.SubElement(p, f'{{{nsA}}}r')
            rPr = etree.SubElement(r, f'{{{nsA}}}rPr', attrib={'lang': 'uz-UZ', 'altLang': 'en-US', 'dirty': '0'})
            if sz:
                rPr.set('sz', str(sz))
            if bold is not None:
                rPr.set('b', '1' if bold else '0')
            if color:
                solidFill = etree.SubElement(rPr, f'{{{nsA}}}solidFill')
                if len(color) == 6 and all(c in '0123456789ABCDEFabcdef' for c in color):
                    srgbClr = etree.SubElement(solidFill, f'{{{nsA}}}srgbClr')
                    srgbClr.set('val', color.upper())
                else:
                    schemeClr = etree.SubElement(solidFill, f'{{{nsA}}}schemeClr')
                    schemeClr.set('val', color)
            t = etree.SubElement(r, f'{{{nsA}}}t')
            t.text = line


def fill_t_oddiy2_slide_1_cover(slide, title, name_surname):
    """1-slayd: Muqova — sarlavha va muallif."""
    shapes = slide.shapes
    # Shape[0]: sarlavha (katta markaziy panel) — KATTA HARFLARDA
    if len(shapes) > 0 and shapes[0].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[0], title.upper() if title else '', algn='l')
    # Shape[1]: muallif (pastda)
    if len(shapes) > 1 and shapes[1].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[1], name_surname, sz=2400, algn='l')


def fill_t_oddiy2_slide_2_plan(slide, plan):
    """2-slayd: Reja — tartiblangan ro'yxat."""
    import re
    shapes = slide.shapes
    if isinstance(plan, dict):
        plan_items = plan.get('content', [])
    else:
        plan_items = plan if isinstance(plan, list) else []

    # Shape[0]: sarlavha "Reja"
    if len(shapes) > 0 and shapes[0].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[0], 'Reja', algn='l')

    # Shape[1]: ro'yxat
    if len(shapes) > 1 and shapes[1].has_text_frame:
        lines = []
        for i, item in enumerate(plan_items, 1):
            clean = re.sub(r'^\d+[\.\)]\s*', '', str(item)).strip()
            lines.append(f'{i}. {clean}')
        text = '\n'.join(lines)
        _t_oddiy2_clear_and_write(shapes[1], text, algn='l')


def fill_t_oddiy2_slide_3_quote(slide, title, content):
    """3-slayd: Quote/Image style — sarlavha + katta matn."""
    shapes = slide.shapes
    body_text = ''
    if isinstance(content, dict):
        body_text = content.get('content', content.get('col1', ''))
    elif isinstance(content, str):
        body_text = content

    if len(shapes) > 0 and shapes[0].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[0], title, algn='l')
    if len(shapes) > 1 and shapes[1].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[1], body_text, algn='l')


def fill_t_oddiy2_slide_4_single(slide, title, content):
    """4-slayd: Single body — sarlavha + bitta katta matn bloki."""
    shapes = slide.shapes
    body_text = ''
    if isinstance(content, dict):
        body_text = content.get('content', content.get('col1', ''))
    elif isinstance(content, str):
        body_text = content

    if len(shapes) > 0 and shapes[0].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[0], title, algn='l')
    if len(shapes) > 1 and shapes[1].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[1], body_text, algn='l')


def fill_t_oddiy2_slide_5_two_col(slide, title, content):
    """5-slayd: Two-column (pastel fon) — sarlavha + 2 ustun."""
    shapes = slide.shapes
    col1_text = ''
    col2_text = ''
    if isinstance(content, dict):
        col1_text = content.get('col1', '')
        col2_text = content.get('col2', '')
        # col1/col2 yo'q bo'lsa content ni ikki teng bo'lib yozish
        if not col1_text and not col2_text:
            full = content.get('content', '')
            if isinstance(full, list):
                mid = len(full) // 2
                col1_text = '\n'.join(str(x) for x in full[:mid])
                col2_text = '\n'.join(str(x) for x in full[mid:])
            elif isinstance(full, str):
                sentences = full.split('. ')
                mid = max(1, len(sentences) // 2)
                col1_text = '. '.join(sentences[:mid]) + ('.' if sentences[:mid] else '')
                col2_text = '. '.join(sentences[mid:])
    elif isinstance(content, str):
        sentences = content.split('. ')
        mid = max(1, len(sentences) // 2)
        col1_text = '. '.join(sentences[:mid]) + '.'
        col2_text = '. '.join(sentences[mid:])

    if len(shapes) > 0 and shapes[0].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[0], title, algn='l')
    if len(shapes) > 1 and shapes[1].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[1], col1_text, algn='l')
    if len(shapes) > 2 and shapes[2].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[2], col2_text, algn='l')


def fill_t_oddiy2_slide_6_two_col_dark(slide, title, content):
    """6-slayd: Two-column dark (to'q yashil fon) — sarlavha + 2 ustun."""
    shapes = slide.shapes
    col1_text = ''
    col2_text = ''
    if isinstance(content, dict):
        col1_text = content.get('col1', '')
        col2_text = content.get('col2', '')
        if not col1_text and not col2_text:
            full = content.get('content', '')
            if isinstance(full, list):
                mid = len(full) // 2
                col1_text = '\n'.join(str(x) for x in full[:mid])
                col2_text = '\n'.join(str(x) for x in full[mid:])
            elif isinstance(full, str):
                sentences = full.split('. ')
                mid = max(1, len(sentences) // 2)
                col1_text = '. '.join(sentences[:mid]) + ('.' if sentences[:mid] else '')
                col2_text = '. '.join(sentences[mid:])
    elif isinstance(content, str):
        sentences = content.split('. ')
        mid = max(1, len(sentences) // 2)
        col1_text = '. '.join(sentences[:mid]) + '.'
        col2_text = '. '.join(sentences[mid:])

    if len(shapes) > 0 and shapes[0].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[0], title, algn='l')
    if len(shapes) > 1 and shapes[1].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[1], col1_text, algn='l')
    if len(shapes) > 2 and shapes[2].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[2], col2_text, algn='l')


def fill_t_oddiy2_slide_7_two_col_dark2(slide, title, content):
    """7-slayd: Two-column dark2 (to'q yashil fon) — sarlavha + 2 ustun."""
    shapes = slide.shapes
    col1_text = ''
    col2_text = ''
    if isinstance(content, dict):
        col1_text = content.get('col1', '')
        col2_text = content.get('col2', '')
        if not col1_text and not col2_text:
            full = content.get('content', '')
            if isinstance(full, list):
                mid = len(full) // 2
                col1_text = '\n'.join(str(x) for x in full[:mid])
                col2_text = '\n'.join(str(x) for x in full[mid:])
            elif isinstance(full, str):
                sentences = full.split('. ')
                mid = max(1, len(sentences) // 2)
                col1_text = '. '.join(sentences[:mid]) + ('.' if sentences[:mid] else '')
                col2_text = '. '.join(sentences[mid:])
    elif isinstance(content, str):
        sentences = content.split('. ')
        mid = max(1, len(sentences) // 2)
        col1_text = '. '.join(sentences[:mid]) + '.'
        col2_text = '. '.join(sentences[mid:])

    if len(shapes) > 0 and shapes[0].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[0], title, algn='l')
    if len(shapes) > 1 and shapes[1].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[1], col1_text, algn='l')
    if len(shapes) > 2 and shapes[2].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[2], col2_text, algn='l')


def fill_t_oddiy2_slide_8_conclusion(slide, title):
    """8-slayd: Xulosa — markaziy matn."""
    shapes = slide.shapes
    if len(shapes) > 0 and shapes[0].has_text_frame:
        _t_oddiy2_clear_and_write(shapes[0], "E'tiboringiz uchun rahmat!", algn='ctr')


def build_slide_structure_oddiy2(content_data_list, plan, slide_count):
    """oddiy2 uchun slayd tuzilmasini aniqlash."""
    # Slayd turlari: 3=quote, 4=single, 5=two_col, 6=two_col_dark, 7=two_col_dark2
    content_slide_types = [3, 4, 5, 6, 7]
    slides = []
    for i, item in enumerate(content_data_list):
        stype = content_slide_types[i % len(content_slide_types)]
        slides.append({'type': stype, 'data': item})
    return slides


def generate_template_oddiy2_presentation(prs, topic, requested_slide_count, language, name_surname, plan, content_data_list, user_images=None):
    """oddiy2 shablon asosida taqdimot yaratish."""
    import io

    slides = prs.slides

    # 1-slayd: muqova
    fill_t_oddiy2_slide_1_cover(slides[0], topic, name_surname)

    # 2-slayd: reja
    fill_t_oddiy2_slide_2_plan(slides[1], plan)

    # Kontent slaydlari (3-7 slaydlar: index 2-6)
    content_fill_funcs = [
        fill_t_oddiy2_slide_3_quote,
        fill_t_oddiy2_slide_4_single,
        fill_t_oddiy2_slide_5_two_col,
        fill_t_oddiy2_slide_6_two_col_dark,
        fill_t_oddiy2_slide_7_two_col_dark2,
    ]

    for i, item in enumerate(content_data_list[:5]):
        slide_idx = 2 + i  # 2, 3, 4, 5, 6
        if slide_idx < len(slides):
            title = item.get('title', '') if isinstance(item, dict) else ''
            content_fill_funcs[i](slides[slide_idx], title, item)

    # 8-slayd: xulosa (index 7)
    if len(slides) > 7:
        fill_t_oddiy2_slide_8_conclusion(slides[7], topic)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ============================================================
# PLATINUM SHABLON (template_num=37) — Gamma Premium uslubi
# Slaydlar: 1-muqova, 2-reja, 3-4-ustunli, 4-3-ustunli,
#           5-chap rasm+o'ng matn, 6-4-blokli, 7-xulosa, 8-outro
# ============================================================

SLIDE_TYPE_NAMES_PLATINUM = {
    0: 'single_body',
    1: 'four_columns',
    2: 'three_columns',
    3: 'image_left',
    4: 'four_blocks',
}


def _pt_clear_write(shape, text, sz=None, bold=None, color=None, align=None):
    """Platinum shablon uchun shape ga matn yozish."""
    from lxml import etree

    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    txBody = tf._txBody
    nsA = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Barcha paragraflarni tozalash
    for p in txBody.findall(f'{{{nsA}}}p'):
        txBody.remove(p)

    # Yangi paragraf yaratish
    p_elem = etree.SubElement(txBody, f'{{{nsA}}}p')
    r_elem = etree.SubElement(p_elem, f'{{{nsA}}}r')
    rPr = etree.SubElement(r_elem, f'{{{nsA}}}rPr', attrib={'lang': 'uz-UZ', 'dirty': '0'})

    if sz:
        rPr.set('sz', str(int(sz * 100)))
    if bold is not None:
        rPr.set('b', '1' if bold else '0')
    if color:
        solidFill = etree.SubElement(rPr, f'{{{nsA}}}solidFill')
        etree.SubElement(solidFill, f'{{{nsA}}}srgbClr', attrib={'val': color.replace('#', '')})

    t_elem = etree.SubElement(r_elem, f'{{{nsA}}}t')
    t_elem.text = str(text) if text else ''

    if align:
        pPr = etree.Element(f'{{{nsA}}}pPr')
        algn_map = {'center': 'ctr', 'left': 'l', 'right': 'r', 'ctr': 'ctr', 'l': 'l'}
        pPr.set('algn', algn_map.get(align, 'l'))
        p_elem.insert(0, pPr)


def _pt_short(text, max_chars=45):
    """Matnni qisqartirish — sarlavha uchun."""
    if not text:
        return ''
    text = str(text).strip()
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rsplit(' ', 1)[0] + '...'


def _pt_body(text, max_chars=200):
    """Matnni qisqartirish — body uchun."""
    if not text:
        return ''
    text = str(text).strip()
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rsplit(' ', 1)[0] + '...'


def _pt_replace_blip(shape, img_path):
    """Shape ichidagi blip (rasm) ni yangi rasm bilan almashtirish."""
    try:
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        blips = shape._element.findall('.//a:blip', {'a': ns_a})
        if not blips:
            return False
        blip = blips[0]
        slide_part = shape.part
        _, new_rId = slide_part.get_or_add_image_part(img_path)
        blip.set(f'{{{ns_r}}}embed', new_rId)
        return True
    except Exception as e:
        logging.warning(f'[PT] blip almashtirish xatoligi: {e}')
    return False


def _pt_fetch_and_replace(slide, image_name, image_query):
    """Pixabay dan rasm yuklab, slide dagi image_name nomli shapega joylashtirish."""
    try:
        img_path = fetch_image(image_query)
        if not img_path:
            return
        for s in slide.shapes:
            if s.name == image_name:
                _pt_replace_blip(s, img_path)
                try:
                    import os
                    os.remove(img_path)
                except Exception:
                    pass
                return
    except Exception as e:
        logging.warning(f'[PT] rasm yuklash xatoligi ({image_query}): {e}')


def fill_platinum_slide_1_cover(slide, topic, name_surname, image_query=None):
    """1-slayd: Muqova — sarlavha (Text 0) + muallif (Text 1) + rasm (Image 0)."""
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        if s.name == 'Text 0':
            _pt_clear_write(s, topic, sz=30, bold=True, color='3B4540')
        elif s.name == 'Text 1':
            desc = name_surname if name_surname else 'Taqdimot'
            _pt_clear_write(s, desc, sz=13, bold=False, color='405449')
    if image_query:
        _pt_fetch_and_replace(slide, 'Image 0', image_query)


def fill_platinum_slide_2_plan(slide, plan):
    """2-slayd: Reja.
    - Text 23 x4: reja matni (raqamsiz, faqat matn)
    - Text 6 (=1), Text 11 (=2), Text 16 (=3), Text 21 (=4): yashil raqam bloklari
    - Shape 3/4/5, Shape 9/10, Shape 14/15, Shape 19/20: bezak shakllari
    Reja soni < 4 bo'lsa, ortiqcha yashil raqam bloklari va bezak shakllari yashiriladi.
    """
    import re
    from lxml import etree
    # plan dict yoki list bo'lishi mumkin
    if isinstance(plan, dict):
        plan_items = plan.get('content', plan.get('items', []))
        if isinstance(plan_items, str):
            plan_items = [plan_items]
    elif isinstance(plan, list):
        plan_items = plan
    else:
        plan_items = []
    # Raqamli prefikslarni tozalash
    clean_items = []
    for item in plan_items:
        if isinstance(item, str):
            item_clean = re.sub(r'^\s*\d+[\.)\-]\s*', '', item.strip())
            if item_clean:
                clean_items.append(item_clean)
    if not clean_items:
        clean_items = ['Reja mavjud emas']
    # Maksimal 3 ta element (3 tadan ortiq bo'lsa qisqartirish)
    clean_items = clean_items[:3]
    plan_count = len(clean_items)

    # Text 23 shapes ni top koordinatasi bo'yicha saralash
    text23_shapes = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.name == 'Text 23'],
        key=lambda s: s.top
    )
    for i, ts in enumerate(text23_shapes):
        if i < plan_count:
            # Faqat matn yoz, raqam qo'shma (yashil raqam bloki allaqachon shablon ichida bor)
            _pt_clear_write(ts, _pt_body(clean_items[i], 120), sz=13, bold=False, color='3B4540')
        else:
            # Bo'sh qil
            _pt_clear_write(ts, '', sz=13, bold=False, color='3B4540')

    # Yashil raqam bloklari va bezak shakllarini reja soniga moslashtirish
    # Har bir reja elementi uchun shape guruhlari:
    # 1-element: Text 6, Shape 3, Shape 4, Shape 5
    # 2-element: Text 11, Shape 9, Shape 10
    # 3-element: Text 16, Shape 14, Shape 15
    # 4-element: Text 21, Shape 19, Shape 20
    plan_shape_groups = [
        ['Text 6', 'Shape 3', 'Shape 4', 'Shape 5'],
        ['Text 11', 'Shape 9', 'Shape 10'],
        ['Text 16', 'Shape 14', 'Shape 15'],
        ['Text 21', 'Shape 19', 'Shape 20'],
    ]
    # Reja sonidan ortiqcha shape larni XML dan o'chirish
    shapes_to_remove = []
    for group_idx, group_names in enumerate(plan_shape_groups):
        if group_idx >= plan_count:
            for shape in slide.shapes:
                if shape.name in group_names:
                    shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception as e:
            import logging
            logging.warning(f'[PT plan] shape o\'chirish xatoligi {shape.name}: {e}')


def fill_platinum_slide_3_four_col(slide, title, content_data):
    """3-slayd: Text 2 sarlavha; (Text 3+Text 4), (Text 5+Text 6), (Text 7+Text 8), (Text 9+Text 10) — 4 ta ustun.
    Har ustunda: odd = sarlavha, even = tavsif.
    Shablon shape nomlari: Text 3(sarlavha), Text 4(tavsif), Text 5(sarlavha), Text 6(tavsif),
                           Text 7(sarlavha), Text 8(tavsif), Text 9(sarlavha), Text 10(tavsif).
    Rasmlar: Image 0-3 (4 ta rasm, har ustun uchun 1 ta).
    """
    items = _pt_extract_items_pairs(content_data, 4)
    name_map = {
        'Text 2': title,
        'Text 3': _pt_short(items[0][0], 35),
        'Text 4': _pt_body(items[0][1], 150),
        'Text 5': _pt_short(items[1][0], 35),
        'Text 6': _pt_body(items[1][1], 150),
        'Text 7': _pt_short(items[2][0], 35),
        'Text 8': _pt_body(items[2][1], 150),
        'Text 9': _pt_short(items[3][0], 35),
        'Text 10': _pt_body(items[3][1], 150),
    }
    title_shapes = {'Text 2', 'Text 3', 'Text 5', 'Text 7', 'Text 9'}
    for s in slide.shapes:
        if s.has_text_frame and s.name in name_map:
            bold = s.name in title_shapes
            sz = 20 if s.name == 'Text 2' else (12 if bold else 11)
            _pt_clear_write(s, name_map[s.name], sz=sz, bold=bold, color='3B4540')
    # Rasmlarni almashtirish
    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    if image_query:
        for img_name in ['Image 0', 'Image 1', 'Image 2', 'Image 3']:
            _pt_fetch_and_replace(slide, img_name, image_query)


def fill_platinum_slide_4_three_col(slide, title, content_data):
    """4-slayd: Text 0 sarlavha; (Text 2+Text 3), (Text 4+Text 5), (Text 6+Text 7) — 3 ta ustun.
    Rasmlar: Image 0-2 (3 ta rasm).
    """
    items = _pt_extract_items_pairs(content_data, 3)
    name_map = {
        'Text 0': title,
        'Text 2': _pt_short(items[0][0], 35),
        'Text 3': _pt_body(items[0][1], 180),
        'Text 4': _pt_short(items[1][0], 35),
        'Text 5': _pt_body(items[1][1], 180),
        'Text 6': _pt_short(items[2][0], 35),
        'Text 7': _pt_body(items[2][1], 180),
    }
    title_shapes = {'Text 0', 'Text 2', 'Text 4', 'Text 6'}
    for s in slide.shapes:
        if s.has_text_frame and s.name in name_map:
            bold = s.name in title_shapes
            sz = 22 if s.name == 'Text 0' else (13 if bold else 11)
            _pt_clear_write(s, name_map[s.name], sz=sz, bold=bold, color='3B4540')
    # Rasmlarni almashtirish
    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    if image_query:
        for img_name in ['Image 0', 'Image 1', 'Image 2']:
            _pt_fetch_and_replace(slide, img_name, image_query)


def fill_platinum_slide_5_image_left(slide, title, content_data):
    """5-slayd: chap rasm, o'ng tomonda Text 0 sarlavha;
    (Text 2+Text 3), (Text 4+Text 5), (Text 6+Text 7) — 3 ta band.
    Rasmlar: Image 0 (asosiy chap rasm), Image 1-3 (kichik rasmlar).
    """
    items = _pt_extract_items_pairs(content_data, 3)
    name_map = {
        'Text 0': title,
        'Text 2': _pt_short(items[0][0], 35),
        'Text 3': _pt_body(items[0][1], 120),
        'Text 4': _pt_short(items[1][0], 35),
        'Text 5': _pt_body(items[1][1], 120),
        'Text 6': _pt_short(items[2][0], 35),
        'Text 7': _pt_body(items[2][1], 120),
    }
    title_shapes = {'Text 0', 'Text 2', 'Text 4', 'Text 6'}
    for s in slide.shapes:
        if s.has_text_frame and s.name in name_map:
            bold = s.name in title_shapes
            sz = 20 if s.name == 'Text 0' else (12 if bold else 11)
            _pt_clear_write(s, name_map[s.name], sz=sz, bold=bold, color='3B4540')
    # Rasmlarni almashtirish
    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    if image_query:
        for img_name in ['Image 0', 'Image 1', 'Image 2', 'Image 3']:
            _pt_fetch_and_replace(slide, img_name, image_query)


def fill_platinum_slide_6_four_blocks(slide, title, content_data):
    """6-slayd: Text 0 sarlavha; (Text 2+Text 3), (Text 5+Text 6), (Text 8+Text 9), (Text 11+Text 12) — 4 ta blok.
    Rasmlar: Image 0-3 (4 ta rasm).
    """
    items = _pt_extract_items_pairs(content_data, 4)
    name_map = {
        'Text 0': title,
        'Text 2': _pt_short(items[0][0], 35),
        'Text 3': _pt_body(items[0][1], 150),
        'Text 5': _pt_short(items[1][0], 35),
        'Text 6': _pt_body(items[1][1], 150),
        'Text 8': _pt_short(items[2][0], 35),
        'Text 9': _pt_body(items[2][1], 150),
        'Text 11': _pt_short(items[3][0], 35),
        'Text 12': _pt_body(items[3][1], 150),
    }
    title_shapes = {'Text 0', 'Text 2', 'Text 5', 'Text 8', 'Text 11'}
    for s in slide.shapes:
        if s.has_text_frame and s.name in name_map:
            bold = s.name in title_shapes
            sz = 20 if s.name == 'Text 0' else (12 if bold else 11)
            _pt_clear_write(s, name_map[s.name], sz=sz, bold=bold, color='3B4540')
    # Rasmlarni almashtirish
    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    if image_query:
        for img_name in ['Image 0', 'Image 1', 'Image 2', 'Image 3']:
            _pt_fetch_and_replace(slide, img_name, image_query)


def fill_platinum_slide_7_conclusion(slide, title, content_data):
    """7-slayd: chap rasm, o'ng tomonda Text 0 sarlavha;
    (Text 2+Text 3), (Text 5+Text 6) — 2 ta karta blok;
    (Text 8+Text 9) — keng pastki blok; Text 11 — xulosa matni.
    Rasmlar: Image 0 (asosiy chap rasm).
    """
    items = _pt_extract_items_pairs(content_data, 3)
    # xulosa matni
    if isinstance(content_data, dict):
        conclusion = content_data.get('conclusion', content_data.get('content', ''))
        if isinstance(conclusion, list):
            conclusion = ' '.join(str(x) for x in conclusion)
        conclusion = _pt_body(str(conclusion), 200)
    elif isinstance(content_data, str):
        conclusion = _pt_body(content_data, 200)
    else:
        conclusion = title

    name_map = {
        'Text 0': title,
        'Text 2': _pt_short(items[0][0], 35),
        'Text 3': _pt_body(items[0][1], 150),
        'Text 5': _pt_short(items[1][0], 35),
        'Text 6': _pt_body(items[1][1], 150),
        'Text 8': _pt_short(items[2][0], 35),
        'Text 9': _pt_body(items[2][1], 150),
        'Text 11': conclusion,
    }
    title_shapes = {'Text 0', 'Text 2', 'Text 5', 'Text 8'}
    for s in slide.shapes:
        if s.has_text_frame and s.name in name_map:
            bold = s.name in title_shapes
            sz = 22 if s.name == 'Text 0' else (12 if bold else 11)
            _pt_clear_write(s, name_map[s.name], sz=sz, bold=bold, color='3B4540')
    # Rasmlarni almashtirish
    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    if image_query:
        _pt_fetch_and_replace(slide, 'Image 0', image_query)


def fill_platinum_slide_8_outro(slide, topic=None):
    """8-slayd: Outro — Text 1 ga E'tiboringiz uchun rahmat. Rasm: Image 0."""
    for s in slide.shapes:
        if s.has_text_frame and s.name == 'Text 1':
            _pt_clear_write(s, "E'tiboringiz uchun rahmat!", sz=28, bold=True, color='3B4540', align='center')
    if topic:
        _pt_fetch_and_replace(slide, 'Image 0', topic)


def _pt_extract_items_pairs(content_data, count):
    """content_data dan count ta (sarlavha, tavsif) juftligi ro'yxatini olish.
    Sarlavha: matnning birinchi 3-4 so'zi (max 30 belgi).
    Body: to'liq matn (max 150 belgi).
    """
    raw = []
    if isinstance(content_data, dict):
        content = content_data.get('content', content_data.get('col1', ''))
        if isinstance(content, list):
            raw = [str(x) for x in content]
        elif isinstance(content, str):
            raw = [s.strip() for s in content.replace('\n', '. ').split('. ') if s.strip()]
    elif isinstance(content_data, str):
        raw = [s.strip() for s in content_data.replace('\n', '. ').split('. ') if s.strip()]

    pairs = []
    for i in range(count):
        if i < len(raw):
            text = raw[i]
            # Sarlavha: birinchi nuqtaga qadar yoki max 30 belgi
            # Agar matnda nuqta bo'lsa, nuqtadan oldingi qism sarlavha
            if '. ' in text:
                first_sentence = text.split('. ')[0].strip()
                if len(first_sentence) <= 25:
                    title_part = first_sentence
                else:
                    # Birinchi 2-3 so'z, max 25 belgi
                    words = first_sentence.split()
                    title_part = ''
                    for w in words[:4]:
                        candidate = (title_part + ' ' + w).strip()
                        if len(candidate) <= 25:
                            title_part = candidate
                        else:
                            break
                    if not title_part:
                        title_part = words[0][:25]
            else:
                # Birinchi 2-3 so'z, max 25 belgi
                words = text.split()
                title_part = ''
                for w in words[:4]:
                    candidate = (title_part + ' ' + w).strip()
                    if len(candidate) <= 25:
                        title_part = candidate
                    else:
                        break
                if not title_part:
                    title_part = words[0][:25]
            # Body: to'liq matn, max 160 belgi
            body_part = text[:160] if len(text) > 160 else text
            pairs.append((title_part, body_part))
        else:
            # Agar element yo'q bo'lsa, mavjud elementlardan birini qayta ishlatish
            if pairs:
                # Oxirgi mavjud elementni qayta ishlatish
                pairs.append(pairs[-1])
            elif raw:
                # raw dan birinchi elementni ishlatish
                text = raw[0]
                words = text.split()
                title_part = ' '.join(words[:3])[:25] if words else 'Ma\'lumot'
                body_part = text[:160]
                pairs.append((title_part, body_part))
            else:
                pairs.append(('', ''))
    return pairs


def _pt_extract_items(content_data, count):
    """content_data dan count ta element ro'yxatini olish (legacy)."""
    items = []
    if isinstance(content_data, dict):
        content = content_data.get('content', content_data.get('col1', ''))
        if isinstance(content, list):
            items = [str(x) for x in content[:count]]
        elif isinstance(content, str):
            sentences = [s.strip() for s in content.replace('\n', '. ').split('. ') if s.strip()]
            items = sentences[:count]
    elif isinstance(content_data, str):
        sentences = [s.strip() for s in content_data.replace('\n', '. ').split('. ') if s.strip()]
        items = sentences[:count]
    # Bo'sh elementlar uchun fallback: oxirgi mavjud elementni qayta ishlatish
    while len(items) < count:
        if items:
            items.append(items[-1])
        else:
            items.append('Ma\'lumot mavjud emas')
    return items


def generate_template_platinum_presentation(prs, topic, requested_slide_count, language,
                                             name_surname, plan, content_data_list,
                                             user_images=None):
    """Platinum (Gamma uslubi) shablon asosida taqdimot yaratish."""
    import io
    import logging
    logger = logging.getLogger(__name__)

    slides = prs.slides
    n = len(slides)

    # 1-slayd: Muqova (rasm ham yuklanadi)
    if n > 0:
        cover_image_query = topic  # muqova rasmi mavzu bo'yicha
        fill_platinum_slide_1_cover(slides[0], topic, name_surname, image_query=cover_image_query)

    # 2-slayd: Reja
    if n > 1:
        fill_platinum_slide_2_plan(slides[1], plan)

    # Kontent slaydlari (3-7, index 2-6)
    fill_funcs = [
        fill_platinum_slide_3_four_col,    # 0
        fill_platinum_slide_4_three_col,   # 1
        fill_platinum_slide_5_image_left,  # 2
        fill_platinum_slide_6_four_blocks, # 3
        fill_platinum_slide_7_conclusion,  # 4
    ]

    for i, data in enumerate(content_data_list):
        slide_idx = i + 2
        if slide_idx >= n - 1:
            break
        slide = slides[slide_idx]
        func = fill_funcs[i % len(fill_funcs)]
        title = data.get('title', topic) if isinstance(data, dict) else topic
        try:
            func(slide, title, data)
        except Exception as e:
            logger.warning(f"platinum slide {slide_idx} fill xatolik: {e}")

    # 8-slayd: Outro (rasm ham yuklanadi)
    if n > 0:
        fill_platinum_slide_8_outro(slides[-1], topic=topic)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════
# GAMMA2 SHABLON (Stil_gamma2) — 8 slaydli qorong'i uslub
# ═══════════════════════════════════════════════════════════════

SLIDE_TYPE_NAMES_GAMMA2 = {
    0: 'four_blocks_2x2',   # 3-slayd
    1: 'numbered_list_4',   # 4-slayd
    2: 'icon_list_3',       # 5-slayd
    3: 'icon_list_3_large', # 6-slayd
    4: 'two_plus_one',      # 7-slayd
}


def _g2_clear_write(shape, text, sz=None, bold=None, color=None, align=None, max_chars=None, line_spacing=None):
    """Gamma2 shablon uchun shape ga matn yozish."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
    # Matn uzunligini cheklash
    if max_chars and text and len(text) > max_chars:
        text = text[:max_chars].rsplit(' ', 1)[0] + '...'
    tf = shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.clear()
    else:
        p_elem = etree.SubElement(tf._txBody, qn('a:p'))
        p = tf.paragraphs[0]
    p = tf.paragraphs[0]
    p.clear()
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    # Line spacing (qatorlar orasidagi interval)
    if line_spacing is not None:
        pPr = p._pPr
        if pPr is None:
            pPr = etree.SubElement(p._p, qn('a:pPr'))
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = lnSpc.find(qn('a:spcPct'))
        if spcPct is None:
            spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))
    run = p.add_run()
    run.text = text
    if sz:
        run.font.size = Pt(sz)
    if bold is not None:
        run.font.bold = bold
    if color:
        try:
            run.font.color.rgb = RGBColor.from_string(color)
        except Exception:
            pass


def _g2_fetch_and_replace(slide, shape_name, query):
    """Gamma2 shablon uchun shape dagi rasmni Pixabay dan yuklab almashtirish."""
    import logging
    import io
    from pptx.util import Emu
    logger = logging.getLogger(__name__)
    pixabay_key = os.environ.get('PIXABAY_API_KEY', '')
    if not pixabay_key:
        logger.warning("PIXABAY_API_KEY yo'q. Rasm o'tkazib yuborildi.")
        return
    try:
        img_data = fetch_image(query)
        if not img_data:
            logger.warning(f"[G2] Rasm yuklanmadi: {query!r}")
            return
        for shape in slide.shapes:
            if shape.name == shape_name:
                pic = shape._element
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                blip = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                if blip is None:
                    blip = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}blip')
                if blip is None:
                    logger.warning(f"[G2] blip topilmadi: {shape_name}")
                    return
                rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if not rId:
                    logger.warning(f"[G2] rId topilmadi: {shape_name}")
                    return
                part = slide.part
                img_part = part.related_parts[rId]
                img_part._blob = img_data
                logger.info(f"[G2] Rasm almashtirildi: {shape_name}")
                return
    except Exception as e:
        logger.warning(f"[G2] Rasm almashtirish xatoligi {shape_name}: {e}")


def fill_gamma2_slide_1_cover(slide, topic, name_surname, image_query=None):
    """1-slayd: Muqova — Text 0 (sarlavha), Text 1 (tavsif), Image 0 (o'ng rasm)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.name == 'Text 0':
                _g2_clear_write(shape, topic, sz=36, bold=True, color='76B9FF')
            elif shape.name == 'Text 1':
                subtitle = name_surname if name_surname else topic
                _g2_clear_write(shape, subtitle, sz=14, bold=False, color='D6E5EF')
    if image_query:
        _g2_fetch_and_replace(slide, 'Image 0', image_query)


def fill_gamma2_slide_2_plan(slide, plan):
    """2-slayd: Reja — Text 0 (REJA), Text 4/8/12 (reja matnlari), Shape+Text raqamlar.
    Shablon 3 ta reja elementi uchun mo'ljallangan.
    Reja soni 3 dan kam bo'lsa, ortiqcha shape lar o'chiriladi.
    """
    import re
    # plan dict yoki list bo'lishi mumkin
    if isinstance(plan, dict):
        plan_items = plan.get('content', plan.get('items', []))
        if isinstance(plan_items, str):
            plan_items = [plan_items]
    elif isinstance(plan, list):
        plan_items = plan
    else:
        plan_items = []
    # Raqamli prefikslarni tozalash
    clean_items = []
    for item in plan_items:
        if isinstance(item, str):
            item_clean = re.sub(r'^\s*\d+[\.)\-]\s*', '', item.strip())
            if item_clean:
                clean_items.append(item_clean)
    if not clean_items:
        clean_items = ['Reja mavjud emas']
    # Maksimal 3 ta element (3 tadan ortiq bo'lsa qisqartirish)
    clean_items = clean_items[:3]
    plan_count = len(clean_items)

    # Reja matni shape lari: Text 4, Text 8, Text 12 (top bo'yicha saralangan)
    plan_text_shapes = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.name in ('Text 4', 'Text 8', 'Text 12')],
        key=lambda s: s.top
    )
    for i, ts in enumerate(plan_text_shapes):
        if i < plan_count:
            _g2_clear_write(ts, clean_items[i], sz=13, bold=False, color='D6E5EF')
        else:
            _g2_clear_write(ts, '', sz=13, bold=False, color='D6E5EF')

    # Reja raqam bloklari guruhlari: har biri (Shape, Text_raqam) juftligi
    # 1-element: Shape 1, Text 2
    # 2-element: Shape 5, Text 6
    # 3-element: Shape 9, Text 10
    plan_shape_groups = [
        ['Shape 1', 'Text 2'],
        ['Shape 5', 'Text 6'],
        ['Shape 9', 'Text 10'],
    ]
    shapes_to_remove = []
    for group_idx, group_names in enumerate(plan_shape_groups):
        if group_idx >= plan_count:
            for shape in slide.shapes:
                if shape.name in group_names:
                    shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        try:
            shape._element.getparent().remove(shape._element)
        except Exception as e:
            import logging
            logging.warning(f'[G2 plan] shape o\'chirish xatoligi {shape.name}: {e}')


def fill_gamma2_slide_3_four_blocks(slide, title, content_data):
    """3-slayd: 2x2 grid — 4 ta blok.
    Shape 1+Text 2+Text 3, Shape 4+Text 5+Text 6,
    Shape 7+Text 8+Text 9, Shape 10+Text 11+Text 12.
    """
    # Sarlavha
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'Text 0':
            _g2_clear_write(shape, title, sz=22, bold=True, color='76B9FF')

    # 4 ta blok uchun (sarlavha, tavsif) juftliklari
    pairs = _pt_extract_items_pairs(content_data, 4)

    block_map = [
        ('Text 2', 'Text 3'),   # 1-blok
        ('Text 5', 'Text 6'),   # 2-blok
        ('Text 8', 'Text 9'),   # 3-blok
        ('Text 11', 'Text 12'), # 4-blok
    ]
    for i, (title_name, body_name) in enumerate(block_map):
        t, b = pairs[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.name == title_name:
                    _g2_clear_write(shape, t, sz=13, bold=True, color='D6E5EF', max_chars=40)
                elif shape.name == body_name:
                    _g2_clear_write(shape, b, sz=11, bold=False, color='D6E5EF', max_chars=110)

    # Rasm almashtirish
    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    _g2_fetch_and_replace(slide, 'Image 0', image_query)


def fill_gamma2_slide_4_numbered_list(slide, title, content_data):
    """4-slayd: Raqamli ro'yxat — 4 ta element.
    Shablon tahlilidan:
    - Text 5 (sarlavha 1), Text 6 (tavsif 1)
    - Text 10 (sarlavha 2), Text 11 (tavsif 2)
    - Text 15 (sarlavha 3), Text 16 (tavsif 3)
    - Text 19 (sarlavha 4), Text 20 (tavsif 4) — agar mavjud bo'lsa
    Raqamlar (Text 4='1', Text 9='2', Text 14='3', Text 19='4') o'zgartirilmaydi.
    """
    # Sarlavha
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'Text 0':
            _g2_clear_write(shape, title, sz=22, bold=True, color='76B9FF')

    pairs = _pt_extract_items_pairs(content_data, 4)

    # Shablon tahlilidan aniq shape nomlar:
    # Text 5 sarlavha, Text 6 tavsif (1-element)
    # Text 10 sarlavha, Text 11 tavsif (2-element)
    # Text 15 sarlavha, Text 16 tavsif (3-element)
    # Text 19 yoki Text 20 — 4-element uchun (agar mavjud)
    # Lekin shablon faqat 3 ta element uchun mo'ljallangan (Text 5/6, Text 10/11, Text 15/16)
    block_map = [
        ('Text 5', 'Text 6'),
        ('Text 10', 'Text 11'),
        ('Text 15', 'Text 16'),
    ]
    for i, (title_name, body_name) in enumerate(block_map):
        if i >= len(pairs):
            break
        t, b = pairs[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.name == title_name:
                    _g2_clear_write(shape, t, sz=13, bold=True, color='D6E5EF')
                elif shape.name == body_name:
                    _g2_clear_write(shape, b, sz=11, bold=False, color='D6E5EF')

    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    _g2_fetch_and_replace(slide, 'Image 0', image_query)


def fill_gamma2_slide_5_icon_list(slide, title, content_data):
    """5-slayd: 3-ikonali ro'yxat — Image 1/2/3 (ikonlar), Text 1+Text 2, Text 3+Text 4, Text 5+Text 6."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'Text 0':
            _g2_clear_write(shape, title, sz=22, bold=True, color='76B9FF')

    pairs = _pt_extract_items_pairs(content_data, 3)

    block_map = [
        ('Text 1', 'Text 2'),
        ('Text 3', 'Text 4'),
        ('Text 5', 'Text 6'),
    ]
    for i, (title_name, body_name) in enumerate(block_map):
        t, b = pairs[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.name == title_name:
                    _g2_clear_write(shape, t, sz=13, bold=True, color='D6E5EF', max_chars=45)
                elif shape.name == body_name:
                    _g2_clear_write(shape, b, sz=11, bold=False, color='D6E5EF', max_chars=90)

    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    _g2_fetch_and_replace(slide, 'Image 0', image_query)


def fill_gamma2_slide_6_icon_list_large(slide, title, content_data):
    """6-slayd: 3-ikonali katta ro'yxat — Image 1/2/3 (katta ikonlar), Text 1+Text 2, Text 3+Text 4, Text 5+Text 6."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'Text 0':
            _g2_clear_write(shape, title, sz=22, bold=True, color='76B9FF')

    pairs = _pt_extract_items_pairs(content_data, 3)

    block_map = [
        ('Text 1', 'Text 2'),
        ('Text 3', 'Text 4'),
        ('Text 5', 'Text 6'),
    ]
    for i, (title_name, body_name) in enumerate(block_map):
        t, b = pairs[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.name == title_name:
                    _g2_clear_write(shape, t, sz=13, bold=True, color='D6E5EF', max_chars=45)
                elif shape.name == body_name:
                    _g2_clear_write(shape, b, sz=11, bold=False, color='D6E5EF', max_chars=90)

    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    _g2_fetch_and_replace(slide, 'Image 0', image_query)


def fill_gamma2_slide_7_two_plus_one(slide, title, content_data):
    """7-slayd: 2+1 blokli — Image 0 (chap), Text 0 (sarlavha o'ng),
    Shape 1+Text 2+Text 3, Shape 4+Text 5+Text 6 (yuqori 2 blok),
    Shape 7+Text 8+Text 9 (keng pastki blok).
    """
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'Text 0':
            _g2_clear_write(shape, title, sz=22, bold=True, color='76B9FF')

    pairs = _pt_extract_items_pairs(content_data, 3)

    block_map = [
        ('Text 2', 'Text 3'),
        ('Text 5', 'Text 6'),
        ('Text 8', 'Text 9'),
    ]
    for i, (title_name, body_name) in enumerate(block_map):
        t, b = pairs[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.name == title_name:
                    _g2_clear_write(shape, t, sz=13, bold=True, color='D6E5EF', max_chars=50, line_spacing=1.0)
                elif shape.name == body_name:
                    _g2_clear_write(shape, b, sz=11, bold=False, color='D6E5EF', max_chars=130, line_spacing=1.0)

    image_query = content_data.get('image_query', title) if isinstance(content_data, dict) else title
    _g2_fetch_and_replace(slide, 'Image 0', image_query)


def fill_gamma2_slide_8_outro(slide, topic=None):
    """8-slayd: Outro — Text 0 ga E'tiboringiz uchun rahmat. Image 0 rasm."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'Text 0':
            tf = shape.text_frame
            tf.word_wrap = True
            # 2 qatorli matn: "E'tiboringiz uchun" va "rahmat!"
            from pptx.dml.color import RGBColor
            from pptx.util import Pt
            from pptx.enum.text import PP_ALIGN
            for para in tf.paragraphs:
                para.clear()
            # Birinchi paragraf
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.LEFT
            r0 = p0.add_run()
            r0.text = "E'tiboringiz uchun"
            r0.font.size = Pt(36)
            r0.font.bold = True
            r0.font.color.rgb = RGBColor.from_string('76B9FF')
            # Ikkinchi paragraf
            from pptx.oxml.ns import qn
            from lxml import etree
            p1_elem = etree.SubElement(tf._txBody, qn('a:p'))
            # Yangi paragraf qo'shish
            from pptx.text.text import _Paragraph
            p1 = tf.paragraphs[-1]
            p1.alignment = PP_ALIGN.LEFT
            r1 = p1.add_run()
            r1.text = "rahmat!"
            r1.font.size = Pt(36)
            r1.font.bold = True
            r1.font.color.rgb = RGBColor.from_string('76B9FF')
    if topic:
        _g2_fetch_and_replace(slide, 'Image 0', topic)


def generate_template_gamma2_presentation(prs, topic, requested_slide_count, language,
                                          name_surname, plan, content_data_list,
                                          user_images=None):
    """Gamma2 (Stil_gamma2) shablon asosida taqdimot yaratish."""
    import io
    import logging
    logger = logging.getLogger(__name__)

    slides = prs.slides
    n = len(slides)

    # 1-slayd: Muqova
    if n > 0:
        fill_gamma2_slide_1_cover(slides[0], topic, name_surname, image_query=topic)

    # 2-slayd: Reja
    if n > 1:
        fill_gamma2_slide_2_plan(slides[1], plan)

    # Kontent slaydlari (3-7, index 2-6)
    fill_funcs = [
        fill_gamma2_slide_3_four_blocks,    # 0
        fill_gamma2_slide_4_numbered_list,  # 1
        fill_gamma2_slide_5_icon_list,      # 2
        fill_gamma2_slide_6_icon_list_large,# 3
        fill_gamma2_slide_7_two_plus_one,   # 4
    ]

    for i, data in enumerate(content_data_list):
        slide_idx = i + 2
        if slide_idx >= n - 1:
            break
        slide = slides[slide_idx]
        func = fill_funcs[i % len(fill_funcs)]
        title = data.get('title', topic) if isinstance(data, dict) else topic
        try:
            func(slide, title, data)
        except Exception as e:
            logger.warning(f"gamma2 slide {slide_idx} fill xatolik: {e}")

    # 8-slayd: Outro
    if n > 0:
        fill_gamma2_slide_8_outro(slides[-1], topic=topic)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
