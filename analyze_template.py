
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_pptx_template(template_path):
    prs = Presentation(template_path)
    print(f"Analyzing template: {template_path}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        print(f"Layout: {slide.slide_layout.name}")
        
        # Analyze shapes (text boxes, placeholders, etc.)
        for shape in slide.shapes:
            print(f"  Shape ID: {shape.shape_id}, Name: {shape.name}, Type: {shape.shape_type}")
            
            if shape.has_text_frame:
                text_frame = shape.text_frame
                print(f"    Text Frame: {text_frame.text.strip()}")
                print(f"    Number of paragraphs: {len(text_frame.paragraphs)}")
                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    print(f"      Paragraph {para_idx+1}: {paragraph.text.strip()}")
            
            if shape.is_placeholder:
                print(f"    Placeholder: {shape.placeholder_format.idx}, Type: {shape.placeholder_format.type}")
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"    Image: {shape.image.filename}")
            
            if shape.has_table:
                print(f"    Table: {shape.table.rows} rows, {shape.table.columns} columns")

# Path to the uploaded template
template_file = "/home/ubuntu/upload/pasted_file_CgwVm7_1.pptx"
analyze_pptx_template(template_file)
