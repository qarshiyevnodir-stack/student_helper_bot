"""Together retry va Gamma rasm almashtirishining qat'iyligini API chaqiruvisiz sinaydi."""
import base64
import hashlib
import os
from pathlib import Path
import sys
import tempfile
import zipfile

from PIL import Image
from pptx import Presentation
import requests

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
os.environ["TOGETHER_API_KEY_1"] = "test-key-one"
os.environ["TOGETHER_API_KEY_2"] = "test-key-two"
os.environ["TOGETHER_API_KEY_3"] = "test-key-three"
import utils

# utils importi .env faylini yuklashi mumkin; test navbati faqat xavfsiz test
# qiymatlari bilan sinovdan o'tishi uchun ularni importdan keyin qayta o'rnatamiz.
os.environ["TOGETHER_API_KEY_1"] = "test-key-one"
os.environ["TOGETHER_API_KEY_2"] = "test-key-two"
os.environ["TOGETHER_API_KEY_3"] = "test-key-three"


class FakeResponse:
    def __init__(self, status_code, image_bytes=b""):
        self.status_code = status_code
        self.headers = {}
        self._image_bytes = image_bytes

    def raise_for_status(self):
        if self.status_code >= 400:
            exc = requests.HTTPError(f"{self.status_code} test error")
            exc.response = self
            raise exc

    def json(self):
        return {"data": [{"b64_json": base64.b64encode(self._image_bytes).decode()}]}


def make_test_image_bytes():
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
        temp_path = Path(temp_file.name)
    try:
        Image.new("RGB", (64, 48), (210, 30, 80)).save(temp_path)
        return temp_path.read_bytes()
    finally:
        temp_path.unlink(missing_ok=True)


def test_retry_then_success(image_bytes):
    calls = []
    responses = [FakeResponse(429), FakeResponse(503), FakeResponse(200, image_bytes)]
    original_post = utils.requests.post
    original_sleep = utils.time.sleep
    original_slot = utils._wait_for_together_slot
    original_next_key = utils._get_next_together_key
    # Fetch funksiyasi kalit tanlashini to'liq mock qilamiz; test hech qachon
    # haqiqiy Together kalitini ishlatmaydi yoki tarmoqqa so'rov yubormaydi.
    keys = iter(["test-key-one", "test-key-two", "test-key-three"])
    utils._get_next_together_key = lambda: next(keys)
    utils.requests.post = lambda _url, headers, **_kwargs: calls.append(headers["Authorization"]) or responses.pop(0)
    utils.time.sleep = lambda _seconds: None
    utils._wait_for_together_slot = lambda: None
    try:
        image_path = utils.fetch_image_together("test topic", style="illustration", topic="test topic")
        assert image_path and Path(image_path).exists()
        assert Path(image_path).read_bytes() == image_bytes
        # 429 va 503 dan so'ng uchinchi urinishda rasm olindi.
        # Kalit qiymatlari ataylab tekshiruv chiqishiga yozilmaydi.
        assert len(calls) == 3
        assert all(header.startswith("Bearer ") for header in calls)
        Path(image_path).unlink(missing_ok=True)
    finally:
        utils.requests.post = original_post
        utils.time.sleep = original_sleep
        utils._wait_for_together_slot = original_slot
        utils._get_next_together_key = original_next_key


def test_gamma_blob_replace_and_fail_closed(image_bytes):
    template = ROOT / "templates" / "shablonlar" / "gamma2.pptx"
    original_fetch = utils.fetch_image_together
    created = []

    def fake_success(*_args, **_kwargs):
        temp_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        temp_file.write(image_bytes)
        temp_file.close()
        created.append(temp_file.name)
        return temp_file.name

    try:
        utils.fetch_image_together = fake_success
        presentation = Presentation(template)
        utils._gamma_place_image(presentation.slides[0], "Image 0", "test topic", "cover", "illustration")
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as output_file:
            output_path = Path(output_file.name)
        presentation.save(output_path)
        with zipfile.ZipFile(output_path) as archive:
            output_images = [archive.read(name) for name in archive.namelist() if name.startswith("ppt/media/")]
        assert hashlib.sha256(image_bytes).hexdigest() in {hashlib.sha256(item).hexdigest() for item in output_images}
        output_path.unlink(missing_ok=True)

        utils.fetch_image_together = lambda *_args, **_kwargs: None
        failing_presentation = Presentation(template)
        try:
            utils._gamma_place_image(failing_presentation.slides[0], "Image 0", "test topic", "cover", "illustration")
            raise AssertionError("GammaImageGenerationError kutilgan edi")
        except utils.GammaImageGenerationError:
            pass

        # To'liq Gamma2 generatori ham eski shablon rasmlari bilan bytes qaytarmasligi kerak.
        try:
            utils.generate_template_gamma2_presentation(
                prs=Presentation(template), topic="test topic", requested_slide_count=5,
                language="uz", name_surname="Test", plan=["Birinchi reja"],
                content_data_list=[{"title": "Test bo'limi", "content": "Test matni"}],
            )
            raise AssertionError("Gamma2 generatori AI rasm xatosida to'xtashi kerak edi")
        except utils.GammaImageGenerationError:
            pass
    finally:
        utils.fetch_image_together = original_fetch
        for path in created:
            Path(path).unlink(missing_ok=True)


image_bytes = make_test_image_bytes()
test_retry_then_success(image_bytes)
test_gamma_blob_replace_and_fail_closed(image_bytes)
print("GAMMA_IMAGE_RELIABILITY_TEST_OK")
