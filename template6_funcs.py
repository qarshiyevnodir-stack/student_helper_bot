CONTENT_SLIDE_TEMPLATE_INDICES_6 = [2, 3, 4, 5, 6]

def build_slide_structure_6(prs, requested_content_count):
    """
    6-shablon uchun slayd tuzilmasini quradi.
    duplicate_slide va move_slide ishlatadi.
    """
    import logging
    from utils import duplicate_slide, move_slide
    
    full_repeats = max(1, round(requested_content_count / 5))
    total_content_slides = full_repeats * 5
    logging.info(f"[T6] Kontent slaydlari: {requested_content_count} so'raldi, "
                 f"{full_repeats} marta takrorlanadi ({total_content_slides} ta kontent slayd)")
    extra_sets_needed = full_repeats - 1
    for set_num in range(extra_sets_needed):
        for slide_template_idx in CONTENT_SLIDE_TEMPLATE_INDICES_6:
            duplicate_slide(prs, slide_template_idx)
        logging.info(f"  [T6] {set_num + 2}-to'plam qo'shildi. Jami slaydlar: {len(prs.slides)}")
    
    # Xulosa slaydini (index 7) oxiriga ko'chirish
    conclusion_current_index = 7
    last_index = len(prs.slides) - 1
    move_slide(prs, conclusion_current_index, last_index)
    logging.info(f"[T6] Yakuniy tuzilma: {len(prs.slides)} ta slayd")
    return total_content_slides

def fill_t6_slide_1_cover(slide, topic, name_surname):
    """
    6-Shablon Slayd 1 — Muqova.
    Shape[1]: Sarlavha (F3F0DF rang)
    Shape[2]: Subtitle (CDD6E2 rang)
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = topic.upper()
        run.font.size = Pt(52.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xF3, 0xF0, 0xDF)
        
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name_surname
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0xCD, 0xD6, 0xE2)

def fill_t6_slide_2_plan(slide, plan_data):
    """
    6-Shablon Slayd 2 — Reja.
    Shape[1]: "Reja:"
    Shape[2]: Reja bandlari
    """
    from lxml import etree
    from utils import calc_body_font_pt
    
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    titles = plan_data.get("content", plan_data.get("titles", []))
    
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.word_wrap = True
        total_chars = sum(len(t) for t in titles)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=22)
        
        txBody = tf._txBody
        for p_elem in txBody.findall(f'{{{ns_a}}}p'):
            txBody.remove(p_elem)
            
        for title in titles:
            safe_title = title.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p_xml = (
                f'<a:p xmlns:a="{ns_a}">' 
                f'<a:pPr algn="l"><a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="uz-UZ" sz="{int(font_pt*100)}" dirty="0">'
                f'<a:solidFill><a:srgbClr val="334155"/></a:solidFill></a:rPr>'
                f'<a:t>{safe_title}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)

def fill_t6_slide_3_text(slide, content_data):
    """
    6-Shablon Slayd 3.
    Shape[3]: Matn
    Shape[4]: Sarlavha
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from utils import calc_body_font_pt
    
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        tf = slide.shapes[4].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        tf = slide.shapes[3].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=24)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

def fill_t6_slide_4_text(slide, content_data):
    """
    6-Shablon Slayd 4.
    Shape[1]: Sarlavha
    Shape[2]: Matn
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from utils import calc_body_font_pt
    
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(33.75)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=24)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

def fill_t6_slide_5_three_steps(slide, content_data):
    """
    6-Shablon Slayd 5.
    Shape[1], [2], [3]: 3 ta qadam
    Shape[4]: Sarlavha
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from utils import calc_body_font_pt
    
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    
    if len(slide.shapes) > 4 and slide.shapes[4].has_text_frame:
        tf = slide.shapes[4].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    for i in range(3):
        shape_idx = i + 1
        if len(slide.shapes) > shape_idx and slide.shapes[shape_idx].has_text_frame:
            text = items[i] if i < len(items) else ""
            tf = slide.shapes[shape_idx].text_frame
            tf.clear()
            tf.word_wrap = True
            font_pt = calc_body_font_pt(len(text), base_pt=14, min_pt=10, max_pt=18)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

def fill_t6_slide_6_two_cols(slide, content_data):
    """
    6-Shablon Slayd 6.
    Shape[1], [2]: 2 ta matn
    Shape[3]: Sarlavha
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from utils import calc_body_font_pt
    
    title = content_data.get("title", "")
    col1 = content_data.get("col1", [])
    col2 = content_data.get("col2", [])
    if not col1 and not col2:
        items = content_data.get("content", [])
        mid = len(items) // 2
        col1 = items[:mid] if mid > 0 else items
        col2 = items[mid:] if mid > 0 else []
        
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        tf = slide.shapes[3].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    for i, col_data in enumerate([col1, col2]):
        shape_idx = i + 1
        if len(slide.shapes) > shape_idx and slide.shapes[shape_idx].has_text_frame:
            tf = slide.shapes[shape_idx].text_frame
            tf.clear()
            tf.word_wrap = True
            text = "\n".join(col_data)
            font_pt = calc_body_font_pt(len(text), base_pt=14, min_pt=10, max_pt=18)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

def fill_t6_slide_7_text(slide, content_data):
    """
    6-Shablon Slayd 7.
    Shape[1]: Sarlavha
    Shape[3]: Matn
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from utils import calc_body_font_pt
    
    title = content_data.get("title", "")
    items = content_data.get("content", [])
    
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x50, 0x88)
        
    if len(slide.shapes) > 3 and slide.shapes[3].has_text_frame:
        tf = slide.shapes[3].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=14, min_pt=10, max_pt=18)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

def fill_t6_slide_8_conclusion(slide, conclusion_data):
    """
    6-Shablon Slayd 8 — Xulosa.
    Shape[1]: "XULOSA" (katta matn)
    Shape[2]: Xulosa matni
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from utils import calc_body_font_pt
    
    items = conclusion_data.get("content", [])
    if not items:
        items = conclusion_data.get("col1", []) + conclusion_data.get("col2", [])
        
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "XULOSA"
        run.font.size = Pt(60)
        run.font.color.rgb = RGBColor(0xCD, 0xD6, 0xE2)
        
    if len(slide.shapes) > 2 and slide.shapes[2].has_text_frame:
        tf = slide.shapes[2].text_frame
        tf.clear()
        tf.word_wrap = True
        total_chars = sum(len(s) for s in items)
        font_pt = calc_body_font_pt(total_chars, base_pt=18, min_pt=12, max_pt=24)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0xCD, 0xD6, 0xE2)

def generate_template_6_presentation(prs, topic, requested_slide_count, language,
                                      name_surname, plan, content_data_list):
    """
    6-shablon asosida to'liq prezentatsiya yaratadi.
    """
    import io
    import logging
    from utils import generate_slide_content
    
    total_content_slides = build_slide_structure_6(prs, requested_slide_count)
    plan_dict = plan if isinstance(plan, dict) else {}
    
    # Slayd 1 — Muqova
    fill_t6_slide_1_cover(prs.slides[0], topic, name_surname)
    
    # Slayd 2 — Reja
    fill_t6_slide_2_plan(prs.slides[1], plan_dict)
    
    # Kontent slaydlari (3-dan boshlab)
    for i in range(total_content_slides):
        slide_index = i + 2
        if slide_index >= len(prs.slides) - 1:
            break
            
        slide = prs.slides[slide_index]
        data = content_data_list[i] if i < len(content_data_list) else {}
        
        slide_type = i % 5
        if slide_type == 0:
            fill_t6_slide_3_text(slide, data)
        elif slide_type == 1:
            fill_t6_slide_4_text(slide, data)
        elif slide_type == 2:
            fill_t6_slide_5_three_steps(slide, data)
        elif slide_type == 3:
            fill_t6_slide_6_two_cols(slide, data)
        elif slide_type == 4:
            fill_t6_slide_7_text(slide, data)
            
        logging.info(f"  [T6] Slayd {slide_index + 1} to'ldirildi (tur {slide_type}): {data.get('title', '')}")
        
    # Xulosa slayd
    conclusion_slide = prs.slides[-1]
    conclusion_data = generate_slide_content(topic, requested_slide_count, requested_slide_count, language, is_conclusion=True)
    if not conclusion_data:
        conclusion_data = {"title": "Xulosa", "content": ["Asosiy xulosalar", "Tavsiyalar"]}
    fill_t6_slide_8_conclusion(conclusion_slide, conclusion_data)
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
